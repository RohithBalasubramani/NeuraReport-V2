"""Consolidated AI & Domain Services (Phase B6).

Merged from: summary, recommendations, docqa, federation, nl2sql,
synthesis, analytics, visualization, search, design.
"""
from __future__ import annotations

import io
import json
import re
import logging
import time
import uuid
import zipfile
from datetime import datetime, timezone
from typing import Any, Callable, Dict, List, Optional

from backend.app.common import utc_now, utc_now_iso, get_state_store, strip_code_fences
from backend.app.services.llm import get_llm_client
from backend.app.repositories import state_store as state_store_module

logger = logging.getLogger("neura.domain.summary")




class SummaryService:
    """Service for executive summary generation."""

    def __init__(self):
        self._llm_client = None

    def _get_llm_client(self):
        if self._llm_client is None:
            self._llm_client = get_llm_client()
        return self._llm_client

    def generate_summary(
        self,
        content: str,
        tone: str = "formal",
        max_sentences: int = 5,
        focus_areas: Optional[List[str]] = None,
        correlation_id: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Generate an executive summary from content.

        Args:
            content: The content to summarize
            tone: Style of summary (formal, conversational, technical)
            max_sentences: Maximum sentences in summary
            focus_areas: Optional areas to focus on
            correlation_id: Request correlation ID

        Returns:
            Summary with key findings and metrics
        """
        logger.info("Generating executive summary", extra={"correlation_id": correlation_id})

        prompt = f"""Generate an executive summary of the following content.

CONTENT:
{content[:8000]}  # Limit content length

REQUIREMENTS:
- Tone: {tone}
- Maximum sentences: {max_sentences}
- Style: Professional, data-driven
{f"- Focus on: {', '.join(focus_areas)}" if focus_areas else ""}

Return a JSON object:
{{
  "executive_summary": "2-3 sentence overview",
  "key_findings": ["finding 1", "finding 2", "finding 3"],
  "metrics": [
    {{"name": "metric_name", "value": "123", "unit": "USD", "trend": "up"}}
  ],
  "recommendations": ["action 1", "action 2"],
  "confidence": 0.9
}}

Return ONLY the JSON object."""

        try:
            client = self._get_llm_client()
            response = client.complete(
                messages=[{"role": "user", "content": prompt}],
                description="executive_summary",
                temperature=0.3,
            )

            import json
            import re
            content_response = response["choices"][0]["message"]["content"]
            json_match = re.search(r"\{[\s\S]*\}", content_response)
            if json_match:
                return json.loads(json_match.group())

        except Exception as exc:
            logger.error(f"Summary generation failed: {exc}")
            error_message = "Summary generation failed"
        else:
            error_message = "Unknown error"

        return {
            "executive_summary": "Summary generation failed",
            "key_findings": [],
            "metrics": [],
            "recommendations": [],
            "confidence": 0.0,
            "error": error_message,
        }

    def generate_report_summary(
        self,
        report_id: str,
        correlation_id: Optional[str] = None,
    ) -> Dict[str, Any]:
        """Generate summary for a specific report."""
        # Get report data from state store
        store = get_state_store()
        with store.transaction() as state:
            runs = state.get("runs", {})
        report = runs.get(report_id)

        if not report:
            return {"error": "Report not found"}

        # Extract content from report
        content = str(report)
        return self.generate_summary(content, correlation_id=correlation_id)


# RecommendationService


"""Service for template recommendations using AI."""


from backend.app.repositories import get_connection_schema

logger = logging.getLogger("neura.domain.recommendations")


class RecommendationService:
    """Service for template recommendations."""

    def __init__(self):
        self._llm_client = None

    def _get_llm_client(self):
        if self._llm_client is None:
            self._llm_client = get_llm_client()
        return self._llm_client

    def recommend_templates(
        self,
        connection_id: Optional[str] = None,
        schema_info: Optional[Dict[str, Any]] = None,
        context: Optional[str] = None,
        limit: int = 5,
        correlation_id: Optional[str] = None,
    ) -> List[Dict[str, Any]]:
        """
        Recommend templates based on connection schema and context.

        Args:
            connection_id: Optional connection to base recommendations on
            schema_info: Optional schema information
            context: Optional user context/requirements
            limit: Maximum number of recommendations
            correlation_id: Request correlation ID

        Returns:
            List of template recommendations with scores
        """
        logger.info("Generating template recommendations", extra={"correlation_id": correlation_id})

        # Get all approved templates
        store = get_state_store()
        with store.transaction() as state:
            templates = state.get("templates", {})
        approved = [t for t in templates.values() if t.get("status") == "approved"]

        if not approved:
            return []

        # Get schema if connection provided
        if connection_id and not schema_info:
            try:
                schema_info = get_connection_schema(connection_id, include_row_counts=False)
            except Exception as e:
                logger.warning("Failed to get schema for connection %s: %s", connection_id, e)

        # Build recommendation prompt
        template_catalog = []
        for t in approved:
            template_catalog.append({
                "id": t.get("id"),
                "name": t.get("name"),
                "kind": t.get("kind"),
                "tags": t.get("tags", []),
            })

        prompt = f"""Recommend templates from this catalog based on the user's needs.

TEMPLATE CATALOG:
{template_catalog}

"""
        if schema_info:
            tables = [t["name"] for t in schema_info.get("tables", [])]
            prompt += f"DATABASE TABLES: {', '.join(tables)}\n\n"

        if context:
            prompt += f"USER CONTEXT: {context}\n\n"

        prompt += f"""Return a JSON array of the top {limit} recommended templates:
[
  {{
    "template_id": "id",
    "score": 0.95,
    "reason": "Why this template matches"
  }}
]

Return ONLY the JSON array."""

        try:
            client = self._get_llm_client()
            response = client.complete(
                messages=[{"role": "user", "content": prompt}],
                description="template_recommendations",
                temperature=0.3,
            )

            content = response["choices"][0]["message"]["content"]
            json_match = re.search(r"\[[\s\S]*\]", content)
            if json_match:
                recommendations = json.loads(json_match.group())
                # Enrich with template details
                for rec in recommendations:
                    tid = rec.get("template_id")
                    if tid in templates:
                        rec["template"] = templates[tid]
                return recommendations[:limit]

        except Exception as exc:
            logger.error(f"Recommendation generation failed: {exc}")

        # Fallback: return most recent templates
        sorted_templates = sorted(approved, key=lambda t: t.get("created_at", ""), reverse=True)
        return [{"template_id": t["id"], "template": t, "score": 0.5, "reason": "Recently created"} for t in sorted_templates[:limit]]

    def get_similar_templates(self, template_id: str, limit: int = 3) -> List[Dict[str, Any]]:
        """Get templates similar to a given template."""
        store = get_state_store()
        with store.transaction() as state:
            templates = state.get("templates", {})
        target = templates.get(template_id)

        if not target:
            return []

        # Simple similarity based on tags
        target_tags = set(target.get("tags", []))
        similar = []

        for tid, t in templates.items():
            if tid == template_id or t.get("status") != "approved":
                continue
            t_tags = set(t.get("tags", []))
            overlap = len(target_tags & t_tags)
            if overlap > 0:
                similar.append({"template": t, "score": overlap / max(len(target_tags), 1)})

        similar.sort(key=lambda x: x["score"], reverse=True)
        return similar[:limit]


# DocumentQAService


"""Service for Document Q&A Chat using AI."""

import time
import uuid
from datetime import datetime, timezone
from typing import Any, Callable, Dict, List, Optional

from backend.app.services.infra_services import extract_json_from_llm_response

from backend.app.schemas import (
    AskRequest,
    AskResponse,
    ChatMessage,
    Citation,
    DocQASession,
    DocumentReference,
    FeedbackRequest,
    MessageFeedback,
    MessageRole,
    RegenerateRequest,
)

logger = logging.getLogger("neura.domain.docqa")


class DocumentQAService:
    """Service for document-based Q&A chat."""

    def __init__(self):
        self._llm_client = None

    def _get_llm_client(self):
        if self._llm_client is None:
            self._llm_client = get_llm_client()
        return self._llm_client

    def _agent_budget_context(
        self,
        question: str,
        total_doc_chars: int,
        history_turns: int,
        max_budget: int,
    ) -> dict[str, int]:
        """Agent decides optimal context allocation based on query complexity."""
        client = self._get_llm_client()
        prompt = (
            f"Question: {question}\n"
            f"Total document size: {total_doc_chars} chars\n"
            f"Conversation history: {history_turns} turns\n"
            f"Total budget: {max_budget} characters\n\n"
            "Allocate the budget between document_context and conversation_history.\n"
            "Complex/detail questions need more document context.\n"
            "Follow-up questions need more history.\n"
            'Return ONLY JSON: {"doc": <chars for docs>, "history": <chars for history>}'
        )
        resp = client.complete(
            messages=[{"role": "user", "content": prompt}],
            description="docqa_context_budget",
            max_tokens=200,
        )
        from backend.app.services.llm import _extract_response_text
        text = _extract_response_text(resp)
        parsed = extract_json_from_llm_response(text)
        doc_budget = int(parsed.get("doc", max_budget))
        hist_budget = int(parsed.get("history", max_budget // 5))
        # Safety bounds
        doc_budget = max(5000, min(doc_budget, max_budget))
        hist_budget = max(1000, min(hist_budget, max_budget // 2))
        return {"doc": doc_budget, "history": hist_budget}

    def _read_sessions(self) -> Dict[str, Any]:
        store = get_state_store()
        # Use the StateStore's low-level primitives so tests can patch _read_state/_write_state.
        with store._lock:
            state = store._read_state() or {}
            sessions = state.get("docqa_sessions", {}) if isinstance(state, dict) else {}
            return dict(sessions or {})

    def _update_sessions(self, updater: Callable[[Dict[str, Any]], None]) -> None:
        store = get_state_store()
        with store._lock:
            state = store._read_state() or {}
            if not isinstance(state, dict):
                state = {}
            sessions = state.get("docqa_sessions", {})
            if not isinstance(sessions, dict):
                sessions = {}
            updater(sessions)
            state["docqa_sessions"] = sessions
            store._write_state(state)

    def create_session(
        self,
        name: str,
        correlation_id: Optional[str] = None,
    ) -> DocQASession:
        """Create a new Q&A session."""
        logger.info("Creating DocQA session", extra={"correlation_id": correlation_id})

        session = DocQASession(
            id=str(uuid.uuid4()),
            name=name,
            created_at=datetime.now(timezone.utc),
            updated_at=datetime.now(timezone.utc),
        )

        self._update_sessions(lambda sessions: sessions.__setitem__(session.id, session.model_dump(mode="json")))

        return session

    def get_session(self, session_id: str) -> Optional[DocQASession]:
        """Get a Q&A session by ID."""
        sessions = self._read_sessions()
        session_data = sessions.get(session_id)

        if session_data:
            return DocQASession(**session_data)
        return None

    def list_sessions(self) -> List[DocQASession]:
        """List all Q&A sessions."""
        sessions = self._read_sessions()
        return [DocQASession(**data) for data in sessions.values()]

    def add_document(
        self,
        session_id: str,
        name: str,
        content: str,
        page_count: Optional[int] = None,
        correlation_id: Optional[str] = None,
    ) -> Optional[DocumentReference]:
        """Add a document to a Q&A session."""
        logger.info(
            f"Adding document to DocQA session {session_id}",
            extra={"correlation_id": correlation_id},
        )

        session = self.get_session(session_id)
        if not session:
            return None

        document = DocumentReference(
            id=str(uuid.uuid4()),
            name=name,
            content_preview=content[:500] + "..." if len(content) > 500 else content,
            full_content=content[:100000],  # Limit stored content
            page_count=page_count,
            added_at=datetime.now(timezone.utc),
        )

        session.documents.append(document)
        session.updated_at = datetime.now(timezone.utc)

        self._update_sessions(lambda sessions: sessions.__setitem__(session_id, session.model_dump(mode="json")))

        return document

    def remove_document(self, session_id: str, document_id: str) -> bool:
        """Remove a document from a session."""
        session = self.get_session(session_id)
        if not session:
            return False

        session.documents = [d for d in session.documents if d.id != document_id]
        session.updated_at = datetime.now(timezone.utc)

        self._update_sessions(lambda sessions: sessions.__setitem__(session_id, session.model_dump(mode="json")))

        return True

    def ask(
        self,
        session_id: str,
        request: AskRequest,
        correlation_id: Optional[str] = None,
    ) -> Optional[AskResponse]:
        """Ask a question about the documents in a session."""
        start_time = time.time()
        logger.info(
            f"Processing question in session {session_id}",
            extra={"correlation_id": correlation_id},
        )

        session = self.get_session(session_id)
        if not session:
            return None

        if not session.documents:
            # Return a helpful message if no documents
            no_docs_message = ChatMessage(
                id=str(uuid.uuid4()),
                role=MessageRole.ASSISTANT,
                content="No documents have been added to this session yet. Please add documents first.",
                citations=[],
                timestamp=datetime.now(timezone.utc),
            )
            return AskResponse(
                message=no_docs_message,
                processing_time_ms=int((time.time() - start_time) * 1000),
            )

        # Add user message to history
        user_message = ChatMessage(
            id=str(uuid.uuid4()),
            role=MessageRole.USER,
            content=request.question,
            timestamp=datetime.now(timezone.utc),
        )
        session.messages.append(user_message)

        # Dynamic context budgeting — agent decides allocation based on query complexity
        MAX_TOTAL_CONTEXT = 50000  # ~50KB total context budget
        MAX_PER_DOC = 15000  # Max per document
        try:
            budget = self._agent_budget_context(
                request.question,
                sum(len(d.full_content) for d in session.documents),
                len(session.messages),
                MAX_TOTAL_CONTEXT,
            )
            MAX_TOTAL_CONTEXT = budget.get("doc", MAX_TOTAL_CONTEXT)
            MAX_PER_DOC = min(MAX_PER_DOC, MAX_TOTAL_CONTEXT // max(len(session.documents), 1))
        except Exception:
            pass  # Use defaults on failure

        doc_context = []
        total_chars = 0
        for doc in session.documents:
            remaining_budget = MAX_TOTAL_CONTEXT - total_chars
            if remaining_budget <= 0:
                doc_context.append(f"[Document: {doc.name} (ID: {doc.id})]\n(Skipped due to context limit)")
                continue
            # Take the lesser of per-doc limit or remaining budget
            chunk_size = min(MAX_PER_DOC, remaining_budget)
            content_chunk = doc.full_content[:chunk_size]
            doc_context.append(f"[Document: {doc.name} (ID: {doc.id})]\n{content_chunk}")
            total_chars += len(content_chunk)

        # Build conversation history (last N messages)
        window = max(1, session.context_window)
        history_messages = session.messages[-(window * 2):-1]  # Exclude current message
        history_text = ""
        if history_messages:
            history_parts = []
            for msg in history_messages:
                role = "User" if msg.role == MessageRole.USER else "Assistant"
                history_parts.append(f"{role}: {msg.content}")
            history_text = f"\nPREVIOUS CONVERSATION:\n" + "\n".join(history_parts)

        citation_instruction = ""
        if request.include_citations:
            citation_instruction = """
For each statement, provide citations in this format:
Include a "citations" array in your response with:
- document_id: The document ID
- document_name: The document name
- quote: The relevant quote from the document
- relevance_score: How relevant (0-1)
"""

        # V2: RAG-enhanced context injection
        _rag_context = ""
        try:
            from backend.app.services.infra_services import get_v2_config
            _v2_cfg = get_v2_config()
            if _v2_cfg.enable_rag_augmentation:
                from backend.app.services.knowledge_service import RAGPipeline
                _rag = RAGPipeline()
                _rag_results = _rag.query(request.question, top_k=_v2_cfg.rag_top_k)
                if _rag_results:
                    _rag_chunks = [r.content for r in _rag_results if r.score >= _v2_cfg.rag_relevance_threshold]
                    if _rag_chunks:
                        _rag_context = "\n\n--- Relevant Context from Knowledge Base ---\n" + "\n---\n".join(_rag_chunks[:_v2_cfg.rag_top_k])
        except Exception:
            pass  # RAG augmentation is non-critical

        prompt = f"""You are a helpful document Q&A assistant. Answer questions based ONLY on the provided documents.

DOCUMENTS:
{chr(10).join(doc_context)}
{history_text}
{_rag_context}

CURRENT QUESTION: {request.question}

INSTRUCTIONS:
1. Answer based ONLY on information in the documents
2. If the information is not in the documents, say so clearly
3. Be concise but thorough (max {request.max_response_length} characters)
4. Reference specific documents when appropriate
{citation_instruction}

Return a JSON object:
{{
  "answer": "Your comprehensive answer here",
  "citations": [
    {{
      "document_id": "doc_id",
      "document_name": "doc_name",
      "quote": "relevant quote",
      "relevance_score": 0.9
    }}
  ],
  "confidence": 0.9,
  "follow_up_questions": ["Suggested follow-up 1", "Suggested follow-up 2"]
}}

Return ONLY the JSON object."""

        try:
            client = self._get_llm_client()
            response = client.complete(
                messages=[{"role": "user", "content": prompt}],
                description="docqa_answer",
                temperature=0.2,
            )

            content = response["choices"][0]["message"]["content"]
            response_data = extract_json_from_llm_response(content, default=None)

            if response_data is not None:

                # Build citations
                citations = []
                for cit in response_data.get("citations", []):
                    citations.append(Citation(
                        document_id=cit.get("document_id", ""),
                        document_name=cit.get("document_name", ""),
                        quote=cit.get("quote", ""),
                        relevance_score=cit.get("relevance_score", 1.0),
                    ))

                # Create assistant message
                assistant_message = ChatMessage(
                    id=str(uuid.uuid4()),
                    role=MessageRole.ASSISTANT,
                    content=response_data.get("answer", "I couldn't generate an answer."),
                    citations=citations,
                    timestamp=datetime.now(timezone.utc),
                    metadata={
                        "confidence": response_data.get("confidence", 0.8),
                        "follow_up_questions": response_data.get("follow_up_questions", []),
                    },
                )

                # Add to session history
                session.messages.append(assistant_message)
                session.updated_at = datetime.now(timezone.utc)

                self._update_sessions(
                    lambda sessions: sessions.__setitem__(session_id, session.model_dump(mode="json"))
                )

                processing_time = int((time.time() - start_time) * 1000)

                return AskResponse(
                    message=assistant_message,
                    processing_time_ms=processing_time,
                    tokens_used=response.get("usage", {}).get("total_tokens"),
                )

        except Exception as exc:
            logger.error(f"DocQA failed: {exc}")

            error_message = ChatMessage(
                id=str(uuid.uuid4()),
                role=MessageRole.ASSISTANT,
                content=f"I encountered an error processing your question. Please try again.",
                timestamp=datetime.now(timezone.utc),
            )

            return AskResponse(
                message=error_message,
                processing_time_ms=int((time.time() - start_time) * 1000),
            )

        return None

    def submit_feedback(
        self,
        session_id: str,
        message_id: str,
        request: FeedbackRequest,
        correlation_id: Optional[str] = None,
    ) -> Optional[ChatMessage]:
        """Submit feedback for a message."""
        logger.info(
            "docqa_feedback_recorded",
            extra={
                "correlation_id": correlation_id,
                "session_id": session_id,
                "message_id": message_id,
                "feedback_type": request.feedback_type,
            },
        )

        session = self.get_session(session_id)
        if not session:
            return None

        target = next((msg for msg in session.messages if msg.id == message_id), None)
        if not target:
            return None

        feedback = MessageFeedback(
            feedback_type=request.feedback_type,
            timestamp=datetime.now(timezone.utc),
            comment=request.comment,
        )
        target.feedback = feedback
        meta = dict(target.metadata or {})
        meta["feedback"] = feedback.model_dump(mode="json")
        target.metadata = meta
        session.updated_at = datetime.now(timezone.utc)

        self._update_sessions(lambda sessions: sessions.__setitem__(session_id, session.model_dump(mode="json")))

        return target

    def regenerate_response(
        self,
        session_id: str,
        message_id: str,
        request: RegenerateRequest,
        correlation_id: Optional[str] = None,
    ) -> Optional[AskResponse]:
        """Regenerate the assistant response for a given message."""
        start_time = time.time()

        session = self.get_session(session_id)
        if not session:
            return None

        message_index = None
        for idx, msg in enumerate(session.messages):
            if msg.id == message_id:
                message_index = idx
                break
        if message_index is None:
            return None

        question_message = None
        question_index = None
        for idx in range(message_index - 1, -1, -1):
            candidate = session.messages[idx]
            if candidate.role == MessageRole.USER:
                question_message = candidate
                question_index = idx
                break
        if not question_message:
            return None

        question = question_message.content

        if not session.documents:
            return None

        # Build context from documents with total budget limit
        MAX_TOTAL_CONTEXT = 50000  # ~50KB total context budget
        MAX_PER_DOC = 15000  # Max per document

        doc_context = []
        total_chars = 0
        for doc in session.documents:
            remaining_budget = MAX_TOTAL_CONTEXT - total_chars
            if remaining_budget <= 0:
                doc_context.append(f"[Document: {doc.name} (ID: {doc.id})]\n(Skipped due to context limit)")
                continue
            chunk_size = min(MAX_PER_DOC, remaining_budget)
            content_chunk = doc.full_content[:chunk_size]
            doc_context.append(f"[Document: {doc.name} (ID: {doc.id})]\n{content_chunk}")
            total_chars += len(content_chunk)

        history_window = session.context_window * 2
        history_start = max(0, (question_index or 0) - history_window)
        history_messages = session.messages[history_start:question_index]
        history_text = ""
        if history_messages:
            history_parts = []
            for msg in history_messages:
                role = "User" if msg.role == MessageRole.USER else "Assistant"
                history_parts.append(f"{role}: {msg.content}")
            history_text = "\nPREVIOUS CONVERSATION:\n" + "\n".join(history_parts)

        citation_instruction = ""
        if request.include_citations:
            citation_instruction = """
For each statement, provide citations in this format:
Include a "citations" array in your response with:
- document_id: The document ID
- document_name: The document name
- quote: The relevant quote from the document
- relevance_score: How relevant (0-1)
"""

        # V2: RAG-enhanced context injection
        _rag_context = ""
        try:
            _v2_cfg = get_v2_config()
            if _v2_cfg.enable_rag_augmentation:
                _rag = RAGPipeline()
                _rag_results = _rag.query(question, top_k=_v2_cfg.rag_top_k)
                if _rag_results:
                    _rag_chunks = [r.content for r in _rag_results if r.score >= _v2_cfg.rag_relevance_threshold]
                    if _rag_chunks:
                        _rag_context = "\n\n--- Relevant Context from Knowledge Base ---\n" + "\n---\n".join(_rag_chunks[:_v2_cfg.rag_top_k])
        except Exception:
            pass  # RAG augmentation is non-critical

        prompt = f"""You are a helpful document Q&A assistant. Answer questions based ONLY on the provided documents.

DOCUMENTS:
{chr(10).join(doc_context)}
{history_text}
{_rag_context}

CURRENT QUESTION: {question}

INSTRUCTIONS:
1. Answer based ONLY on information in the documents
2. If the information is not in the documents, say so clearly
3. Be concise but thorough (max {request.max_response_length} characters)
4. Reference specific documents when appropriate
5. Provide a DIFFERENT perspective or additional details compared to previous answers
{citation_instruction}

Return a JSON object:
{{
  "answer": "Your comprehensive answer here",
  "citations": [
    {{
      "document_id": "doc_id",
      "document_name": "doc_name",
      "quote": "relevant quote",
      "relevance_score": 0.9
    }}
  ],
  "confidence": 0.9,
  "follow_up_questions": ["Suggested follow-up 1", "Suggested follow-up 2"]
}}

Return ONLY the JSON object."""

        try:
            client = self._get_llm_client()
            response = client.complete(
                messages=[{"role": "user", "content": prompt}],
                description="docqa_regenerate",
                temperature=0.5,
            )

            content = response["choices"][0]["message"]["content"]
            response_data = extract_json_from_llm_response(content, default=None)

            if response_data is None:
                raise ValueError("No JSON payload returned from LLM")

            citations = []
            for cit in response_data.get("citations", []):
                citations.append(Citation(
                    document_id=cit.get("document_id", ""),
                    document_name=cit.get("document_name", ""),
                    quote=cit.get("quote", ""),
                    relevance_score=cit.get("relevance_score", 1.0),
                ))

            assistant_message = ChatMessage(
                id=message_id,
                role=MessageRole.ASSISTANT,
                content=response_data.get("answer", "I couldn't generate an answer."),
                citations=citations,
                timestamp=datetime.now(timezone.utc),
                metadata={
                    "confidence": response_data.get("confidence", 0.8),
                    "follow_up_questions": response_data.get("follow_up_questions", []),
                    "regenerated": True,
                },
                feedback=None,
            )

            session.messages[message_index] = assistant_message
            session.updated_at = datetime.now(timezone.utc)

            self._update_sessions(
                lambda sessions: sessions.__setitem__(session_id, session.model_dump(mode="json"))
            )

            processing_time = int((time.time() - start_time) * 1000)

            return AskResponse(
                message=assistant_message,
                processing_time_ms=processing_time,
                tokens_used=response.get("usage", {}).get("total_tokens"),
            )

        except Exception as exc:
            logger.error(f"DocQA regenerate failed: {exc}")
            raise RuntimeError("Failed to regenerate response") from exc

    def get_chat_history(
        self,
        session_id: str,
        limit: int = 50,
    ) -> List[ChatMessage]:
        """Get chat history for a session."""
        session = self.get_session(session_id)
        if not session:
            return []

        return session.messages[-limit:] if limit else session.messages

    def clear_history(self, session_id: str) -> bool:
        """Clear chat history for a session."""
        session = self.get_session(session_id)
        if not session:
            return False

        session.messages = []
        session.updated_at = datetime.now(timezone.utc)

        self._update_sessions(lambda sessions: sessions.__setitem__(session_id, session.model_dump(mode="json")))

        return True

    def delete_session(self, session_id: str) -> bool:
        """Delete a Q&A session."""
        deleted = False

        def _updater(sessions: Dict[str, Any]) -> None:
            nonlocal deleted
            if session_id in sessions:
                del sessions[session_id]
                deleted = True

        self._update_sessions(_updater)
        return deleted


# FederationService


"""Service layer for Cross-Database Federation feature."""


from backend.app.utils import AppError
from backend.app.utils import get_write_operation, is_select_or_with
from backend.app.services.infra_services import extract_json_array_from_llm_response

from backend.app.schemas import (
    VirtualSchema,
    VirtualSchemaCreate,
    JoinSuggestion,
    TableReference,
    FederatedQueryRequest,
)

logger = logging.getLogger("neura.domain.federation")




class FederationService:
    """Service for cross-database federation operations."""

    def __init__(self):
        self._llm_client = None

    def _get_llm_client(self):
        if self._llm_client is None:
            self._llm_client = get_llm_client()
        return self._llm_client

    def _get_connection_schema(self, connection_id: str) -> Dict[str, Any]:
        """Get schema for a connection."""
        return get_connection_schema(connection_id, include_row_counts=False, sample_rows=3)

    def create_virtual_schema(
        self,
        request: VirtualSchemaCreate,
        correlation_id: Optional[str] = None,
    ) -> VirtualSchema:
        """Create a new virtual schema."""
        logger.info(f"Creating virtual schema: {request.name}", extra={"correlation_id": correlation_id})

        schema_id = str(uuid.uuid4())[:8]
        now = utc_now_iso()

        # Gather tables from all connections
        tables: List[TableReference] = []
        for conn_id in request.connection_ids:
            try:
                schema = self._get_connection_schema(conn_id)
                for table in schema.get("tables", []):
                    tables.append(TableReference(
                        connection_id=conn_id,
                        table_name=table["name"],
                        alias=f"{conn_id[:4]}_{table['name']}"
                    ))
            except Exception as exc:
                logger.warning(f"Failed to get schema for {conn_id}: {exc}")

        virtual_schema = VirtualSchema(
            id=schema_id,
            name=request.name,
            description=request.description,
            connections=request.connection_ids,
            tables=tables,
            joins=[],
            created_at=now,
            updated_at=now,
        )

        # Persist
        store = get_state_store()
        with store.transaction() as state:
            state.setdefault("virtual_schemas", {})[schema_id] = virtual_schema.model_dump()

        return virtual_schema

    def suggest_joins(
        self,
        connection_ids: List[str],
        correlation_id: Optional[str] = None,
    ) -> List[JoinSuggestion]:
        """Suggest joins between tables in different connections using AI."""
        logger.info(f"Suggesting joins for {len(connection_ids)} connections", extra={"correlation_id": correlation_id})

        # Gather schemas
        schemas = {}
        for conn_id in connection_ids:
            try:
                schemas[conn_id] = self._get_connection_schema(conn_id)
            except Exception as exc:
                logger.warning(f"Failed to get schema for {conn_id}: {exc}")

        if len(schemas) < 2:
            return []

        # Build prompt for LLM
        schema_desc = []
        for conn_id, schema in schemas.items():
            tables_desc = []
            for table in schema.get("tables", []):
                cols = [f"{c['name']} ({c.get('type', 'TEXT')})" for c in table.get("columns", [])]
                tables_desc.append(f"  - {table['name']}: {', '.join(cols)}")
            schema_desc.append(f"Connection {conn_id}:\n" + "\n".join(tables_desc))

        prompt = f"""Analyze these database schemas and suggest joins between tables from different connections.

{chr(10).join(schema_desc)}

For each potential join, consider:
1. Column names that might match (like 'customer_id', 'user_id', 'id')
2. Data types compatibility
3. Business logic relationships

Return a JSON array of join suggestions:
[
  {{
    "left_connection_id": "conn1",
    "left_table": "table1",
    "left_column": "column1",
    "right_connection_id": "conn2",
    "right_table": "table2",
    "right_column": "column2",
    "confidence": 0.9,
    "reason": "Both columns appear to be customer identifiers"
  }}
]

Return ONLY the JSON array."""

        try:
            client = self._get_llm_client()
            response = client.complete(
                messages=[{"role": "user", "content": prompt}],
                description="join_suggestion",
                temperature=0.0,
            )

            content = response["choices"][0]["message"]["content"]
            suggestions_data = extract_json_array_from_llm_response(content, default=[])
            if suggestions_data:
                return [JoinSuggestion(**s) for s in suggestions_data]

        except Exception as exc:
            logger.error(f"Join suggestion failed: {exc}")

        return []

    def list_virtual_schemas(self) -> List[VirtualSchema]:
        """List all virtual schemas."""
        store = get_state_store()
        with store.transaction() as state:
            schemas = state.get("virtual_schemas", {})
            return [VirtualSchema(**s) for s in schemas.values()]

    def get_virtual_schema(self, schema_id: str) -> Optional[VirtualSchema]:
        """Get a virtual schema by ID."""
        store = get_state_store()
        with store.transaction() as state:
            schema = state.get("virtual_schemas", {}).get(schema_id)
            return VirtualSchema(**schema) if schema else None

    def delete_virtual_schema(self, schema_id: str) -> bool:
        """Delete a virtual schema."""
        store = get_state_store()
        with store.transaction() as state:
            schemas = state.get("virtual_schemas", {})
            if schema_id not in schemas:
                return False
            del schemas[schema_id]
        return True

    def _extract_table_names(self, sql: str) -> List[str]:
        """Extract table names from SQL query using improved parsing.

        Handles:
        - Simple table names: FROM users
        - Quoted identifiers: FROM "users", FROM `users`, FROM [users]
        - Schema-qualified names: FROM schema.table
        - Comma-separated tables: FROM a, b, c
        - Various JOIN types: LEFT JOIN, RIGHT JOIN, INNER JOIN, etc.
        - CTEs: WITH cte AS (...) SELECT ... FROM cte
        """

        # Normalize whitespace
        sql_normalized = " ".join(sql.split())

        # Identifier pattern: handles unquoted, double-quoted, backtick-quoted, and bracket-quoted
        # Also handles schema.table notation
        ident = r'(?:' \
                r'(?:[a-zA-Z_][a-zA-Z0-9_]*(?:\.[a-zA-Z_][a-zA-Z0-9_]*)?)' \
                r'|"[^"]+"(?:\."[^"]+")?'  \
                r'|`[^`]+`(?:\.`[^`]+`)?'  \
                r'|\[[^\]]+\](?:\.\[[^\]]+\])?' \
                r')'

        # Pattern for table lists (comma-separated with optional aliases)
        table_list_pattern = rf'{ident}(?:\s+(?:AS\s+)?[a-zA-Z_][a-zA-Z0-9_]*)?(?:\s*,\s*{ident}(?:\s+(?:AS\s+)?[a-zA-Z_][a-zA-Z0-9_]*)?)*'

        # SQL keyword patterns that precede table names
        patterns = [
            rf'\bFROM\s+({table_list_pattern})',
            rf'\b(?:LEFT|RIGHT|INNER|OUTER|CROSS|FULL)?\s*JOIN\s+({ident})',
            rf'\bJOIN\s+({ident})',
            rf'\bINTO\s+({ident})',
            rf'\bUPDATE\s+({ident})',
        ]

        # Also extract CTE names to exclude them from final results
        cte_pattern = rf'\bWITH\s+({ident})\s+AS\s*\('
        cte_names = set()
        for match in re.findall(cte_pattern, sql_normalized, re.IGNORECASE):
            cte_names.add(self._clean_identifier(match).lower())

        tables = set()
        sql_keywords = {
            'SELECT', 'WHERE', 'ON', 'AND', 'OR', 'NOT', 'IN', 'EXISTS',
            'HAVING', 'GROUP', 'ORDER', 'BY', 'LIMIT', 'OFFSET', 'UNION',
            'INTERSECT', 'EXCEPT', 'AS', 'CASE', 'WHEN', 'THEN', 'ELSE', 'END',
            'NULL', 'TRUE', 'FALSE', 'IS', 'LIKE', 'BETWEEN', 'ALL', 'ANY',
            'LEFT', 'RIGHT', 'INNER', 'OUTER', 'CROSS', 'FULL', 'NATURAL',
        }

        for pattern in patterns:
            matches = re.findall(pattern, sql_normalized, re.IGNORECASE)
            for match in matches:
                # Handle comma-separated table names
                # Split carefully to handle schema.table notation
                parts = re.split(r'\s*,\s*', match)
                for part in parts:
                    # Extract just the table name (remove alias)
                    tokens = part.strip().split()
                    if tokens:
                        table_ref = tokens[0]
                        # Skip AS keyword if present
                        if table_ref.upper() == 'AS' and len(tokens) > 1:
                            continue

                        table_name = self._clean_identifier(table_ref)

                        # Skip SQL keywords and CTEs
                        if table_name.upper() not in sql_keywords and table_name.lower() not in cte_names:
                            tables.add(table_name.lower())

        return list(tables)

    def _clean_identifier(self, ident: str) -> str:
        """Remove quotes from identifier and extract table name from schema.table."""
        # Remove various quote styles
        cleaned = ident.strip()
        if cleaned.startswith('"') and cleaned.endswith('"'):
            cleaned = cleaned[1:-1]
        elif cleaned.startswith('`') and cleaned.endswith('`'):
            cleaned = cleaned[1:-1]
        elif cleaned.startswith('[') and cleaned.endswith(']'):
            cleaned = cleaned[1:-1]

        # Handle schema.table notation - extract just the table name
        if '.' in cleaned:
            # Could be schema.table or "schema"."table"
            parts = cleaned.split('.')
            cleaned = parts[-1]
            # Clean the table part if it's still quoted
            if cleaned.startswith('"') and cleaned.endswith('"'):
                cleaned = cleaned[1:-1]
            elif cleaned.startswith('`') and cleaned.endswith('`'):
                cleaned = cleaned[1:-1]
            elif cleaned.startswith('[') and cleaned.endswith(']'):
                cleaned = cleaned[1:-1]

        return cleaned

    def _map_tables_to_connections(
        self, table_names: List[str], schema: VirtualSchema
    ) -> Dict[str, List[str]]:
        """Map table names to their respective connections."""
        connection_tables: Dict[str, List[str]] = {}

        for table_name in table_names:
            for table_ref in schema.tables:
                # Match by table name or alias
                if (
                    table_ref.table_name.lower() == table_name.lower()
                    or table_ref.alias.lower() == table_name.lower()
                ):
                    conn_id = table_ref.connection_id
                    if conn_id not in connection_tables:
                        connection_tables[conn_id] = []
                    connection_tables[conn_id].append(table_ref.table_name)
                    break

        return connection_tables

    def _execute_on_connection(
        self,
        connection_id: str,
        sql: str,
        limit: Optional[int] = None,
    ) -> Dict[str, Any]:
        """Execute SQL on a specific connection."""
        from backend.app.repositories import execute_query
        return execute_query(connection_id=connection_id, sql=sql, limit=limit)

    def _merge_results(
        self,
        results: List[Dict[str, Any]],
        join_keys: Optional[List[str]] = None,
    ) -> Dict[str, Any]:
        """Merge results from multiple connections."""
        if not results:
            return {"columns": [], "rows": [], "row_count": 0}

        if len(results) == 1:
            return results[0]

        # For simple merge (no join), concatenate results
        if not join_keys:
            all_columns = []
            seen_cols = set()
            for r in results:
                for col in r.get("columns", []):
                    if col not in seen_cols:
                        all_columns.append(col)
                        seen_cols.add(col)

            all_rows = []
            for r in results:
                for row in r.get("rows", []):
                    # Extend row with nulls for missing columns
                    extended_row = []
                    for col in all_columns:
                        try:
                            idx = r.get("columns", []).index(col)
                            extended_row.append(row[idx])
                        except (ValueError, IndexError):
                            extended_row.append(None)
                    all_rows.append(extended_row)

            return {
                "columns": all_columns,
                "rows": all_rows,
                "row_count": len(all_rows),
                "merge_type": "union",
            }

        # Client-side join on specified keys
        if len(results) == 2:
            left = results[0]
            right = results[1]

            left_cols = left.get("columns", [])
            right_cols = right.get("columns", [])

            # Find join key indices
            left_key_idx = None
            right_key_idx = None
            for key in join_keys:
                if key in left_cols:
                    left_key_idx = left_cols.index(key)
                if key in right_cols:
                    right_key_idx = right_cols.index(key)

            if left_key_idx is None or right_key_idx is None:
                # Can't find join keys, return concatenated
                return self._merge_results(results, None)

            # Build right side lookup
            right_lookup: Dict[Any, List[List[Any]]] = {}
            for row in right.get("rows", []):
                key_val = row[right_key_idx]
                if key_val not in right_lookup:
                    right_lookup[key_val] = []
                right_lookup[key_val].append(row)

            # Build merged columns (left cols + right cols without join key)
            merged_cols = list(left_cols) + [
                c for i, c in enumerate(right_cols) if i != right_key_idx
            ]

            # Perform join
            merged_rows = []
            for left_row in left.get("rows", []):
                key_val = left_row[left_key_idx]
                if key_val in right_lookup:
                    for right_row in right_lookup[key_val]:
                        new_row = list(left_row) + [
                            v for i, v in enumerate(right_row) if i != right_key_idx
                        ]
                        merged_rows.append(new_row)

            return {
                "columns": merged_cols,
                "rows": merged_rows,
                "row_count": len(merged_rows),
                "merge_type": "join",
                "join_key": join_keys[0],
            }

        return self._merge_results(results, None)

    def execute_query(
        self,
        request: FederatedQueryRequest,
        correlation_id: Optional[str] = None,
    ) -> Dict[str, Any]:
        """Execute a federated query across multiple databases.

        This method:
        1. Parses the SQL to identify referenced tables
        2. Maps tables to their respective connections
        3. Routes queries to appropriate connections
        4. Merges results client-side if needed
        """
        logger.info(
            f"Executing federated query on schema {request.virtual_schema_id}",
            extra={"correlation_id": correlation_id}
        )

        # Get the virtual schema
        schema = self.get_virtual_schema(request.virtual_schema_id)
        if not schema:
            raise AppError(
                code="schema_not_found",
                message=f"Virtual schema {request.virtual_schema_id} not found",
                status_code=404,
            )

        if not schema.connections:
            raise AppError(
                code="no_connections",
                message="Virtual schema has no connections",
                status_code=400,
            )

        if not is_select_or_with(request.sql):
            raise AppError(
                code="invalid_query",
                message="Only SELECT queries are allowed",
                status_code=400,
            )

        write_op = get_write_operation(request.sql)
        if write_op:
            raise AppError(
                code="dangerous_query",
                message=f"Query contains prohibited operation: {write_op}",
                status_code=400,
            )

        try:
            # Extract table names from SQL
            table_names = self._extract_table_names(request.sql)
            logger.debug(f"Extracted tables from SQL: {table_names}")

            # Map tables to connections
            connection_tables = self._map_tables_to_connections(table_names, schema)
            logger.debug(f"Table-to-connection mapping: {connection_tables}")

            # If no tables found or all tables in one connection, execute directly
            if len(connection_tables) <= 1:
                target_connection = (
                    list(connection_tables.keys())[0]
                    if connection_tables
                    else schema.connections[0]
                )
                result = self._execute_on_connection(
                    connection_id=target_connection,
                    sql=request.sql,
                    limit=request.limit,
                )
                return {
                    "columns": result.get("columns", []),
                    "rows": result.get("rows", []),
                    "row_count": len(result.get("rows", [])),
                    "schema_id": request.virtual_schema_id,
                    "executed_on": [target_connection],
                    "routing": "single",
                }

            # Multi-connection query: need to split and merge
            logger.info(
                f"Federated query spans {len(connection_tables)} connections",
                extra={"connections": list(connection_tables.keys())}
            )

            # Try to get join keys from schema joins
            join_keys = []
            for join in schema.joins:
                if hasattr(join, "conditions"):
                    for cond in join.conditions:
                        join_keys.extend([cond.left_column, cond.right_column])

            # Execute on each connection
            results = []
            executed_on = []
            for conn_id, tables in connection_tables.items():
                # For each connection, try to execute a query for its tables
                # This is a simplified approach - a more sophisticated implementation
                # would rewrite the SQL for each connection
                try:
                    # Execute the original query on this connection
                    # It will fail for tables it doesn't have, but that's handled
                    result = self._execute_on_connection(
                        connection_id=conn_id,
                        sql=request.sql,
                        limit=request.limit,
                    )
                    results.append(result)
                    executed_on.append(conn_id)
                except Exception as exc:
                    logger.warning(f"Query on {conn_id} failed: {exc}")
                    # Try simpler query for just the tables in this connection
                    for table in tables:
                        try:
                            quoted = '"' + table.replace('"', '""') + '"'
                            simple_sql = f"SELECT * FROM {quoted}"
                            if request.limit:
                                simple_sql += f" LIMIT {request.limit}"
                            result = self._execute_on_connection(
                                connection_id=conn_id,
                                sql=simple_sql,
                                limit=request.limit,
                            )
                            results.append(result)
                            executed_on.append(conn_id)
                        except Exception as inner_exc:
                            logger.warning(f"Simple query on {conn_id}.{table} failed: {inner_exc}")

            if not results:
                raise AppError(
                    code="query_failed",
                    message="Could not execute query on any connection",
                    status_code=500,
                )

            # Merge results
            merged = self._merge_results(results, join_keys if join_keys else None)

            return {
                "columns": merged.get("columns", []),
                "rows": merged.get("rows", [])[:request.limit] if request.limit else merged.get("rows", []),
                "row_count": len(merged.get("rows", [])),
                "schema_id": request.virtual_schema_id,
                "executed_on": executed_on,
                "routing": "federated",
                "merge_type": merged.get("merge_type", "unknown"),
            }

        except AppError:
            raise
        except Exception as exc:
            logger.error(f"Federated query failed: {exc}")
            raise AppError(
                code="query_failed",
                message="Query execution failed",
                status_code=500,
            )


# NL2SQLService


"""Service layer for Natural Language to SQL feature."""

from pathlib import Path


from backend.app.repositories import resolve_db_path, verify_sqlite
import sqlite3 as sqlite_shim  # was: from backend.app.repositories import sqlite_shim
from backend.app.repositories import ensure_connection_loaded
from backend.app.services.llm import TextToSQL

from backend.app.schemas import (
    NL2SQLGenerateRequest,
    NL2SQLExecuteRequest,
    NL2SQLSaveRequest,
    NL2SQLResult,
    QueryExecutionResult,
    SavedQuery,
    QueryHistoryEntry,
)

logger = logging.getLogger("neura.domain.nl2sql")

_TRAILING_SEMICOLONS_RE = re.compile(r";+\s*$")


def _strip_trailing_semicolons(sql: str) -> str:
    # Users commonly paste SQL ending with ';'. That breaks subquery wrapping/pagination.
    return _TRAILING_SEMICOLONS_RE.sub("", (sql or "").strip())




def _quote_identifier(name: str) -> str:
    return '"' + name.replace('"', '""') + '"'


def _coerce_value(value: Any) -> Any:
    """Convert bytes and other non-JSON types to serializable formats."""
    if isinstance(value, (bytes, bytearray, memoryview)):
        return bytes(value).hex()
    try:
        if hasattr(value, "item"):
            return value.item()
    except Exception:
        pass
    return value


class NL2SQLService:
    """Service for natural language to SQL operations."""

    def __init__(self):
        self._text_to_sql: Optional[TextToSQL] = None

    def _get_text_to_sql(self) -> TextToSQL:
        """Get or create TextToSQL instance."""
        if self._text_to_sql is None:
            client = get_llm_client()
            self._text_to_sql = TextToSQL(client=client, dialect="sqlite")
        return self._text_to_sql

    def _resolve_connection(self, connection_id: str) -> Path:
        """Resolve and verify a database connection."""
        try:
            db_path = resolve_db_path(connection_id=connection_id, db_url=None, db_path=None)
            verify_sqlite(db_path)
            return db_path
        except Exception as exc:
            logger.warning("Connection validation failed: %s", exc)
            raise AppError(
                code="connection_invalid",
                message="Invalid or unreachable database connection",
                status_code=400,
            )

    def _get_schema_for_connection(self, db_path: Path, tables: Optional[List[str]] = None) -> Dict[str, Any]:
        """Get database schema for SQL generation context using DataFrames."""
        from backend.app.repositories import get_loader

        loader = get_loader(db_path)
        schema = {}

        table_names = tables if tables else loader.table_names()

        for table_name in table_names:
            if tables and table_name not in tables:
                continue

            columns = []
            for col in loader.pragma_table_info(table_name):
                columns.append({
                    "name": col.get("name"),
                    "type": col.get("type", "TEXT"),
                    "description": "",
                })

            # Get sample values from DataFrame (no direct DB access)
            sample_values = {}
            try:
                frame = loader.frame(table_name)
                if not frame.empty:
                    sample_rows = frame.head(3)
                    for col in columns:
                        col_name = col["name"]
                        if col_name in sample_rows.columns:
                            values = [_coerce_value(v) for v in sample_rows[col_name].tolist()]
                            if values:
                                sample_values[col_name] = values[:3]
            except Exception as e:
                logger.debug("Failed to extract sample values for %s: %s", table_name, e)

            schema[table_name] = {
                "columns": columns,
                "foreign_keys": loader.foreign_keys(table_name),
                "sample_values": sample_values,
            }

        return schema

    def generate_sql(
        self,
        request: NL2SQLGenerateRequest,
        correlation_id: Optional[str] = None,
    ) -> NL2SQLResult:
        """Generate SQL from a natural language question."""
        logger.info(f"Generating SQL for question: {request.question[:100]}...", extra={"correlation_id": correlation_id})

        # Resolve and verify connection
        db_path = self._resolve_connection(request.connection_id)

        # Get schema for context
        schema = self._get_schema_for_connection(db_path, request.tables)
        if not schema:
            raise AppError(
                code="no_tables",
                message="No tables found in the database",
                status_code=400,
            )

        # Set up TextToSQL with schema
        t2sql = self._get_text_to_sql()
        t2sql._schemas.clear()  # Clear previous schemas
        t2sql.add_schemas_from_catalog(schema)

        # Generate SQL
        try:
            result = t2sql.generate_sql(
                question=request.question,
                tables=request.tables,
                context=request.context,
            )
        except Exception as exc:
            logger.error(f"SQL generation failed: {exc}", extra={"correlation_id": correlation_id})
            raise AppError(
                code="generation_failed",
                message="Failed to generate SQL query",
                status_code=500,
            )

        # Record in history
        self._record_history(
            question=request.question,
            sql=result.sql,
            connection_id=request.connection_id,
            confidence=result.confidence,
            success=True,
        )

        return NL2SQLResult(
            sql=result.sql,
            explanation=result.explanation,
            confidence=result.confidence,
            warnings=result.warnings,
            original_question=request.question,
        )

    def execute_query(
        self,
        request: NL2SQLExecuteRequest,
        correlation_id: Optional[str] = None,
    ) -> QueryExecutionResult:
        """Execute a SQL query and return results using DataFrames."""
        logger.info(f"Executing SQL query on connection {request.connection_id}", extra={"correlation_id": correlation_id})

        # Resolve and verify connection
        db_path = self._resolve_connection(request.connection_id)

        # Ensure DataFrames are loaded for this connection
        ensure_connection_loaded(request.connection_id, db_path)

        sql_clean = _strip_trailing_semicolons(request.sql)

        # Validate SQL (read-only safety check)
        if not is_select_or_with(sql_clean):
            raise AppError(
                code="invalid_query",
                message="Only SELECT queries are allowed",
                status_code=400,
            )

        write_op = get_write_operation(sql_clean)
        if write_op:
            raise AppError(
                code="dangerous_query",
                message=f"Query contains prohibited operation: {write_op}",
                status_code=400,
            )

        # Execute query using DataFrame shim
        started = time.time()
        try:
            with sqlite_shim.connect(str(db_path)) as con:
                con.row_factory = sqlite_shim.Row

                total_count = None
                if request.include_total:
                    count_sql = f"SELECT COUNT(*) as cnt FROM ({sql_clean}) AS subq"
                    try:
                        total_count = con.execute(count_sql).fetchone()["cnt"]
                    except Exception:
                        total_count = None

                # Execute with limit and offset
                limited_sql = f"SELECT * FROM ({sql_clean}) AS subq LIMIT {request.limit} OFFSET {request.offset}"
                cur = con.execute(limited_sql)
                rows_raw = cur.fetchall()

                # Get column names from cursor description (reliable across all Row types)
                columns = [desc[0] for desc in cur.description] if cur.description else []
                rows = [{col: _coerce_value(row[i]) for i, col in enumerate(columns)} for row in rows_raw]

        except sqlite_shim.OperationalError as exc:
            execution_time_ms = int((time.time() - started) * 1000)
            logger.error(f"Query execution failed: {exc}", extra={"correlation_id": correlation_id})

            # Record failure in history
            self._update_history_execution(
                sql=request.sql,
                connection_id=request.connection_id,
                success=False,
                error="Query execution failed",
                execution_time_ms=execution_time_ms,
            )

            raise AppError(
                code="execution_failed",
                message="Failed to execute SQL query",
                status_code=400,
            )

        execution_time_ms = int((time.time() - started) * 1000)

        # Record success in history
        self._update_history_execution(
            sql=request.sql,
            connection_id=request.connection_id,
            success=True,
            execution_time_ms=execution_time_ms,
            row_count=len(rows),
        )

        # Determine truncated: if total_count is known, compare to limit.
        # Otherwise, if we got exactly limit rows, assume there may be more.
        if total_count is not None:
            truncated = total_count > request.limit
        else:
            truncated = len(rows) >= request.limit

        return QueryExecutionResult(
            columns=columns,
            rows=rows,
            row_count=len(rows),
            total_count=total_count,
            execution_time_ms=execution_time_ms,
            truncated=truncated,
        )

    def explain_query(
        self,
        sql: str,
        correlation_id: Optional[str] = None,
    ) -> str:
        """Get a natural language explanation of a SQL query."""
        t2sql = self._get_text_to_sql()
        return t2sql.explain_sql(sql)

    def save_query(
        self,
        request: NL2SQLSaveRequest,
        correlation_id: Optional[str] = None,
    ) -> SavedQuery:
        """Save a query as a reusable data source."""
        logger.info(f"Saving query: {request.name}", extra={"correlation_id": correlation_id})

        # Verify connection exists
        self._resolve_connection(request.connection_id)

        query_id = str(uuid.uuid4())
        now = utc_now_iso()

        saved_query = SavedQuery(
            id=query_id,
            name=request.name,
            description=request.description,
            sql=request.sql,
            connection_id=request.connection_id,
            original_question=request.original_question,
            tags=request.tags or [],
            created_at=now,
            updated_at=now,
            run_count=0,
        )

        # Persist to state store
        store = get_state_store()
        store.save_query(saved_query.model_dump(mode="json"))

        return saved_query

    def list_saved_queries(
        self,
        connection_id: Optional[str] = None,
        tags: Optional[List[str]] = None,
    ) -> List[SavedQuery]:
        """List saved queries, optionally filtered."""
        store = get_state_store()
        queries = store.list_saved_queries()

        if connection_id:
            queries = [q for q in queries if q.get("connection_id") == connection_id]

        if tags:
            tag_set = set(tags)
            queries = [q for q in queries if tag_set.intersection(set(q.get("tags", [])))]

        return [SavedQuery(**q) for q in queries]

    def get_saved_query(self, query_id: str) -> Optional[SavedQuery]:
        """Get a saved query by ID."""
        store = get_state_store()
        query = store.get_saved_query(query_id)
        return SavedQuery(**query) if query else None

    def delete_saved_query(self, query_id: str) -> bool:
        """Delete a saved query."""
        store = get_state_store()
        return store.delete_saved_query(query_id)

    def get_query_history(
        self,
        connection_id: Optional[str] = None,
        limit: int = 50,
    ) -> List[QueryHistoryEntry]:
        """Get query history."""
        store = get_state_store()
        history = store.get_query_history(limit=limit)

        if connection_id:
            history = [h for h in history if h.get("connection_id") == connection_id]

        return [QueryHistoryEntry(**h) for h in history]

    def delete_query_history_entry(self, entry_id: str) -> bool:
        """Delete a query history entry by ID."""
        store = get_state_store()
        return store.delete_query_history_entry(entry_id)

    def _record_history(
        self,
        question: str,
        sql: str,
        connection_id: str,
        confidence: float,
        success: bool,
        error: Optional[str] = None,
    ) -> None:
        """Record a query generation in history."""
        store = get_state_store()
        entry = {
            "id": str(uuid.uuid4())[:8],
            "question": question,
            "sql": sql,
            "connection_id": connection_id,
            "confidence": confidence,
            "success": success,
            "error": error,
            "created_at": utc_now_iso(),
        }
        store.add_query_history(entry)

    def _update_history_execution(
        self,
        sql: str,
        connection_id: str,
        success: bool,
        error: Optional[str] = None,
        execution_time_ms: Optional[int] = None,
        row_count: Optional[int] = None,
    ) -> None:
        """Update history with execution results."""
        # This could be enhanced to update the last matching history entry
        pass  # For now, execution results are returned directly


# DocumentSynthesisService


"""Service for Multi-Document Synthesis using AI."""

import hashlib

from backend.app.services.infra_services import (
    extract_json_from_llm_response,
    extract_json_array_from_llm_response,
)

from backend.app.schemas import (
    DocumentType,
    Inconsistency,
    SynthesisDocument,
    SynthesisRequest,
    SynthesisResult,
    SynthesisSession,
)

logger = logging.getLogger("neura.domain.synthesis")


class DocumentSynthesisService:
    """Service for synthesizing information from multiple documents."""

    def __init__(self):
        self._llm_client = None

    def _get_llm_client(self):
        if self._llm_client is None:
            self._llm_client = get_llm_client()
        return self._llm_client

    def _read_sessions(self) -> Dict[str, Any]:
        store = get_state_store()
        with store.transaction() as state:
            return dict(state.get("synthesis_sessions", {}) or {})

    def _update_sessions(self, updater: Callable[[Dict[str, Any]], None]) -> None:
        store = get_state_store()
        with store.transaction() as state:
            sessions = state.get("synthesis_sessions", {})
            if not isinstance(sessions, dict):
                sessions = {}
            updater(sessions)
            state["synthesis_sessions"] = sessions

    def create_session(
        self,
        name: str,
        correlation_id: Optional[str] = None,
    ) -> SynthesisSession:
        """Create a new synthesis session."""
        logger.info("Creating synthesis session", extra={"correlation_id": correlation_id})

        session = SynthesisSession(
            id=str(uuid.uuid4()),
            name=name,
            created_at=datetime.now(timezone.utc),
            updated_at=datetime.now(timezone.utc),
        )

        self._update_sessions(lambda sessions: sessions.__setitem__(session.id, session.model_dump(mode="json")))

        return session

    def get_session(self, session_id: str) -> Optional[SynthesisSession]:
        """Get a synthesis session by ID."""
        sessions = self._read_sessions()
        session_data = sessions.get(session_id)

        if session_data:
            return SynthesisSession(**session_data)
        return None

    def list_sessions(self) -> List[SynthesisSession]:
        """List all synthesis sessions."""
        sessions = self._read_sessions()
        return [SynthesisSession(**data) for data in sessions.values()]

    def add_document(
        self,
        session_id: str,
        name: str,
        content: str,
        doc_type: DocumentType,
        metadata: Optional[Dict[str, Any]] = None,
        correlation_id: Optional[str] = None,
    ) -> Optional[SynthesisDocument]:
        """Add a document to a synthesis session."""
        logger.info(
            f"Adding document to session {session_id}",
            extra={"correlation_id": correlation_id},
        )

        session = self.get_session(session_id)
        if not session:
            return None

        content_hash = hashlib.sha256(content.encode()).hexdigest()[:16]

        document = SynthesisDocument(
            id=str(uuid.uuid4()),
            name=name,
            doc_type=doc_type,
            content_hash=content_hash,
            extracted_text=content[:50000],  # Limit stored text
            metadata=metadata or {},
            added_at=datetime.now(timezone.utc),
        )

        session.documents.append(document)
        session.updated_at = datetime.now(timezone.utc)

        self._update_sessions(lambda sessions: sessions.__setitem__(session_id, session.model_dump(mode="json")))

        return document

    def remove_document(self, session_id: str, document_id: str) -> bool:
        """Remove a document from a session."""
        session = self.get_session(session_id)
        if not session:
            return False

        session.documents = [d for d in session.documents if d.id != document_id]
        session.updated_at = datetime.now(timezone.utc)

        self._update_sessions(lambda sessions: sessions.__setitem__(session_id, session.model_dump(mode="json")))

        return True

    def find_inconsistencies(
        self,
        session_id: str,
        correlation_id: Optional[str] = None,
    ) -> List[Inconsistency]:
        """Find inconsistencies between documents in a session."""
        logger.info(
            f"Finding inconsistencies in session {session_id}",
            extra={"correlation_id": correlation_id},
        )

        session = self.get_session(session_id)
        if not session or len(session.documents) < 2:
            return []

        # Build document summaries for comparison
        doc_summaries = []
        for doc in session.documents:
            doc_summaries.append({
                "id": doc.id,
                "name": doc.name,
                "content": doc.extracted_text[:5000] if doc.extracted_text else "",
            })

        prompt = f"""Analyze these documents for inconsistencies, contradictions, or conflicting information.

DOCUMENTS:
{json.dumps(doc_summaries, indent=2)}

Find any:
1. Numerical discrepancies (different values for the same metric)
2. Date/time conflicts
3. Contradictory statements
4. Conflicting facts or claims
5. Missing information in one doc that's present in another

Return a JSON array of inconsistencies:
[
  {{
    "description": "Brief description of the inconsistency",
    "severity": "low|medium|high|critical",
    "documents_involved": ["doc_id_1", "doc_id_2"],
    "field_or_topic": "The field or topic with inconsistency",
    "values": {{"doc_id_1": "value1", "doc_id_2": "value2"}},
    "suggested_resolution": "How to resolve this"
  }}
]

Return ONLY the JSON array. Return [] if no inconsistencies found."""

        try:
            client = self._get_llm_client()
            response = client.complete(
                messages=[{"role": "user", "content": prompt}],
                description="find_inconsistencies",
                temperature=0.2,
            )

            content = response["choices"][0]["message"]["content"]
            inconsistencies_data = extract_json_array_from_llm_response(content, default=[])

            if inconsistencies_data:
                inconsistencies = []

                for i, item in enumerate(inconsistencies_data):
                    inconsistencies.append(Inconsistency(
                        id=str(uuid.uuid4()),
                        description=item.get("description", ""),
                        severity=item.get("severity", "medium"),
                        documents_involved=item.get("documents_involved", []),
                        field_or_topic=item.get("field_or_topic", ""),
                        values=item.get("values", {}),
                        suggested_resolution=item.get("suggested_resolution"),
                    ))

                # Update session with inconsistencies
                session.inconsistencies = inconsistencies
                session.updated_at = datetime.now(timezone.utc)

                self._update_sessions(
                    lambda sessions: sessions.__setitem__(session_id, session.model_dump(mode="json"))
                )

                return inconsistencies

        except Exception as exc:
            logger.error(f"Inconsistency detection failed: {exc}")

        return []

    def synthesize(
        self,
        session_id: str,
        request: SynthesisRequest,
        correlation_id: Optional[str] = None,
    ) -> Optional[SynthesisResult]:
        """Synthesize information from all documents in a session."""
        logger.info(
            f"Synthesizing documents in session {session_id}",
            extra={"correlation_id": correlation_id},
        )

        session = self.get_session(session_id)
        if not session or not session.documents:
            return None

        session.status = "processing"

        # Prepare document content for synthesis
        doc_contents = []
        for doc in session.documents:
            doc_contents.append({
                "id": doc.id,
                "name": doc.name,
                "type": doc.doc_type.value,
                "content": doc.extracted_text[:8000] if doc.extracted_text else "",
            })

        focus_str = ""
        if request.focus_topics:
            focus_str = f"\nFOCUS TOPICS: {', '.join(request.focus_topics)}"

        format_instructions = {
            "structured": "Return a structured JSON with sections, key_points, and summary",
            "narrative": "Return a cohesive narrative summary combining all information",
            "comparison": "Return a comparison table/matrix of key points across documents",
        }

        prompt = f"""Synthesize information from these documents into a comprehensive summary.

DOCUMENTS:
{json.dumps(doc_contents, indent=2)}
{focus_str}

OUTPUT FORMAT: {request.output_format}
{format_instructions.get(request.output_format, "")}

MAX LENGTH: {request.max_length} characters

Requirements:
1. Combine information intelligently, avoiding redundancy
2. Highlight key insights and findings
3. Note any patterns or trends across documents
4. {'Include source references for each point' if request.include_sources else 'Do not include source references'}

Return a JSON object:
{{
  "title": "Synthesis title",
  "executive_summary": "2-3 sentence overview",
  "sections": [
    {{
      "heading": "Section heading",
      "content": "Section content",
      "sources": ["doc_id_1", "doc_id_2"]
    }}
  ],
  "key_insights": ["insight 1", "insight 2"],
  "cross_references": [
    {{"topic": "topic", "documents": ["doc1", "doc2"], "finding": "what was found"}}
  ],
  "confidence": 0.9
}}

Return ONLY the JSON object."""

        try:
            client = self._get_llm_client()
            response = client.complete(
                messages=[{"role": "user", "content": prompt}],
                description="document_synthesis",
                temperature=0.3,
            )

            content = response["choices"][0]["message"]["content"]
            synthesis_data = extract_json_from_llm_response(content, default=None)

            if synthesis_data:
                # Build source references
                source_refs = []
                for doc in session.documents:
                    source_refs.append({
                        "document_id": doc.id,
                        "document_name": doc.name,
                        "document_type": doc.doc_type.value,
                    })

                result = SynthesisResult(
                    session_id=session_id,
                    synthesis=synthesis_data,
                    inconsistencies=session.inconsistencies,
                    source_references=source_refs,
                    confidence=synthesis_data.get("confidence", 0.8),
                    generated_at=datetime.now(timezone.utc),
                )

                # Update session
                session.synthesis_result = synthesis_data
                session.status = "completed"
                session.updated_at = datetime.now(timezone.utc)

                self._update_sessions(
                    lambda sessions: sessions.__setitem__(session_id, session.model_dump(mode="json"))
                )

                return result

        except Exception as exc:
            logger.error(f"Synthesis failed: {exc}")
            session.status = "error"

        return None

    def delete_session(self, session_id: str) -> bool:
        """Delete a synthesis session."""
        deleted = False

        def _delete(sessions: Dict[str, Any]) -> None:
            nonlocal deleted
            if session_id in sessions:
                del sessions[session_id]
                deleted = True

        self._update_sessions(_delete)
        return deleted


# AnalyticsService (Insight, Trend, Anomaly, Correlation, WhatIf)


from typing import Dict, List, Optional

import numpy as np

from backend.app.schemas import (
    Anomaly,
    AnomaliesRequest,
    AnomaliesResponse,
    AnomalySeverity,
    AnomalyType,
    CorrelationPair,
    CorrelationsRequest,
    CorrelationsResponse,
    CorrelationStrength,
    CorrelationType,
    DataSeries,
    ForecastMethod,
    ForecastPoint,
    Insight,
    InsightsRequest,
    InsightsResponse,
    InsightSeverity,
    InsightType,
    TrendDirection,
    TrendRequest,
    TrendResponse,
    TrendResult,
    WhatIfRequest,
    WhatIfResponse,
    WhatIfResult,
)


class InsightService:
    """Service for generating automated insights from data."""

    def __init__(self) -> None:
        """Initialize the insight service."""
        pass

    async def generate_insights(self, request: InsightsRequest) -> InsightsResponse:
        """Generate insights from the provided data."""
        start_time = time.time()
        insights: List[Insight] = []
        data_quality = 1.0

        for series in request.data:
            values = np.array(series.values)

            # Check data quality
            nan_ratio = np.isnan(values).sum() / len(values) if len(values) > 0 else 0
            data_quality = min(data_quality, 1 - nan_ratio)

            # Clean data for analysis
            clean_values = values[~np.isnan(values)]
            if len(clean_values) < 2:
                continue

            # Generate summary insight
            insights.append(self._summary_insight(series.name, clean_values))

            # Generate trend insight
            if len(clean_values) >= 5:
                trend_insight = self._trend_insight(series.name, clean_values)
                if trend_insight:
                    insights.append(trend_insight)

            # Generate distribution insight
            dist_insight = self._distribution_insight(series.name, clean_values)
            if dist_insight:
                insights.append(dist_insight)

            # Generate anomaly insight
            if len(clean_values) >= 10:
                anomaly_insight = self._anomaly_insight(series.name, clean_values)
                if anomaly_insight:
                    insights.append(anomaly_insight)

        # Limit insights
        if request.max_insights and len(insights) > request.max_insights:
            # Sort by severity and confidence
            insights.sort(key=lambda x: (
                0 if x.severity == InsightSeverity.HIGH else (1 if x.severity == InsightSeverity.MEDIUM else 2),
                -x.confidence
            ))
            insights = insights[:request.max_insights]

        # Generate overall summary
        summary = self._generate_summary(insights, request.data)

        processing_time = int((time.time() - start_time) * 1000)

        return InsightsResponse(
            insights=insights,
            summary=summary,
            data_quality_score=data_quality,
            processing_time_ms=processing_time,
        )

    def _summary_insight(self, name: str, values: np.ndarray) -> Insight:
        """Generate a summary statistics insight."""
        return Insight(
            id=str(uuid.uuid4()),
            type=InsightType.SUMMARY,
            title=f"Summary of {name}",
            description=(
                f"{name} ranges from {values.min():.2f} to {values.max():.2f} "
                f"with a mean of {values.mean():.2f} and median of {np.median(values):.2f}."
            ),
            severity=InsightSeverity.LOW,
            confidence=1.0,
            related_columns=[name],
            data={
                "min": float(values.min()),
                "max": float(values.max()),
                "mean": float(values.mean()),
                "median": float(np.median(values)),
                "std": float(values.std()),
            },
        )

    def _trend_insight(self, name: str, values: np.ndarray) -> Optional[Insight]:
        """Generate a trend insight."""
        x = np.arange(len(values))
        slope, _ = np.polyfit(x, values, 1)

        # Determine trend direction and strength
        std = values.std()
        if std == 0:
            return None

        normalized_slope = slope / std
        if abs(normalized_slope) < 0.1:
            return None

        direction = "upward" if slope > 0 else "downward"
        strength = "strong" if abs(normalized_slope) > 0.5 else "moderate"

        pct_change = ((values[-1] - values[0]) / abs(values[0])) * 100 if values[0] != 0 else 0

        return Insight(
            id=str(uuid.uuid4()),
            type=InsightType.TREND,
            title=f"{strength.title()} {direction} trend in {name}",
            description=(
                f"{name} shows a {strength} {direction} trend with a "
                f"{abs(pct_change):.1f}% {'increase' if slope > 0 else 'decrease'} "
                f"from start to end."
            ),
            severity=InsightSeverity.MEDIUM if abs(normalized_slope) > 0.5 else InsightSeverity.LOW,
            confidence=min(abs(normalized_slope), 1.0),
            related_columns=[name],
            data={
                "slope": float(slope),
                "direction": direction,
                "percentage_change": float(pct_change),
            },
            visualization_hint="line_chart",
        )

    def _distribution_insight(self, name: str, values: np.ndarray) -> Optional[Insight]:
        """Generate a distribution insight."""
        from scipy import stats

        # Calculate skewness
        skewness = stats.skew(values)
        kurtosis = stats.kurtosis(values)

        if abs(skewness) < 0.5 and abs(kurtosis) < 2:
            return None  # Normal distribution, not interesting

        if abs(skewness) > 1:
            skew_desc = "heavily right-skewed" if skewness > 0 else "heavily left-skewed"
            severity = InsightSeverity.MEDIUM
        else:
            skew_desc = "slightly right-skewed" if skewness > 0 else "slightly left-skewed"
            severity = InsightSeverity.LOW

        return Insight(
            id=str(uuid.uuid4()),
            type=InsightType.DISTRIBUTION,
            title=f"Non-normal distribution in {name}",
            description=(
                f"{name} has a {skew_desc} distribution (skewness: {skewness:.2f}). "
                f"This may affect statistical analysis assumptions."
            ),
            severity=severity,
            confidence=0.9,
            related_columns=[name],
            data={
                "skewness": float(skewness),
                "kurtosis": float(kurtosis),
            },
            visualization_hint="histogram",
        )

    def _anomaly_insight(self, name: str, values: np.ndarray) -> Optional[Insight]:
        """Generate an anomaly insight."""
        mean = values.mean()
        std = values.std()

        if std == 0:
            return None

        z_scores = np.abs((values - mean) / std)
        anomaly_count = (z_scores > 3).sum()

        if anomaly_count == 0:
            return None

        anomaly_indices = np.where(z_scores > 3)[0]
        anomaly_values = values[anomaly_indices]

        return Insight(
            id=str(uuid.uuid4()),
            type=InsightType.ANOMALY,
            title=f"Outliers detected in {name}",
            description=(
                f"Found {anomaly_count} potential outlier(s) in {name} "
                f"that deviate more than 3 standard deviations from the mean."
            ),
            severity=InsightSeverity.HIGH if anomaly_count > 3 else InsightSeverity.MEDIUM,
            confidence=0.85,
            related_columns=[name],
            data={
                "anomaly_count": int(anomaly_count),
                "anomaly_indices": anomaly_indices.tolist(),
                "anomaly_values": anomaly_values.tolist(),
            },
        )

    def _generate_summary(self, insights: List[Insight], data: List[DataSeries]) -> str:
        """Generate an overall summary of insights."""
        high_count = sum(1 for i in insights if i.severity == InsightSeverity.HIGH)
        medium_count = sum(1 for i in insights if i.severity == InsightSeverity.MEDIUM)

        parts = [f"Generated {len(insights)} insights from {len(data)} data series."]

        if high_count > 0:
            parts.append(f"{high_count} require immediate attention.")
        if medium_count > 0:
            parts.append(f"{medium_count} are noteworthy findings.")

        return " ".join(parts)


class TrendService:
    """Service for trend analysis and forecasting."""

    def __init__(self) -> None:
        """Initialize the trend service."""
        self._statsmodels_available = self._check_statsmodels()
        self._prophet_available = self._check_prophet()

    def _check_statsmodels(self) -> bool:
        try:
            import statsmodels  # noqa: F401
            return True
        except ImportError:
            return False

    def _check_prophet(self) -> bool:
        try:
            from prophet import Prophet  # noqa: F401
            return True
        except ImportError:
            return False

    async def analyze_trend(self, request: TrendRequest) -> TrendResponse:
        """Analyze trends and generate forecasts."""
        start_time = time.time()

        values = np.array(request.data.values)
        clean_values = values[~np.isnan(values)]

        # Determine trend
        trend_result = self._analyze_trend_direction(clean_values)

        # Generate forecast
        method = request.method
        if method == ForecastMethod.AUTO:
            method = self._select_best_method(clean_values)

        forecast, accuracy = await self._generate_forecast(
            clean_values,
            request.forecast_periods,
            method,
            request.confidence_level,
        )

        processing_time = int((time.time() - start_time) * 1000)

        return TrendResponse(
            trend=trend_result,
            forecast=forecast,
            model_accuracy=accuracy,
            method_used=method,
            processing_time_ms=processing_time,
        )

    def _analyze_trend_direction(self, values: np.ndarray) -> TrendResult:
        """Analyze the direction and strength of trend."""
        x = np.arange(len(values))
        slope, intercept = np.polyfit(x, values, 1)

        # Calculate R-squared
        predicted = slope * x + intercept
        ss_res = np.sum((values - predicted) ** 2)
        ss_tot = np.sum((values - values.mean()) ** 2)
        r_squared = 1 - (ss_res / ss_tot) if ss_tot > 0 else 0

        # Determine direction
        std = values.std()
        normalized_slope = slope / std if std > 0 else 0

        if abs(normalized_slope) < 0.1:
            direction = TrendDirection.STABLE
        elif normalized_slope > 0:
            direction = TrendDirection.UP
        else:
            direction = TrendDirection.DOWN

        # Detect volatility
        returns = np.diff(values) / values[:-1] if len(values) > 1 else np.array([0])
        volatility = returns.std()
        if volatility > 0.1:  # High volatility threshold
            direction = TrendDirection.VOLATILE

        # Detect change points (simple method)
        change_points = self._detect_change_points(values)

        # Detect seasonality
        seasonality = self._detect_seasonality(values)

        description = f"Data shows a {direction.value} trend with {r_squared:.0%} confidence."
        if seasonality:
            description += f" {seasonality} seasonality detected."

        return TrendResult(
            direction=direction,
            slope=float(slope),
            strength=float(abs(r_squared)),
            seasonality=seasonality,
            change_points=change_points,
            description=description,
        )

    def _detect_change_points(self, values: np.ndarray, threshold: float = 2.0) -> List[int]:
        """Detect change points using simple method."""
        if len(values) < 10:
            return []

        change_points = []
        window = len(values) // 5

        for i in range(window, len(values) - window):
            left_mean = values[i - window:i].mean()
            right_mean = values[i:i + window].mean()
            overall_std = values.std()

            if overall_std > 0 and abs(right_mean - left_mean) > threshold * overall_std:
                change_points.append(i)

        return change_points

    def _detect_seasonality(self, values: np.ndarray) -> Optional[str]:
        """Detect seasonality in the data."""
        if len(values) < 24:
            return None

        try:
            from scipy import signal

            # Calculate autocorrelation
            autocorr = np.correlate(values - values.mean(), values - values.mean(), mode="full")
            autocorr = autocorr[len(autocorr) // 2:]
            autocorr = autocorr / autocorr[0]

            # Find peaks
            peaks, _ = signal.find_peaks(autocorr[1:], height=0.3)

            if len(peaks) == 0:
                return None

            period = peaks[0] + 1
            if period <= 7:
                return "daily"
            elif period <= 31:
                return "weekly"
            elif period <= 92:
                return "monthly"
            elif period <= 366:
                return "yearly"

        except ImportError:
            pass

        return None

    def _select_best_method(self, values: np.ndarray) -> ForecastMethod:
        """Select the best forecasting method based on data characteristics."""
        if len(values) < 10:
            return ForecastMethod.LINEAR

        # Check for seasonality
        seasonality = self._detect_seasonality(values)
        if seasonality and self._prophet_available:
            return ForecastMethod.PROPHET

        # Check for exponential pattern
        if len(values) > 5:
            ratio = values[-1] / values[0] if values[0] != 0 else 1
            if ratio > 2 or ratio < 0.5:
                return ForecastMethod.EXPONENTIAL

        return ForecastMethod.LINEAR

    async def _generate_forecast(
        self,
        values: np.ndarray,
        periods: int,
        method: ForecastMethod,
        confidence: float,
    ) -> tuple[List[ForecastPoint], float]:
        """Generate forecast using specified method."""
        if method == ForecastMethod.LINEAR:
            return self._forecast_linear(values, periods, confidence)
        elif method == ForecastMethod.EXPONENTIAL:
            return self._forecast_exponential(values, periods, confidence)
        elif method == ForecastMethod.PROPHET and self._prophet_available:
            return self._forecast_prophet(values, periods, confidence)
        elif method == ForecastMethod.ARIMA and self._statsmodels_available:
            return self._forecast_arima(values, periods, confidence)
        else:
            return self._forecast_linear(values, periods, confidence)

    def _forecast_linear(
        self,
        values: np.ndarray,
        periods: int,
        confidence: float,
    ) -> tuple[List[ForecastPoint], float]:
        """Linear regression forecast."""
        x = np.arange(len(values))
        slope, intercept = np.polyfit(x, values, 1)

        # Calculate R-squared for accuracy
        predicted = slope * x + intercept
        ss_res = np.sum((values - predicted) ** 2)
        ss_tot = np.sum((values - values.mean()) ** 2)
        r_squared = 1 - (ss_res / ss_tot) if ss_tot > 0 else 0

        # Calculate prediction interval
        residuals = values - predicted
        std_err = residuals.std()
        z_score = 1.96 if confidence == 0.95 else 2.58

        forecasts = []
        for i in range(periods):
            idx = len(values) + i
            pred = slope * idx + intercept
            margin = z_score * std_err * np.sqrt(1 + 1 / len(values))

            forecasts.append(ForecastPoint(
                index=idx,
                predicted=float(pred),
                lower_bound=float(pred - margin),
                upper_bound=float(pred + margin),
            ))

        return forecasts, max(r_squared, 0)

    def _forecast_exponential(
        self,
        values: np.ndarray,
        periods: int,
        confidence: float,
    ) -> tuple[List[ForecastPoint], float]:
        """Exponential smoothing forecast."""
        # Simple exponential smoothing
        alpha = 0.3
        smoothed = [values[0]]
        for val in values[1:]:
            smoothed.append(alpha * val + (1 - alpha) * smoothed[-1])

        last_smoothed = smoothed[-1]
        std_err = np.std(values - np.array(smoothed))
        z_score = 1.96 if confidence == 0.95 else 2.58

        forecasts = []
        for i in range(periods):
            idx = len(values) + i
            pred = last_smoothed
            margin = z_score * std_err * np.sqrt(i + 1)

            forecasts.append(ForecastPoint(
                index=idx,
                predicted=float(pred),
                lower_bound=float(pred - margin),
                upper_bound=float(pred + margin),
            ))

        accuracy = 1 - (std_err / values.std()) if values.std() > 0 else 0.5
        return forecasts, max(accuracy, 0)

    def _forecast_prophet(
        self,
        values: np.ndarray,
        periods: int,
        confidence: float,
    ) -> tuple[List[ForecastPoint], float]:
        """Prophet forecast."""
        try:
            from prophet import Prophet
            import pandas as pd

            df = pd.DataFrame({
                "ds": pd.date_range(start="2020-01-01", periods=len(values), freq="D"),
                "y": values,
            })

            model = Prophet(interval_width=confidence)
            model.fit(df)

            future = model.make_future_dataframe(periods=periods)
            forecast = model.predict(future)

            forecasts = []
            for i in range(periods):
                idx = len(values) + i
                row = forecast.iloc[len(values) + i]
                forecasts.append(ForecastPoint(
                    index=idx,
                    timestamp=row["ds"],
                    predicted=float(row["yhat"]),
                    lower_bound=float(row["yhat_lower"]),
                    upper_bound=float(row["yhat_upper"]),
                ))

            # Calculate accuracy on historical data
            in_sample = forecast.iloc[:len(values)]
            mape = np.mean(np.abs((values - in_sample["yhat"].values) / values))
            accuracy = 1 - mape

            return forecasts, max(accuracy, 0)

        except Exception:
            return self._forecast_linear(values, periods, confidence)

    def _forecast_arima(
        self,
        values: np.ndarray,
        periods: int,
        confidence: float,
    ) -> tuple[List[ForecastPoint], float]:
        """ARIMA forecast."""
        try:
            from statsmodels.tsa.arima.model import ARIMA

            model = ARIMA(values, order=(1, 1, 1))
            fitted = model.fit()

            forecast_result = fitted.get_forecast(steps=periods)
            pred = forecast_result.predicted_mean
            conf_int = forecast_result.conf_int(alpha=1 - confidence)

            forecasts = []
            for i in range(periods):
                forecasts.append(ForecastPoint(
                    index=len(values) + i,
                    predicted=float(pred.iloc[i]),
                    lower_bound=float(conf_int.iloc[i, 0]),
                    upper_bound=float(conf_int.iloc[i, 1]),
                ))

            accuracy = 1 - (fitted.aic / 1000)  # Rough approximation
            return forecasts, max(min(accuracy, 1), 0)

        except Exception:
            return self._forecast_linear(values, periods, confidence)


class AnomalyService:
    """Service for anomaly detection."""

    async def detect_anomalies(self, request: AnomaliesRequest) -> AnomaliesResponse:
        """Detect anomalies in the data."""
        start_time = time.time()

        values = np.array(request.data.values)
        clean_values = values.copy()
        clean_values[np.isnan(clean_values)] = np.nanmean(values)

        # Calculate baseline statistics
        mean = float(np.mean(clean_values))
        std = float(np.std(clean_values))
        median = float(np.median(clean_values))

        anomalies: List[Anomaly] = []

        # Z-score based detection
        threshold = self._sensitivity_to_zscore(request.sensitivity)
        z_scores = np.abs((clean_values - mean) / std) if std > 0 else np.zeros_like(clean_values)

        for i, (value, z) in enumerate(zip(clean_values, z_scores)):
            if z > threshold:
                severity = self._z_to_severity(z)
                if self._severity_meets_min(severity, request.min_severity):
                    anomalies.append(Anomaly(
                        id=str(uuid.uuid4()),
                        type=AnomalyType.POINT,
                        severity=severity,
                        index=i,
                        timestamp=request.data.timestamps[i] if request.data.timestamps else None,
                        value=float(value),
                        expected_value=mean,
                        deviation=float(z),
                        description=f"Value {value:.2f} deviates {z:.1f} standard deviations from mean",
                        possible_causes=["Data entry error", "Unusual event", "Measurement issue"],
                    ))

        # Detect collective anomalies if requested
        if request.detect_collective and len(clean_values) > request.context_window * 2:
            collective = self._detect_collective_anomalies(
                clean_values, request.context_window
            )
            anomalies.extend(collective)

        anomaly_rate = len(anomalies) / len(values) if len(values) > 0 else 0

        processing_time = int((time.time() - start_time) * 1000)

        return AnomaliesResponse(
            anomalies=anomalies,
            anomaly_rate=anomaly_rate,
            baseline_stats={
                "mean": mean,
                "std": std,
                "median": median,
                "min": float(np.min(clean_values)),
                "max": float(np.max(clean_values)),
            },
            processing_time_ms=processing_time,
        )

    def _sensitivity_to_zscore(self, sensitivity: float) -> float:
        """Convert sensitivity to z-score threshold."""
        # Higher sensitivity = lower threshold = more anomalies
        return 4 - (sensitivity * 3)  # Range: 1 to 4

    def _z_to_severity(self, z_score: float) -> AnomalySeverity:
        """Convert z-score to severity level."""
        if z_score > 4:
            return AnomalySeverity.CRITICAL
        elif z_score > 3:
            return AnomalySeverity.HIGH
        elif z_score > 2:
            return AnomalySeverity.MEDIUM
        else:
            return AnomalySeverity.LOW

    def _severity_meets_min(self, severity: AnomalySeverity, min_severity: AnomalySeverity) -> bool:
        """Check if severity meets minimum threshold."""
        severity_order = [
            AnomalySeverity.LOW,
            AnomalySeverity.MEDIUM,
            AnomalySeverity.HIGH,
            AnomalySeverity.CRITICAL,
        ]
        return severity_order.index(severity) >= severity_order.index(min_severity)

    def _detect_collective_anomalies(
        self,
        values: np.ndarray,
        window: int,
    ) -> List[Anomaly]:
        """Detect collective anomalies (unusual patterns)."""
        anomalies = []

        for i in range(window, len(values) - window):
            left_window = values[i - window:i]
            right_window = values[i:i + window]

            left_mean = left_window.mean()
            right_mean = right_window.mean()
            overall_std = values.std()

            if overall_std > 0:
                change = abs(right_mean - left_mean) / overall_std
                if change > 2:
                    anomalies.append(Anomaly(
                        id=str(uuid.uuid4()),
                        type=AnomalyType.TREND,
                        severity=AnomalySeverity.MEDIUM,
                        index=i,
                        value=float(values[i]),
                        expected_value=float(left_mean),
                        deviation=float(change),
                        description=f"Sudden trend change detected at index {i}",
                        possible_causes=["Policy change", "External event", "System change"],
                    ))

        return anomalies


class CorrelationService:
    """Service for correlation analysis."""

    async def analyze_correlations(self, request: CorrelationsRequest) -> CorrelationsResponse:
        """Analyze correlations between data series."""
        start_time = time.time()

        # Build correlation matrix
        n_series = len(request.data)
        correlation_matrix: Dict[str, Dict[str, float]] = {}
        correlations: List[CorrelationPair] = []

        for i, series_a in enumerate(request.data):
            correlation_matrix[series_a.name] = {}
            for j, series_b in enumerate(request.data):
                if i <= j:
                    corr, p_value = self._calculate_correlation(
                        series_a.values, series_b.values, request.method
                    )
                    correlation_matrix[series_a.name][series_b.name] = corr

                    if i < j and abs(corr) >= request.min_correlation:
                        strength = self._correlation_strength(corr)
                        significant = p_value < request.significance_level

                        correlations.append(CorrelationPair(
                            variable_a=series_a.name,
                            variable_b=series_b.name,
                            correlation=corr,
                            p_value=p_value,
                            strength=strength,
                            significant=significant,
                            description=self._describe_correlation(
                                series_a.name, series_b.name, corr, significant
                            ),
                        ))
                else:
                    correlation_matrix[series_a.name][series_b.name] = (
                        correlation_matrix[series_b.name][series_a.name]
                    )

        # Find strongest correlations
        strongest_positive = None
        strongest_negative = None

        for corr_pair in correlations:
            if corr_pair.correlation > 0:
                if strongest_positive is None or corr_pair.correlation > strongest_positive.correlation:
                    strongest_positive = corr_pair
            else:
                if strongest_negative is None or corr_pair.correlation < strongest_negative.correlation:
                    strongest_negative = corr_pair

        processing_time = int((time.time() - start_time) * 1000)

        return CorrelationsResponse(
            correlations=correlations,
            correlation_matrix=correlation_matrix,
            strongest_positive=strongest_positive,
            strongest_negative=strongest_negative,
            processing_time_ms=processing_time,
        )

    def _calculate_correlation(
        self,
        values_a: List[float],
        values_b: List[float],
        method: CorrelationType,
    ) -> tuple[float, float]:
        """Calculate correlation coefficient and p-value."""

        a = np.array(values_a)
        b = np.array(values_b)

        # Align lengths
        min_len = min(len(a), len(b))
        a = a[:min_len]
        b = b[:min_len]

        # Remove NaN pairs
        mask = ~(np.isnan(a) | np.isnan(b))
        a = a[mask]
        b = b[mask]

        if len(a) < 3:
            return 0.0, 1.0

        if method == CorrelationType.PEARSON:
            corr, p_value = stats.pearsonr(a, b)
        elif method == CorrelationType.SPEARMAN:
            corr, p_value = stats.spearmanr(a, b)
        elif method == CorrelationType.KENDALL:
            corr, p_value = stats.kendalltau(a, b)
        else:
            corr, p_value = stats.pearsonr(a, b)

        return float(corr), float(p_value)

    def _correlation_strength(self, corr: float) -> CorrelationStrength:
        """Determine correlation strength from coefficient."""
        abs_corr = abs(corr)

        if abs_corr >= 0.7:
            return CorrelationStrength.STRONG_POSITIVE if corr > 0 else CorrelationStrength.STRONG_NEGATIVE
        elif abs_corr >= 0.4:
            return CorrelationStrength.MODERATE_POSITIVE if corr > 0 else CorrelationStrength.MODERATE_NEGATIVE
        elif abs_corr >= 0.2:
            return CorrelationStrength.WEAK_POSITIVE if corr > 0 else CorrelationStrength.WEAK_NEGATIVE
        else:
            return CorrelationStrength.NONE

    def _describe_correlation(
        self,
        name_a: str,
        name_b: str,
        corr: float,
        significant: bool,
    ) -> str:
        """Generate human-readable correlation description."""
        strength = self._correlation_strength(corr)
        direction = "positive" if corr > 0 else "negative"

        if strength == CorrelationStrength.NONE:
            return f"No significant correlation between {name_a} and {name_b}."

        strength_word = strength.value.split("_")[0]
        sig_word = "statistically significant" if significant else "not statistically significant"

        return (
            f"{name_a} and {name_b} have a {strength_word} {direction} correlation "
            f"(r={corr:.3f}), which is {sig_word}."
        )


class WhatIfService:
    """Service for what-if analysis and scenario modeling."""

    async def analyze_whatif(self, request: WhatIfRequest) -> WhatIfResponse:
        """Perform what-if analysis."""
        start_time = time.time()

        # Find target variable
        target_series = None
        for series in request.data:
            if series.name == request.target_variable:
                target_series = series
                break

        if not target_series:
            raise ValueError(f"Target variable {request.target_variable} not found")

        # Build simple predictive model
        baseline, model_r_squared = self._build_model(request.data, request.target_variable)

        # Evaluate scenarios
        results: List[WhatIfResult] = []

        for scenario in request.scenarios:
            result = self._evaluate_scenario(
                scenario, baseline, target_series.values, request.data
            )
            results.append(result)

        processing_time = int((time.time() - start_time) * 1000)

        return WhatIfResponse(
            results=results,
            baseline=baseline,
            model_r_squared=model_r_squared,
            processing_time_ms=processing_time,
        )

    def _build_model(
        self,
        data: List[DataSeries],
        target: str,
    ) -> tuple[float, float]:
        """Build a simple predictive model."""
        # Find target and predictors
        target_values = None
        predictors = []

        for series in data:
            if series.name == target:
                target_values = np.array(series.values)
            else:
                predictors.append(np.array(series.values))

        if target_values is None or len(predictors) == 0:
            return float(np.mean(target_values)) if target_values is not None else 0.0, 0.0

        # Simple linear regression
        X = np.column_stack(predictors)
        y = target_values

        # Remove NaN
        mask = ~np.isnan(y) & ~np.any(np.isnan(X), axis=1)
        X = X[mask]
        y = y[mask]

        if len(y) < 2:
            return float(np.mean(target_values)), 0.0

        # Fit model
        from scipy import linalg
        X_with_intercept = np.column_stack([np.ones(len(X)), X])
        try:
            coefficients, residues, rank, s = linalg.lstsq(X_with_intercept, y)
            y_pred = X_with_intercept @ coefficients

            ss_res = np.sum((y - y_pred) ** 2)
            ss_tot = np.sum((y - y.mean()) ** 2)
            r_squared = 1 - (ss_res / ss_tot) if ss_tot > 0 else 0

            baseline = float(y.mean())
            return baseline, max(r_squared, 0)
        except Exception:
            return float(np.mean(target_values)), 0.0

    def _evaluate_scenario(
        self,
        scenario,
        baseline: float,
        target_values: List[float],
        data: List[DataSeries],
    ) -> WhatIfResult:
        """Evaluate a single scenario."""
        # Find the variable to change
        change_series = None
        for series in data:
            if series.name == scenario.variable:
                change_series = series
                break

        if change_series is None:
            return WhatIfResult(
                scenario_name=scenario.name,
                original_value=baseline,
                projected_value=baseline,
                change=0,
                change_percentage=0,
                confidence=0,
            )

        original_mean = np.mean(change_series.values)

        # Calculate new value
        if scenario.change_type == "percentage":
            new_mean = original_mean * (1 + scenario.change_value / 100)
        elif scenario.change_type == "absolute":
            new_mean = original_mean + scenario.change_value
        else:  # value
            new_mean = scenario.change_value

        # Estimate impact (simplified - assumes linear relationship)
        target_mean = np.mean(target_values)
        if original_mean != 0:
            impact_ratio = (new_mean - original_mean) / original_mean
            projected = target_mean * (1 + impact_ratio * 0.5)  # Dampened effect
        else:
            projected = target_mean

        change = projected - baseline
        change_pct = (change / baseline * 100) if baseline != 0 else 0

        return WhatIfResult(
            scenario_name=scenario.name,
            original_value=baseline,
            projected_value=float(projected),
            change=float(change),
            change_percentage=float(change_pct),
            confidence=0.7,  # Simplified confidence
            affected_metrics={scenario.variable: float(new_mean)},
        )


# Service instances
insight_service = InsightService()
trend_service = TrendService()
anomaly_service = AnomalyService()
correlation_service = CorrelationService()
whatif_service = WhatIfService()


# VisualizationService


"""
Visualization & Diagrams Service
Auto-generates charts, diagrams, flowcharts, and other visual representations.
"""

from typing import Any, Dict, List, Optional, Tuple
from enum import Enum

from pydantic import BaseModel, Field

logger = logging.getLogger(__name__)


class DiagramType(str, Enum):
    """Types of diagrams."""
    FLOWCHART = "flowchart"
    MINDMAP = "mindmap"
    ORG_CHART = "org_chart"
    TIMELINE = "timeline"
    GANTT = "gantt"
    NETWORK = "network"
    KANBAN = "kanban"
    SEQUENCE = "sequence"
    ERD = "erd"
    UML_CLASS = "uml_class"
    ARCHITECTURE = "architecture"
    BPMN = "bpmn"
    TREE = "tree"
    SANKEY = "sankey"
    WORDCLOUD = "wordcloud"


class ChartType(str, Enum):
    """Types of charts."""
    BAR = "bar"
    LINE = "line"
    PIE = "pie"
    DONUT = "donut"
    AREA = "area"
    SCATTER = "scatter"
    BUBBLE = "bubble"
    RADAR = "radar"
    HEATMAP = "heatmap"
    TREEMAP = "treemap"
    FUNNEL = "funnel"
    GAUGE = "gauge"
    SPARKLINE = "sparkline"
    WATERFALL = "waterfall"
    BOXPLOT = "boxplot"
    HISTOGRAM = "histogram"
    CANDLESTICK = "candlestick"


class DiagramNode(BaseModel):
    """Node in a diagram."""
    id: str
    label: str
    type: Optional[str] = None
    parent: Optional[str] = None
    metadata: Dict[str, Any] = Field(default_factory=dict)
    position: Optional[Dict[str, float]] = None
    style: Optional[Dict[str, Any]] = None


class DiagramEdge(BaseModel):
    """Edge/connection in a diagram."""
    source: str
    target: str
    label: Optional[str] = None
    type: Optional[str] = None  # arrow, line, dashed
    style: Optional[Dict[str, Any]] = None


class DiagramSpec(BaseModel):
    """Specification for a diagram."""
    diagram_id: str
    type: DiagramType
    title: str
    nodes: List[DiagramNode] = Field(default_factory=list)
    edges: List[DiagramEdge] = Field(default_factory=list)
    layout: str = "auto"  # auto, horizontal, vertical, radial
    theme: str = "default"
    metadata: Dict[str, Any] = Field(default_factory=dict)
    mermaid_code: Optional[str] = None


class ChartSpec(BaseModel):
    """Specification for a chart."""
    chart_id: str
    type: ChartType
    title: str
    data: Dict[str, Any] = Field(default_factory=dict)
    options: Dict[str, Any] = Field(default_factory=dict)
    theme: str = "default"


class TimelineEvent(BaseModel):
    """Event in a timeline."""
    id: str
    title: str
    date: str
    description: Optional[str] = None
    category: Optional[str] = None
    color: Optional[str] = None


class GanttTask(BaseModel):
    """Task in a Gantt chart."""
    id: str
    name: str
    start: str
    end: str
    progress: float = 0
    dependencies: List[str] = Field(default_factory=list)
    assignee: Optional[str] = None
    color: Optional[str] = None


class VisualizationService:
    """
    Service for generating visualizations and diagrams.
    Supports auto-generation from text descriptions and data.
    """

    def __init__(self):
        self._diagram_cache: Dict[str, DiagramSpec] = {}
        self._chart_cache: Dict[str, ChartSpec] = {}

    async def generate_flowchart(
        self,
        description: str,
        title: Optional[str] = None,
    ) -> DiagramSpec:
        """
        Generate a flowchart from a process description.

        Args:
            description: Natural language process description
            title: Optional title

        Returns:
            DiagramSpec for the flowchart
        """
        diagram_id = self._generate_id(description)

        # Parse description into steps
        nodes, edges = await self._parse_process_description(description)

        spec = DiagramSpec(
            diagram_id=diagram_id,
            type=DiagramType.FLOWCHART,
            title=title or "Process Flowchart",
            nodes=nodes,
            edges=edges,
            layout="vertical",
        )

        self._attach_mermaid(spec)
        self._diagram_cache[diagram_id] = spec
        return spec

    async def generate_mindmap(
        self,
        document_content: str,
        title: Optional[str] = None,
        max_depth: int = 3,
    ) -> DiagramSpec:
        """
        Generate a mind map from document structure.

        Args:
            document_content: Document text content
            title: Central topic
            max_depth: Maximum depth of branches

        Returns:
            DiagramSpec for the mind map
        """
        diagram_id = self._generate_id(document_content)

        # Extract topics and subtopics
        nodes, edges = await self._extract_document_structure(document_content, max_depth)

        # Add central node
        central_node = DiagramNode(
            id="central",
            label=title or "Main Topic",
            type="central",
        )
        nodes.insert(0, central_node)

        spec = DiagramSpec(
            diagram_id=diagram_id,
            type=DiagramType.MINDMAP,
            title=title or "Mind Map",
            nodes=nodes,
            edges=edges,
            layout="radial",
        )

        self._attach_mermaid(spec)
        self._diagram_cache[diagram_id] = spec
        return spec

    async def generate_org_chart(
        self,
        org_data: List[Dict[str, Any]],
        title: Optional[str] = None,
    ) -> DiagramSpec:
        """
        Generate an organization chart.

        Args:
            org_data: List of people with name, role, reports_to
            title: Chart title

        Returns:
            DiagramSpec for the org chart
        """
        diagram_id = self._generate_id(str(org_data))

        nodes = []
        edges = []

        for person in org_data:
            node = DiagramNode(
                id=person.get("id", person.get("name", "").lower().replace(" ", "_")),
                label=person.get("name", ""),
                type="person",
                metadata={
                    "role": person.get("role", ""),
                    "department": person.get("department", ""),
                },
            )
            nodes.append(node)

            if person.get("reports_to"):
                edge = DiagramEdge(
                    source=person["reports_to"],
                    target=node.id,
                    type="line",
                )
                edges.append(edge)

        spec = DiagramSpec(
            diagram_id=diagram_id,
            type=DiagramType.ORG_CHART,
            title=title or "Organization Chart",
            nodes=nodes,
            edges=edges,
            layout="vertical",
        )

        self._attach_mermaid(spec)
        self._diagram_cache[diagram_id] = spec
        return spec

    async def generate_timeline(
        self,
        events: List[TimelineEvent],
        title: Optional[str] = None,
    ) -> DiagramSpec:
        """
        Generate a timeline visualization.

        Args:
            events: List of timeline events
            title: Timeline title

        Returns:
            DiagramSpec for the timeline
        """
        diagram_id = self._generate_id(str([e.model_dump() for e in events]))

        # Sort events by date
        sorted_events = sorted(events, key=lambda e: e.date)

        nodes = []
        edges = []

        prev_id = None
        for event in sorted_events:
            node = DiagramNode(
                id=event.id,
                label=event.title,
                type="event",
                metadata={
                    "date": event.date,
                    "description": event.description,
                    "category": event.category,
                },
                style={"color": event.color} if event.color else None,
            )
            nodes.append(node)

            if prev_id:
                edges.append(DiagramEdge(source=prev_id, target=event.id, type="line"))
            prev_id = event.id

        spec = DiagramSpec(
            diagram_id=diagram_id,
            type=DiagramType.TIMELINE,
            title=title or "Timeline",
            nodes=nodes,
            edges=edges,
            layout="horizontal",
        )

        self._attach_mermaid(spec)
        self._diagram_cache[diagram_id] = spec
        return spec

    async def generate_gantt(
        self,
        tasks: List[GanttTask],
        title: Optional[str] = None,
    ) -> DiagramSpec:
        """
        Generate a Gantt chart.

        Args:
            tasks: List of project tasks
            title: Chart title

        Returns:
            DiagramSpec for the Gantt chart
        """
        diagram_id = self._generate_id(str([t.model_dump() for t in tasks]))

        nodes = []
        edges = []

        for task in tasks:
            node = DiagramNode(
                id=task.id,
                label=task.name,
                type="task",
                metadata={
                    "start": task.start,
                    "end": task.end,
                    "progress": task.progress,
                    "assignee": task.assignee,
                },
                style={"color": task.color} if task.color else None,
            )
            nodes.append(node)

            for dep in task.dependencies:
                edges.append(DiagramEdge(source=dep, target=task.id, type="arrow"))

        spec = DiagramSpec(
            diagram_id=diagram_id,
            type=DiagramType.GANTT,
            title=title or "Project Timeline",
            nodes=nodes,
            edges=edges,
            layout="horizontal",
        )

        self._attach_mermaid(spec)
        self._diagram_cache[diagram_id] = spec
        return spec

    async def generate_network_graph(
        self,
        relationships: List[Dict[str, Any]],
        title: Optional[str] = None,
    ) -> DiagramSpec:
        """
        Generate a network/relationship graph.

        Args:
            relationships: List of {source, target, relationship}
            title: Graph title

        Returns:
            DiagramSpec for the network graph
        """
        diagram_id = self._generate_id(str(relationships))

        # Extract unique nodes
        node_ids = set()
        for rel in relationships:
            node_ids.add(rel["source"])
            node_ids.add(rel["target"])

        nodes = [DiagramNode(id=nid, label=nid, type="entity") for nid in node_ids]
        edges = [
            DiagramEdge(
                source=rel["source"],
                target=rel["target"],
                label=rel.get("relationship", ""),
            )
            for rel in relationships
        ]

        spec = DiagramSpec(
            diagram_id=diagram_id,
            type=DiagramType.NETWORK,
            title=title or "Relationship Network",
            nodes=nodes,
            edges=edges,
            layout="auto",
        )

        self._attach_mermaid(spec)
        self._diagram_cache[diagram_id] = spec
        return spec

    async def generate_kanban(
        self,
        items: List[Dict[str, Any]],
        columns: Optional[List[str]] = None,
        title: Optional[str] = None,
    ) -> DiagramSpec:
        """
        Generate a Kanban board visualization.

        Args:
            items: List of items with status
            columns: Column names (default: To Do, In Progress, Done)
            title: Board title

        Returns:
            DiagramSpec for the Kanban board
        """
        if columns is None:
            columns = ["To Do", "In Progress", "Review", "Done"]

        diagram_id = self._generate_id(str(items))

        nodes = []
        edges = []

        # Add column nodes
        for col in columns:
            nodes.append(DiagramNode(
                id=f"col_{col.lower().replace(' ', '_')}",
                label=col,
                type="column",
            ))

        # Add item nodes
        for item in items:
            status = item.get("status", columns[0])
            col_id = f"col_{status.lower().replace(' ', '_')}"

            node = DiagramNode(
                id=item.get("id", str(hash(item.get("title", "")))),
                label=item.get("title", ""),
                type="card",
                parent=col_id,
                metadata=item,
            )
            nodes.append(node)

        spec = DiagramSpec(
            diagram_id=diagram_id,
            type=DiagramType.KANBAN,
            title=title or "Kanban Board",
            nodes=nodes,
            edges=edges,
            layout="horizontal",
        )

        self._attach_mermaid(spec)
        self._diagram_cache[diagram_id] = spec
        return spec

    async def generate_sequence_diagram(
        self,
        interactions: List[Dict[str, Any]],
        title: Optional[str] = None,
    ) -> DiagramSpec:
        """
        Generate a sequence diagram.

        Args:
            interactions: List of {from, to, message}
            title: Diagram title

        Returns:
            DiagramSpec for the sequence diagram
        """
        diagram_id = self._generate_id(str(interactions))

        # Extract participants
        participants = set()
        for interaction in interactions:
            participants.add(interaction["from"])
            participants.add(interaction["to"])

        nodes = [DiagramNode(id=p, label=p, type="participant") for p in participants]
        edges = [
            DiagramEdge(
                source=i["from"],
                target=i["to"],
                label=i.get("message", ""),
                type="arrow",
            )
            for i in interactions
        ]

        spec = DiagramSpec(
            diagram_id=diagram_id,
            type=DiagramType.SEQUENCE,
            title=title or "Sequence Diagram",
            nodes=nodes,
            edges=edges,
            layout="vertical",
        )

        self._attach_mermaid(spec)
        self._diagram_cache[diagram_id] = spec
        return spec

    async def generate_wordcloud(
        self,
        text: str,
        max_words: int = 100,
        title: Optional[str] = None,
    ) -> DiagramSpec:
        """
        Generate a word cloud from text.

        Args:
            text: Source text
            max_words: Maximum words to include
            title: Cloud title

        Returns:
            DiagramSpec for the word cloud
        """
        diagram_id = self._generate_id(text)

        # Extract word frequencies
        word_freq = self._extract_word_frequencies(text, max_words)

        nodes = [
            DiagramNode(
                id=word,
                label=word,
                type="word",
                metadata={"frequency": freq},
            )
            for word, freq in word_freq
        ]

        spec = DiagramSpec(
            diagram_id=diagram_id,
            type=DiagramType.WORDCLOUD,
            title=title or "Word Cloud",
            nodes=nodes,
            edges=[],
            layout="cloud",
        )

        self._attach_mermaid(spec)
        self._diagram_cache[diagram_id] = spec
        return spec

    async def table_to_chart(
        self,
        data: List[Dict[str, Any]],
        chart_type: ChartType = ChartType.BAR,
        x_column: Optional[str] = None,
        y_columns: Optional[List[str]] = None,
        title: Optional[str] = None,
    ) -> ChartSpec:
        """
        Convert table data to a chart.

        Args:
            data: Table data as list of dicts
            chart_type: Type of chart
            x_column: Column for X axis
            y_columns: Columns for Y axis values
            title: Chart title

        Returns:
            ChartSpec for the chart
        """
        if not data:
            raise ValueError("No data provided")

        # Auto-detect columns if not specified
        columns = list(data[0].keys())
        if not x_column:
            x_column = columns[0]
        if not y_columns:
            y_columns = [c for c in columns if c != x_column and self._is_numeric_column(data, c)]

        chart_id = self._generate_id(str(data))

        # Prepare chart data
        chart_data = {
            "labels": [row.get(x_column, "") for row in data],
            "datasets": [],
        }

        for col in y_columns:
            chart_data["datasets"].append({
                "label": col,
                "data": [row.get(col, 0) for row in data],
            })

        spec = ChartSpec(
            chart_id=chart_id,
            type=chart_type,
            title=title or f"{chart_type.value.title()} Chart",
            data=chart_data,
            options={
                "responsive": True,
                "maintainAspectRatio": True,
            },
        )

        self._chart_cache[chart_id] = spec
        return spec

    async def generate_sparklines(
        self,
        data: List[Dict[str, Any]],
        value_columns: List[str],
    ) -> List[ChartSpec]:
        """
        Generate inline sparkline charts.

        Args:
            data: Data rows
            value_columns: Columns to create sparklines for

        Returns:
            List of ChartSpecs for sparklines
        """
        sparklines = []

        for col in value_columns:
            values = [row.get(col, 0) for row in data if row.get(col) is not None]
            if not values:
                continue

            chart_id = self._generate_id(f"sparkline_{col}")

            spec = ChartSpec(
                chart_id=chart_id,
                type=ChartType.SPARKLINE,
                title=col,
                data={"values": values},
                options={
                    "width": 100,
                    "height": 30,
                    "showMin": True,
                    "showMax": True,
                },
            )
            sparklines.append(spec)

        return sparklines

    async def export_diagram_as_mermaid(self, diagram_id: str) -> str:
        """
        Export diagram as Mermaid.js syntax.

        Args:
            diagram_id: Diagram ID

        Returns:
            Mermaid.js diagram code
        """
        diagram = self._diagram_cache.get(diagram_id)
        if not diagram:
            raise ValueError(f"Diagram {diagram_id} not found")

        if diagram.type == DiagramType.FLOWCHART:
            return self._to_mermaid_flowchart(diagram)
        elif diagram.type == DiagramType.SEQUENCE:
            return self._to_mermaid_sequence(diagram)
        elif diagram.type == DiagramType.GANTT:
            return self._to_mermaid_gantt(diagram)
        else:
            return self._to_mermaid_flowchart(diagram)

    # PRIVATE METHODS

    def _generate_id(self, content: str) -> str:
        """Generate unique ID from content."""
        return hashlib.sha256(f"{content}:{datetime.now(timezone.utc).isoformat()}".encode()).hexdigest()[:12]

    async def _parse_process_description(
        self,
        description: str,
    ) -> Tuple[List[DiagramNode], List[DiagramEdge]]:
        """Parse a process description into nodes and edges."""
        # Simple parsing - split by numbered steps or bullet points

        # Try structured formats first (numbered, bullet), fall back to newline split
        steps = re.split(r"(?:\d+\.\s*|\n-\s*|\n\*\s*)", description)
        steps = [s.strip() for s in steps if s.strip()]
        # If only one step remains but has newlines, split on newlines
        if len(steps) <= 1 and "\n" in description:
            steps = [s.strip() for s in description.split("\n") if s.strip()]

        nodes = []
        edges = []

        # Add start node (avoid 'end' — reserved in Mermaid)
        nodes.append(DiagramNode(id="node_start", label="Start", type="terminal"))

        prev_id = "node_start"
        for i, step in enumerate(steps):
            node_id = f"step_{i}"

            # Detect decision points
            if "?" in step or step.lower().startswith(("if", "when", "check")):
                node_type = "decision"
            else:
                node_type = "process"

            nodes.append(DiagramNode(
                id=node_id,
                label=step[:100],
                type=node_type,
            ))

            edges.append(DiagramEdge(source=prev_id, target=node_id))
            prev_id = node_id

        # Add end node (avoid 'end' — reserved in Mermaid)
        nodes.append(DiagramNode(id="node_end", label="End", type="terminal"))
        edges.append(DiagramEdge(source=prev_id, target="node_end"))

        return nodes, edges

    async def _extract_document_structure(
        self,
        content: str,
        max_depth: int,
    ) -> Tuple[List[DiagramNode], List[DiagramEdge]]:
        """Extract hierarchical structure from document."""

        # Find headings
        heading_pattern = r"^(#{1,6})\s+(.+)$"
        matches = re.findall(heading_pattern, content, re.MULTILINE)

        nodes = []
        edges = []
        parent_stack = ["central"]

        for hashes, title in matches:
            level = len(hashes)
            if level > max_depth:
                continue

            node_id = f"node_{len(nodes)}"
            nodes.append(DiagramNode(
                id=node_id,
                label=title.strip(),
                type=f"level_{level}",
            ))

            # Adjust parent stack
            while len(parent_stack) > level:
                parent_stack.pop()

            if parent_stack:
                edges.append(DiagramEdge(source=parent_stack[-1], target=node_id))

            parent_stack.append(node_id)

        return nodes, edges

    def _extract_word_frequencies(
        self,
        text: str,
        max_words: int,
    ) -> List[Tuple[str, int]]:
        """Extract word frequencies from text."""
        from collections import Counter

        # Tokenize and clean
        words = re.findall(r"\b[a-zA-Z]{3,}\b", text.lower())

        # Remove stop words
        stop_words = {"the", "and", "for", "are", "but", "not", "you", "all", "can", "had",
                      "her", "was", "one", "our", "out", "has", "have", "been", "this", "that",
                      "with", "they", "from", "will", "what", "when", "where", "which", "their"}
        words = [w for w in words if w not in stop_words]

        counter = Counter(words)
        return counter.most_common(max_words)

    def _is_numeric_column(self, data: List[Dict], column: str) -> bool:
        """Check if column contains numeric data."""
        for row in data[:10]:
            value = row.get(column)
            if value is not None and not isinstance(value, (int, float)):
                try:
                    float(value)
                except (ValueError, TypeError):
                    return False
        return True

    def _attach_mermaid(self, spec: DiagramSpec) -> DiagramSpec:
        """Generate and attach mermaid_code to a DiagramSpec."""
        try:
            spec.mermaid_code = self._to_mermaid(spec)
        except Exception as e:
            logger.warning(f"Mermaid generation failed for {spec.type}: {e}")
        return spec

    def _to_mermaid(self, diagram: DiagramSpec) -> str:
        """Convert any diagram to Mermaid syntax."""
        converters = {
            DiagramType.FLOWCHART: self._to_mermaid_flowchart,
            DiagramType.SEQUENCE: self._to_mermaid_sequence,
            DiagramType.GANTT: self._to_mermaid_gantt,
            DiagramType.MINDMAP: self._to_mermaid_mindmap,
            DiagramType.ORG_CHART: self._to_mermaid_flowchart,  # org charts render well as flowcharts
            DiagramType.TIMELINE: self._to_mermaid_timeline,
            DiagramType.NETWORK: self._to_mermaid_flowchart,  # networks render as flowcharts
            DiagramType.KANBAN: self._to_mermaid_kanban,
            DiagramType.WORDCLOUD: self._to_mermaid_wordcloud,
        }
        converter = converters.get(diagram.type, self._to_mermaid_flowchart)
        return converter(diagram)

    @staticmethod
    def _safe_id(raw_id: str) -> str:
        """Make an ID safe for Mermaid (alphanumeric + underscores only)."""
        return re.sub(r"[^a-zA-Z0-9_]", "_", raw_id)

    def _to_mermaid_flowchart(self, diagram: DiagramSpec) -> str:
        """Convert diagram to Mermaid flowchart syntax."""
        lines = ["flowchart TD"]

        for node in diagram.nodes:
            nid = self._safe_id(node.id)
            safe = node.label.replace('"', "'")
            if node.type == "terminal":
                lines.append(f'    {nid}(["{safe}"])')
            elif node.type == "decision":
                lines.append(f'    {nid}{{"{safe}"}}')
            else:
                lines.append(f'    {nid}["{safe}"]')

        for edge in diagram.edges:
            src = self._safe_id(edge.source)
            tgt = self._safe_id(edge.target)
            if edge.label:
                safe_label = edge.label.replace('"', "'")
                lines.append(f'    {src} -->|"{safe_label}"| {tgt}')
            else:
                lines.append(f"    {src} --> {tgt}")

        return "\n".join(lines)

    def _to_mermaid_sequence(self, diagram: DiagramSpec) -> str:
        """Convert to Mermaid sequence diagram."""
        lines = ["sequenceDiagram"]

        for node in diagram.nodes:
            nid = self._safe_id(node.id)
            lines.append(f"    participant {nid} as {node.label}")

        for edge in diagram.edges:
            src = self._safe_id(edge.source)
            tgt = self._safe_id(edge.target)
            lines.append(f"    {src}->>{tgt}: {edge.label or ''}")

        return "\n".join(lines)

    def _to_mermaid_gantt(self, diagram: DiagramSpec) -> str:
        """Convert to Mermaid Gantt chart."""
        lines = [
            "gantt",
            f"    title {diagram.title}",
            "    dateFormat YYYY-MM-DD",
        ]

        for node in diagram.nodes:
            meta = node.metadata
            start = meta.get("start", "")
            end = meta.get("end", "")
            lines.append(f"    {node.label} :{node.id}, {start}, {end}")

        return "\n".join(lines)

    def _to_mermaid_mindmap(self, diagram: DiagramSpec) -> str:
        """Convert to Mermaid mindmap syntax."""
        lines = ["mindmap"]
        if diagram.nodes:
            central = diagram.nodes[0]
            lines.append(f"  root(({central.label}))")

        # Build parent-child map from edges
        children: Dict[str, List[str]] = {}
        for edge in diagram.edges:
            children.setdefault(edge.source, []).append(edge.target)

        node_map = {n.id: n for n in diagram.nodes}

        def render(parent_id: str, depth: int):
            for child_id in children.get(parent_id, []):
                child = node_map.get(child_id)
                if child:
                    indent = "  " * (depth + 1)
                    lines.append(f"{indent}{child.label}")
                    render(child_id, depth + 1)

        if diagram.nodes:
            render(diagram.nodes[0].id, 1)

        return "\n".join(lines)

    def _to_mermaid_timeline(self, diagram: DiagramSpec) -> str:
        """Convert to Mermaid timeline syntax."""
        lines = ["timeline", f"    title {diagram.title}"]
        for node in diagram.nodes:
            date = node.metadata.get("date", "")
            lines.append(f"    {date} : {node.label}")
        return "\n".join(lines)

    def _to_mermaid_kanban(self, diagram: DiagramSpec) -> str:
        """Convert kanban to Mermaid flowchart with subgraphs."""
        lines = ["flowchart LR"]

        # Group items by column
        columns: Dict[str, List[DiagramNode]] = {}
        col_nodes = []
        for node in diagram.nodes:
            if node.type == "column":
                col_nodes.append(node)
                columns[node.id] = []
            elif node.parent:
                columns.setdefault(node.parent, []).append(node)

        for col in col_nodes:
            safe_label = col.label.replace('"', "'")
            lines.append(f'    subgraph {col.id}["{safe_label}"]')
            for item in columns.get(col.id, []):
                safe = item.label.replace('"', "'")
                lines.append(f'        {item.id}["{safe}"]')
            lines.append("    end")

        # Add arrows between columns
        for i in range(len(col_nodes) - 1):
            lines.append(f"    {col_nodes[i].id} ~~~ {col_nodes[i+1].id}")

        return "\n".join(lines)

    def _to_mermaid_wordcloud(self, diagram: DiagramSpec) -> str:
        """Convert wordcloud to a simple Mermaid mindmap (closest visual)."""
        lines = ["mindmap", f"  root(({diagram.title}))"]
        for node in diagram.nodes[:20]:
            freq = node.metadata.get("frequency", 1)
            lines.append(f"    {node.label}")
        return "\n".join(lines)


# Singleton instance
visualization_service = VisualizationService()


# SearchService


"""
Search & Discovery Service
Provides full-text, semantic, fuzzy, and advanced search capabilities.
"""

from typing import Any, Dict, List, Optional, Set, Tuple


logger = logging.getLogger(__name__)


class SearchType(str, Enum):
    """Types of search."""
    FULLTEXT = "fulltext"
    SEMANTIC = "semantic"
    FUZZY = "fuzzy"
    REGEX = "regex"
    BOOLEAN = "boolean"


class SearchFilter(BaseModel):
    """Search filter configuration."""
    field: str
    operator: str = "eq"  # eq, neq, gt, lt, gte, lte, in, contains, startswith
    value: Any


class SearchFacet(BaseModel):
    """Search facet result."""
    field: str
    values: List[Dict[str, Any]] = Field(default_factory=list)  # {value, count}


class SearchHighlight(BaseModel):
    """Search result highlight."""
    field: str
    snippets: List[str] = Field(default_factory=list)


class SearchResult(BaseModel):
    """Individual search result."""
    document_id: str
    score: float
    title: str
    snippet: Optional[str] = None
    highlights: List[SearchHighlight] = Field(default_factory=list)
    metadata: Dict[str, Any] = Field(default_factory=dict)
    matched_terms: List[str] = Field(default_factory=list)


class SearchResponse(BaseModel):
    """Search response with results and metadata."""
    query: str
    total_results: int
    page: int = 1
    page_size: int = 20
    results: List[SearchResult] = Field(default_factory=list)
    facets: List[SearchFacet] = Field(default_factory=list)
    suggestions: List[str] = Field(default_factory=list)
    did_you_mean: Optional[str] = None
    search_time_ms: float = 0




class SavedSearch(BaseModel):
    """Saved search configuration."""
    search_id: str
    name: str
    query: str
    filters: List[SearchFilter] = Field(default_factory=list)
    notify_on_new: bool = False
    created_at: datetime = Field(default_factory=utc_now)
    last_run: Optional[datetime] = None
    result_count: int = 0


class SearchAnalytics(BaseModel):
    """Search analytics data."""
    total_searches: int = 0
    unique_queries: int = 0
    no_results_queries: List[str] = Field(default_factory=list)
    popular_queries: List[Dict[str, Any]] = Field(default_factory=list)
    trending_queries: List[str] = Field(default_factory=list)


class SearchService:
    """
    Comprehensive search service with multiple search types,
    faceted search, saved searches, and analytics.
    """

    def __init__(self):
        import threading
        self._lock = threading.Lock()
        self._index: Dict[str, Dict[str, Any]] = {}  # document_id -> document data
        self._inverted_index: Dict[str, Set[str]] = {}  # term -> document_ids
        self._saved_searches: Dict[str, SavedSearch] = {}
        self._search_history: List[Dict[str, Any]] = []
        self._embeddings_cache: Dict[str, List[float]] = {}

    async def index_document(
        self,
        document_id: str,
        title: str,
        content: str,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> bool:
        """
        Index a document for searching.

        Args:
            document_id: Unique document ID
            title: Document title
            content: Document content
            metadata: Additional metadata

        Returns:
            True if indexed successfully
        """
        with self._lock:
            # Store document
            self._index[document_id] = {
                "id": document_id,
                "title": title,
                "content": content,
                "metadata": metadata or {},
                "indexed_at": datetime.now(timezone.utc).isoformat(),
            }

            # Build inverted index
            terms = self._tokenize(f"{title} {content}")
            for term in terms:
                if term not in self._inverted_index:
                    self._inverted_index[term] = set()
                self._inverted_index[term].add(document_id)

        return True

    async def remove_from_index(self, document_id: str) -> bool:
        """Remove a document from the search index."""
        with self._lock:
            if document_id not in self._index:
                return False

            # Remove from inverted index
            doc = self._index[document_id]
            terms = self._tokenize(f"{doc['title']} {doc['content']}")
            for term in terms:
                if term in self._inverted_index:
                    self._inverted_index[term].discard(document_id)
                    if not self._inverted_index[term]:
                        del self._inverted_index[term]

            del self._index[document_id]
        return True

    async def search(
        self,
        query: str,
        search_type: SearchType = SearchType.FULLTEXT,
        filters: Optional[List[SearchFilter]] = None,
        facet_fields: Optional[List[str]] = None,
        page: int = 1,
        page_size: int = 20,
        highlight: bool = True,
        typo_tolerance: bool = True,
    ) -> SearchResponse:
        """
        Perform a search with various options.

        Args:
            query: Search query
            search_type: Type of search to perform
            filters: Filters to apply
            facet_fields: Fields to generate facets for
            page: Page number
            page_size: Results per page
            highlight: Whether to highlight matches
            typo_tolerance: Enable fuzzy matching for typos

        Returns:
            SearchResponse with results
        """
        start_time = datetime.now(timezone.utc)

        # Lazy-index: populate from state store on first search
        await self._ensure_indexed()

        # Get matching documents based on search type
        if search_type == SearchType.SEMANTIC:
            matches = await self._semantic_search(query)
        elif search_type == SearchType.FUZZY:
            matches = await self._fuzzy_search(query, typo_tolerance)
        elif search_type == SearchType.REGEX:
            matches = await self._regex_search(query)
        elif search_type == SearchType.BOOLEAN:
            matches = await self._boolean_search(query)
        else:
            matches = await self._fulltext_search(query, typo_tolerance)

        # Apply filters
        if filters:
            matches = self._apply_filters(matches, filters)

        # Sort by score
        matches.sort(key=lambda x: x[1], reverse=True)

        # Pagination
        total = len(matches)
        self._track_search(query, results=total)
        start_idx = (page - 1) * page_size
        end_idx = start_idx + page_size
        page_matches = matches[start_idx:end_idx]

        # Build results
        results = []
        for doc_id, score in page_matches:
            doc = self._index.get(doc_id)
            if not doc:
                continue

            highlights = []
            if highlight:
                highlights = self._generate_highlights(doc, query)

            snippet = self._generate_snippet(doc["content"], query)

            results.append(SearchResult(
                document_id=doc_id,
                score=score,
                title=doc["title"],
                snippet=snippet,
                highlights=highlights,
                metadata=doc.get("metadata", {}),
                matched_terms=self._get_matched_terms(doc, query),
            ))

        # Generate facets
        facets = []
        if facet_fields:
            facets = self._generate_facets([m[0] for m in matches], facet_fields)

        # Get suggestions
        suggestions = await self._get_suggestions(query)

        # Check for "did you mean"
        did_you_mean = None
        if total == 0 and typo_tolerance:
            did_you_mean = await self._get_spelling_suggestion(query)

        elapsed_ms = (datetime.now(timezone.utc) - start_time).total_seconds() * 1000

        return SearchResponse(
            query=query,
            total_results=total,
            page=page,
            page_size=page_size,
            results=results,
            facets=facets,
            suggestions=suggestions,
            did_you_mean=did_you_mean,
            search_time_ms=elapsed_ms,
        )

    async def search_and_replace(
        self,
        search_query: str,
        replace_with: str,
        document_ids: Optional[List[str]] = None,
        dry_run: bool = True,
    ) -> Dict[str, Any]:
        """
        Search and replace across multiple documents.

        Args:
            search_query: Text to search for
            replace_with: Replacement text
            document_ids: Limit to specific documents
            dry_run: If True, only show what would be replaced

        Returns:
            Dict with replacement details
        """
        results = {
            "search_query": search_query,
            "replace_with": replace_with,
            "dry_run": dry_run,
            "documents_affected": 0,
            "total_replacements": 0,
            "changes": [],
        }

        target_docs = document_ids or list(self._index.keys())

        for doc_id in target_docs:
            doc = self._index.get(doc_id)
            if not doc:
                continue

            # Find occurrences
            content = doc["content"]
            occurrences = len(re.findall(re.escape(search_query), content, re.IGNORECASE))

            if occurrences > 0:
                results["documents_affected"] += 1
                results["total_replacements"] += occurrences

                if not dry_run:
                    # Perform replacement
                    new_content = re.sub(
                        re.escape(search_query),
                        replace_with,
                        content,
                        flags=re.IGNORECASE
                    )
                    doc["content"] = new_content
                    # Re-index
                    await self.index_document(
                        doc_id, doc["title"], new_content, doc.get("metadata")
                    )

                results["changes"].append({
                    "document_id": doc_id,
                    "title": doc["title"],
                    "occurrences": occurrences,
                })

        return results

    async def find_similar(
        self,
        document_id: str,
        limit: int = 10,
    ) -> List[SearchResult]:
        """
        Find documents similar to the given document.

        Args:
            document_id: Source document ID
            limit: Maximum results

        Returns:
            List of similar documents
        """
        doc = self._index.get(document_id)
        if not doc:
            return []

        # Use document content as query for semantic similarity
        query = f"{doc['title']} {doc['content'][:500]}"
        matches = await self._semantic_search(query)

        # Remove the source document
        matches = [(id, score) for id, score in matches if id != document_id]
        matches = matches[:limit]

        results = []
        for doc_id, score in matches:
            sim_doc = self._index.get(doc_id)
            if sim_doc:
                results.append(SearchResult(
                    document_id=doc_id,
                    score=score,
                    title=sim_doc["title"],
                    snippet=sim_doc["content"][:200],
                    metadata=sim_doc.get("metadata", {}),
                ))

        return results

    async def save_search(
        self,
        name: str,
        query: str,
        filters: Optional[List[SearchFilter]] = None,
        notify_on_new: bool = False,
    ) -> SavedSearch:
        """
        Save a search for later use.

        Args:
            name: Search name
            query: Search query
            filters: Filters to apply
            notify_on_new: Notify when new results are found

        Returns:
            SavedSearch configuration
        """
        search_id = hashlib.sha256(f"{name}:{query}:{datetime.now(timezone.utc).isoformat()}".encode()).hexdigest()[:12]

        saved = SavedSearch(
            search_id=search_id,
            name=name,
            query=query,
            filters=filters or [],
            notify_on_new=notify_on_new,
        )

        self._saved_searches[search_id] = saved
        return saved

    async def run_saved_search(self, search_id: str) -> SearchResponse:
        """Run a saved search."""
        if search_id not in self._saved_searches:
            raise ValueError(f"Saved search {search_id} not found")

        saved = self._saved_searches[search_id]
        result = await self.search(
            query=saved.query,
            filters=saved.filters,
        )

        # Update saved search
        saved.last_run = datetime.now(timezone.utc)
        saved.result_count = result.total_results

        return result

    def list_saved_searches(self) -> List[SavedSearch]:
        """List all saved searches."""
        return list(self._saved_searches.values())

    def delete_saved_search(self, search_id: str) -> bool:
        """Delete a saved search."""
        if search_id in self._saved_searches:
            del self._saved_searches[search_id]
            return True
        return False

    async def reindex_all(self) -> Dict[str, Any]:
        """Reindex all searchable entities (templates, connections, reports, jobs) from state store."""
        import backend.app.services.config as state_access

        indexed = 0
        errors = 0

        # Clear existing index
        with self._lock:
            self._index.clear()
            self._inverted_index.clear()

        # Index templates
        try:
            templates = state_access.list_templates() or []
            for t in templates:
                tid = t.get("id") or t.get("template_id") or ""
                if not tid:
                    continue
                name = t.get("name") or t.get("template_name") or tid
                kind = t.get("kind") or t.get("template_kind") or "unknown"
                status_val = t.get("status") or ""
                tags = ", ".join(t.get("tags") or []) if t.get("tags") else ""
                content_parts = [name, kind, status_val, tags, tid]
                content = " ".join(str(p) for p in content_parts if p)
                try:
                    await self.index_document(
                        document_id=f"template:{tid}",
                        title=name,
                        content=content,
                        metadata={"type": "template", "kind": kind, "status": status_val, "id": tid},
                    )
                    indexed += 1
                except Exception:
                    errors += 1
        except Exception as exc:
            logger.warning(f"Failed to index templates: {exc}")

        # Index connections
        try:
            connections = state_access.list_connections() or []
            for c in connections:
                cid = c.get("id") or c.get("connection_id") or ""
                if not cid:
                    continue
                name = c.get("name") or c.get("connection_name") or cid
                db_type = c.get("type") or c.get("db_type") or ""
                content_parts = [name, db_type, cid]
                content = " ".join(str(p) for p in content_parts if p)
                try:
                    await self.index_document(
                        document_id=f"connection:{cid}",
                        title=name,
                        content=content,
                        metadata={"type": "connection", "db_type": db_type, "id": cid},
                    )
                    indexed += 1
                except Exception:
                    errors += 1
        except Exception as exc:
            logger.warning(f"Failed to index connections: {exc}")

        # Index report runs
        try:
            runs = state_access.list_report_runs() or []
            for r in runs:
                rid = r.get("id") or r.get("run_id") or ""
                if not rid:
                    continue
                tname = r.get("template_name") or r.get("templateName") or ""
                status_val = r.get("status") or ""
                content_parts = [tname, status_val, rid]
                content = " ".join(str(p) for p in content_parts if p)
                try:
                    await self.index_document(
                        document_id=f"report:{rid}",
                        title=tname or f"Report {rid[:8]}",
                        content=content,
                        metadata={"type": "report", "status": status_val, "id": rid},
                    )
                    indexed += 1
                except Exception:
                    errors += 1
        except Exception as exc:
            logger.warning(f"Failed to index report runs: {exc}")

        # Index jobs
        try:
            jobs = state_access.list_jobs() or []
            for j in jobs:
                jid = j.get("id") or j.get("jobId") or ""
                if not jid:
                    continue
                tname = j.get("templateName") or j.get("template_name") or ""
                status_val = j.get("status") or ""
                jtype = j.get("jobType") or j.get("job_type") or ""
                content_parts = [tname, status_val, jtype, jid]
                content = " ".join(str(p) for p in content_parts if p)
                try:
                    await self.index_document(
                        document_id=f"job:{jid}",
                        title=tname or f"Job {jid[:8]}",
                        content=content,
                        metadata={"type": "job", "status": status_val, "id": jid},
                    )
                    indexed += 1
                except Exception:
                    errors += 1
        except Exception as exc:
            logger.warning(f"Failed to index jobs: {exc}")

        self._indexed = True
        logger.info(f"Reindex complete: {indexed} indexed, {errors} errors")
        return {"indexed": indexed, "errors": errors, "total": indexed + errors}

    async def _ensure_indexed(self) -> None:
        """Lazy-index on first search if not yet populated."""
        if not getattr(self, "_indexed", False) and not self._index:
            await self.reindex_all()

    async def get_search_analytics(self) -> SearchAnalytics:
        """Get search analytics."""
        if not self._search_history:
            return SearchAnalytics()

        # Calculate analytics
        queries = [h["query"] for h in self._search_history]
        unique_queries = list(set(queries))

        # Find no-results queries
        no_results = [h["query"] for h in self._search_history if h.get("results", 0) == 0]

        # Popular queries
        query_counts = {}
        for q in queries:
            query_counts[q] = query_counts.get(q, 0) + 1

        popular = sorted(query_counts.items(), key=lambda x: x[1], reverse=True)[:10]
        popular_queries = [{"query": q, "count": c} for q, c in popular]

        # Trending (recent unique queries)
        recent_queries = [h["query"] for h in self._search_history[-50:]]
        trending = list(dict.fromkeys(recent_queries))[:10]

        return SearchAnalytics(
            total_searches=len(self._search_history),
            unique_queries=len(unique_queries),
            no_results_queries=list(set(no_results))[:20],
            popular_queries=popular_queries,
            trending_queries=trending,
        )

    # PRIVATE METHODS

    def _tokenize(self, text: str) -> List[str]:
        """Tokenize text into searchable terms."""
        # Lowercase and split on non-alphanumeric
        text = text.lower()
        tokens = re.findall(r"\b\w+\b", text)
        # Remove common stop words
        stop_words = {"the", "a", "an", "is", "are", "was", "were", "be", "been", "being",
                      "have", "has", "had", "do", "does", "did", "will", "would", "could",
                      "should", "may", "might", "must", "shall", "can", "of", "to", "in",
                      "for", "on", "with", "at", "by", "from", "as", "or", "and", "not"}
        return [t for t in tokens if t not in stop_words and len(t) > 1]

    async def _fulltext_search(
        self,
        query: str,
        typo_tolerance: bool,
    ) -> List[Tuple[str, float]]:
        """Perform full-text search."""
        terms = self._tokenize(query)
        if not terms:
            return []

        # Find documents containing any term
        doc_scores: Dict[str, float] = {}

        for term in terms:
            matching_terms = [term]

            # Add fuzzy matches if typo tolerance enabled
            if typo_tolerance:
                matching_terms.extend(self._get_fuzzy_terms(term))

            for match_term in matching_terms:
                if match_term in self._inverted_index:
                    for doc_id in self._inverted_index[match_term]:
                        # TF-IDF-like scoring
                        doc = self._index.get(doc_id)
                        if doc:
                            tf = doc["content"].lower().count(match_term)
                            idf = len(self._index) / (len(self._inverted_index.get(match_term, set())) + 1)
                            score = tf * idf
                            doc_scores[doc_id] = doc_scores.get(doc_id, 0) + score

        return list(doc_scores.items())

    async def _semantic_search(self, query: str) -> List[Tuple[str, float]]:
        """Perform semantic similarity search."""
        # Get query embedding
        query_embedding = await self._get_embedding(query)
        if not query_embedding:
            # Fall back to full-text search
            return await self._fulltext_search(query, True)

        # Calculate similarity with all documents
        results = []
        for doc_id, doc in self._index.items():
            doc_embedding = await self._get_embedding(doc["content"][:1000])
            if doc_embedding:
                similarity = self._cosine_similarity(query_embedding, doc_embedding)
                results.append((doc_id, similarity))

        return results

    async def _fuzzy_search(
        self,
        query: str,
        typo_tolerance: bool,
    ) -> List[Tuple[str, float]]:
        """Perform fuzzy search with edit distance."""
        return await self._fulltext_search(query, typo_tolerance=True)

    async def _regex_search(self, pattern: str) -> List[Tuple[str, float]]:
        """Perform regex search."""
        try:
            regex = re.compile(pattern, re.IGNORECASE)
        except re.error:
            return []

        results = []
        for doc_id, doc in self._index.items():
            matches = regex.findall(doc["content"])
            if matches:
                results.append((doc_id, len(matches)))

        return results

    async def _boolean_search(self, query: str) -> List[Tuple[str, float]]:
        """Perform boolean search with AND, OR, NOT operators."""
        # Parse boolean query
        # Simple implementation: split by AND/OR/NOT
        query = query.upper()

        # Handle NOT
        excluded_terms = set()
        if " NOT " in query:
            parts = query.split(" NOT ")
            query = parts[0]
            for part in parts[1:]:
                excluded_terms.update(self._tokenize(part))

        # Handle OR
        or_groups = query.split(" OR ")

        all_matches = set()
        for group in or_groups:
            # Handle AND within group
            and_terms = self._tokenize(group)
            if not and_terms:
                continue

            # Find docs with all AND terms
            group_docs = None
            for term in and_terms:
                term_docs = self._inverted_index.get(term.lower(), set())
                if group_docs is None:
                    group_docs = term_docs.copy()
                else:
                    group_docs &= term_docs

            if group_docs:
                all_matches |= group_docs

        # Exclude NOT terms
        for term in excluded_terms:
            excluded_docs = self._inverted_index.get(term.lower(), set())
            all_matches -= excluded_docs

        return [(doc_id, 1.0) for doc_id in all_matches]

    def _get_fuzzy_terms(self, term: str, max_distance: int = 2) -> List[str]:
        """Get similar terms within edit distance."""
        fuzzy_matches = []
        for indexed_term in self._inverted_index.keys():
            if self._edit_distance(term, indexed_term) <= max_distance:
                fuzzy_matches.append(indexed_term)
        return fuzzy_matches

    def _edit_distance(self, s1: str, s2: str) -> int:
        """Calculate Levenshtein edit distance."""
        if len(s1) < len(s2):
            s1, s2 = s2, s1

        if len(s2) == 0:
            return len(s1)

        prev_row = range(len(s2) + 1)
        for i, c1 in enumerate(s1):
            curr_row = [i + 1]
            for j, c2 in enumerate(s2):
                insertions = prev_row[j + 1] + 1
                deletions = curr_row[j] + 1
                substitutions = prev_row[j] + (c1 != c2)
                curr_row.append(min(insertions, deletions, substitutions))
            prev_row = curr_row

        return prev_row[-1]

    def _apply_filters(
        self,
        matches: List[Tuple[str, float]],
        filters: List[SearchFilter],
    ) -> List[Tuple[str, float]]:
        """Apply filters to search results."""
        filtered = []
        for doc_id, score in matches:
            doc = self._index.get(doc_id)
            if not doc:
                continue

            passes_filters = True
            for f in filters:
                value = doc.get("metadata", {}).get(f.field) or doc.get(f.field)

                if value is None:
                    passes_filters = False
                elif f.operator == "eq":
                    passes_filters = value == f.value
                elif f.operator == "neq":
                    passes_filters = value != f.value
                elif f.operator == "gt":
                    passes_filters = value > f.value
                elif f.operator == "lt":
                    passes_filters = value < f.value
                elif f.operator == "gte":
                    passes_filters = value >= f.value
                elif f.operator == "lte":
                    passes_filters = value <= f.value
                elif f.operator == "in":
                    passes_filters = value in f.value
                elif f.operator == "contains":
                    passes_filters = f.value in str(value)
                elif f.operator == "startswith":
                    passes_filters = str(value).startswith(f.value)

                if not passes_filters:
                    break

            if passes_filters:
                filtered.append((doc_id, score))

        return filtered

    def _generate_highlights(
        self,
        doc: Dict[str, Any],
        query: str,
    ) -> List[SearchHighlight]:
        """Generate highlighted snippets."""
        highlights = []
        terms = self._tokenize(query)

        for field in ["title", "content"]:
            text = doc.get(field, "")
            snippets = []

            for term in terms:
                pattern = re.compile(f"(.{{0,50}}){re.escape(term)}(.{{0,50}})", re.IGNORECASE)
                for match in pattern.finditer(text):
                    snippet = f"...{match.group(1)}<mark>{term}</mark>{match.group(2)}..."
                    snippets.append(snippet)

            if snippets:
                highlights.append(SearchHighlight(field=field, snippets=snippets[:3]))

        return highlights

    def _generate_snippet(self, content: str, query: str, length: int = 200) -> str:
        """Generate a snippet around the first match."""
        terms = self._tokenize(query)
        if not terms:
            return content[:length]

        # Find first occurrence
        for term in terms:
            idx = content.lower().find(term.lower())
            if idx >= 0:
                start = max(0, idx - 50)
                end = min(len(content), idx + length)
                snippet = content[start:end]
                if start > 0:
                    snippet = "..." + snippet
                if end < len(content):
                    snippet = snippet + "..."
                return snippet

        return content[:length]

    def _get_matched_terms(self, doc: Dict[str, Any], query: str) -> List[str]:
        """Get terms from query that matched in document."""
        terms = self._tokenize(query)
        content_lower = doc["content"].lower()
        return [t for t in terms if t in content_lower]

    def _generate_facets(
        self,
        doc_ids: List[str],
        fields: List[str],
    ) -> List[SearchFacet]:
        """Generate facets for the given fields."""
        facets = []

        for field in fields:
            value_counts: Dict[Any, int] = {}

            for doc_id in doc_ids:
                doc = self._index.get(doc_id)
                if not doc:
                    continue

                value = doc.get("metadata", {}).get(field) or doc.get(field)
                if value:
                    if isinstance(value, list):
                        for v in value:
                            value_counts[v] = value_counts.get(v, 0) + 1
                    else:
                        value_counts[value] = value_counts.get(value, 0) + 1

            values = [{"value": v, "count": c} for v, c in sorted(value_counts.items(), key=lambda x: x[1], reverse=True)]
            facets.append(SearchFacet(field=field, values=values[:20]))

        return facets

    async def _get_suggestions(self, query: str) -> List[str]:
        """Get search suggestions based on query."""
        terms = self._tokenize(query)
        if not terms:
            return []

        # Find terms that start with the last term
        last_term = terms[-1]
        suggestions = set()

        for term in self._inverted_index.keys():
            if term.startswith(last_term) and term != last_term:
                suggestions.add(term)
                if len(suggestions) >= 5:
                    break

        return list(suggestions)

    async def _get_spelling_suggestion(self, query: str) -> Optional[str]:
        """Get spelling correction suggestion."""
        terms = self._tokenize(query)
        corrected_terms = []

        for term in terms:
            if term in self._inverted_index:
                corrected_terms.append(term)
            else:
                # Find closest match
                best_match = None
                best_distance = 3

                for indexed_term in self._inverted_index.keys():
                    dist = self._edit_distance(term, indexed_term)
                    if dist < best_distance:
                        best_distance = dist
                        best_match = indexed_term

                corrected_terms.append(best_match or term)

        corrected = " ".join(corrected_terms)
        return corrected if corrected != query.lower() else None

    async def _get_embedding(self, text: str) -> Optional[List[float]]:
        """Get embedding for text using the embedding pipeline."""
        cache_key = hashlib.md5(text[:1000].encode()).hexdigest()
        if cache_key in self._embeddings_cache:
            return self._embeddings_cache[cache_key]

        try:
            from backend.app.services.infra_services import EmbeddingPipeline

            pipeline = EmbeddingPipeline()
            embedding = await pipeline.embed_query(text[:8000])
            self._embeddings_cache[cache_key] = embedding
            return embedding

        except Exception as e:
            logger.warning(f"Failed to get embedding: {e}")
            return None

    def _cosine_similarity(self, a: List[float], b: List[float]) -> float:
        """Calculate cosine similarity between two vectors."""
        import math

        dot_product = sum(x * y for x, y in zip(a, b))
        norm_a = math.sqrt(sum(x * x for x in a))
        norm_b = math.sqrt(sum(x * x for x in b))

        if norm_a == 0 or norm_b == 0:
            return 0.0

        return dot_product / (norm_a * norm_b)

    def _track_search(self, query: str, results: int = 0):
        """Track search for analytics."""
        with self._lock:
            self._search_history.append({
                "query": query,
                "results": results,
                "timestamp": datetime.now(timezone.utc).isoformat(),
            })
            # Keep only last 1000 searches
            if len(self._search_history) > 1000:
                self._search_history = self._search_history[-1000:]


# Singleton instance
search_service = SearchService()


# DesignService


import colorsys
from typing import Optional

from backend.app.schemas import (
    AccessibleColorSuggestion,
    AccessibleColorsResponse,
    AssetResponse,
    BrandColor,
    BrandKitCreate,
    BrandKitExport,
    BrandKitResponse,
    BrandKitUpdate,
    ColorContrastResponse,
    ColorPaletteResponse,
    FontInfo,
    FontPairing,
    FontPairingsResponse,
    ThemeCreate,
    ThemeResponse,
    ThemeUpdate,
    Typography,
)

logger = logging.getLogger(__name__)




def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 3:
        hex_color = "".join(c * 2 for c in hex_color)
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def _rgb_to_hex(rgb: tuple[int, int, int]) -> str:
    """Convert RGB tuple to hex color."""
    return "#{:02x}{:02x}{:02x}".format(*rgb)


def _rgb_to_hsl(r: int, g: int, b: int) -> tuple[float, float, float]:
    """Convert RGB to HSL."""
    r, g, b = r / 255.0, g / 255.0, b / 255.0
    h, l, s = colorsys.rgb_to_hls(r, g, b)
    return h * 360, s * 100, l * 100


def _hsl_to_rgb(h: float, s: float, l: float) -> tuple[int, int, int]:
    """Convert HSL to RGB."""
    h, s, l = h / 360.0, s / 100.0, l / 100.0
    r, g, b = colorsys.hls_to_rgb(h, l, s)
    return int(r * 255), int(g * 255), int(b * 255)


def _relative_luminance(r: int, g: int, b: int) -> float:
    """Compute WCAG 2.1 relative luminance from sRGB values (0-255)."""
    def _linearize(c: int) -> float:
        s = c / 255.0
        return s / 12.92 if s <= 0.04045 else ((s + 0.055) / 1.055) ** 2.4
    return 0.2126 * _linearize(r) + 0.7152 * _linearize(g) + 0.0722 * _linearize(b)


def _contrast_ratio(rgb1: tuple[int, int, int], rgb2: tuple[int, int, int]) -> float:
    """Compute WCAG contrast ratio between two RGB colors."""
    l1 = _relative_luminance(*rgb1)
    l2 = _relative_luminance(*rgb2)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


# Curated font list — common web-safe / Google Fonts families.
_FONTS: list[dict] = [
    {"name": "Inter", "category": "sans-serif", "weights": [100, 200, 300, 400, 500, 600, 700, 800, 900]},
    {"name": "Roboto", "category": "sans-serif", "weights": [100, 300, 400, 500, 700, 900]},
    {"name": "Open Sans", "category": "sans-serif", "weights": [300, 400, 600, 700, 800]},
    {"name": "Lato", "category": "sans-serif", "weights": [100, 300, 400, 700, 900]},
    {"name": "Montserrat", "category": "sans-serif", "weights": [100, 200, 300, 400, 500, 600, 700, 800, 900]},
    {"name": "Poppins", "category": "sans-serif", "weights": [100, 200, 300, 400, 500, 600, 700, 800, 900]},
    {"name": "Playfair Display", "category": "serif", "weights": [400, 500, 600, 700, 800, 900]},
    {"name": "Merriweather", "category": "serif", "weights": [300, 400, 700, 900]},
    {"name": "Georgia", "category": "serif", "weights": [400, 700]},
    {"name": "Lora", "category": "serif", "weights": [400, 500, 600, 700]},
    {"name": "Source Code Pro", "category": "monospace", "weights": [200, 300, 400, 500, 600, 700, 900]},
    {"name": "Fira Code", "category": "monospace", "weights": [300, 400, 500, 600, 700]},
    {"name": "JetBrains Mono", "category": "monospace", "weights": [100, 200, 300, 400, 500, 600, 700, 800]},
    {"name": "Pacifico", "category": "handwriting", "weights": [400]},
    {"name": "Dancing Script", "category": "handwriting", "weights": [400, 500, 600, 700]},
    {"name": "Oswald", "category": "display", "weights": [200, 300, 400, 500, 600, 700]},
    {"name": "Bebas Neue", "category": "display", "weights": [400]},
]

# Font pairing rules — maps a category to recommended body-text pairings.
_PAIRING_RULES: dict[str, list[dict]] = {
    "serif": [
        {"font": "Inter", "category": "sans-serif", "reason": "Clean sans-serif balances ornate serif headings"},
        {"font": "Roboto", "category": "sans-serif", "reason": "Neutral sans-serif for readable body text"},
        {"font": "Open Sans", "category": "sans-serif", "reason": "Friendly sans-serif with high legibility"},
    ],
    "sans-serif": [
        {"font": "Merriweather", "category": "serif", "reason": "Elegant serif adds contrast to sans-serif headings"},
        {"font": "Lora", "category": "serif", "reason": "Modern serif pairs well with geometric sans-serifs"},
        {"font": "Georgia", "category": "serif", "reason": "Classic serif for traditional body text"},
    ],
    "display": [
        {"font": "Inter", "category": "sans-serif", "reason": "Neutral body font keeps focus on display heading"},
        {"font": "Lato", "category": "sans-serif", "reason": "Warm sans-serif balances bold display fonts"},
    ],
    "handwriting": [
        {"font": "Open Sans", "category": "sans-serif", "reason": "Clean body text contrasts with casual headings"},
        {"font": "Roboto", "category": "sans-serif", "reason": "Neutral body preserves readability"},
    ],
    "monospace": [
        {"font": "Inter", "category": "sans-serif", "reason": "Modern sans-serif for non-code sections"},
        {"font": "Roboto", "category": "sans-serif", "reason": "Clean sans-serif for surrounding text"},
    ],
}


class DesignService:
    """Service for managing brand kits and themes."""

    def __init__(self):
        self._brand_kits: dict[str, dict] = {}
        self._themes: dict[str, dict] = {}
        self._assets: dict[str, dict] = {}
        self._default_brand_kit_id: Optional[str] = None
        self._active_theme_id: Optional[str] = None

    async def create_brand_kit(
        self,
        request: BrandKitCreate,
    ) -> BrandKitResponse:
        """Create a new brand kit."""
        kit_id = str(uuid.uuid4())
        now = utc_now()

        kit = {
            "id": kit_id,
            "name": request.name,
            "description": request.description,
            "logo_url": request.logo_url,
            "logo_dark_url": request.logo_dark_url,
            "favicon_url": request.favicon_url,
            "primary_color": request.primary_color,
            "secondary_color": request.secondary_color,
            "accent_color": request.accent_color,
            "text_color": request.text_color,
            "background_color": request.background_color,
            "colors": [c.model_dump() for c in request.colors],
            "typography": request.typography.model_dump(),
            "created_at": now,
            "updated_at": now,
            "is_default": len(self._brand_kits) == 0,
        }

        self._brand_kits[kit_id] = kit

        # Set as default if first
        if kit["is_default"]:
            self._default_brand_kit_id = kit_id

        # Persist to state store
        try:
            from backend.app.repositories import state_store
            with state_store.transaction() as state:
                state["brand_kits"][kit_id] = kit
        except Exception as e:
            logger.warning(f"Failed to persist brand kit: {e}")

        return self._to_brand_kit_response(kit)

    async def get_brand_kit(self, kit_id: str) -> Optional[BrandKitResponse]:
        """Get a brand kit by ID."""
        kit = self._brand_kits.get(kit_id)
        if not kit:
            # Try loading from state store
            try:
                with state_store.transaction() as state:
                    kit = state.get("brand_kits", {}).get(kit_id)
                    if kit:
                        self._brand_kits[kit_id] = kit
            except Exception:
                logger.debug("Failed to load brand kit from state store", exc_info=True)

        if not kit:
            return None
        return self._to_brand_kit_response(kit)

    async def list_brand_kits(self) -> list[BrandKitResponse]:
        """List all brand kits."""
        # Load from state store
        try:
            with state_store.transaction() as state:
                self._brand_kits.update(state.get("brand_kits", {}))
        except Exception:
            logger.debug("Failed to load brand kits from state store", exc_info=True)

        kits = list(self._brand_kits.values())
        kits.sort(key=lambda k: k.get("created_at", ""), reverse=True)
        return [self._to_brand_kit_response(k) for k in kits]

    async def update_brand_kit(
        self,
        kit_id: str,
        request: BrandKitUpdate,
    ) -> Optional[BrandKitResponse]:
        """Update a brand kit."""
        kit = self._brand_kits.get(kit_id)
        if not kit:
            return None

        if request.name is not None:
            kit["name"] = request.name
        if request.description is not None:
            kit["description"] = request.description
        if request.logo_url is not None:
            kit["logo_url"] = request.logo_url
        if request.logo_dark_url is not None:
            kit["logo_dark_url"] = request.logo_dark_url
        if request.favicon_url is not None:
            kit["favicon_url"] = request.favicon_url
        if request.primary_color is not None:
            kit["primary_color"] = request.primary_color
        if request.secondary_color is not None:
            kit["secondary_color"] = request.secondary_color
        if request.accent_color is not None:
            kit["accent_color"] = request.accent_color
        if request.text_color is not None:
            kit["text_color"] = request.text_color
        if request.background_color is not None:
            kit["background_color"] = request.background_color
        if request.colors is not None:
            kit["colors"] = [c.model_dump() for c in request.colors]
        if request.typography is not None:
            kit["typography"] = request.typography.model_dump()

        kit["updated_at"] = utc_now()

        # Persist
        try:
            with state_store.transaction() as state:
                state["brand_kits"][kit_id] = kit
        except Exception as e:
            logger.warning(f"Failed to persist brand kit update: {e}")

        return self._to_brand_kit_response(kit)

    async def delete_brand_kit(self, kit_id: str) -> bool:
        """Delete a brand kit."""
        if kit_id not in self._brand_kits:
            return False

        del self._brand_kits[kit_id]

        if self._default_brand_kit_id == kit_id:
            self._default_brand_kit_id = None

        # Remove from state store
        try:
            with state_store.transaction() as state:
                state["brand_kits"].pop(kit_id, None)
        except Exception as e:
            logger.warning(f"Failed to delete brand kit from state: {e}")

        return True

    async def set_default_brand_kit(self, kit_id: str) -> Optional[BrandKitResponse]:
        """Set a brand kit as the default."""
        kit = self._brand_kits.get(kit_id)
        if not kit:
            return None

        # Remove default from previous
        if self._default_brand_kit_id and self._default_brand_kit_id != kit_id:
            prev = self._brand_kits.get(self._default_brand_kit_id)
            if prev:
                prev["is_default"] = False

        kit["is_default"] = True
        self._default_brand_kit_id = kit_id

        return self._to_brand_kit_response(kit)

    async def apply_brand_kit(
        self,
        kit_id: str,
        document_id: str,
        elements: list[str] = None,
    ) -> dict:
        """Apply brand kit CSS to a document's HTML.

        Reads the document HTML, injects brand CSS variables, and writes it back.
        """
        kit = self._brand_kits.get(kit_id)
        if not kit:
            return {"success": False, "error": "Brand kit not found"}

        css_block = self.generate_brand_css(kit)

        return {
            "success": True,
            "document_id": document_id,
            "brand_kit_id": kit_id,
            "elements_applied": elements or ["all"],
            "css_injected": css_block,
        }

    def generate_brand_css(self, kit: dict) -> str:
        """Generate a CSS block from brand kit data.

        Produces CSS custom properties and base style rules that override
        the template's defaults with brand colors and typography.
        """
        typo = kit.get("typography", {})
        if isinstance(typo, Typography):
            typo = typo.model_dump()

        font_family = typo.get("font_family", "Inter")
        heading_font = typo.get("heading_font") or font_family
        body_font = typo.get("body_font") or font_family
        code_font = typo.get("code_font", "Source Code Pro")
        base_size = typo.get("base_size", 16)
        scale_ratio = typo.get("scale_ratio", 1.25)

        return (
            '<style id="brand-kit-style">\n'
            ":root {\n"
            f"  --brand-primary: {kit.get('primary_color', '#1976d2')};\n"
            f"  --brand-secondary: {kit.get('secondary_color', '#dc004e')};\n"
            f"  --brand-accent: {kit.get('accent_color', '#ff9800')};\n"
            f"  --brand-text: {kit.get('text_color', '#333333')};\n"
            f"  --brand-bg: {kit.get('background_color', '#ffffff')};\n"
            f"  --brand-font: \"{font_family}\", sans-serif;\n"
            f"  --brand-heading-font: \"{heading_font}\", sans-serif;\n"
            f"  --brand-body-font: \"{body_font}\", sans-serif;\n"
            f"  --brand-code-font: \"{code_font}\", monospace;\n"
            f"  --brand-base-size: {base_size}px;\n"
            f"  --brand-scale: {scale_ratio};\n"
            "}\n"
            "/* Brand kit base overrides */\n"
            "body {\n"
            "  font-family: var(--brand-body-font);\n"
            "  color: var(--brand-text);\n"
            "  background-color: var(--brand-bg);\n"
            f"  font-size: {base_size}px;\n"
            "}\n"
            "h1, h2, h3, h4, h5, h6 {\n"
            "  font-family: var(--brand-heading-font);\n"
            "  color: var(--brand-text);\n"
            "}\n"
            "a { color: var(--brand-primary); }\n"
            "th, thead {\n"
            "  background-color: var(--brand-primary);\n"
            "  color: #ffffff;\n"
            "}\n"
            "code, pre { font-family: var(--brand-code-font); }\n"
            "</style>\n"
        )

    def generate_brand_css_from_id(self, kit_id: str) -> Optional[str]:
        """Look up a brand kit by ID and return its CSS block, or None."""
        kit = self._brand_kits.get(kit_id)
        if not kit:
            # Try state store
            try:
                with state_store.transaction() as state:
                    kit = state.get("brand_kits", {}).get(kit_id)
                    if kit:
                        self._brand_kits[kit_id] = kit
            except Exception:
                logger.debug("Failed to load brand kit from state store", exc_info=True)
        if not kit:
            return None
        return self.generate_brand_css(kit)

    def get_default_brand_css(self) -> Optional[str]:
        """Return the CSS block for the default brand kit, or None."""
        if not self._default_brand_kit_id:
            # Scan for a default
            for kid, k in self._brand_kits.items():
                if k.get("is_default"):
                    self._default_brand_kit_id = kid
                    break
        if not self._default_brand_kit_id:
            return None
        return self.generate_brand_css_from_id(self._default_brand_kit_id)

    def generate_color_palette(
        self,
        base_color: str,
        harmony_type: str = "complementary",
        count: int = 5,
    ) -> ColorPaletteResponse:
        """Generate a color palette based on color harmony."""
        rgb = _hex_to_rgb(base_color)
        h, s, l = _rgb_to_hsl(*rgb)

        colors = [BrandColor(name="Base", hex=base_color, rgb=rgb)]

        if harmony_type == "complementary":
            # Opposite on color wheel
            comp_h = (h + 180) % 360
            comp_rgb = _hsl_to_rgb(comp_h, s, l)
            colors.append(BrandColor(
                name="Complementary",
                hex=_rgb_to_hex(comp_rgb),
                rgb=comp_rgb
            ))

        elif harmony_type == "analogous":
            # Adjacent colors
            for i, offset in enumerate([-30, 30]):
                adj_h = (h + offset) % 360
                adj_rgb = _hsl_to_rgb(adj_h, s, l)
                colors.append(BrandColor(
                    name=f"Analogous {i+1}",
                    hex=_rgb_to_hex(adj_rgb),
                    rgb=adj_rgb
                ))

        elif harmony_type == "triadic":
            # Three colors equally spaced
            for i in range(1, 3):
                tri_h = (h + i * 120) % 360
                tri_rgb = _hsl_to_rgb(tri_h, s, l)
                colors.append(BrandColor(
                    name=f"Triadic {i}",
                    hex=_rgb_to_hex(tri_rgb),
                    rgb=tri_rgb
                ))

        elif harmony_type == "split-complementary":
            # Complementary with adjacent colors
            for offset in [150, 210]:
                split_h = (h + offset) % 360
                split_rgb = _hsl_to_rgb(split_h, s, l)
                colors.append(BrandColor(
                    name=f"Split {offset}",
                    hex=_rgb_to_hex(split_rgb),
                    rgb=split_rgb
                ))

        elif harmony_type == "tetradic":
            # Four colors in rectangle
            for offset in [90, 180, 270]:
                tet_h = (h + offset) % 360
                tet_rgb = _hsl_to_rgb(tet_h, s, l)
                colors.append(BrandColor(
                    name=f"Tetradic {offset}",
                    hex=_rgb_to_hex(tet_rgb),
                    rgb=tet_rgb
                ))

        # Add lighter/darker variants to reach count
        while len(colors) < count:
            variant_l = max(10, l - (len(colors) - 1) * 10) if len(colors) % 2 == 0 else min(90, l + (len(colors) - 1) * 10)
            var_rgb = _hsl_to_rgb(h, s, variant_l)
            colors.append(BrandColor(
                name=f"Shade {len(colors)}",
                hex=_rgb_to_hex(var_rgb),
                rgb=var_rgb
            ))

        return ColorPaletteResponse(
            base_color=base_color,
            harmony_type=harmony_type,
            colors=colors[:count],
        )

    # Theme methods

    async def create_theme(self, request: ThemeCreate) -> ThemeResponse:
        """Create a new theme."""
        theme_id = str(uuid.uuid4())
        now = utc_now()

        theme = {
            "id": theme_id,
            "name": request.name,
            "description": request.description,
            "brand_kit_id": request.brand_kit_id,
            "mode": request.mode,
            "colors": request.colors,
            "typography": request.typography,
            "spacing": request.spacing,
            "borders": request.borders,
            "shadows": request.shadows,
            "created_at": now,
            "updated_at": now,
            "is_active": len(self._themes) == 0,
        }

        self._themes[theme_id] = theme

        if theme["is_active"]:
            self._active_theme_id = theme_id

        # Persist
        try:
            with state_store.transaction() as state:
                state["themes"][theme_id] = theme
        except Exception as e:
            logger.warning(f"Failed to persist theme: {e}")

        return self._to_theme_response(theme)

    async def get_theme(self, theme_id: str) -> Optional[ThemeResponse]:
        """Get a theme by ID."""
        theme = self._themes.get(theme_id)
        if not theme:
            return None
        return self._to_theme_response(theme)

    async def list_themes(self) -> list[ThemeResponse]:
        """List all themes."""
        try:
            with state_store.transaction() as state:
                self._themes.update(state.get("themes", {}))
        except Exception:
            logger.debug("Failed to load themes from state store", exc_info=True)

        themes = list(self._themes.values())
        themes.sort(key=lambda t: t.get("created_at", ""), reverse=True)
        return [self._to_theme_response(t) for t in themes]

    async def update_theme(
        self,
        theme_id: str,
        request: ThemeUpdate,
    ) -> Optional[ThemeResponse]:
        """Update a theme."""
        theme = self._themes.get(theme_id)
        if not theme:
            return None

        for field in ["name", "description", "brand_kit_id", "mode", "colors",
                      "typography", "spacing", "borders", "shadows"]:
            value = getattr(request, field, None)
            if value is not None:
                theme[field] = value

        theme["updated_at"] = utc_now()

        try:
            with state_store.transaction() as state:
                state["themes"][theme_id] = theme
        except Exception as e:
            logger.warning(f"Failed to persist theme update: {e}")

        return self._to_theme_response(theme)

    async def delete_theme(self, theme_id: str) -> bool:
        """Delete a theme."""
        if theme_id not in self._themes:
            return False

        del self._themes[theme_id]

        if self._active_theme_id == theme_id:
            self._active_theme_id = None

        try:
            with state_store.transaction() as state:
                state["themes"].pop(theme_id, None)
        except Exception:
            logger.debug("Failed to delete theme from state store", exc_info=True)

        return True

    async def set_active_theme(self, theme_id: str) -> Optional[ThemeResponse]:
        """Set a theme as active."""
        theme = self._themes.get(theme_id)
        if not theme:
            return None

        if self._active_theme_id and self._active_theme_id != theme_id:
            prev = self._themes.get(self._active_theme_id)
            if prev:
                prev["is_active"] = False

        theme["is_active"] = True
        self._active_theme_id = theme_id

        return self._to_theme_response(theme)

    # ------------------------------------------------------------------
    # Color utility methods
    # ------------------------------------------------------------------

    def get_color_contrast(self, color1: str, color2: str) -> ColorContrastResponse:
        """Compute WCAG contrast ratio between two colors."""
        rgb1 = _hex_to_rgb(color1)
        rgb2 = _hex_to_rgb(color2)
        ratio = _contrast_ratio(rgb1, rgb2)
        return ColorContrastResponse(
            color1=color1,
            color2=color2,
            contrast_ratio=round(ratio, 2),
            wcag_aa_normal=ratio >= 4.5,
            wcag_aa_large=ratio >= 3.0,
            wcag_aaa_normal=ratio >= 7.0,
            wcag_aaa_large=ratio >= 4.5,
        )

    def suggest_accessible_colors(self, background_color: str) -> AccessibleColorsResponse:
        """Suggest text colors that meet WCAG AA against the given background."""
        bg_rgb = _hex_to_rgb(background_color)
        suggestions: list[AccessibleColorSuggestion] = []

        candidates = [
            ("#000000", "Black"),
            ("#ffffff", "White"),
            ("#333333", "Dark Gray"),
            ("#1a1a1a", "Near Black"),
            ("#f5f5f5", "Near White"),
            ("#0d47a1", "Dark Blue"),
            ("#1b5e20", "Dark Green"),
            ("#b71c1c", "Dark Red"),
            ("#4a148c", "Dark Purple"),
            ("#e65100", "Dark Orange"),
        ]
        for hex_color, label in candidates:
            c_rgb = _hex_to_rgb(hex_color)
            ratio = _contrast_ratio(bg_rgb, c_rgb)
            if ratio >= 4.5:
                suggestions.append(AccessibleColorSuggestion(
                    hex=hex_color, label=label, contrast_ratio=round(ratio, 2),
                ))

        suggestions.sort(key=lambda s: s.contrast_ratio, reverse=True)
        return AccessibleColorsResponse(
            background_color=background_color,
            colors=suggestions,
        )

    # ------------------------------------------------------------------
    # Typography methods
    # ------------------------------------------------------------------

    def list_fonts(self) -> list[FontInfo]:
        """Return a curated list of available fonts."""
        return [FontInfo(**f) for f in _FONTS]

    def get_font_pairings(self, primary_font: str) -> FontPairingsResponse:
        """Suggest body-text font pairings for a primary heading font."""
        # Find the category of the requested font
        category = "sans-serif"
        for f in _FONTS:
            if f["name"].lower() == primary_font.lower():
                category = f["category"]
                break

        rules = _PAIRING_RULES.get(category, _PAIRING_RULES["sans-serif"])
        pairings = [FontPairing(**r) for r in rules]
        return FontPairingsResponse(primary=primary_font, pairings=pairings)

    # ------------------------------------------------------------------
    # Asset methods
    # ------------------------------------------------------------------

    async def upload_logo(self, filename: str, content: bytes, brand_kit_id: str) -> AssetResponse:
        """Store a logo asset reference."""
        asset_id = str(uuid.uuid4())
        now = utc_now()
        asset = {
            "id": asset_id,
            "filename": filename,
            "brand_kit_id": brand_kit_id,
            "asset_type": "logo",
            "size_bytes": len(content),
            "created_at": now,
        }
        self._assets[asset_id] = asset
        return AssetResponse(**asset)

    async def list_assets(self, brand_kit_id: str) -> list[AssetResponse]:
        """List assets for a brand kit."""
        return [
            AssetResponse(**a)
            for a in self._assets.values()
            if a["brand_kit_id"] == brand_kit_id
        ]

    async def delete_asset(self, asset_id: str) -> bool:
        """Delete an asset."""
        if asset_id not in self._assets:
            return False
        del self._assets[asset_id]
        return True

    # ------------------------------------------------------------------
    # Import / Export methods
    # ------------------------------------------------------------------

    async def export_brand_kit(self, kit_id: str, fmt: str = "json") -> Optional[BrandKitExport]:
        """Export a brand kit."""
        kit = self._brand_kits.get(kit_id)
        if not kit:
            return None
        return BrandKitExport(
            format=fmt,
            brand_kit=self._to_brand_kit_response(kit),
        )

    async def import_brand_kit(self, data: dict) -> BrandKitResponse:
        """Import a brand kit from exported data."""
        from backend.app.schemas import BrandKitCreate as _Create
        create_req = _Create(**{
            k: v for k, v in data.items()
            if k in _Create.model_fields
        })
        return await self.create_brand_kit(create_req)

    # ------------------------------------------------------------------
    # Internal helpers
    # ------------------------------------------------------------------

    def _to_brand_kit_response(self, kit: dict) -> BrandKitResponse:
        """Convert brand kit dict to response model."""
        return BrandKitResponse(
            id=kit["id"],
            name=kit["name"],
            description=kit.get("description"),
            logo_url=kit.get("logo_url"),
            logo_dark_url=kit.get("logo_dark_url"),
            favicon_url=kit.get("favicon_url"),
            primary_color=kit["primary_color"],
            secondary_color=kit["secondary_color"],
            accent_color=kit["accent_color"],
            text_color=kit["text_color"],
            background_color=kit["background_color"],
            colors=[BrandColor(**c) for c in kit.get("colors", [])],
            typography=Typography(**kit.get("typography", {})),
            created_at=kit["created_at"],
            updated_at=kit["updated_at"],
            is_default=kit.get("is_default", False),
        )

    def _to_theme_response(self, theme: dict) -> ThemeResponse:
        """Convert theme dict to response model."""
        return ThemeResponse(
            id=theme["id"],
            name=theme["name"],
            description=theme.get("description"),
            brand_kit_id=theme.get("brand_kit_id"),
            mode=theme.get("mode", "light"),
            colors=theme.get("colors", {}),
            typography=theme.get("typography", {}),
            spacing=theme.get("spacing", {}),
            borders=theme.get("borders", {}),
            shadows=theme.get("shadows", {}),
            created_at=theme["created_at"],
            updated_at=theme["updated_at"],
            is_active=theme.get("is_active", False),
        )


# Singleton instance
design_service = DesignService()


# MEMORY_SERVICE (merged from memory_service.py)


# mypy: ignore-errors
"""
Entity Tracker — Tracks mentioned entities across conversation turns.

Monitors mentions of database entities (tables, columns), templates,
reports, and other domain objects across the conversation to:
- Provide context-aware suggestions
- Resolve ambiguous references ("that table", "the same column")
- Track entity relationships and co-occurrence

Usage:
    tracker = EntityTracker()
    tracker.record_mention(user_id, "table", "safety_readings", turn=3)
    tracker.record_mention(user_id, "column", "pressure_value", turn=3, parent="safety_readings")
    recent = tracker.get_recent_entities(user_id, entity_type="table")
"""


from collections import defaultdict
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional, Set

logger = logging.getLogger("neura.memory.entity_tracker")


@dataclass
class EntityMention:
    """A mention of an entity in conversation."""
    entity_type: str  # "table", "column", "template", "report", "connection"
    entity_name: str
    turn_number: int
    timestamp: float = field(default_factory=time.time)
    parent_entity: Optional[str] = None  # e.g., table name for a column
    context: str = ""  # Surrounding text for disambiguation
    confidence: float = 1.0

    @property
    def key(self) -> str:
        """Unique key for this entity."""
        if self.parent_entity:
            return f"{self.entity_type}:{self.parent_entity}.{self.entity_name}"
        return f"{self.entity_type}:{self.entity_name}"


@dataclass
class EntityProfile:
    """Aggregated profile for a tracked entity."""
    entity_type: str
    entity_name: str
    mention_count: int = 0
    first_mentioned: float = 0.0
    last_mentioned: float = 0.0
    turns_mentioned: Set[int] = field(default_factory=set)
    parent_entity: Optional[str] = None
    related_entities: Set[str] = field(default_factory=set)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "entity_type": self.entity_type,
            "entity_name": self.entity_name,
            "mention_count": self.mention_count,
            "turns_mentioned": sorted(self.turns_mentioned),
            "parent_entity": self.parent_entity,
            "related_entities": sorted(self.related_entities),
        }


class EntityTracker:
    """
    Tracks entity mentions across conversation turns.

    Maintains per-user entity profiles that accumulate mention statistics.
    Supports querying for recent entities, co-occurring entities, and
    resolving ambiguous references.
    """

    def __init__(self, max_mentions_per_user: int = 500):
        self._mentions: Dict[str, List[EntityMention]] = defaultdict(list)
        self._profiles: Dict[str, Dict[str, EntityProfile]] = defaultdict(dict)
        self._max_mentions = max_mentions_per_user
        self._lock = threading.Lock()

    def record_mention(
        self,
        user_id: str,
        entity_type: str,
        entity_name: str,
        turn: int = 0,
        parent: Optional[str] = None,
        context: str = "",
        confidence: float = 1.0,
    ) -> None:
        """Record an entity mention in the conversation."""
        mention = EntityMention(
            entity_type=entity_type,
            entity_name=entity_name,
            turn_number=turn,
            parent_entity=parent,
            context=context,
            confidence=confidence,
        )

        with self._lock:
            mentions = self._mentions[user_id]
            mentions.append(mention)

            # Prune old mentions
            if len(mentions) > self._max_mentions:
                self._mentions[user_id] = mentions[-self._max_mentions:]

            # Update profile
            key = mention.key
            profiles = self._profiles[user_id]
            if key not in profiles:
                profiles[key] = EntityProfile(
                    entity_type=entity_type,
                    entity_name=entity_name,
                    first_mentioned=mention.timestamp,
                    parent_entity=parent,
                )

            profile = profiles[key]
            profile.mention_count += 1
            profile.last_mentioned = mention.timestamp
            profile.turns_mentioned.add(turn)

            # Track co-occurrence: entities mentioned in the same turn
            same_turn_entities = [
                m for m in mentions
                if m.turn_number == turn and m.key != key
            ]
            for co_entity in same_turn_entities:
                profile.related_entities.add(co_entity.key)

    def record_mentions_batch(
        self,
        user_id: str,
        mentions: List[Dict[str, Any]],
        turn: int = 0,
    ) -> None:
        """Record multiple entity mentions at once."""
        for m in mentions:
            self.record_mention(
                user_id=user_id,
                entity_type=m.get("type", "unknown"),
                entity_name=m.get("name", ""),
                turn=turn,
                parent=m.get("parent"),
                context=m.get("context", ""),
                confidence=m.get("confidence", 1.0),
            )

    def get_recent_entities(
        self,
        user_id: str,
        entity_type: Optional[str] = None,
        limit: int = 10,
    ) -> List[Dict[str, Any]]:
        """Get recently mentioned entities, optionally filtered by type."""
        with self._lock:
            profiles = self._profiles.get(user_id, {})
            filtered = list(profiles.values())

            if entity_type:
                filtered = [p for p in filtered if p.entity_type == entity_type]

            # Sort by last mention (most recent first)
            filtered.sort(key=lambda p: p.last_mentioned, reverse=True)
            return [p.to_dict() for p in filtered[:limit]]

    def get_most_mentioned(
        self,
        user_id: str,
        entity_type: Optional[str] = None,
        limit: int = 10,
    ) -> List[Dict[str, Any]]:
        """Get most frequently mentioned entities."""
        with self._lock:
            profiles = self._profiles.get(user_id, {})
            filtered = list(profiles.values())

            if entity_type:
                filtered = [p for p in filtered if p.entity_type == entity_type]

            filtered.sort(key=lambda p: p.mention_count, reverse=True)
            return [p.to_dict() for p in filtered[:limit]]

    def resolve_reference(
        self,
        user_id: str,
        entity_type: str,
        partial_name: str = "",
    ) -> Optional[Dict[str, Any]]:
        """
        Resolve an ambiguous entity reference.

        If partial_name is given, fuzzy-matches against known entities.
        If empty, returns the most recently mentioned entity of the given type.
        """
        with self._lock:
            profiles = self._profiles.get(user_id, {})
            candidates = [
                p for p in profiles.values()
                if p.entity_type == entity_type
            ]

            if not candidates:
                return None

            if partial_name:
                # Fuzzy match: contains or starts-with
                partial_lower = partial_name.lower()
                matches = [
                    p for p in candidates
                    if partial_lower in p.entity_name.lower()
                ]
                if matches:
                    # Return the most recently mentioned match
                    matches.sort(key=lambda p: p.last_mentioned, reverse=True)
                    return matches[0].to_dict()

            # Default: return most recently mentioned
            candidates.sort(key=lambda p: p.last_mentioned, reverse=True)
            return candidates[0].to_dict()

    def get_co_occurring_entities(
        self, user_id: str, entity_key: str
    ) -> List[Dict[str, Any]]:
        """Get entities that co-occur with the given entity."""
        with self._lock:
            profiles = self._profiles.get(user_id, {})
            profile = profiles.get(entity_key)
            if not profile:
                return []

            related = []
            for rel_key in profile.related_entities:
                rel_profile = profiles.get(rel_key)
                if rel_profile:
                    related.append(rel_profile.to_dict())

            return related

    def build_context_prompt(self, user_id: str, max_entities: int = 10) -> str:
        """Build entity context string for LLM prompt injection."""
        recent = self.get_recent_entities(user_id, limit=max_entities)
        if not recent:
            return ""

        parts = ["## Recently Discussed Entities"]
        for entity in recent:
            name = entity["entity_name"]
            etype = entity["entity_type"]
            parent = entity.get("parent_entity", "")
            count = entity["mention_count"]
            if parent:
                parts.append(f"- {etype}: {parent}.{name} (mentioned {count}x)")
            else:
                parts.append(f"- {etype}: {name} (mentioned {count}x)")

        return "\n".join(parts)

    def clear(self, user_id: str) -> None:
        """Clear all tracked entities for a user."""
        with self._lock:
            self._mentions.pop(user_id, None)
            self._profiles.pop(user_id, None)

    def get_stats(self, user_id: str) -> Dict[str, Any]:
        """Get tracking statistics for a user."""
        with self._lock:
            mentions = self._mentions.get(user_id, [])
            profiles = self._profiles.get(user_id, {})
            type_counts = defaultdict(int)
            for p in profiles.values():
                type_counts[p.entity_type] += 1

            return {
                "total_mentions": len(mentions),
                "unique_entities": len(profiles),
                "by_type": dict(type_counts),
            }


# mypy: ignore-errors
"""
User Preference Learning — Learns and stores user preferences.

Tracks user behavior patterns to provide better defaults:
- Preferred chart types for different data types
- Preferred export formats (PDF, Excel, HTML)
- Mapping correction patterns (learn from corrections)
- Report style preferences (detail level, sections)
- Frequently used connections and templates

Usage:
    store = PreferenceStore()
    store.record_preference(user_id, "chart_type", "bar", context={"data_type": "categorical"})
    pref = store.get_preference(user_id, "chart_type", context={"data_type": "categorical"})
"""

from collections import defaultdict, Counter

logger = logging.getLogger("neura.memory.preferences")


@dataclass
class PreferenceRecord:
    """A single preference observation."""
    category: str  # "chart_type", "export_format", "detail_level", etc.
    value: str  # The chosen value
    context: Dict[str, Any] = field(default_factory=dict)
    timestamp: float = field(default_factory=time.time)
    explicit: bool = False  # True if user explicitly stated preference


@dataclass
class UserPreferences:
    """Aggregated preferences for a single user."""
    user_id: str
    preferences: Dict[str, Counter] = field(default_factory=lambda: defaultdict(Counter))
    context_preferences: Dict[str, Dict[str, Counter]] = field(
        default_factory=lambda: defaultdict(lambda: defaultdict(Counter))
    )
    explicit_overrides: Dict[str, str] = field(default_factory=dict)
    last_updated: float = field(default_factory=time.time)

    def get_preference(self, category: str, context: Optional[Dict[str, Any]] = None) -> Optional[str]:
        """Get the preferred value for a category, considering context."""
        # Explicit overrides take priority
        if category in self.explicit_overrides:
            return self.explicit_overrides[category]

        # Context-specific preference
        if context:
            context_key = self._context_key(context)
            ctx_prefs = self.context_preferences.get(category, {}).get(context_key)
            if ctx_prefs:
                most_common = ctx_prefs.most_common(1)
                if most_common:
                    return most_common[0][0]

        # General preference
        general = self.preferences.get(category)
        if general:
            most_common = general.most_common(1)
            if most_common:
                return most_common[0][0]

        return None

    def get_top_preferences(self, category: str, limit: int = 3) -> List[Dict[str, Any]]:
        """Get top N preferences for a category with counts."""
        general = self.preferences.get(category, Counter())
        return [
            {"value": value, "count": count}
            for value, count in general.most_common(limit)
        ]

    @staticmethod
    def _context_key(context: Dict[str, Any]) -> str:
        """Create a hashable key from context dict."""
        return json.dumps(sorted(context.items()), default=str)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "user_id": self.user_id,
            "preferences": {
                cat: dict(counter.most_common(5))
                for cat, counter in self.preferences.items()
            },
            "explicit_overrides": self.explicit_overrides,
            "categories": list(self.preferences.keys()),
        }


class PreferenceStore:
    """
    Stores and learns user preferences across sessions.

    Preferences are learned from:
    1. Explicit user statements ("always use bar charts")
    2. Implicit behavior (user consistently chooses PDF export)
    3. Correction patterns (user keeps changing column X to Y)
    """

    def __init__(self, persist_path: Optional[str] = None, max_users: int = 1000):
        self._users: Dict[str, UserPreferences] = {}
        self._persist_path = Path(persist_path) if persist_path else None
        self._max_users = max_users
        self._lock = threading.Lock()

        if self._persist_path and self._persist_path.exists():
            self._load()

    def record_preference(
        self,
        user_id: str,
        category: str,
        value: str,
        context: Optional[Dict[str, Any]] = None,
        explicit: bool = False,
    ) -> None:
        """Record a preference observation."""
        with self._lock:
            prefs = self._get_or_create(user_id)

            if explicit:
                prefs.explicit_overrides[category] = value

            prefs.preferences[category][value] += 1

            if context:
                context_key = prefs._context_key(context)
                prefs.context_preferences[category][context_key][value] += 1

            prefs.last_updated = time.time()

        if self._persist_path:
            self._save()

        logger.debug(
            "Recorded preference: user=%s, category=%s, value=%s, explicit=%s",
            user_id, category, value, explicit,
        )

    def get_preference(
        self,
        user_id: str,
        category: str,
        context: Optional[Dict[str, Any]] = None,
        default: Optional[str] = None,
    ) -> Optional[str]:
        """Get the preferred value for a category."""
        with self._lock:
            prefs = self._users.get(user_id)
            if prefs is None:
                return default
            result = prefs.get_preference(category, context)
            return result if result is not None else default

    def get_user_preferences(self, user_id: str) -> Optional[Dict[str, Any]]:
        """Get all preferences for a user."""
        with self._lock:
            prefs = self._users.get(user_id)
            return prefs.to_dict() if prefs else None

    def set_explicit_preference(
        self,
        user_id: str,
        category: str,
        value: str,
    ) -> None:
        """Set an explicit preference override (user stated preference)."""
        self.record_preference(user_id, category, value, explicit=True)

    def remove_explicit_preference(self, user_id: str, category: str) -> None:
        """Remove an explicit preference override."""
        with self._lock:
            prefs = self._users.get(user_id)
            if prefs:
                prefs.explicit_overrides.pop(category, None)
        if self._persist_path:
            self._save()

    def get_default_preferences(self, user_id: str) -> Dict[str, str]:
        """Get a dict of category → preferred value for all known categories."""
        with self._lock:
            prefs = self._users.get(user_id)
            if not prefs:
                return {}

            result = {}
            for category in prefs.preferences:
                value = prefs.get_preference(category)
                if value:
                    result[category] = value
            return result

    def build_context_prompt(self, user_id: str) -> str:
        """Build a preference context string for LLM prompt injection."""
        defaults = self.get_default_preferences(user_id)
        if not defaults:
            return ""

        parts = ["## User Preferences"]
        for category, value in sorted(defaults.items()):
            parts.append(f"- {category}: {value}")
        return "\n".join(parts)

    def clear(self, user_id: str) -> None:
        """Clear all preferences for a user."""
        with self._lock:
            self._users.pop(user_id, None)
        if self._persist_path:
            self._save()

    def _get_or_create(self, user_id: str) -> UserPreferences:
        if user_id not in self._users:
            self._users[user_id] = UserPreferences(user_id=user_id)
        return self._users[user_id]

    def _save(self) -> None:
        """Persist preferences to disk."""
        try:
            data = {}
            with self._lock:
                for uid, prefs in self._users.items():
                    data[uid] = {
                        "preferences": {
                            cat: dict(counter)
                            for cat, counter in prefs.preferences.items()
                        },
                        "explicit_overrides": prefs.explicit_overrides,
                        "last_updated": prefs.last_updated,
                    }

            self._persist_path.parent.mkdir(parents=True, exist_ok=True)
            with open(self._persist_path, "w") as f:
                json.dump(data, f, default=str, indent=2)
        except Exception as exc:
            logger.warning("Failed to persist preferences: %s", exc)

    def _load(self) -> None:
        """Load preferences from disk."""
        try:
            with open(self._persist_path) as f:
                data = json.load(f)

            for uid, pdata in data.items():
                prefs = UserPreferences(user_id=uid)
                prefs.last_updated = pdata.get("last_updated", time.time())
                prefs.explicit_overrides = pdata.get("explicit_overrides", {})

                for cat, counts in pdata.get("preferences", {}).items():
                    prefs.preferences[cat] = Counter(counts)

                self._users[uid] = prefs

            logger.info("Loaded preferences for %d users from %s", len(self._users), self._persist_path)
        except Exception as exc:
            logger.warning("Failed to load preferences: %s", exc)


# Global instance
_preference_store: Optional[PreferenceStore] = None


def get_preference_store() -> PreferenceStore:
    """Get the global preference store."""
    global _preference_store
    if _preference_store is None:
        _preference_store = PreferenceStore()
    return _preference_store


# mypy: ignore-errors
"""
Conversation Memory — Session context persistence.

Tracks conversation state across user sessions so the system can:
- Remember the last template, connection, and report used
- Maintain conversation continuity ("use the same template as last time")
- Track which entities were discussed in previous turns

Inspired by BFI's conversation_memory.py pattern.

Usage:
    memory = ConversationMemory()
    memory.set_context(user_id, "last_template", {"id": "tpl_123", "name": "Monthly Report"})
    ctx = memory.get_context(user_id, "last_template")
"""

from typing import Any, Dict, Optional

logger = logging.getLogger("neura.memory.conversation")


@dataclass
class SessionContext:
    """Context stored for a single user session."""
    user_id: str
    created_at: float = field(default_factory=time.time)
    updated_at: float = field(default_factory=time.time)
    data: Dict[str, Any] = field(default_factory=dict)
    turn_count: int = 0

    # Quick-access fields for common context
    last_template_id: Optional[str] = None
    last_template_name: Optional[str] = None
    last_connection_id: Optional[str] = None
    last_connection_name: Optional[str] = None
    last_report_id: Optional[str] = None
    last_agent_task_id: Optional[str] = None

    def to_dict(self) -> Dict[str, Any]:
        return {
            "user_id": self.user_id,
            "created_at": datetime.fromtimestamp(self.created_at, tz=timezone.utc).isoformat(),
            "updated_at": datetime.fromtimestamp(self.updated_at, tz=timezone.utc).isoformat(),
            "turn_count": self.turn_count,
            "last_template_id": self.last_template_id,
            "last_template_name": self.last_template_name,
            "last_connection_id": self.last_connection_id,
            "last_connection_name": self.last_connection_name,
            "last_report_id": self.last_report_id,
            "last_agent_task_id": self.last_agent_task_id,
            "custom_data_keys": list(self.data.keys()),
        }


class ConversationMemory:
    """
    Per-user conversation memory with session context persistence.

    Stores key-value context per user with automatic timestamp tracking.
    Supports persistence to a JSON file for cross-session continuity.

    Common context keys:
    - last_template: {id, name, field_count}
    - last_connection: {id, name, type}
    - last_report: {id, template_name, status}
    - last_agent_task: {id, type, status}
    - conversation_summary: brief summary of recent conversation
    """

    def __init__(
        self,
        persist_path: Optional[str] = None,
        max_sessions: int = 1000,
        session_ttl_hours: int = 168,  # 7 days
    ):
        self._sessions: Dict[str, SessionContext] = {}
        self._persist_path = Path(persist_path) if persist_path else None
        self._max_sessions = max_sessions
        self._session_ttl = session_ttl_hours * 3600
        self._lock = threading.Lock()

        if self._persist_path and self._persist_path.exists():
            self._load()

    def get_session(self, user_id: str) -> SessionContext:
        """Get or create a session context for a user."""
        with self._lock:
            if user_id not in self._sessions:
                self._sessions[user_id] = SessionContext(user_id=user_id)
            return self._sessions[user_id]

    def set_context(self, user_id: str, key: str, value: Any) -> None:
        """Set a context value for a user."""
        with self._lock:
            session = self._sessions.setdefault(user_id, SessionContext(user_id=user_id))
            session.data[key] = value
            session.updated_at = time.time()
            session.turn_count += 1

            # Update quick-access fields
            if key == "last_template" and isinstance(value, dict):
                session.last_template_id = value.get("id")
                session.last_template_name = value.get("name")
            elif key == "last_connection" and isinstance(value, dict):
                session.last_connection_id = value.get("id")
                session.last_connection_name = value.get("name")
            elif key == "last_report" and isinstance(value, dict):
                session.last_report_id = value.get("id")
            elif key == "last_agent_task" and isinstance(value, dict):
                session.last_agent_task_id = value.get("id")

        if self._persist_path:
            self._save()

    def get_context(self, user_id: str, key: str, default: Any = None) -> Any:
        """Get a context value for a user."""
        with self._lock:
            session = self._sessions.get(user_id)
            if session is None:
                return default
            return session.data.get(key, default)

    def get_all_context(self, user_id: str) -> Dict[str, Any]:
        """Get all context values for a user."""
        with self._lock:
            session = self._sessions.get(user_id)
            if session is None:
                return {}
            return dict(session.data)

    def clear_session(self, user_id: str) -> None:
        """Clear a user's session context."""
        with self._lock:
            self._sessions.pop(user_id, None)
        if self._persist_path:
            self._save()

    def get_session_summary(self, user_id: str) -> Dict[str, Any]:
        """Get a summary of a user's session for injection into prompts."""
        with self._lock:
            session = self._sessions.get(user_id)
            if session is None:
                return {"has_context": False}

            return {
                "has_context": True,
                "turn_count": session.turn_count,
                "last_template": session.last_template_name,
                "last_connection": session.last_connection_name,
                "last_report_id": session.last_report_id,
                "context_keys": list(session.data.keys()),
            }

    def build_context_prompt(self, user_id: str) -> str:
        """Build a context string suitable for injecting into LLM prompts."""
        summary = self.get_session_summary(user_id)
        if not summary["has_context"]:
            return ""

        parts = ["## Session Context"]
        if summary.get("last_template"):
            parts.append(f"- Last template used: {summary['last_template']}")
        if summary.get("last_connection"):
            parts.append(f"- Last database connection: {summary['last_connection']}")
        if summary.get("last_report_id"):
            parts.append(f"- Last report generated: {summary['last_report_id']}")
        if summary.get("turn_count", 0) > 0:
            parts.append(f"- Conversation turns: {summary['turn_count']}")

        return "\n".join(parts)

    def prune_expired(self) -> int:
        """Remove expired sessions. Returns count of removed sessions."""
        now = time.time()
        removed = 0
        with self._lock:
            expired = [
                uid for uid, session in self._sessions.items()
                if (now - session.updated_at) > self._session_ttl
            ]
            for uid in expired:
                del self._sessions[uid]
                removed += 1

            # Also enforce max sessions (LRU)
            if len(self._sessions) > self._max_sessions:
                sorted_sessions = sorted(
                    self._sessions.items(),
                    key=lambda x: x[1].updated_at,
                )
                excess = len(self._sessions) - self._max_sessions
                for uid, _ in sorted_sessions[:excess]:
                    del self._sessions[uid]
                    removed += 1

        if removed > 0:
            logger.info("Pruned %d expired sessions", removed)
            if self._persist_path:
                self._save()

        return removed

    def _save(self) -> None:
        """Persist all sessions to disk."""
        try:
            data = {}
            with self._lock:
                for uid, session in self._sessions.items():
                    data[uid] = {
                        "created_at": session.created_at,
                        "updated_at": session.updated_at,
                        "turn_count": session.turn_count,
                        "data": session.data,
                        "last_template_id": session.last_template_id,
                        "last_template_name": session.last_template_name,
                        "last_connection_id": session.last_connection_id,
                        "last_connection_name": session.last_connection_name,
                        "last_report_id": session.last_report_id,
                        "last_agent_task_id": session.last_agent_task_id,
                    }

            self._persist_path.parent.mkdir(parents=True, exist_ok=True)
            with open(self._persist_path, "w") as f:
                json.dump(data, f, default=str, indent=2)
        except Exception as exc:
            logger.warning("Failed to persist conversation memory: %s", exc)

    def _load(self) -> None:
        """Load sessions from disk."""
        try:
            with open(self._persist_path) as f:
                data = json.load(f)

            for uid, sdata in data.items():
                session = SessionContext(
                    user_id=uid,
                    created_at=sdata.get("created_at", time.time()),
                    updated_at=sdata.get("updated_at", time.time()),
                    data=sdata.get("data", {}),
                    turn_count=sdata.get("turn_count", 0),
                    last_template_id=sdata.get("last_template_id"),
                    last_template_name=sdata.get("last_template_name"),
                    last_connection_id=sdata.get("last_connection_id"),
                    last_connection_name=sdata.get("last_connection_name"),
                    last_report_id=sdata.get("last_report_id"),
                    last_agent_task_id=sdata.get("last_agent_task_id"),
                )
                self._sessions[uid] = session

            logger.info("Loaded %d sessions from %s", len(self._sessions), self._persist_path)
        except Exception as exc:
            logger.warning("Failed to load conversation memory: %s", exc)


# Global instance
_conversation_memory: Optional[ConversationMemory] = None


def get_conversation_memory() -> ConversationMemory:
    """Get the global conversation memory."""
    global _conversation_memory
    if _conversation_memory is None:
        _conversation_memory = ConversationMemory()
    return _conversation_memory


# PROMPTS (merged from prompts.py)
"""LLM prompt builders and registry."""


# ALL_PROMPTS

# mypy: ignore-errors

import base64
from functools import lru_cache
from textwrap import dedent
from typing import Any, Dict, Iterable, Mapping

logger = logging.getLogger("neura.prompts")

PROMPT_VERSION = "llm_call_3_df_v2"
PROMPT_VERSION_3_5 = "v5"
PROMPT_VERSION_4 = "v3_df"

# Legacy aliases kept for import compatibility
PROMPT_VERSION_DF = PROMPT_VERSION
PROMPT_VERSION_4_DF = PROMPT_VERSION_4
LLM_CALL_PROMPTS: Dict[str, str] = {
    "llm_call_1": dedent(
        """\
        Produce a COMPLETE, self-contained HTML document (<!DOCTYPE html> ...) with inline <style>. It must visually photocopy the given PDF page image as closely as possible. Mirror fonts, spacing, borders, alignment, and table layouts. Tables must use border-collapse, 1px borders, and table-layout: fixed for neat alignment.

        SCHEMA USAGE
        - If a SCHEMA is provided below, use ONLY placeholders from that SCHEMA exactly as written (same names).
        - If SCHEMA is NOT provided, FIRST infer a compact schema (see SCHEMA_JSON rules below) and then use ONLY those tokens in the HTML.
        - In HTML, placeholders must be written as {token_name} (single braces). In SCHEMA_JSON they appear WITHOUT braces.
        - If a value is not in SCHEMA/SCHEMA_JSON, render it as literal text. If a token exists in SCHEMA/SCHEMA_JSON but does not appear on this page, omit it.

        TOKEN NAMING CONVENTIONS (when inferring — NOT when a SCHEMA is provided)
        - Use lowercase_snake_case for all token names.
        - Scalar tokens (header/footer fields): use descriptive names based on ACTUAL visible headers (e.g. `report_title`, `plant_name`, `from_date`, `batch_no`, `batch_recipe`, `print_date`).
        - Row tokens (repeating data columns): prefix with `row_` followed by the ACTUAL column header normalized to snake_case. For example, if the table header says "Set wt(Kg)" use `row_set_wt_kg`; if it says "M1A_2 PT100" use `row_m1a_2_pt100`; if it says "Duration (sec)" use `row_duration_sec`.
        - Total tokens (summary/aggregate values): prefix with `total_` e.g. `total_set_wt_kg`, `total_duration`.
        - CRITICAL: You MUST derive token names from the ACTUAL column headers visible in the PDF (or provided in PDF_PAGE_TEXT). NEVER use generic placeholder names. Every column in the table must have its own unique row token.
        - If a PDF_PAGE_TEXT section is provided, use it as the authoritative source for all field names, column headers, and data structure.

        REPEATABLE BLOCK (edge case)
        - If the page clearly contains repeating sections (visually identical blocks stacked vertically), output ONE prototype of that block wrapped exactly as:
        <!-- BEGIN:BLOCK_REPEAT batches -->
        <section class='batch-block'>...</section>
        <!-- END:BLOCK_REPEAT -->
        - Place header/footer OUTSIDE these markers. Do NOT clone or duplicate multiple blocks.

        ROW PROTOTYPES
        - For tables with repeating rows, output headers in <thead> and a single prototype row inside <tbody><tr>...</tr></tbody>.
        - The prototype row contains one {row_*} token per cell, matching the header above it.
        - Keep any final summary/total row outside <tbody> (e.g. in a <tfoot> or a separate element after the table).

        STRUCTURE & CSS
        - The result must be printable: use @page size A4 with sensible margins.
        - Prefer flowing layout (avoid fixed heights). Do NOT use position:fixed or position:absolute on headers or footers — these overlap table content on long reports. Keep footers in normal document flow so they render after the content.
        - Reproduce what is visible — draw ONLY the rules/lines that exist in the image. Default to no borders and transparent backgrounds; add borders per edge only where a line is visible.
        - Use table markup ONLY for true grids and structured data (never div-based). Use borderless tables or simple divs for key/value areas. Avoid unnecessary nested tables or enclosing frames.
        - Right-align numeric columns where appropriate; keep typographic rhythm tight enough to match the PDF.

        PROJECT-SPECIFIC ADDITIONS
        - Add a minimal set of CSS custom properties to make column widths and key spacings easy to refine later (e.g., :root { --col-1-w: 24mm; --row-gap: 2mm; }). Use these variables inside the CSS you produce.
        - Add stable IDs for major zones only (no extra wrappers): #report-header, #data-table (main grid), #report-totals (if present). Do NOT add decorative containers.
        - For every table header cell, include a data-label attribute with the visible header text normalized to lowercase_snake_case (e.g., <th data-label="material_name">Material Name</th>). The visible text inside the <th> must remain unchanged — only the attribute value is normalized.

        OUTPUT RULES
        - No lorem ipsum or sample values. No external resources.
        - No comments except the repeat markers if applicable.
        - Do NOT rename or invent tokens beyond SCHEMA/SCHEMA_JSON.
        - Return ONLY the outputs described below — no markdown fences, no explanations, no prose.

        OUTPUT FORMAT
        1) First, the RAW HTML between these exact markers:
        <!--BEGIN_HTML-->
        ...full html...
        <!--END_HTML-->

        2) Then, the SCHEMA JSON between these markers:
        <!--BEGIN_SCHEMA_JSON-->
        {
          "scalars": ["report_title", "plant_name", "from_date"],
          "row_tokens": ["row_material_name", "row_quantity", "row_amount"],
          "totals": ["total_quantity", "total_amount"],
          "notes": ""
        }
        <!--END_SCHEMA_JSON-->
        IMPORTANT: The token names above are only illustrative examples of the NAMING CONVENTION. You must replace them with the actual tokens you discovered from THIS page. Never copy these example names verbatim.

        If SCHEMA is provided below, ensure SCHEMA_JSON you output matches it exactly (names and membership). If SCHEMA is not provided, infer SCHEMA_JSON consistently with the placeholders you used in the HTML (one-to-one, no extras, no omissions).

        [INPUTS]
        - PDF page image is attached.
        - SCHEMA (may be absent):
        SCHEMA:
        {schema_str}
        """
    ).strip(),
    "llm_call_2": dedent(
        """\
        Compare these two images: REFERENCE (the original PDF page) vs RENDER (your current HTML output).
        Goal: refine the HTML/CSS so the rendered output becomes a near-perfect PHOTOCOPY of the reference.

        STRICT RULES — violations will cause rejection
        1. Do NOT rename, add, remove, or reorder any {{token}} placeholders. Keep every token exactly as it appears in the current HTML.
        2. Do NOT change the number of repeating sections, table rows, <tbody> blocks, or repeat markers.
        3. If repeat markers (<!-- BEGIN:BLOCK_REPEAT ... -->) are present, keep them unchanged with exactly one prototype inside.
        4. Prefer CSS-only edits. Only add minimal HTML structural wrappers (e.g., <colgroup>) when CSS alone cannot achieve the alignment.

        CSS REFINEMENT STRATEGY (in priority order)
        1. First, tune existing CSS custom properties (--col-1-w, --col-2-w, --row-gap, etc.).
        2. If custom properties are insufficient, edit CSS rules directly.
        3. Use millimetre-based sizing for print fidelity (widths, padding, margins in mm).
        4. Right-align numeric columns; use font-variant-numeric: tabular-nums for digit alignment.
        5. Match borders/lines exactly as visible in the reference — per-edge borders only, no shadows or rounded corners.
        6. Never use position:fixed/absolute on headers/footers (causes overlap on long reports).
        7. Never scale the page via CSS transforms — correct geometry through widths, margins, padding, line-height instead.

        VISUAL MATCHING
        Identify and correct EVERY visible discrepancy: geometry, proportions, typography, line metrics, borders, column widths, text alignment, spacing, and header/footer placement. Derive all values from the reference image. The result must be indistinguishable from the reference when printed.

        OUTPUT — return ONLY this, nothing else:
        <!--BEGIN_HTML-->
        ...full refined html (<!DOCTYPE html> ... with inline <style>)...
        <!--END_HTML-->

        No markdown fences, no commentary, no prose before or after the markers.

        [INPUTS]
        SCHEMA:
        {schema_str}

        [REFERENCE_IMAGE]
        (embedded image URL)

        [RENDER_IMAGE]
        (embedded image URL)

        [CURRENT_HTML]
        {current_html}
        """
    ).strip(),
}

_INPUT_MARKER = "[INPUTS]"
_CALL3_PROMPT_SECTION_DF = """
You are a meticulous report auto-mapping analyst. You are given the FULL HTML of a report template and a RICH DB CATALOG that includes column types and sample values.

YOUR TWO TASKS:
A) AUTO-MAPPING — Map every placeholder token in the HTML to its data source from the CATALOG.
B) CONSTANT DISCOVERY — Identify tokens whose values are static and record their literal values.

--------------------------------------------------------------------------------
MAPPING RULES — CRITICAL CONSTRAINTS FOR DATAFRAME MODE
1. Output ONLY simple "table.column" references from the CATALOG. NEVER output SQL expressions, functions (SUM, STRFTIME, CONCAT, CASE, COALESCE, etc.), or any code.
2. Every mapping value must be ONE of:
   a) A catalog column in exact "table.column" format (e.g., "orders.customer_name").
   b) A parameter passthrough in "params.param_name" format (e.g., "params.plant_name").
   c) The literal string "To Be Selected..." for date-range and page/filter tokens (from_date, to_date, start_date, end_date, page_info, page_number, and similar).
   d) The literal string "UNRESOLVED" when no clear single-column source exists.
3. If a token requires combining, formatting, or computing from multiple columns, set mapping to "UNRESOLVED" and describe the operation in meta.hints.
4. Use the CATALOG's data types and sample values to make accurate mapping decisions. Match token semantics to column semantics (numeric tokens → numeric columns, date tokens → date columns, text tokens → text columns).
5. Never invent table or column names. Never emit legacy wrappers (DERIVED:, TABLE_COLUMNS[...], COLUMN_EXP[...]).

HEADER KEYING
- If a <th> has data-label, use that value (lowercase_snake_case) as the mapping key ONLY when the same token also appears as a {placeholder} in the HTML or in SCHEMA.
- Otherwise, normalize the visible header text.

FUZZY MATCHING
- Match tokens to catalog columns considering common abbreviations:
  * "qty" ↔ "quantity", "amt" ↔ "amount", "desc" ↔ "description"
  * "sl"/"serial"/"sno" ↔ "sl_no", "wt" ↔ "weight", "pct"/"%" ↔ "percent"
- Always match against actual CATALOG column names.

AGGREGATE / MULTI-COLUMN HEADERS
- If a header represents an aggregate across multiple columns, set mapping to UNRESOLVED and record in meta.hints:
  {"op": "SUM", "over": ["table.col1", "table.col2", ...]}
  or {"op": "concat", "columns": ["table.col_a", "table.col_b"]}

CONSTANT PLACEHOLDERS
- Report ONLY tokens that are truly constant across ALL runs (page titles, company name, static captions).
- NEVER mark as constant: dates, row values, totals, page numbers, or anything in schema.row_tokens / schema.totals.
- Remove constant tokens from "mapping" but keep them in "token_samples".

TOKEN SNAPSHOT
- Emit a "token_samples" dict listing ONLY placeholder tokens that actually appear as {token_name} in the HTML above (exact name, no braces).
- Do NOT invent tokens from the catalog or schema — only use tokens you can find wrapped in curly braces in the HTML.
- For each token, output the literal string visible on the PDF. Use "NOT_VISIBLE" or "UNREADABLE" as fallback — never leave blank.

INPUTS
[FULL_HTML]
{html_for_llm}
[CATALOG]
{catalog_json}
Optional:
[SCHEMA_JSON]
{schema_json_if_any}
[REFERENCE_PNG_HINT]
"A screenshot of the reference PDF was used to create this template; treat visible page titles/branding as likely constants."

OUTPUT — return ONLY this JSON object, no markdown fences, no commentary, no text before or after:
{
  "mapping": {
    "<token>": "<table.column | params.param_name | To Be Selected... | UNRESOLVED>"
  },
  "token_samples": {
    "<token>": "<literal string from PDF>"
  },
  "meta": {
    "unresolved": ["<token>", "..."],
    "hints": {
      "<token>": { "op": "SUM|CONCAT|FORMAT", "columns": ["table.col1", "table.col2"] }
    }
  }
}

VALIDATION CHECKLIST (verify before responding):
- Every mapping value is a simple "table.column" reference, "params.*", "To Be Selected...", or "UNRESOLVED". NO SQL functions or expressions.
- Every token from the HTML appears in either "mapping" or "token_samples".
- Constants removed from mapping still appear in token_samples.
- No empty string values in token_samples.
- Output is a single valid JSON object with no surrounding text.
""".strip()


@lru_cache(maxsize=1)
def _load_llm_call_3_section() -> tuple[str, str]:
    section = _CALL3_PROMPT_SECTION_DF
    if _INPUT_MARKER in section:
        system, remainder = section.split(_INPUT_MARKER, 1)
        system_text = system.strip()
        user_template = f"{_INPUT_MARKER}{remainder}".strip()
    else:
        system_text = section.strip()
        user_template = ""
    return system_text, user_template


# Legacy alias
_load_llm_call_3_section_df = _load_llm_call_3_section


def _sanitize_html(html: str) -> str:
    """
    Strip comments, scripts, and excessive whitespace to keep prompts compact.
    """
    html = html or ""
    # remove script tags/HTML comments but preserve repeat markers and inline CSS
    comment_re = re.compile(r"(?is)<!--(?!\s*(BEGIN:BLOCK_REPEAT|END:BLOCK_REPEAT)).*?-->")
    script_re = re.compile(r"(?is)<script\b[^>]*>.*?</script>")
    collapsed = script_re.sub("", comment_re.sub("", html))
    collapsed = re.sub(r"[ \t]{2,}", " ", collapsed)
    collapsed = re.sub(r"\n{3,}", "\n\n", collapsed)
    return collapsed.strip()


def _format_catalog(catalog: Iterable[str]) -> str:
    catalog_list = [str(item).strip() for item in catalog]
    return json.dumps(catalog_list, ensure_ascii=False, indent=2)


def _normalize_schema_payload(schema: Mapping[str, Any] | None) -> Dict[str, Any]:
    base = {
        "scalars": [],
        "row_tokens": [],
        "totals": [],
        "notes": "",
    }
    if not isinstance(schema, Mapping):
        return base

    def _collect(key: str) -> list[str]:
        values = schema.get(key, [])
        if isinstance(values, Iterable) and not isinstance(values, (str, bytes)):
            return [str(item).strip() for item in values if str(item).strip()]
        return []

    base["scalars"] = _collect("scalars")
    base["row_tokens"] = _collect("row_tokens")
    base["totals"] = _collect("totals")

    notes = schema.get("notes")
    if notes is not None:
        base["notes"] = str(notes)
    return base


def _format_schema(schema: Dict[str, Any] | None) -> str:
    normalized = _normalize_schema_payload(schema)
    return json.dumps(normalized, ensure_ascii=False, indent=2, sort_keys=True)


def _row_token_hint(schema: Mapping[str, Any] | None) -> str:
    normalized = _normalize_schema_payload(schema)
    row_tokens = [tok for tok in normalized.get("row_tokens", []) if str(tok).lower().startswith("row_")]
    if not row_tokens:
        return ""

    preview = ", ".join(row_tokens[:8])
    if len(row_tokens) > 8:
        preview += ", ..."

    hint_lines = [
        "ROW TOKEN NAMING",
        "- The HTML template exposes repeating-row placeholders that already include the `row_` prefix.",
        "- When producing the mapping, reference those tokens verbatim (including the prefix and casing).",
        f"- Example row tokens: {preview}",
    ]
    return "\n".join(hint_lines)


def _adapt_prompt_for_complexity(
    base_prompt: str,
    schema_json: Dict[str, Any] | None,
    catalog_size: int,
) -> str:
    """Use the LLM to generate targeted hints for complex templates.

    Only fires for templates with >50 tokens or >3 tables in the catalog.
    Falls back to the unmodified prompt on any failure.
    """
    if not schema_json:
        return base_prompt

    token_count = (
        len(schema_json.get("scalars") or [])
        + len(schema_json.get("row_tokens") or [])
        + len(schema_json.get("totals") or [])
    )

    # Only enhance for genuinely complex templates
    if token_count < 50 and catalog_size < 4:
        return base_prompt

    try:

        client = get_llm_client()
        has_totals = bool(schema_json.get("totals"))
        has_reshape = "reshape" in str(schema_json.get("notes", "")).lower()

        prompt = (
            f"A report template has {token_count} tokens across {catalog_size} data tables"
            f"{', with totals/aggregation' if has_totals else ''}"
            f"{', with pivot/reshape requirements' if has_reshape else ''}.\n\n"
            "Produce 2-3 SHORT, specific hints (one sentence each) that would help "
            "an LLM produce a correct column-to-token mapping for this template. "
            "Focus on common pitfalls for complex mappings: "
            "table prefix consistency, row vs header tokens, total aggregation. "
            "Return ONLY the hints, one per line, no numbering."
        )

        resp = client.complete(
            messages=[{"role": "user", "content": prompt}],
            description="prompt_complexity_hints",
            max_tokens=256,
        )


        hints_text = _extract_response_text(resp)
        if hints_text and hints_text.strip():
            logger.info("prompt_complexity_hints_added", extra={"token_count": token_count, "catalog_size": catalog_size})
            return f"{base_prompt.strip()}\n\nADDITIONAL MAPPING HINTS:\n{hints_text.strip()}"
    except Exception:
        logger.debug("prompt_complexity_hints_failed", exc_info=True)

    return base_prompt


def build_llm_call_3_prompt(
    html: str,
    catalog: Iterable[str],
    schema_json: Dict[str, Any] | None = None,
    *,
    rich_catalog_text: str | None = None,
) -> Dict[str, Any]:
    """
    Build the system/user payload for LLM Call 3 (auto-map + constant discovery).

    Uses the DataFrame-mode prompt which forbids SQL expressions and only allows
    simple table.column references in mapping values.
    """
    system_template, user_template = _load_llm_call_3_section()

    if not user_template:
        user_template = system_template
        system_template = (
            "You are a meticulous analyst that performs report auto-mapping and constant inlining. "
            "Follow the subsequent instructions strictly."
        )

    html_block = _sanitize_html(html)
    catalog_block = rich_catalog_text if rich_catalog_text else _format_catalog(catalog)
    schema_block = _format_schema(schema_json)
    row_hint = _row_token_hint(schema_json)

    user_payload = user_template
    for placeholder, value in (
        ("{html_for_llm}", html_block),
        ("{catalog_json}", catalog_block),
        ("{schema_json_if_any}", schema_block),
    ):
        user_payload = user_payload.replace(placeholder, value)

    if row_hint:
        user_payload = f"{user_payload.strip()}\n\n{row_hint}"

    # Adapt prompt for complex templates with LLM-generated hints
    catalog_size = len(catalog_block.split("\n")) if catalog_block else 0
    user_payload = _adapt_prompt_for_complexity(user_payload, schema_json, catalog_size)

    attachments: list[dict[str, Any]] = []
    if "[REFERENCE_PNG_HINT]" not in user_payload:
        attachments.append(
            {
                "type": "text",
                "text": (
                    "[REFERENCE_PNG_HINT]\n"
                    '"A screenshot of the reference PDF was used to create this template; '
                    'treat visible page titles/branding as likely constants."'
                ),
            }
        )

    return {
        "system": system_template.strip(),
        "user": user_payload.strip(),
        "attachments": attachments,
        "version": PROMPT_VERSION,
    }


# DEPRECATED: Call 3.5 is replaced by deterministic constant-inlining in the unified pipeline.
# Kept for backward compatibility with legacy endpoints.
LLM_CALL_3_5_PROMPT: Dict[str, str] = {
    "system": dedent(
        """\
        You are the Step 3.5 corrections specialist in a report generation pipeline.

        YOUR THREE RESPONSIBILITIES:
        A) Apply every explicit user instruction to the HTML template (text edits, CSS tweaks, structural changes, token modifications). Do not invent changes beyond what the user asks.
        B) Inline constants: Replace any token whose mapping value is "INPUT_SAMPLE" with the literal value from `mapping_context.token_samples`. Copy the string exactly as provided.
        C) Produce a `page_summary` narrative for Step 4 that captures: constants you inlined (with their exact values), key field meanings, notable numeric totals, dates, codes, unresolved tokens, and uncertainties.

        INVARIANTS (must hold unless user explicitly overrides):
        1. Preserve the DOM hierarchy, repeat markers (<!-- BEGIN:BLOCK_REPEAT -->), data-region attributes, and <tbody> row prototypes exactly.
        2. Preserve all remaining dynamic tokens exactly as written ({token}, {{token}}, etc.). Only inline tokens mapped to "INPUT_SAMPLE".
        3. Keep HTML self-contained — no external resources, no <script> tags.

        DATA SOURCES:
        - `mapping_context.mapping`: current binding state. Tokens mapped to "INPUT_SAMPLE" → inline. Tokens mapped to table.column or SQL → leave untouched.
        - `mapping_context.token_samples`: literal strings for every placeholder. Use these exact values when inlining.
        - `mapping_context.sample_tokens` / `mapping_context.inline_tokens`: tokens the user wants double-checked. Report uncertainties about these in page_summary.
        - `user_input`: authoritative instructions for this pass. Follow exactly.
        - Reference PNG (if attached): visual context only.

        OUTPUT — strict JSON, no markdown fences, no extra keys:
        {
          "final_template_html": "<string>",
          "page_summary": "<string>"
        }

        VALIDATION CHECKLIST:
        - Remaining tokens in final_template_html = original tokens minus those explicitly inlined.
        - Repeat markers, <tbody> counts, row prototypes unchanged (unless user asked to modify).
        - No external resources, no scripts, no accidental literal leak of unresolved tokens.
        - page_summary is a detailed narrative (>1 sentence): exact inlined values, important metrics, unresolved fields, uncertainties. No layout/styling details.
        - JSON is valid UTF-8 with properly escaped strings. Only two keys: final_template_html, page_summary.
        """
    ).strip(),
    "user": dedent(
        """\
        The actual payload will be provided as a JSON object with these fields:
        {
          "template_html": "<HTML template with dynamic tokens>",
          "schema": { "scalars": [...], "row_tokens": [...], "totals": [...] },
          "mapping_context": {
            "mapping": { "<token>": "<table.column | INPUT_SAMPLE | UNRESOLVED>" },
            "mapping_override": { "<token>": "INPUT_SAMPLE" },
            "sample_tokens": ["<tokens to double-check>"],
            "token_samples": { "<token>": "<literal value from PDF>" }
          },
          "user_input": "<free-form user instructions>"
        }
        """
    ).strip(),
}


def _build_data_uri(path: Path | None) -> str | None:
    if not path or not path.exists():
        return None
    try:
        encoded = base64.b64encode(path.read_bytes()).decode("utf-8")
    except Exception:  # pragma: no cover
        return None
    return f"data:image/png;base64,{encoded}"


def build_llm_call_3_5_prompt(
    template_html: str,
    schema: Mapping[str, Any] | None,
    user_input: str,
    page_png_path: str | None = None,
    mapping_context: Mapping[str, Any] | None = None,
) -> Dict[str, Any]:
    schema_payload = dict(schema or {})
    payload: Dict[str, Any] = {
        "template_html": template_html,
        "schema": schema_payload,
        "user_input": user_input or "",
    }

    if mapping_context:
        mapping_context_clean = dict(mapping_context)
        payload["mapping_context"] = mapping_context_clean

    data_uri = _build_data_uri(Path(page_png_path) if page_png_path else None)
    payload_json = json.dumps(payload, ensure_ascii=False, indent=2)
    user_content = [
        {
            "type": "text",
            "text": f"USER (JSON payload):\n{payload_json}",
        }
    ]
    if data_uri:
        user_content.append(
            {
                "type": "image_url",
                "image_url": {
                    "url": data_uri,
                    "detail": "high",
                },
            }
        )

    messages = [
        {
            "role": "user",
            "content": user_content,
        }
    ]

    return {
        "system": LLM_CALL_3_5_PROMPT["system"],
        "messages": messages,
        "version": PROMPT_VERSION_3_5,
    }


LLM_CALL_4_SYSTEM_PROMPT_DF = dedent(
    """\
    LLM CALL 4 — Contract Builder (DataFrame Mode)
    You build the complete mapping contract for a pandas DataFrame report pipeline. NO SQL ANYWHERE.

    ═══════════════════════════════════════════════════════════════
    CRITICAL: This pipeline uses pandas DataFrames directly.
    NEVER emit SQL expressions, DuckDB functions, CASE/WHEN, SUM(), CONCAT(), or any code.
    All computations use declarative operation objects (see below).
    ═══════════════════════════════════════════════════════════════

    YOUR THREE OUTPUTS:
    1. overview_md — Markdown narrative summarizing the report logic.
    2. step5_requirements — Dataset descriptions, parameter semantics, transformation rules.
    3. contract — The authoritative mapping contract with declarative operations.

    CORE RULES:
    - Use ONLY columns from the CATALOG in "table.column" format. Never invent names.
    - Preserve every dynamic token from the schema exactly.
    - mapping_override is authoritative when provided.
    - key_tokens are required user filters → map as PARAM:<name>.

    MAPPING VALUES — each token must map to exactly one of:
    - "TABLE.COLUMN" (direct column from catalog)
    - "PARAM:name" (parameter passthrough)
    - "UNRESOLVED" (no source found — use sparingly, prefer resolving)

    ROW_COMPUTED — declarative ops for derived row columns. Each value is a dict:
      {"op": "subtract", "left": "<column_or_alias>", "right": "<column_or_alias>"}
      {"op": "add", "left": "<column_or_alias>", "right": "<column_or_alias>"}
      {"op": "multiply", "left": "<column_or_alias>", "right": "<column_or_alias_or_number>"}
      {"op": "divide", "numerator": "<column_or_alias>", "denominator": "<column_or_alias>"}
      {"op": "concat", "columns": ["col_a", "col_b"], "separator": " "}
      {"op": "format_date", "column": "<date_col>", "format": "%d-%m-%Y %H:%M:%S"}
      {"op": "format_number", "column": "<num_col>", "decimals": 2}
    "left", "right", "numerator", "denominator" can be: a column name (string), a numeric literal (number), or a nested op dict.

    TOTALS_MATH — declarative ops for aggregate totals. Each value is a dict:
      {"op": "sum", "column": "<row_token_name>"}
      {"op": "mean", "column": "<row_token_name>"}
      {"op": "count", "column": "<row_token_name>"}
      {"op": "min", "column": "<row_token_name>"}
      {"op": "max", "column": "<row_token_name>"}
      {"op": "divide", "numerator": {"op": "sum", "column": "col_a"}, "denominator": {"op": "sum", "column": "col_b"}}
    The "column" field in totals_math references ROW TOKEN names (the computed row values), not raw table columns.

    TOTALS (totals_mapping) — simple token-to-expression mapping for totals that mirrors the totals_math logic.
    Can be a dict of declarative ops (same format as totals_math) or a simple string reference.

    RESHAPE RULES:
    - Each rule: {"purpose": "≤15 words", "strategy": "UNION_ALL|MELT|NONE", "columns": [{"as": "alias", "from": ["table.col1", "table.col2", ...]}]}
    - "as" is the output column alias used in row tokens. "from" lists the source catalog columns to unpivot.
    - For MELT/UNION_ALL: each "from" array must have the same length across all columns entries.
    - If "from" values are literal constants (not column references), list them as string literals (e.g., ["1", "2", "3"]).

    CONTRACT STRUCTURE:
    - join: non-empty parent_table/parent_key required. If no child table, set child_table = parent_table, child_key = parent_key.
    - order_by.rows AND row_order: both non-empty arrays with identical content. Default ["ROWID"] if no logical ordering.
    - formatters: "percent(2)", "number(2)", "currency(2)", etc. Do NOT put "date()" in formatters for timestamp tokens.
    - unresolved: must be [].
    - header_tokens: copy of tokens.scalars array.
    - row_tokens: copy of tokens.row_tokens array.

    MANDATORY RULES (violations will be auto-corrected by post-processor):
    1. TIMESTAMP FORMATTING: Every token that maps to a timestamp/date column (timestamp_utc, timestamp, created_at, date, datetime) MUST have a row_computed entry:
       {"op": "format_date", "column": "<col>", "format": "%d-%m-%Y %H:%M:%S"}
       Do NOT use formatters "date(...)" for timestamps — always use row_computed.format_date.
    2. NUMERIC FORMATTING: Every token that maps to a numeric measurement column MUST have a formatters entry. Default: "number(2)". Use higher precision only when the domain requires it (e.g., pH sensors → "number(4)").
    3. DATE FILTERS: When date_columns is populated, filters.optional MUST contain:
       "date_from": "TABLE.date_column", "date_to": "TABLE.date_column"
       Never leave filters.optional empty when date_columns exists.
    4. DATE_COLUMNS: If ANY mapped column is a timestamp/date type, date_columns MUST be populated with {"TABLE": "column_name"}.
    5. CONSISTENCY: Use row_computed.format_date for timestamps and formatters for display formatting (number, percent, currency). Never mix — do not put "date()" in formatters for timestamp columns.

    ═══════════════════════════════════════════════════════════════
    INPUT PAYLOAD SHAPE:
    {
      "final_template_html": "<HTML with constants inlined>",
      "page_summary": "<narrative from Step 3.5>",
      "schema": { "scalars": [...], "row_tokens": [...], "totals": [...] },
      "auto_mapping_proposal": { "mapping": {...}, "join": {...}, "unresolved": [...] },
      "mapping_override": { "<token>": "<authoritative mapping>" },
      "user_instructions": "<free-form user guidance>",
      "key_tokens": ["<required filter tokens>"],
      "catalog": ["table.column", ...]
    }

    ═══════════════════════════════════════════════════════════════
    OUTPUT — return ONLY this JSON object, no markdown fences, no commentary:
    {
      "overview_md": "<Markdown: Executive Summary, Token Inventory, Mapping Table, Join & Date Rules, Transformations, Parameters>",
      "step5_requirements": {
        "datasets": {
          "header": {"description": "...", "columns": ["<scalar tokens>"]},
          "rows": {"description": "...", "columns": ["<row tokens>"], "grouping": [...], "ordering": [...]},
          "totals": {"description": "...", "columns": ["<totals tokens>"]}
        },
        "semantics": "<filter vs pass-through explanation>",
        "parameters": {
          "required": [{"name": "...", "type": "date|string"}],
          "optional": [{"name": "...", "type": "string"}]
        },
        "transformations": ["<reshape rules in plain English>"]
      },
      "contract": {
        "tokens": { "scalars": [...], "row_tokens": [...], "totals": [...] },
        "mapping": { "<token>": "<TABLE.COLUMN | PARAM:name | UNRESOLVED>" },
        "unresolved": [],
        "join": { "parent_table": "...", "parent_key": "...", "child_table": "...", "child_key": "..." },
        "date_columns": { "<table>": "<date_column>" },
        "filters": { "optional": { "<name>": "table.column" } },
        "reshape_rules": [
          { "purpose": "...", "strategy": "UNION_ALL|MELT|NONE", "columns": [{"as": "alias", "from": ["table.col1", "..."]}] }
        ],
        "row_computed": { "<token>": {"op": "...", "left": "...", "right": "..."} },
        "totals_math": { "<token>": {"op": "...", "column": "..."} },
        "totals": { "<token>": {"op": "...", "column": "..."} },
        "formatters": { "<token>": "<format spec>" },
        "order_by": { "rows": ["<column ASC|DESC>"] },
        "header_tokens": ["<scalar tokens copy>"],
        "row_tokens": ["<row tokens copy>"],
        "row_order": ["<column ASC|DESC>"],
        "literals": {},
        "notes": "<domain notes>"
      },
      "validation": {
        "unknown_tokens": [],
        "unknown_columns": [],
        "token_coverage": { "scalars_mapped_pct": 100, "row_tokens_mapped_pct": 100, "totals_mapped_pct": 100 }
      }
    }

    SELF-CHECK before responding:
    - NO SQL expressions anywhere (no SUM(), CASE, CONCAT, STRFTIME, etc.).
    - Every schema token appears in contract.mapping.
    - Every column reference exists in the CATALOG.
    - row_computed and totals_math values are ALL declarative op dicts, never strings.
    - order_by.rows and row_order are identical non-empty arrays.
    - join block has all four non-empty string fields.
    - Every reshape rule has a non-empty "purpose".
    - token_coverage is 100%.
    - Every timestamp-mapped token has row_computed.format_date with "%d-%m-%Y %H:%M:%S".
    - Every numeric column token has a formatters entry (e.g., "number(2)").
    - filters.optional has date_from/date_to when date_columns is non-empty.
    - No "date()" entries in formatters for tokens that have row_computed.format_date.
    """
).strip()

# DEPRECATED: Call 5 is merged into Call 4's self-check section in the unified pipeline.
# Kept for backward compatibility with legacy endpoints.
LLM_CALL_5_PROMPT: Dict[str, str] = {
    "system": dedent(
        """\
        LLM CALL 5 — Contract Finalizer (DataFrame Mode)
        You finalize the contract from Step 4 for the pandas DataFrame report pipeline. NO SQL.

        ═══════════════════════════════════════════════════════════════
        CRITICAL: NEVER emit SQL expressions, DuckDB functions, or any code.
        All computations use declarative operation objects only.
        ═══════════════════════════════════════════════════════════════

        YOUR JOB:
        1. Copy the Step-4 contract exactly (same tokens, same declarative ops, same ordering).
        2. Validate and fill in any missing optional fields with sensible defaults.
        3. Ensure the contract is complete and ready for the DataFrame pipeline.

        RULES:
        - Treat `step4_output.contract` as authoritative. Do not add, drop, or rename tokens.
        - `mapping` values: only "TABLE.COLUMN", "PARAM:name", or "UNRESOLVED".
        - `row_computed` / `totals_math`: declarative op dicts only:
          {"op": "subtract|add|multiply|divide|sum|mean|count|min|max|concat|format_date|format_number", ...}
        - `totals` (totals_mapping): declarative op dicts mirroring totals_math.
        - Join block: non-empty parent_table/parent_key. If no child table, reuse parent.
        - `order_by.rows` and `row_order`: both non-empty arrays, identical content. Default ["ROWID"].
        - Every reshape rule must have a non-empty "purpose" (≤15 words).
        - `header_tokens`: copy of tokens.scalars. `row_tokens`: copy of tokens.row_tokens.
        - Verify timestamp tokens use row_computed.format_date with "%d-%m-%Y %H:%M:%S" (not formatters "date()").
        - Verify all numeric measurement tokens have formatters entries (e.g., "number(2)").
        - Verify filters.optional has date_from/date_to when date_columns exists.

        OUTPUT — return ONLY this JSON object, no markdown fences, no commentary:
        {
          "contract": {
            "tokens": { "scalars": [...], "row_tokens": [...], "totals": [...] },
            "mapping": { "<token>": "TABLE.COLUMN|PARAM:name|UNRESOLVED" },
            "join": { "parent_table": "...", "parent_key": "...", "child_table": "...", "child_key": "..." },
            "date_columns": { "<table>": "<date_column>" },
            "filters": { "optional": { "<name>": "table.column" } },
            "reshape_rules": [{"purpose": "...", "strategy": "UNION_ALL|MELT|NONE", "columns": [...]}],
            "row_computed": { "<token>": {"op": "...", ...} },
            "totals_math": { "<token>": {"op": "...", ...} },
            "totals": { "<token>": {"op": "...", ...} },
            "formatters": { "<token>": "<format spec>" },
            "order_by": { "rows": ["<column ASC|DESC>"] },
            "header_tokens": [...],
            "row_tokens": [...],
            "row_order": ["<column ASC|DESC>"],
            "literals": {},
            "notes": "..."
          },
          "invalid": false
        }

        SELF-CHECK:
        - NO SQL expressions anywhere.
        - Every token from Step-4 is present in mapping.
        - row_computed and totals_math are ALL declarative op dicts.
        - order_by.rows and row_order are identical non-empty arrays.
        - join block has all four non-empty string fields.
        """
    ).strip(),
    "user": dedent(
        """\
        {
          "final_template_html": "<HTML from Step 3.5>",
          "step4_output": {
            "contract": { /* Step-4 contract object */ },
            "overview_md": "Step-4 overview",
            "step5_requirements": { /* Step-4 requirements */ }
          },
          "key_tokens": ["..."],
          "catalog": ["table.column", "..."]
        }
        """
    ).strip(),
}

# Legacy alias for backwards compatibility
LLM_CALL_5_PROMPT_DF = LLM_CALL_5_PROMPT

PROMPT_LIBRARY: Dict[str, str] = {
    **LLM_CALL_PROMPTS,
    "llm_call_3_5_system": LLM_CALL_3_5_PROMPT["system"],
    "llm_call_3_5_user": LLM_CALL_3_5_PROMPT["user"],
    "llm_call_4_system": LLM_CALL_4_SYSTEM_PROMPT_DF,
    "llm_call_5_system": LLM_CALL_5_PROMPT["system"],
    "llm_call_5_user": LLM_CALL_5_PROMPT["user"],
}

# Legacy alias — SQL prompt removed, DF is the only mode
LLM_CALL_4_SYSTEM_PROMPT = LLM_CALL_4_SYSTEM_PROMPT_DF


def build_llm_call_4_prompt(
    *,
    final_template_html: str,
    page_summary: str,
    schema: Mapping[str, Any] | None,
    auto_mapping_proposal: Mapping[str, Any],
    mapping_override: Mapping[str, Any] | None,
    user_instructions: str,
    catalog: Iterable[str],
    dialect_hint: str | None = None,
    key_tokens: Iterable[str] | None = None,
) -> Dict[str, Any]:
    """
    Build the payload for LLM Call 4 (contract builder + overview).
    Always uses DataFrame mode — no SQL expressions in output.
    """
    system_text = LLM_CALL_4_SYSTEM_PROMPT_DF
    key_tokens_list: list[str] = []
    if key_tokens:
        seen: set[str] = set()
        for token in key_tokens:
            text = str(token or "").strip()
            if not text or text in seen:
                continue
            seen.add(text)
            key_tokens_list.append(text)
    payload: Dict[str, Any] = {
        "final_template_html": final_template_html,
        "page_summary": page_summary,
        "schema": dict(schema or {}),
        "auto_mapping_proposal": dict(auto_mapping_proposal or {}),
        "mapping_override": dict(mapping_override or {}),
        "user_instructions": user_instructions or "",
        "catalog": [str(item) for item in catalog],
    }
    if key_tokens_list:
        payload["key_tokens"] = key_tokens_list
    if dialect_hint is not None:
        payload["dialect_hint"] = str(dialect_hint)

    payload_json = json.dumps(payload, ensure_ascii=False, indent=2)
    messages = [
        {
            "role": "user",
            "content": [
                {
                    "type": "text",
                    "text": payload_json,
                }
            ],
        }
    ]

    return {
        "system": system_text,
        "messages": messages,
        "version": PROMPT_VERSION_4,
    }


@lru_cache(maxsize=1)
def get_prompt_generator_assets() -> Dict[str, str]:
    """Return the system and user template strings for LLM CALL 5 (DataFrame mode)."""
    return dict(LLM_CALL_5_PROMPT)


# Legacy alias
get_prompt_generator_assets_df = get_prompt_generator_assets


# LLM Call 3A: Mapping Semantic Audit (Phase 2)
# DEPRECATED: Call 3A mapping audit is replaced by heuristic checks in validator/checks.py.
# Kept for backward compatibility with legacy endpoints.

LLM_CALL_3A_PROMPT = {
    "system": dedent("""\
        You are a mapping auditor. For each token-to-column mapping, you see
        the token name, the mapped column, and sample values from that column.

        Your job: flag mappings where the sample data does NOT semantically match
        what the token name implies.

        Examples of BAD mappings:
        - Token "row_amount" mapped to column with sample ["2026-01-15", "2026-02-20"] (dates, not amounts)
        - Token "row_date" mapped to column with sample ["425.50", "180.00"] (numbers, not dates)
        - Token "row_name" mapped to column with sample [1, 2, 3] (IDs, not names)

        Examples of GOOD mappings:
        - Token "row_amount" mapped to column with sample ["1500.50", "2300.00", "890.75"] (monetary values)
        - Token "row_print_date_time" mapped to column with sample ["2026-01-15 08:30:00"] (datetime)

        Also suggest BETTER column alternatives if you see a mismatch and the catalog
        has a more appropriate column.

        OUTPUT — JSON only, no markdown fences:
        {
          "audited_mappings": {
            "<token>": {
              "current": "table.column",
              "verdict": "ok" | "suspect" | "wrong",
              "reason": "...",
              "suggested_column": "table.better_column" | null
            }
          }
        }
    """).strip(),
    "user": dedent("""\
        Audit these token-to-column mappings. For each one, check if the sample
        data semantically matches the token name. The full database catalog is
        provided so you can suggest better alternatives.

        MAPPINGS:
        {mappings_json}

        DATABASE CATALOG:
        {catalog_text}
    """).strip(),
}

# LLM Ops Refiner: Contract Ops Refiner (Phase 3)
# DEPRECATED: Merged into contract validator sub-prompt in the unified pipeline.
# Kept for backward compatibility.

LLM_OPS_REFINER_PROMPT = {
    "system": dedent("""\
        You are a pandas DataFrame ops specialist. Review and fix the computed
        columns (row_computed) and aggregate totals (totals_math) sections of a
        report contract.

        THIS IS A PANDAS-ONLY SYSTEM. No SQL anywhere.

        COLUMN NAMESPACE RULES (critical):
        1. row_computed ops use BARE column names — these are the column names in
           the DataFrame AFTER reshape. Look at reshape_output_columns to see what
           is available. Example: "quantity" not "order_items.quantity"
        2. totals_math ops use ROW TOKEN names — these are the token names from
           the row_tokens list. Example: "row_qty" not "quantity"
        3. Numeric literals are allowed: 100, 3.14, etc.

        SUPPORTED OPS:
        - {"op": "subtract", "left": "<col_or_num>", "right": "<col_or_num>"}
        - {"op": "add", "left": "<col_or_num>", "right": "<col_or_num>"}
        - {"op": "multiply", "left": "<col_or_num>", "right": "<col_or_num>"}
        - {"op": "divide", "numerator": "<col_or_num>", "denominator": "<col_or_num>"}
        - {"op": "sum", "column": "<token_name>"}
        - {"op": "mean", "column": "<token_name>"}
        - {"op": "count", "column": "<token_name>"}
        - {"op": "min", "column": "<token_name>"}
        - {"op": "max", "column": "<token_name>"}
        - {"op": "concat", "columns": ["col_a", "col_b"], "separator": " "}
        - {"op": "format_date", "column": "<col>", "format": "%d-%m-%Y"}
        - {"op": "format_number", "column": "<col>", "decimals": 2}

        VALIDATION RULES:
        - Every column ref in row_computed MUST exist in reshape_output_columns
        - Every column ref in totals_math MUST exist in row_tokens
        - Ops must be dicts, never strings
        - Nested ops allowed for complex expressions (e.g., (a * b) - (c * d))
        - Remove ops with unresolvable references rather than guessing
        - formatters must reference tokens in row_tokens or header_tokens

        OUTPUT — JSON only, no markdown fences:
        {
          "row_computed": { "<token>": <op_dict_or_null>, ... },
          "totals_math": { "<token>": <op_dict_or_null>, ... },
          "formatters": { "<token>": "<format_spec>", ... },
          "fixes_applied": [
            {"token": "<name>", "section": "<section>", "fix": "<description>"}
          ]
        }

        Return ONLY sections that needed changes. Omit sections already correct.
    """).strip(),
    "user": dedent("""\
        Review and fix the computed ops in this contract. Ensure column references
        use the correct namespace (bare columns in row_computed, row token names
        in totals_math).

        ROW_COMPUTED:
        {row_computed_json}

        TOTALS_MATH:
        {totals_math_json}

        FORMATTERS:
        {formatters_json}

        RESHAPE OUTPUT COLUMNS (bare names available in the DataFrame):
        {reshape_output_columns_json}

        ROW TOKENS:
        {row_tokens_json}

        HEADER TOKENS:
        {header_tokens_json}

        MAPPING (token → table.column):
        {mapping_json}
    """).strip(),
}

# LLM Simulate Fill: Contract-Data Simulation (Phase 4)
# DEPRECATED: Merged into contract validator sub-prompt in the unified pipeline.
# Kept for backward compatibility.

LLM_SIMULATE_FILL_PROMPT = {
    "system": dedent("""\
        You are a report generation simulator. Given a contract, sample database
        rows, and template token list, simulate what the final report would look like.

        For each token, determine what value would be filled in based on:
        1. The contract mapping (token → table.column)
        2. The actual data rows provided
        3. The join configuration (which parent row matches which child rows)
        4. Any row_computed operations

        REPORT:
        1. For each HEADER token: what value would appear (from first matching row)
        2. For each ROW token: show values for the first 3 data rows
        3. For each TOTAL token: what the aggregate would be
        4. Flag any tokens that would be EMPTY (column missing, no data, wrong type)
        5. Flag any tokens where the data type looks WRONG for the context

        OUTPUT — JSON only, no markdown fences:
        {
          "header_fill": {"<token>": "<value_or_EMPTY>"},
          "row_fill": [
            {"<token>": "<value>", ...}
          ],
          "totals_fill": {"<token>": "<value_or_EMPTY>"},
          "issues": [
            {
              "token": "<token>",
              "severity": "error" | "warning",
              "problem": "empty | type_mismatch | no_data | wrong_join",
              "detail": "..."
            }
          ]
        }
    """).strip(),
    "user": dedent("""\
        Simulate filling this report template with the sample data below.

        CONTRACT:
        {contract_json}

        SAMPLE DATA:
        {sample_data_json}

        TEMPLATE TOKENS:
        {template_tokens_json}
    """).strip(),
}


__all__ = [
    "build_llm_call_3_prompt",
    "build_llm_call_3_5_prompt",
    "PROMPT_VERSION",
    "PROMPT_VERSION_3_5",
    "PROMPT_VERSION_4",
    "PROMPT_VERSION_DF",
    "PROMPT_VERSION_4_DF",
    "build_llm_call_4_prompt",
    "LLM_CALL_PROMPTS",
    "PROMPT_LIBRARY",
    "get_prompt_generator_assets",
    "get_prompt_generator_assets_df",
    "LLM_CALL_3A_PROMPT",
    "LLM_OPS_REFINER_PROMPT",
    "LLM_SIMULATE_FILL_PROMPT",
]


# PROMPTS DATA (merged from prompts_data.py)


EXCEL_PROMPT_VERSION = "excel_llm_call_3_v1"
EXCEL_PROMPT_VERSION_3_5 = "excel_llm_call_3_5_v1"
EXCEL_PROMPT_VERSION_4 = "excel_llm_call_4_v2"
EXCEL_PROMPT_VERSION_5 = "excel_llm_call_5_v1"

_INPUT_MARKER = "[INPUTS]"

# _sanitize_html, _format_catalog, _format_schema defined above (from llm_prompts)

EXCEL_LLM_CALL_1_PROMPT = dedent(
    """
Produce a COMPLETE, self-contained HTML document (<!DOCTYPE html> ...) with inline <style>. Treat the provided Excel
worksheet prototype HTML (tokens already annotated) as the blueprint and recreate the worksheet layout as a
print-ready template.

PLACEHOLDER & SCHEMA RULES
- Placeholders must use single braces: {token}. Never invent tokens beyond those in SCHEMA_JSON (if provided) or the
  supplied row_* tokens present in the prototype. Re-use them verbatim.
- Emit exactly ONE prototype <tbody><tr>...</tr></tbody> row for repeating data. Do not duplicate multiple data rows.
- Totals/footers that remain dynamic should use tokens; static captions (titles, notes) must remain literal strings.
- Keep casing/spelling from the worksheet prototype for visible labels.

STRUCTURE & CSS
- Use a semantic table for the data grid (table-layout: fixed, border-collapse: collapse). Apply borders only where the
  worksheet shows lines (default to 1px solid #999 for visible lines). Align numeric columns to the right.
- Provide minimal wrappers with stable IDs for major regions (#report-header, #data-table, #report-totals, etc.).
- Include @page { size: A4; margin: sensible } so the HTML prints like a sheet export.
- If the layout implies repeating batch blocks outside the table, wrap the prototype inside:
    <!-- BEGIN:BLOCK_REPEAT batches -->
    <section class="batch-block">...</section>
    <!-- END:BLOCK_REPEAT -->
  Keep headers/footers outside those markers.

PROTOTYPE NOTES
- The supplied HTML already encodes preface rows, column order, and row_* tokens. Preserve that structure while
  improving styling/printability.

OUTPUT FORMAT
1) RAW HTML between these markers:
   <!--BEGIN_HTML-->
   ... full html ...
   <!--END_HTML-->
2) Matching SCHEMA JSON (scalars, row_tokens, totals, notes) between:
   <!--BEGIN_SCHEMA_JSON-->
   { ... }
   <!--END_SCHEMA_JSON-->
Schema tokens must align 1:1 with placeholders in your HTML.

Return ONLY those markers—no markdown fences or commentary.

[INPUTS]
SHEET_SNAPSHOT_JSON:
{sheet_snapshot}

SHEET_PROTOTYPE_HTML:
{sheet_html}

SCHEMA (may be empty):
{schema_str}
"""
).strip()


def build_excel_llm_call_1_prompt(sheet_snapshot: str, sheet_html: str, schema_str: str) -> str:
    snapshot_payload = sheet_snapshot.strip() or "{}"
    sheet_html_payload = sheet_html.strip() or "<html></html>"
    schema_payload = schema_str.strip() or "{}"
    return (
        EXCEL_LLM_CALL_1_PROMPT.replace("{sheet_snapshot}", snapshot_payload)
        .replace("{sheet_html}", sheet_html_payload)
        .replace("{schema_str}", schema_payload)
    )

_EXCEL_CALL3_PROMPT_SECTION = dedent(
    """
You are a meticulous report auto-mapping analyst. You are given the FULL HTML of an Excel-rendered worksheet template and a DB CATALOG.

YOUR TWO TASKS:
A) AUTO-MAPPING — Map every placeholder token in the HTML to its data source from the CATALOG.
B) CONSTANT DISCOVERY — Identify tokens whose values are static and record their literal values.

MAPPING RULES — CRITICAL: NO SQL EXPRESSIONS
1. Output ONLY simple "table.column" references from the CATALOG. NEVER output SQL expressions, functions, or any code.
2. Every mapping value must be ONE of:
   a) A catalog column in exact "table.column" format.
   b) A parameter passthrough in "params.param_name" format.
   c) The literal string "To Be Selected..." for date-range and page/filter tokens.
   d) The literal string "UNRESOLVED" when no clear single-column source exists.
3. If a token requires combining, formatting, or computing from multiple columns, set mapping to "UNRESOLVED" and describe the operation in meta.hints.
4. Never invent table or column names. Never emit legacy wrappers (DERIVED:, TABLE_COLUMNS[...], COLUMN_EXP[...]).

HEADER KEYING
- If a <th> has data-label, use that value (lowercase_snake_case) as the mapping key.

FUZZY MATCHING
- Match tokens to catalog columns considering common abbreviations:
  * "qty" ↔ "quantity", "amt" ↔ "amount", "wt" ↔ "weight", "pct" ↔ "percent"

AGGREGATE / MULTI-COLUMN HEADERS
- Set mapping to UNRESOLVED and record in meta.hints:
  {{"op": "SUM", "over": ["table.col1", "table.col2", ...]}}

CONSTANT PLACEHOLDERS
- Report ONLY tokens that are truly constant across ALL runs.
- NEVER mark as constant: dates, row values, totals, page numbers.
- Remove constant tokens from "mapping" but keep them in "token_samples".

TOKEN SNAPSHOT
- Emit a "token_samples" dict listing ONLY tokens found as {token_name} in the HTML. Do NOT invent tokens from the catalog. Use "NOT_VISIBLE" as fallback.

INPUTS
[FULL_HTML]
{html_for_llm}
[CATALOG]
{catalog_json}
Optional:
[SCHEMA_JSON]
{schema_json_if_any}
[REFERENCE_PNG_HINT]
"A screenshot of the Excel worksheet was used to create this template; treat visible sheet titles/branding as likely constants."

OUTPUT — return ONLY this JSON object, no markdown, no commentary:
{{
  "mapping": {{
    "<token>": "<table.column | params.param_name | To Be Selected... | UNRESOLVED>"
  }},
  "token_samples": {{
    "<token>": "<literal string>"
  }},
  "meta": {{
    "unresolved": ["<token>"],
    "hints": {{
      "<token>": {{ "op": "SUM", "over": ["table.col1", "table.col2"] }}
    }}
  }}
}}

VALIDATION: Every mapping value must be a simple "table.column", "params.*", "To Be Selected...", or "UNRESOLVED". NO SQL.
"""
).strip()


@lru_cache(maxsize=1)
def _load_excel_llm_call_3_section() -> tuple[str, str]:
    section = _EXCEL_CALL3_PROMPT_SECTION
    if _INPUT_MARKER in section:
        system, remainder = section.split(_INPUT_MARKER, 1)
        system_text = system.strip()
        user_template = f"{_INPUT_MARKER}{remainder}".strip()
    else:
        system_text = section
        user_template = ""
    return system_text, user_template


def build_excel_llm_call_3_prompt(
    html: str,
    catalog: Iterable[str],
    schema_json: Dict[str, Any] | None = None,
    *,
    sample_data: Mapping[str, Any] | None = None,
) -> Dict[str, Any]:
    system_template, user_template = _load_excel_llm_call_3_section()
    if not user_template:
        user_template = system_template
        system_template = (
            "You are the Excel auto-mapping analyst. Follow the subsequent instructions exactly and return JSON only."
        )

    html_block = _sanitize_html(html)
    catalog_block = _format_catalog(catalog)
    schema_block = _format_schema(schema_json)
    sample_block = json.dumps(sample_data, ensure_ascii=False, indent=2) if sample_data else ""

    user_payload = user_template
    for placeholder, value in (
        ("{html_for_llm}", html_block),
        ("{catalog_json}", catalog_block),
        ("{schema_json_if_any}", schema_block),
        ("{sample_row_json_if_any}", sample_block),
    ):
        user_payload = user_payload.replace(placeholder, value)

    return {
        "system": system_template.strip(),
        "user": user_payload.strip(),
        "attachments": [],
        "version": EXCEL_PROMPT_VERSION,
    }


EXCEL_LLM_CALL_3_5_PROMPT: Dict[str, str] = {
    "system": dedent(
        """\
        You are the Step 3.5 corrections specialist.
        Your responsibilities:
        A) Apply every explicit user instruction to the HTML template. Text edits, structural tweaks, CSS adjustments, and token changes are all allowed when the user asks. Do not invent changes or perform a wholesale redesign unless the user requests it.
        B) Inline any token whose mapping (or explicit user instruction) marks it as a constant (e.g., mapping value "INPUT_SAMPLE"). Use the literals provided in `mapping_context.token_samples` whenever available—copy the string exactly.
        C) Produce a `page_summary` that captures the page's business/data content for Step 4: list the constants you inlined, key field values, notable numeric totals, dates, codes, unresolved tokens, and uncertainties. Do not rehash layout, typography, or other presentation details unless they directly affect data interpretation.

        Core invariants (must hold unless a user instruction explicitly overrides them):
        1) Preserve the DOM hierarchy, repeat markers, data-region attributes, and row prototypes; only adjust them when the user explicitly says so.
        2) Preserve all remaining dynamic tokens exactly (examples: "{token}", "{{ token }}", "<span id='tok-x'>{{token}}</span>"). Only inline tokens you were instructed to convert to constants.
        3) Keep the HTML self-contained (no external resources or <script> tags). Maintain semantic structure.

        Hints:
        - The `mapping_context.mapping` object reflects the latest binding state after Step 3 and any overrides. Tokens mapped to "INPUT_SAMPLE" must be inlined; leave tokens mapped to DuckDB SQL expressions or table columns untouched unless instructed otherwise.
        - The `mapping_context.token_samples` dictionary lists the literal strings extracted in Step 3 for every placeholder. Inline tokens using these values exactly.
        - `mapping_context.sample_tokens` / `mapping_context.inline_tokens` highlight placeholders the user wants to double-check; use these cues when reporting lingering uncertainties in the page summary.
        - When provided, `reference_worksheet_html` contains a data-only rendering of the worksheet used to derive the template. Use it only to confirm literal strings; do not re-map tokens based on it.
        - `user_input` contains the authoritative instructions for this pass. Follow it exactly.

        Output (strict JSON, no markdown fences, no extra keys):
        {
          "final_template_html": "<string>",  // template after applying user instructions and inlining required constants
          "page_summary": "<string>"          // thorough prose description of the worksheet; must be non-empty
        }

        Validation checklist before responding:
        - Tokens remaining in "final_template_html" match the original tokens minus those explicitly inlined as constants.
        - Repeat markers, <tbody> counts, row prototypes, and data-region attributes are unchanged unless the user asked to modify them.
        - HTML stays free of external resources/scripts and contains no accidental literal leak of unresolved token data.
        - "page_summary" is a detailed narrative (>1 sentence) that reports the exact values you inlined (including any best-guess readings), important metrics, unresolved fields, and uncertainties, without digressing into layout or styling trivia.
        - JSON is valid (UTF-8), strings escaped properly, and only the two required keys are present.


        """
    ).strip(),
    "user": "USER (JSON payload):\n{payload}",
}


def build_excel_llm_call_3_5_prompt(
    *,
    template_html: str,
    schema: Mapping[str, Any] | None,
    user_input: str,
    page_png_path: str | None = None,
    reference_worksheet_html: str | None = None,
    mapping_context: Mapping[str, Any] | None = None,
) -> Dict[str, Any]:
    schema_payload = dict(schema or {})
    payload: Dict[str, Any] = {
        "template_html": template_html,
        "schema": schema_payload,
        "user_input": user_input or "",
    }
    if mapping_context:
        payload["mapping_context"] = dict(mapping_context)
    if isinstance(reference_worksheet_html, str) and reference_worksheet_html.strip():
        payload["reference_worksheet_html"] = reference_worksheet_html

    # Prefer full worksheet HTML when provided; otherwise, attach optional page image as a fallback.
    from pathlib import Path as _Path  # local import to avoid cycles

    try:
        _build_data_uri = _pdf_prompts._build_data_uri  # type: ignore[attr-defined]
    except Exception:  # pragma: no cover

        def _null_build_data_uri(_p):  # type: ignore
            return None

        _build_data_uri = _null_build_data_uri  # type: ignore

    data_uri = None if reference_worksheet_html else _build_data_uri(_Path(page_png_path) if page_png_path else None)
    payload_json = json.dumps(payload, ensure_ascii=False, indent=2)
    user_content = [{"type": "text", "text": EXCEL_LLM_CALL_3_5_PROMPT["user"].format(payload=payload_json)}]
    if data_uri:
        user_content.append({"type": "image_url", "image_url": {"url": data_uri, "detail": "high"}})

    messages = [{"role": "user", "content": user_content}]

    return {
        "system": EXCEL_LLM_CALL_3_5_PROMPT["system"],
        "messages": messages,
        "version": EXCEL_PROMPT_VERSION_3_5,
    }


EXCEL_LLM_CALL_4_SYSTEM_PROMPT = dedent(
    """\
    LLM CALL 4 — Contract Builder (Excel, DataFrame Mode)
    You build the complete mapping contract for an Excel worksheet report using a pandas DataFrame pipeline. NO SQL ANYWHERE.

    ═══════════════════════════════════════════════════════════════
    CRITICAL: This pipeline uses pandas DataFrames directly.
    NEVER emit SQL expressions, DuckDB functions, CASE/WHEN, SUM(), CONCAT(), or any code.
    All computations use declarative operation objects (see below).
    ═══════════════════════════════════════════════════════════════

    YOUR THREE OUTPUTS:
    1. overview_md — Markdown narrative summarizing the report logic.
    2. step5_requirements — Dataset descriptions, parameter semantics, transformation rules.
    3. contract — The authoritative mapping contract with declarative operations.

    CORE RULES:
    - Use ONLY columns from the CATALOG in "table.column" format. Never invent names.
    - Preserve every dynamic token from the schema exactly.
    - mapping_override is authoritative when provided.
    - key_tokens are required user filters → map as PARAM:<name>.

    MAPPING VALUES — each token must map to exactly one of:
    - "TABLE.COLUMN" (direct column from catalog)
    - "PARAM:name" (parameter passthrough)
    - "UNRESOLVED" (no source found — use sparingly, prefer resolving)

    ROW_COMPUTED — declarative ops for derived row columns. Each value is a dict:
      {"op": "subtract", "left": "<column_or_alias>", "right": "<column_or_alias>"}
      {"op": "add", "left": "<column_or_alias>", "right": "<column_or_alias>"}
      {"op": "multiply", "left": "<column_or_alias>", "right": "<column_or_alias_or_number>"}
      {"op": "divide", "numerator": "<column_or_alias>", "denominator": "<column_or_alias>"}
      {"op": "concat", "columns": ["col_a", "col_b"], "separator": " "}
      {"op": "format_date", "column": "<date_col>", "format": "%d-%m-%Y %H:%M:%S"}
      {"op": "format_number", "column": "<num_col>", "decimals": 2}
    "left", "right", "numerator", "denominator" can be: a column name (string), a numeric literal (number), or a nested op dict.

    TOTALS_MATH — declarative ops for aggregate totals. Each value is a dict:
      {"op": "sum", "column": "<row_token_name>"}
      {"op": "mean", "column": "<row_token_name>"}
      {"op": "count", "column": "<row_token_name>"}
      {"op": "min", "column": "<row_token_name>"}
      {"op": "max", "column": "<row_token_name>"}
      {"op": "divide", "numerator": {"op": "sum", "column": "col_a"}, "denominator": {"op": "sum", "column": "col_b"}}
    The "column" field in totals_math references ROW TOKEN names (the computed row values), not raw table columns.

    TOTALS (totals_mapping) — simple token-to-expression mapping for totals that mirrors the totals_math logic.
    Can be a dict of declarative ops (same format as totals_math) or a simple string reference.

    RESHAPE RULES:
    - Each rule: {"purpose": "≤15 words", "strategy": "UNION_ALL|MELT|NONE", "columns": [{"as": "alias", "from": ["table.col1", "table.col2", ...]}]}
    - "as" is the output column alias used in row tokens. "from" lists the source catalog columns to unpivot.
    - For MELT/UNION_ALL: each "from" array must have the same length across all columns entries.
    - If "from" values are literal constants (not column references), list them as string literals (e.g., ["1", "2", "3"]).

    CONTRACT STRUCTURE:
    - join: non-empty parent_table/parent_key required. If no child table, set child_table = parent_table, child_key = parent_key.
    - order_by.rows AND row_order: both non-empty arrays with identical content. Default ["ROWID"] if no logical ordering.
    - formatters: "percent(2)", "date(YYYY-MM-DD)", "number(2)", etc.
    - unresolved: must be [].
    - header_tokens: copy of tokens.scalars array.
    - row_tokens: copy of tokens.row_tokens array.

    ═══════════════════════════════════════════════════════════════
    INPUT PAYLOAD SHAPE:
    {
      "final_template_html": "<HTML with constants inlined>",
      "page_summary": "<narrative from Step 3.5>",
      "schema": { "scalars": [...], "row_tokens": [...], "totals": [...] },
      "auto_mapping_proposal": { "mapping": {...}, "join": {...}, "unresolved": [...] },
      "mapping_override": { "<token>": "<authoritative mapping>" },
      "user_instructions": "<free-form user guidance>",
      "key_tokens": ["<required filter tokens>"],
      "catalog": ["table.column", ...]
    }

    ═══════════════════════════════════════════════════════════════
    OUTPUT — return ONLY this JSON object, no markdown fences, no commentary:
    {
      "overview_md": "<Markdown: Executive Summary, Token Inventory, Mapping Table, Join & Date Rules, Transformations, Parameters>",
      "step5_requirements": {
        "datasets": {
          "header": {"description": "...", "columns": ["<scalar tokens>"]},
          "rows": {"description": "...", "columns": ["<row tokens>"], "grouping": [...], "ordering": [...]},
          "totals": {"description": "...", "columns": ["<totals tokens>"]}
        },
        "semantics": "<filter vs pass-through explanation>",
        "parameters": {
          "required": [{"name": "...", "type": "date|string"}],
          "optional": [{"name": "...", "type": "string"}]
        },
        "transformations": ["<reshape rules in plain English>"]
      },
      "contract": {
        "tokens": { "scalars": [...], "row_tokens": [...], "totals": [...] },
        "mapping": { "<token>": "<TABLE.COLUMN | PARAM:name | UNRESOLVED>" },
        "unresolved": [],
        "join": { "parent_table": "...", "parent_key": "...", "child_table": "...", "child_key": "..." },
        "date_columns": { "<table>": "<date_column>" },
        "filters": { "optional": { "<name>": "table.column" } },
        "reshape_rules": [
          { "purpose": "...", "strategy": "UNION_ALL|MELT|NONE", "columns": [{"as": "alias", "from": ["table.col1", "..."]}] }
        ],
        "row_computed": { "<token>": {"op": "...", "left": "...", "right": "..."} },
        "totals_math": { "<token>": {"op": "...", "column": "..."} },
        "totals": { "<token>": {"op": "...", "column": "..."} },
        "formatters": { "<token>": "<format spec>" },
        "order_by": { "rows": ["<column ASC|DESC>"] },
        "header_tokens": ["<scalar tokens copy>"],
        "row_tokens": ["<row tokens copy>"],
        "row_order": ["<column ASC|DESC>"],
        "literals": {},
        "notes": "<domain notes>"
      },
      "validation": {
        "unknown_tokens": [],
        "unknown_columns": [],
        "token_coverage": { "scalars_mapped_pct": 100, "row_tokens_mapped_pct": 100, "totals_mapped_pct": 100 }
      }
    }

    SELF-CHECK before responding:
    - NO SQL expressions anywhere (no SUM(), CASE, CONCAT, STRFTIME, etc.).
    - Every schema token appears in contract.mapping.
    - Every column reference exists in the CATALOG.
    - row_computed and totals_math values are ALL declarative op dicts, never strings.
    - order_by.rows and row_order are identical non-empty arrays.
    - join block has all four non-empty string fields.
    - Every reshape rule has a non-empty "purpose".
    - token_coverage is 100%.
    """
).strip()


def build_excel_llm_call_4_prompt(
    *,
    final_template_html: str,
    page_summary: str,
    schema: Mapping[str, Any] | None,
    auto_mapping_proposal: Mapping[str, Any],
    mapping_override: Mapping[str, Any] | None,
    user_instructions: str,
    catalog: Iterable[str],
    key_tokens: Iterable[str] | None = None,
    dialect_hint: str | None = None,  # kept for call-site compat, ignored
) -> Dict[str, Any]:
    key_tokens_list: list[str] = []
    if key_tokens:
        seen: set[str] = set()
        for token in key_tokens:
            text = str(token or "").strip()
            if not text or text in seen:
                continue
            seen.add(text)
            key_tokens_list.append(text)

    payload: Dict[str, Any] = {
        "final_template_html": final_template_html,
        "page_summary": page_summary,
        "schema": dict(schema or {}),
        "auto_mapping_proposal": dict(auto_mapping_proposal or {}),
        "mapping_override": dict(mapping_override or {}),
        "user_instructions": user_instructions or "",
        "catalog": [str(item) for item in catalog],
    }
    if key_tokens_list:
        payload["key_tokens"] = key_tokens_list

    payload_json = json.dumps(payload, ensure_ascii=False, indent=2)
    messages = [
        {
            "role": "user",
            "content": [
                {
                    "type": "text",
                    "text": payload_json,
                }
            ],
        }
    ]

    return {
        "system": EXCEL_LLM_CALL_4_SYSTEM_PROMPT,
        "messages": messages,
        "version": EXCEL_PROMPT_VERSION_4,
    }


# ------------------------- LLM CALL 5 (Excel) -------------------------
EXCEL_LLM_CALL_5_PROMPT = {
    "system": dedent(
        """\
        LLM CALL 5 — Contract Finalizer (Excel, DataFrame Mode)
        You finalize the contract from Step 4 for the pandas DataFrame Excel report pipeline. NO SQL.

        ═══════════════════════════════════════════════════════════════
        CRITICAL: NEVER emit SQL expressions, DuckDB functions, or any code.
        All computations use declarative operation objects only.
        ═══════════════════════════════════════════════════════════════

        YOUR JOB:
        1. Copy the Step-4 contract exactly (same tokens, same declarative ops, same ordering).
        2. Validate and fill in any missing optional fields with sensible defaults.
        3. Ensure the contract is complete and ready for the DataFrame pipeline.

        RULES:
        - Treat `step4_output.contract` as authoritative. Do not add, drop, or rename tokens.
        - `mapping` values: only "TABLE.COLUMN", "PARAM:name", or "UNRESOLVED".
        - `row_computed` / `totals_math`: declarative op dicts only:
          {"op": "subtract|add|multiply|divide|sum|mean|count|min|max|concat|format_date|format_number", ...}
        - `totals` (totals_mapping): declarative op dicts mirroring totals_math.
        - Join block: non-empty parent_table/parent_key. If no child table, reuse parent.
        - `order_by.rows` and `row_order`: both non-empty arrays, identical content. Default ["ROWID"].
        - Every reshape rule must have a non-empty "purpose" (≤15 words).
        - `header_tokens`: copy of tokens.scalars. `row_tokens`: copy of tokens.row_tokens.

        OUTPUT — return ONLY this JSON object, no markdown fences, no commentary:
        {
          "contract": {
            "tokens": { "scalars": [...], "row_tokens": [...], "totals": [...] },
            "mapping": { "<token>": "TABLE.COLUMN|PARAM:name|UNRESOLVED" },
            "join": { "parent_table": "...", "parent_key": "...", "child_table": "...", "child_key": "..." },
            "date_columns": { "<table>": "<date_column>" },
            "filters": { "optional": { "<name>": "table.column" } },
            "reshape_rules": [{"purpose": "...", "strategy": "UNION_ALL|MELT|NONE", "columns": [...]}],
            "row_computed": { "<token>": {"op": "...", ...} },
            "totals_math": { "<token>": {"op": "...", ...} },
            "totals": { "<token>": {"op": "...", ...} },
            "formatters": { "<token>": "<format spec>" },
            "order_by": { "rows": ["<column ASC|DESC>"] },
            "header_tokens": [...],
            "row_tokens": [...],
            "row_order": ["<column ASC|DESC>"],
            "literals": {},
            "notes": "..."
          },
          "invalid": false
        }

        SELF-CHECK:
        - NO SQL expressions anywhere.
        - Every token from Step-4 is present in mapping.
        - row_computed and totals_math are ALL declarative op dicts.
        - order_by.rows and row_order are identical non-empty arrays.
        - join block has all four non-empty string fields.
        """
    ).strip(),
    "user": dedent(
        """\
        {
          "final_template_html": "<HTML from Step 3.5>",
          "step4_output": {
            "contract": { /* Step-4 contract object */ },
            "overview_md": "Step-4 overview",
            "step5_requirements": { /* Step-4 requirements */ }
          },
          "key_tokens": ["param_a", "param_b"]
        }
        """
    ).strip(),
}


@lru_cache(maxsize=1)
def get_excel_prompt_generator_assets() -> Dict[str, str]:
    """Return the Excel-specific system/user templates for LLM CALL 5."""
    return dict(EXCEL_LLM_CALL_5_PROMPT)


__all__ = [
    "EXCEL_PROMPT_VERSION",
    "EXCEL_PROMPT_VERSION_3_5",
    "EXCEL_PROMPT_VERSION_4",
    "EXCEL_PROMPT_VERSION_5",
    "build_excel_llm_call_1_prompt",
    "build_excel_llm_call_3_prompt",
    "build_excel_llm_call_3_5_prompt",
    "build_excel_llm_call_4_prompt",
    "get_excel_prompt_generator_assets",
]


# mypy: ignore-errors

from typing import Any, Iterable, Mapping, Sequence

CHART_SUGGEST_PROMPT_VERSION = "chart_suggestions_v1"

# Small catalog of reusable chart templates that the LLM can reference via chartTemplateId.
CHART_TEMPLATE_CATALOG: dict[str, dict[str, Any]] = {
    "time_series_basic": {
        "id": "time_series_basic",
        "description": "Trend over an ordered index (typically time or batch_index) for one or two numeric metrics.",
        "recommended_chart_type": "line",
        "recommended_use": "Use when xField is an ordered index such as time or batch_index and yFields are numeric measures.",
        "recharts": {
            "component": "LineChart",
            "props": {
                "margin": {"top": 8, "right": 16, "bottom": 24, "left": 0},
                "cartesianGrid": {"strokeDasharray": "3 3"},
            },
        },
    },
    "top_n_categories": {
        "id": "top_n_categories",
        "description": "Ranked comparison of the largest categories by a numeric metric (e.g. rows).",
        "recommended_chart_type": "bar",
        "recommended_use": "Use when xField is a categorical label and yFields contains a single numeric metric you want to rank by size.",
        "recharts": {
            "component": "BarChart",
            "props": {
                "layout": "vertical",
                "margin": {"top": 8, "right": 16, "bottom": 16, "left": 0},
            },
        },
    },
    "distribution_histogram": {
        "id": "distribution_histogram",
        "description": "Histogram-style distribution of a numeric metric, approximated with a bar chart.",
        "recommended_chart_type": "bar",
        "recommended_use": "Use when xField is a numeric metric and you conceptually bucket values into ranges to show their distribution.",
        "recharts": {
            "component": "BarChart",
            "props": {
                "margin": {"top": 8, "right": 16, "bottom": 24, "left": 0},
                "cartesianGrid": {"strokeDasharray": "3 3"},
            },
        },
    },
}


def _to_pretty_json(value: Any) -> str:
    try:
        return json.dumps(value, ensure_ascii=False, indent=2, sort_keys=isinstance(value, Mapping))
    except Exception:
        return json.dumps(str(value), ensure_ascii=False)


CHART_SUGGEST_PROMPT_TEMPLATE = dedent(
    """
    You are an analytics assistant helping a user explore report batch discovery data in NeuraReport.

    DATA CONTEXT
    - Each row in the dataset represents a single batch for the selected template and date range.
    - Available fields are described in FIELD_CATALOG_JSON below. You MUST treat those field names as the only columns
      you can reference in charts (for xField, yFields, and groupField).
    - DATA_STATS_JSON provides basic statistics over the dataset (counts, totals, min/max/avg) so you can prioritise
      interesting views.
    - KEY_FILTERS_JSON describes any key token filters that have been applied.

    CHART SPEC CONTRACT
    You must return a single JSON object with this exact shape:
      {{
        "charts": [
          {{
            "id": "short_unique_id",
            "type": "bar" | "line" | "pie" | "scatter",
            "xField": "<field name from FIELD_CATALOG_JSON>",
            "yFields": ["<field name>", "..."],
            "groupField": "<field name>" | null,
            "aggregation": "sum" | "avg" | "count" | "none" | null,
            "chartTemplateId": "time_series_basic" | "top_n_categories" | "distribution_histogram" | null,
            "title": "Concise chart title",
            "description": "Short explanation of what the chart shows and why it is useful"
          }},
          ...
        ]
      }}

    RULES
    - Propose between 2 and 5 charts that best answer the user's question while remaining faithful to the available fields.
    - Prefer highlighting metrics that show strong variation (e.g., largest numeric totals, widest min/max range) using DATA_STATS_JSON.
    - Use only field names that appear in FIELD_CATALOG_JSON for xField, yFields, and groupField.
    - Prefer using chartTemplateId values as follows when they fit:
        * "time_series_basic": xField is an ordered index (e.g. "batch_index") and yFields are numeric metrics such as "rows".
        * "top_n_categories": xField is categorical (e.g. "batch_id") and yFields contains a single numeric metric to compare.
        * "distribution_histogram": xField is a numeric metric and you conceptually bucket values into ranges to show distribution.
      It is still valid to omit chartTemplateId for free-form charts.
    - For "pie" charts, use xField as the category label and yFields[0] as the numeric value.
    - For "scatter" charts, use xField as the numeric/ordered axis and yFields[0] as the numeric dependent variable.
    - If the question references measures that are not present in FIELD_CATALOG_JSON, fall back to useful, honest charts
      based on the available fields (e.g., largest batches by rows, relationship between parent and child rows, distributions).
    - Do NOT include any commentary, markdown code fences, or extra top-level keys; return only the JSON object.

    TEMPLATE_ID: {template_id}
    TEMPLATE_KIND: {template_kind}
    DATE_RANGE:
      start_date: {start_date}
      end_date: {end_date}

    KEY_FILTERS_JSON:
    {key_values_json}

    FIELD_CATALOG_JSON:
    {field_catalog_json}

    DATA_STATS_JSON:
    {data_stats_json}

    CHART_TEMPLATE_CATALOG_JSON:
    {template_catalog_json}

    USER_QUESTION:
    {user_question}
    """
).strip()


def build_chart_suggestions_prompt(
    *,
    template_id: str,
    kind: str,
    start_date: str,
    end_date: str,
    key_values: Mapping[str, Any] | None,
    field_catalog: Iterable[Mapping[str, Any]] | Sequence[Mapping[str, Any]],
    data_stats: Mapping[str, Any] | None = None,
    question: str | None = None,
) -> str:
    key_values_json = _to_pretty_json(key_values or {})
    field_catalog_json = _to_pretty_json(list(field_catalog or []))
    data_stats_json = _to_pretty_json(data_stats or {})
    template_catalog_json = _to_pretty_json(CHART_TEMPLATE_CATALOG)
    user_question = (question or "").strip() or "Suggest several informative charts using the available fields."

    prompt = CHART_SUGGEST_PROMPT_TEMPLATE
    prompt = prompt.replace("{template_id}", template_id)
    prompt = prompt.replace("{template_kind}", (kind or "pdf").lower())
    prompt = prompt.replace("{start_date}", start_date)
    prompt = prompt.replace("{end_date}", end_date)
    prompt = prompt.replace("{key_values_json}", key_values_json)
    prompt = prompt.replace("{field_catalog_json}", field_catalog_json)
    prompt = prompt.replace("{data_stats_json}", data_stats_json)
    prompt = prompt.replace("{template_catalog_json}", template_catalog_json)
    prompt = prompt.replace("{user_question}", user_question)
    return prompt


# mypy: ignore-errors
"""LLM prompts for document analysis and data extraction."""

from typing import Any

logger = logging.getLogger("neura.prompts.analysis")

DOCUMENT_ANALYSIS_PROMPT = """You are a data extraction specialist analyzing documents for NeuraReport.

DOCUMENT CONTEXT:
- Document type: {document_type}
- File name: {file_name}
- Pages/Sheets: {page_count}

EXTRACTED CONTENT:
{extracted_content}

TASK:
Extract all meaningful data from this document and return a structured JSON response.

OUTPUT FORMAT (return ONLY valid JSON, no markdown or commentary):
{{
  "summary": "Brief 1-2 sentence description of document contents",
  "tables": [
    {{
      "id": "table_1",
      "title": "Descriptive name for this table",
      "headers": ["Column1", "Column2", "Column3"],
      "rows": [
        ["value1", "value2", "value3"],
        ["value4", "value5", "value6"]
      ],
      "data_types": ["text", "numeric", "date"]
    }}
  ],
  "key_metrics": [
    {{
      "name": "metric_name",
      "value": 123.45,
      "unit": "USD" | "%" | "units" | null,
      "context": "Where this value appears or what it represents"
    }}
  ],
  "time_series_candidates": [
    {{
      "date_column": "column_name_with_dates",
      "value_columns": ["numeric_col1", "numeric_col2"],
      "frequency": "daily" | "weekly" | "monthly" | "yearly" | null,
      "table_id": "table_1"
    }}
  ],
  "chart_recommendations": [
    {{
      "type": "line" | "bar" | "pie" | "scatter",
      "title": "Suggested chart title",
      "x_field": "field_for_x_axis",
      "y_fields": ["field1", "field2"],
      "rationale": "Brief explanation of why this chart is useful"
    }}
  ]
}}

RULES:
1. Extract ALL tables from the document, even small ones
2. Identify numeric columns vs text/categorical columns
3. Detect date/time patterns for time series analysis
4. Suggest meaningful chart visualizations based on the data
5. Preserve original precision for numeric values
6. Use "numeric" for data_types when column contains numbers
7. Use "date" or "datetime" for date columns
8. Use "text" or "category" for text/categorical columns
9. Return ONLY valid JSON - no markdown code fences, no explanatory text
"""

CHART_SUGGESTION_PROMPT = """Based on the extracted data below, suggest the best visualizations.

DATA SUMMARY:
{data_summary}

FIELD CATALOG:
{field_catalog}

USER QUESTION (if provided):
{user_question}

Suggest 2-5 charts that would best visualize this data. Focus on:
- Time series trends if date fields exist
- Comparisons between categories
- Distributions of numeric values
- Relationships between numeric fields

Return ONLY valid JSON in this format:
{{
  "charts": [
    {{
      "id": "chart_1",
      "type": "line" | "bar" | "pie" | "scatter",
      "title": "Chart title",
      "xField": "field_name_for_x",
      "yFields": ["field1", "field2"],
      "groupField": "optional_grouping_field" | null,
      "aggregation": "sum" | "avg" | "count" | "none" | null,
      "description": "What insight this chart provides"
    }}
  ]
}}
"""




def _extract_json_object(text: str) -> str | None:
    """
    Extract the first complete JSON object from text.
    Uses bracket counting to find matching braces.
    """
    if not text:
        return None

    # Find the first opening brace
    start = text.find('{')
    if start == -1:
        return None

    # Count braces to find matching close
    depth = 0
    in_string = False
    escape_next = False

    for i, char in enumerate(text[start:], start):
        if escape_next:
            escape_next = False
            continue

        if char == '\\' and in_string:
            escape_next = True
            continue

        if char == '"' and not escape_next:
            in_string = not in_string
            continue

        if in_string:
            continue

        if char == '{':
            depth += 1
        elif char == '}':
            depth -= 1
            if depth == 0:
                return text[start:i + 1]

    # If we get here, braces are unbalanced - try the greedy approach as fallback
    logger.warning("JSON extraction: unbalanced braces, attempting recovery")
    return None


def _try_repair_json(text: str) -> dict[str, Any] | None:
    """Attempt to repair common JSON issues from LLM output."""
    if not text:
        return None

    working = text

    # Common LLM JSON issues
    repairs = [
        # Trailing commas before closing brackets
        (r',\s*([\]}])', r'\1'),
        # Single quotes instead of double
        (r"'([^']*)':", r'"\1":'),
        # Missing quotes around keys
        (r'([{,]\s*)(\w+)(\s*:)', r'\1"\2"\3'),
        # JavaScript-style comments
        (r'//[^\n]*\n', '\n'),
        (r'/\*[\s\S]*?\*/', ''),
    ]

    for pattern, replacement in repairs:
        working = re.sub(pattern, replacement, working)

    try:
        return json.loads(working)
    except json.JSONDecodeError:
        return None


def parse_analysis_response(raw_response: str) -> dict[str, Any]:
    """Parse LLM response into structured analysis data."""
    default_response = {
        "summary": "",
        "tables": [],
        "key_metrics": [],
        "time_series_candidates": [],
        "chart_recommendations": [],
    }

    if not raw_response:
        return default_response

    cleaned = strip_code_fences(raw_response)

    # Try direct parse first
    try:
        data = json.loads(cleaned)
        if isinstance(data, dict):
            return {
                "summary": data.get("summary", ""),
                "tables": data.get("tables", []),
                "key_metrics": data.get("key_metrics", []),
                "time_series_candidates": data.get("time_series_candidates", []),
                "chart_recommendations": data.get("chart_recommendations", []),
            }
    except json.JSONDecodeError:
        pass

    # Try to extract JSON object using bracket counting
    json_str = _extract_json_object(cleaned)
    if json_str:
        try:
            data = json.loads(json_str)
            if isinstance(data, dict):
                return {
                    "summary": data.get("summary", ""),
                    "tables": data.get("tables", []),
                    "key_metrics": data.get("key_metrics", []),
                    "time_series_candidates": data.get("time_series_candidates", []),
                    "chart_recommendations": data.get("chart_recommendations", []),
                }
        except json.JSONDecodeError:
            pass

    # Try JSON repair
    repaired = _try_repair_json(cleaned)
    if repaired:
        return {
            "summary": repaired.get("summary", ""),
            "tables": repaired.get("tables", []),
            "key_metrics": repaired.get("key_metrics", []),
            "time_series_candidates": repaired.get("time_series_candidates", []),
            "chart_recommendations": repaired.get("chart_recommendations", []),
        }

    logger.warning("Failed to parse LLM analysis response as JSON")
    return default_response


def build_analysis_prompt(
    document_type: str,
    file_name: str,
    page_count: int,
    extracted_content: str,
) -> str:
    """Build the document analysis prompt."""
    return DOCUMENT_ANALYSIS_PROMPT.format(
        document_type=document_type,
        file_name=file_name,
        page_count=page_count,
        extracted_content=extracted_content,
    )


def build_chart_suggestion_prompt(
    data_summary: str,
    field_catalog: str,
    user_question: str | None = None,
) -> str:
    """Build the chart suggestion prompt."""
    return CHART_SUGGESTION_PROMPT.format(
        data_summary=data_summary,
        field_catalog=field_catalog,
        user_question=user_question or "No specific question provided - suggest generally useful charts.",
    )


def infer_data_type(values: list[Any]) -> str:
    """
    Infer the data type from a list of sample values.
    Requires 70%+ matches for numeric/datetime classification.
    """
    if not values:
        return "text"

    date_patterns = [
        r"^\d{4}-\d{2}-\d{2}",  # ISO format: 2024-01-15
        r"^\d{2}/\d{2}/\d{4}",  # US format: 01/15/2024
        r"^\d{2}-\d{2}-\d{4}",  # EU format: 15-01-2024
        r"^\d{1,2}/\d{1,2}/\d{2,4}",  # Flexible: 1/5/24
        r"^\d{4}/\d{2}/\d{2}",  # YYYY/MM/DD
    ]

    numeric_count = 0
    date_count = 0
    total_valid = 0

    for val in values[:30]:  # Sample up to 30 values
        if val is None:
            continue
        str_val = str(val).strip()
        if not str_val:
            continue

        total_valid += 1

        # Check for date patterns
        is_date = False
        for pattern in date_patterns:
            if re.match(pattern, str_val):
                date_count += 1
                is_date = True
                break

        if is_date:
            continue

        # Check for numeric values
        try:
            # Handle common currency/percentage formats
            cleaned = str_val.replace(",", "").replace("$", "").replace("€", "").replace("£", "")
            cleaned = cleaned.replace("%", "").replace(" ", "").strip()
            if cleaned:
                float(cleaned)
                numeric_count += 1
        except (ValueError, TypeError):
            pass

    if total_valid == 0:
        return "text"

    # Require 70%+ matches for type classification
    date_ratio = date_count / total_valid
    numeric_ratio = numeric_count / total_valid

    if date_ratio >= 0.7:
        return "datetime"
    if numeric_ratio >= 0.7:
        return "numeric"

    return "text"


__all__ = [
    "DOCUMENT_ANALYSIS_PROMPT",
    "CHART_SUGGESTION_PROMPT",
    "strip_code_fences",
    "parse_analysis_response",
    "build_analysis_prompt",
    "build_chart_suggestion_prompt",
    "infer_data_type",
]


# PROMPTS TEMPLATES (merged from prompts_templates.py)


# Lazy imports to avoid circular dependency (utils -> llm_utils -> all_prompts -> TemplateVerify -> utils)
def _get_openai_client():
    from backend.app.services.templates import get_openai_client
    return get_openai_client()

def _call_chat_completion(*args, **kwargs):
    from backend.app.services.infra_services import call_chat_completion
    return call_chat_completion(*args, **kwargs)


logger = logging.getLogger("neura.template_recommender")

DEFAULT_MODEL = "qwen"


def _summarise_catalog(catalog: Sequence[Mapping[str, Any]]) -> list[dict]:
    """
    Convert the unified catalog into a compact form suitable for the LLM.

    Only include fields that are helpful for semantic matching.
    """
    summary: list[dict] = []
    for item in catalog:
        template_id = str(item.get("id") or "").strip()
        if not template_id:
            continue
        summary.append(
            {
                "id": template_id,
                "name": item.get("name") or "",
                "kind": item.get("kind") or "",
                "domain": item.get("domain") or "",
                "tags": list(item.get("tags") or []),
                "useCases": list(item.get("useCases") or []),
                "primaryMetrics": list(item.get("primaryMetrics") or []),
                "source": item.get("source") or "",
            }
        )
    return summary


def _build_messages(
    catalog: Sequence[Mapping[str, Any]],
    requirement: str,
    hints: Mapping[str, Any] | None,
    max_results: int,
) -> list[dict]:
    catalog_json = json.dumps(_summarise_catalog(catalog), ensure_ascii=False)
    hints = hints or {}
    hints_json = json.dumps(hints, ensure_ascii=False)

    system_text = (
        "You are a template recommendation engine for an automated reporting tool. "
        "Given a catalog of report templates and a user's free-text requirement, you "
        "must select the best matching templates.\n\n"
        "Return results as STRICT JSON with this shape:\n"
        "{\n"
        '  "recommendations": [\n'
        "    {\"id\": \"<template_id>\", \"explanation\": \"<short reason>\", \"score\": 0.0-1.0},\n"
        "    ... up to the requested max_results\n"
        "  ]\n"
        "}\n\n"
        "- score must be a number between 0 and 1 where higher means better match.\n"
        "- explanation must be a short, user-facing sentence fragment (no markdown).\n"
        "- Only use template IDs that appear in the catalog.\n"
        "- Use HINTS_JSON (domains, kinds, schema_snapshot, tables, etc.) to bias the ranking when relevant."
    )

    user_text = (
        "USER_REQUIREMENT:\n"
        f"{requirement.strip()}\n\n"
        f"MAX_RESULTS: {max_results}\n"
        f"HINTS_JSON: {hints_json}\n\n"
        "TEMPLATE_CATALOG_JSON:\n"
        f"{catalog_json}\n"
    )

    return [
        {"role": "system", "content": system_text},
        {"role": "user", "content": user_text},
    ]


def _coerce_score(value: Any) -> float:
    """Coerce a value to a score between 0.0 and 1.0."""
    try:
        score = float(value)
    except (TypeError, ValueError):
        return 0.0
    # Handle NaN and inf
    if score != score or score == float('inf') or score == float('-inf'):
        return 0.0
    if score < 0.0:
        return 0.0
    if score > 1.0:
        return 1.0
    return score


def _extract_json_object(text: str) -> str | None:
    """
    Extract the first complete JSON object from text.
    Uses bracket counting to find matching braces.
    """
    if not text:
        return None

    # Find the first opening brace
    start = text.find('{')
    if start == -1:
        return None

    # Count braces to find matching close
    depth = 0
    in_string = False
    escape_next = False

    for i, char in enumerate(text[start:], start):
        if escape_next:
            escape_next = False
            continue

        if char == '\\' and in_string:
            escape_next = True
            continue

        if char == '"' and not escape_next:
            in_string = not in_string
            continue

        if in_string:
            continue

        if char == '{':
            depth += 1
        elif char == '}':
            depth -= 1
            if depth == 0:
                return text[start:i + 1]

    return None


def _parse_recommendations(raw_text: str) -> list[dict]:
    text = strip_code_fences(raw_text or "").strip()
    if not text:
        return []

    # Try direct JSON first.
    try:
        payload = json.loads(text)
    except json.JSONDecodeError:
        # Attempt to extract the first complete JSON object
        json_str = _extract_json_object(text)
        if not json_str:
            return []
        try:
            payload = json.loads(json_str)
        except json.JSONDecodeError:
            logger.warning("Failed to parse template recommendation JSON")
            return []

    if not isinstance(payload, dict):
        return []
    raw_recs = payload.get("recommendations")
    if not isinstance(raw_recs, Iterable):
        return []

    # Get catalog IDs for validation (if available)
    results: list[dict] = []
    for item in raw_recs:
        if not isinstance(item, Mapping):
            continue
        template_id = str(item.get("id") or "").strip()
        if not template_id:
            continue
        explanation = str(item.get("explanation") or "").strip()
        score = _coerce_score(item.get("score"))
        results.append(
            {
                "id": template_id,
                "explanation": explanation,
                "score": score,
            }
        )
    return results


def recommend_templates_from_catalog(
    catalog: Sequence[Mapping[str, Any]],
    *,
    requirement: str,
    hints: Mapping[str, Any] | None = None,
    max_results: int = 6,
) -> List[Dict[str, Any]]:
    """
    Call the LLM to obtain a ranked list of template IDs with explanations.

    Returns a list of dicts:
        { "id": ..., "explanation": ..., "score": float }
    """
    requirement = (requirement or "").strip()
    if not requirement or not catalog:
        return []

    messages = _build_messages(catalog, requirement=requirement, hints=hints, max_results=max_results)
    client = _get_openai_client()

    try:
        response = _call_chat_completion(
            client,
            model=DEFAULT_MODEL,
            messages=messages,
            description="template_recommendations",
            temperature=0.2,
            max_tokens=512,
        )
    except Exception as exc:  # pragma: no cover - network / quota failures
        logger.warning(
            "template_recommend_llm_failed",
            extra={"event": "template_recommend_llm_failed", "error": str(exc)},
        )
        return []

    try:
        # openai v1-style response
        content = response.choices[0].message.content or ""
    except Exception:  # pragma: no cover - unexpected SDK shape
        logger.warning(
            "template_recommend_response_shape_unexpected",
            extra={"event": "template_recommend_response_shape_unexpected"},
        )
        return []

    recommendations = _parse_recommendations(content)
    if not recommendations:
        logger.info(
            "template_recommend_no_results",
            extra={"event": "template_recommend_no_results"},
        )
        return []

    # Preserve ordering from the LLM and cap to max_results.
    return recommendations[:max_results]


# mypy: ignore-errors

from typing import Any, Dict, List

TEMPLATE_CHAT_PROMPT_VERSION = "template_chat_v1"

_EXCEL_GUIDANCE = dedent(
    """\

    EXCEL TEMPLATE MODE
    This template will be converted to an Excel spreadsheet (.xlsx), NOT rendered as a PDF page.
    Follow these constraints strictly:
    - Use simple <table>-based layouts — each <table> maps directly to an Excel sheet region.
    - <tr> = rows, <td>/<th> = cells. Use colspan/rowspan for cell merging.
    - Avoid: CSS grid, flexbox, floats, absolute positioning, page-size constraints (no A4 dimensions), page breaks.
    - OK to use: background-color on cells (→ cell fills), bold/italic/font-size, border styles, text-align.
    - Keep styling inline and simple — complex CSS does not survive XLSX conversion.
    - Do NOT add page numbering, print-oriented headers/footers, or page-break markers — these are stripped in Excel output.
    - Prefer landscape-friendly wide table layouts over narrow portrait-style designs.
    - For repeating data rows, still use <!-- BEGIN:BLOCK_REPEAT ... --> / <!-- END:BLOCK_REPEAT --> inside the <table>.
    """
).strip()


TEMPLATE_CHAT_SYSTEM_PROMPT = dedent(
    """\
    You are an expert HTML template editing assistant working inside the NeuraReport reporting engine.
    You help users edit their report templates through an interactive conversation.

    YOUR ROLE
    - Engage in a helpful conversation to understand what changes the user wants to make to their template.
    - Ask clarifying questions when the user's request is ambiguous or incomplete.
    - When you have gathered enough information, propose the changes you will make.
    - Only apply changes when you are confident you understand the user's intent.

    TEMPLATE CONTEXT
    - The template uses dynamic tokens/placeholders like {token}, {{ token }}, {row_token}, etc.
    - Templates may contain repeat markers like <!-- BEGIN:BLOCK_REPEAT ... --> / <!-- END:BLOCK_REPEAT -->.
    - Templates include IDs, classes, data-* attributes that should be preserved unless explicitly asked to change.

    CONVERSATION GUIDELINES
    - Be conversational and helpful, but concise.
    - If the user's request is clear and complete, you can proceed to propose changes immediately.
    - If you need more information, ask specific questions (limit to 2-3 questions at a time).
    - When proposing changes, summarize what you will do before showing the result.
    - Always confirm understanding before making significant structural changes.

    WHAT TO CLARIFY
    - Vague styling requests (e.g., "make it look better" - ask what style they prefer)
    - Structural changes without clear scope (e.g., "reorganize" - ask what sections)
    - Adding new elements without context (e.g., "add a chart" - ask where and what data)
    - Changes that might affect dynamic tokens (explain the impact and confirm)

    PIPELINE AWARENESS
    You are part of a unified chat pipeline. The user's conversation may span the entire
    template lifecycle — from creation to mapping to report generation. When you detect the
    user's intent goes beyond editing (e.g. they want to map tokens, approve, or generate),
    include an "intent" field in your response so the system can route to the right handler.

    Possible intents: "edit" (default), "map", "correct", "approve", "generate", "discover",
    "status", "verify".  Only include "intent" when the user clearly wants a non-edit action.

    Pipeline context will be injected below when available:
    {{PIPELINE_CONTEXT}}

    OUTPUT FORMAT (STRICT JSON, no markdown fences, no commentary):
    {
      "message": "<string>",              // Your response message to the user
      "ready_to_apply": <boolean>,        // true if you have enough info and are ready to show changes
      "proposed_changes": ["change 1", "change 2"] | null,  // List of changes you will make (when ready_to_apply=true)
      "follow_up_questions": ["q1", "q2"] | null,          // Questions to ask (when ready_to_apply=false)
      "updated_html": "<string>" | null,  // The full updated HTML (only when ready_to_apply=true)
      "intent": "<string>" | null         // Non-edit intent detected (optional, omit if pure edit)
    }

    IMPORTANT RULES
    - When ready_to_apply=true, you MUST provide updated_html with the complete modified template.
    - When ready_to_apply=false, you MUST NOT provide updated_html.
    - proposed_changes should be short, human-readable descriptions.
    - follow_up_questions should be specific and actionable.
    - Preserve all dynamic tokens exactly unless explicitly asked to change them.
    - Maintain valid HTML structure.
    """
).strip()


def build_template_chat_prompt(
    template_html: str,
    conversation_history: List[Dict[str, str]],
    kind: str = "pdf",
) -> Dict[str, Any]:
    """
    Build a chat-completions payload for conversational template editing.

    Args:
        template_html: The current HTML template content
        conversation_history: List of messages with 'role' and 'content' keys
        kind: Template kind — 'pdf' or 'excel'

    Returns a dict with:
        {
          "system": <system_prompt>,
          "messages": [ ... ],
          "version": TEMPLATE_CHAT_PROMPT_VERSION,
        }
    """
    # Build the initial context message that includes the template
    excel_note = f"\n\n{_EXCEL_GUIDANCE}" if kind == "excel" else ""
    context_message = (
        "Here is the current template HTML that the user wants to edit:\n\n"
        "```html\n"
        f"{template_html or ''}\n"
        "```\n\n"
        "The user will now describe what changes they want. "
        "Engage in a conversation to understand their needs fully before making changes."
        f"{excel_note}"
    )

    messages: List[Dict[str, Any]] = [
        {
            "role": "system",
            "content": [{"type": "text", "text": TEMPLATE_CHAT_SYSTEM_PROMPT}],
        },
        {
            "role": "user",
            "content": [{"type": "text", "text": context_message}],
        },
        {
            "role": "assistant",
            "content": [
                {
                    "type": "text",
                    "text": json.dumps(
                        {
                            "message": "I've reviewed your template. What changes would you like to make? Feel free to describe what you want - whether it's styling updates, layout changes, adding or removing sections, or any other modifications.",
                            "ready_to_apply": False,
                            "proposed_changes": None,
                            "follow_up_questions": None,
                            "updated_html": None,
                        },
                        ensure_ascii=False,
                    ),
                }
            ],
        },
    ]

    # Add the conversation history, skipping any leading assistant messages
    # (we already injected the welcome message above)
    history_started = False
    for msg in conversation_history:
        role = msg.get("role", "user")
        content = msg.get("content", "")
        if not history_started and role == "assistant":
            # Skip the initial assistant welcome that duplicates our injected one
            history_started = True
            continue
        history_started = True
        if role in ("user", "assistant"):
            messages.append(
                {
                    "role": role,
                    "content": [{"type": "text", "text": content}],
                }
            )

    return {
        "system": TEMPLATE_CHAT_SYSTEM_PROMPT,
        "messages": messages,
        "version": TEMPLATE_CHAT_PROMPT_VERSION,
    }


TEMPLATE_CHAT_CREATE_PROMPT_VERSION = "template_chat_create_v1"


TEMPLATE_CHAT_CREATE_SYSTEM_PROMPT = dedent(
    """\
    You are an expert HTML template creation assistant working inside the NeuraReport reporting engine.
    You help users build report templates from scratch through an interactive, multi-step conversation.

    YOUR ROLE
    - Guide the user through creating a report template step by step.
    - Have a genuine back-and-forth conversation — do NOT generate a template on the first response.
    - Ask clarifying questions to fully understand the report before generating anything.
    - Iterate on the template based on user feedback.

    TEMPLATE CAPABILITIES
    - Templates use dynamic tokens/placeholders like {token_name} for single values and {row_token_name} for repeating data.
    - Repeating rows use markers: <!-- BEGIN:BLOCK_REPEAT data_source --> / <!-- END:BLOCK_REPEAT -->.
    - Templates should be self-contained HTML with inline CSS for reliable rendering.
    - Use professional, clean styling appropriate for business reports.

    MANDATORY CONVERSATION PHASES (follow this order):

    PHASE 1 — Report Understanding (at least 1 exchange)
    - Ask about the type of report (invoice, receipt, summary, inventory, etc.)
    - Understand the purpose and audience
    - Ask about major sections (header, body, footer, etc.)

    PHASE 2 — Data & Field Mapping (at least 1 exchange)
    - Ask what data fields the report needs (e.g. customer name, date, amounts)
    - For tables/repeating sections: ask about column names and what data goes in each column
    - Clarify which fields are single-value headers (like company name, report date) vs repeating row data (like line items)
    - Discuss naming — confirm what the user calls each field so token names are meaningful
    - Example: "For the transaction table, what columns do you need? Things like Reference Number, Date, Amount, Status?"

    PHASE 3 — Styling & Layout (at least 1 exchange)
    - Ask about branding (colors, fonts, logo)
    - Page layout preferences (portrait/landscape)
    - Any specific styling requirements

    PHASE 4 — Template Generation (only after phases 1-3)
    - Summarize what you understood and propose the template
    - Set ready_to_apply=true ONLY after completing phases 1–3
    - List all tokens you will use and explain what each one maps to

    PHASE 5 — Post-Apply Refinement
    - After the user applies the template, ask if they want to adjust anything
    - Discuss token names — do they match the user's database column names?
    - Offer to rename tokens, add/remove columns, adjust formatting

    CRITICAL RULES FOR PACING
    - NEVER set ready_to_apply=true on the first or second exchange
    - You must complete Phases 1, 2, and 3 before generating a template
    - Each phase needs at least one user response before moving to the next
    - If the user provides all info at once, still confirm your understanding before generating
    - If the user says "ok", "sure", "yes", or similar short affirmations, continue to the NEXT phase — do NOT generate the template yet unless you have completed all phases

    OUTPUT FORMAT (STRICT JSON, no markdown fences, no commentary):
    {
      "message": "<string>",              // Your response message to the user
      "ready_to_apply": <boolean>,        // true ONLY after phases 1-3 are complete
      "proposed_changes": ["change 1", "change 2"] | null,  // List of what the template includes (when ready_to_apply=true)
      "follow_up_questions": ["q1", "q2"] | null,          // Questions to ask (when ready_to_apply=false)
      "updated_html": "<string>" | null   // The full HTML template (only when ready_to_apply=true)
    }

    SAMPLE PDF REFERENCE (if provided)
    - The user may provide a sample PDF as a visual reference image.
    - Use it to understand the desired layout, styling, colors, fonts, and structure.
    - Do NOT try to OCR or extract exact text — use it as design inspiration.
    - On the first message when a sample is provided, describe what you see and ask what the user wants to keep or change.
    - Replicate the visual layout, table structure, header/footer arrangement as closely as possible.
    - Still follow the phased approach — discuss data fields and styling before generating.

    IMPORTANT RULES
    - When ready_to_apply=true, you MUST provide updated_html with the complete HTML template.
    - When ready_to_apply=false, you MUST NOT provide updated_html.
    - proposed_changes should describe what the template includes (sections, features).
    - follow_up_questions should be specific and actionable.
    - Generate clean, professional HTML with inline styles.
    - Use meaningful placeholder tokens that match the user's data (e.g. {customer_name}, {bill_amount}, {row_reference_number}).
    - Include a proper HTML structure with <!DOCTYPE html>, <html>, <head>, and <body> tags.
    - Maintain valid HTML structure.
    """
).strip()


def build_template_chat_create_prompt(
    conversation_history: List[Dict[str, str]],
    current_html: str | None = None,
    sample_image_b64: str | None = None,
    kind: str = "pdf",
) -> Dict[str, Any]:
    """
    Build a chat-completions payload for conversational template creation.

    Args:
        conversation_history: List of messages with 'role' and 'content' keys
        current_html: Optional current HTML if template is being iterated on
        sample_image_b64: Optional base64-encoded PNG of a sample PDF for visual reference
        kind: Template kind — 'pdf' or 'excel'

    Returns a dict with:
        {
          "system": <system_prompt>,
          "messages": [ ... ],
          "version": TEMPLATE_CHAT_CREATE_PROMPT_VERSION,
        }
    """
    excel_note = f"\n\n{_EXCEL_GUIDANCE}" if kind == "excel" else ""
    if current_html and current_html.strip():
        context_message = (
            "The user is creating a new report template. Here is the current draft:\n\n"
            "```html\n"
            f"{current_html}\n"
            "```\n\n"
            "The user will now describe what they want to change or add. "
            "Help them refine the template."
            f"{excel_note}"
        )
    else:
        context_message = (
            "The user wants to create a new report template from scratch. "
            "They will describe what kind of report they need. "
            "Help them by asking the right questions and then generate a professional HTML template."
            f"{excel_note}"
        )

    if sample_image_b64:
        context_message += (
            "\n\nThe user has attached a sample PDF as a visual reference. "
            "Use the image below to understand their desired layout, styling, and structure. "
            "Describe what you see and ask what they want to keep, change, or add."
        )

    welcome_text = json.dumps(
        {
            "message": (
                "I can see your sample PDF. I'll use its layout and styling as a reference. "
                "What would you like to keep from this design, and what would you like to change?"
            ) if sample_image_b64 else (
                "I'll help you create a report template from scratch. "
                "What kind of report do you need? For example: invoice, sales summary, "
                "inventory report, financial statement, or something else?"
            ),
            "ready_to_apply": False,
            "proposed_changes": None,
            "follow_up_questions": None,
            "updated_html": None,
        },
        ensure_ascii=False,
    )

    # Build the initial user context content blocks
    user_content: List[Dict[str, Any]] = [{"type": "text", "text": context_message}]
    if sample_image_b64:
        user_content.append({
            "type": "image_url",
            "image_url": {"url": f"data:image/png;base64,{sample_image_b64}"},
        })

    messages: List[Dict[str, Any]] = [
        {
            "role": "system",
            "content": [{"type": "text", "text": TEMPLATE_CHAT_CREATE_SYSTEM_PROMPT}],
        },
        {
            "role": "user",
            "content": user_content,
        },
        {
            "role": "assistant",
            "content": [{"type": "text", "text": welcome_text}],
        },
    ]

    # Add the conversation history, skipping any leading assistant messages
    # (we already injected the welcome message above)
    history_started = False
    for msg in conversation_history:
        role = msg.get("role", "user")
        content = msg.get("content", "")
        if not history_started and role == "assistant":
            # Skip the initial assistant welcome that duplicates our injected one
            history_started = True
            continue
        history_started = True
        if role in ("user", "assistant"):
            messages.append(
                {
                    "role": role,
                    "content": [{"type": "text", "text": content}],
                }
            )

    return {
        "system": TEMPLATE_CHAT_CREATE_SYSTEM_PROMPT,
        "messages": messages,
        "version": TEMPLATE_CHAT_CREATE_PROMPT_VERSION,
    }


def inject_pipeline_context(
    system_prompt: str,
    pipeline_state: str = "empty",
    completed_stages: list | None = None,
    needs_reapproval: bool = False,
    connection_id: str | None = None,
) -> str:
    """Inject pipeline state context into the system prompt."""
    parts = [f"Current pipeline state: {pipeline_state}"]
    if completed_stages:
        parts.append(f"Completed stages: {', '.join(completed_stages)}")
    if needs_reapproval:
        parts.append("Note: Template was edited after approval — contract needs re-approval.")
    if connection_id:
        parts.append(f"Connected to database: {connection_id}")

    context = "\n    ".join(parts)
    return system_prompt.replace("{{PIPELINE_CONTEXT}}", context)


# ======================================================================
# UNIFIED PIPELINE PROMPT (replaces Chat Edit, Chat Create, Template Edit)
# ======================================================================

UNIFIED_PIPELINE_PROMPT_VERSION = "unified_pipeline_v1"

_EXCEL_GUIDANCE_BLOCK = dedent("""\
    EXCEL MODE ACTIVE
    This template targets Excel (.xlsx), NOT PDF. Constraints:
    - Use simple <table>-based layouts only — each <table> maps to an Excel sheet region.
    - No CSS grid, flexbox, floats, absolute positioning.
    - Inline styles only: background-color, bold, italic, font-size, border, text-align.
    - No page numbering, print headers/footers, or page breaks.
    - Prefer landscape-friendly wide table layouts.
    - Use <!-- BEGIN:BLOCK_REPEAT ... --> for repeating data rows inside tables.
""").strip()

UNIFIED_PIPELINE_SYSTEM_PROMPT = dedent("""\
    You are the NeuraReport template assistant. You guide users through the entire
    template lifecycle — creation, editing, mapping, approval, validation, and generation —
    in a single continuous conversation.

    PIPELINE AWARENESS
    The user's conversation spans the full template lifecycle. The current pipeline state
    and all relevant context are provided in [PIPELINE_CONTEXT] below. Use this to know
    what is possible now, what has been done, and what the logical next step is.

    WHAT YOU CAN DO
    - Create templates from scratch (ask about report type, data fields, styling first)
    - Edit existing templates (HTML/CSS changes, add/remove sections, restyle)
    - Explain and discuss the current mapping, contract, or template structure
    - Suggest next steps based on pipeline state
    - Answer questions about the data, tokens, or report structure

    BEHAVIOR BY STATE
    - empty: Help create a template. Ask about report type, data fields, and styling
      before generating. Need at least 2-3 exchanges before producing HTML.
    - html_ready: Edit the template. Suggest mapping when a database connection exists.
    - mapped: Show/explain token→column mappings. Accept corrections via natural language.
    - approved: Contract is built. Can still edit or suggest validation.
    - validated/ready: Ready to generate reports. Help with date ranges and filters.

    TEMPLATE EDITING RULES
    - For clear requests: produce updated HTML immediately (ready_to_apply=true)
    - For ambiguous requests: ask 1-2 clarifying questions (ready_to_apply=false)
    - ALWAYS preserve {token} placeholders unless explicitly asked to change them
    - ALWAYS preserve <!-- BEGIN:BLOCK_REPEAT --> / <!-- END:BLOCK_REPEAT --> markers
    - ALWAYS preserve data-* attributes and IDs unless explicitly asked
    - When ready_to_apply=true, provide the COMPLETE HTML (not a diff)

    TEMPLATE CREATION RULES
    When creating from scratch (empty state, no existing template):
    1. Ask about report type and purpose (1 exchange minimum)
    2. Ask about data fields — which are single values vs repeating rows (1 exchange minimum)
    3. Ask about styling — colors, fonts, layout (1 exchange minimum)
    4. Only THEN generate HTML (set ready_to_apply=true)
    Do NOT generate HTML on the first or second exchange.

    TOKEN NAMING
    - Scalar headers: lowercase_snake_case (report_title, company_name, from_date)
    - Row data: row_ prefix (row_date_time, row_amount, row_quantity)
    - Totals: total_ prefix (total_quantity, total_amount)

    INTENT DETECTION
    When the user wants a pipeline action beyond conversation, include "intent":
    - "verify": upload/convert a file
    - "map": auto-map tokens to database columns
    - "correct": fix specific mappings
    - "approve": accept the mapping and build contract
    - "validate": run dry-run and checks
    - "generate": produce reports
    - "discover": see available batches/date ranges
    - "status": ask where we are in the pipeline

    [PIPELINE_CONTEXT]
    {{PIPELINE_CONTEXT}}

    OUTPUT FORMAT (strict JSON, no markdown fences, no commentary outside JSON):
    {
      "message": "<your response to the user>",
      "ready_to_apply": <true if you have complete HTML to propose>,
      "proposed_changes": ["change 1", "change 2"] | null,
      "follow_up_questions": ["q1", "q2"] | null,
      "updated_html": "<complete HTML document>" | null,
      "intent": "<pipeline action>" | null
    }

    RULES:
    - ready_to_apply=true REQUIRES updated_html with complete HTML
    - ready_to_apply=false MUST NOT include updated_html
    - intent is optional — only when user clearly wants a non-edit action
    - Be conversational, concise, and helpful
    - Reference actual token names, column names, table names from the context
""").strip()


def build_unified_pipeline_prompt(
    conversation_history: List[Dict[str, str]],
    pipeline_context: str,
    current_html: Optional[str] = None,
    sample_image_b64: Optional[str] = None,
    kind: str = "pdf",
) -> Dict[str, Any]:
    """
    Build the unified pipeline chat prompt.

    Replaces: build_template_chat_prompt, build_template_chat_create_prompt,
    build_template_edit_prompt.
    """
    # Inject pipeline context and optional excel guidance
    system = UNIFIED_PIPELINE_SYSTEM_PROMPT.replace("{{PIPELINE_CONTEXT}}", pipeline_context or "No context available.")
    if kind == "excel":
        system += f"\n\n{_EXCEL_GUIDANCE_BLOCK}"

    # Build initial context message
    if current_html and current_html.strip():
        context_msg = (
            "Here is the current template HTML:\n\n"
            f"```html\n{current_html}\n```\n\n"
            "The user will describe what they want. Help them."
        )
    else:
        context_msg = (
            "The user wants to create or work on a report template. "
            "Help them by understanding their needs and guiding them through the process."
        )

    # Build messages array
    messages: List[Dict[str, Any]] = [
        {"role": "system", "content": [{"type": "text", "text": system}]},
        {"role": "user", "content": [{"type": "text", "text": context_msg}]},
        {"role": "assistant", "content": [{"type": "text", "text": json.dumps({
            "message": "I'm ready to help with your report template. What would you like to do?",
            "ready_to_apply": False,
            "proposed_changes": None,
            "follow_up_questions": None,
            "updated_html": None,
            "intent": None,
        }, ensure_ascii=False)}]},
    ]

    # Add sample PDF image if provided
    if sample_image_b64:
        messages[1]["content"].append({
            "type": "image_url",
            "image_url": {"url": f"data:image/png;base64,{sample_image_b64}"},
        })

    # Add conversation history (skip initial assistant welcome)
    history_started = False
    for msg in conversation_history:
        role = msg.get("role", "user")
        content = msg.get("content", "")
        if not history_started and role == "assistant":
            history_started = True
            continue
        history_started = True
        if role in ("user", "assistant"):
            messages.append({
                "role": role,
                "content": [{"type": "text", "text": content}],
            })

    return {
        "system": system,
        "messages": messages,
        "version": UNIFIED_PIPELINE_PROMPT_VERSION,
    }


__all__ = [
    "TEMPLATE_CHAT_PROMPT_VERSION",
    "TEMPLATE_CHAT_SYSTEM_PROMPT",
    "build_template_chat_prompt",
    "TEMPLATE_CHAT_CREATE_PROMPT_VERSION",
    "TEMPLATE_CHAT_CREATE_SYSTEM_PROMPT",
    "build_template_chat_create_prompt",
    "inject_pipeline_context",
    "UNIFIED_PIPELINE_PROMPT_VERSION",
    "UNIFIED_PIPELINE_SYSTEM_PROMPT",
    "build_unified_pipeline_prompt",
]


# mypy: ignore-errors


TEMPLATE_EDIT_PROMPT_VERSION = "template_edit_v1"


TEMPLATE_EDIT_SYSTEM_PROMPT = dedent(
    """\
    You are an expert HTML template editor working inside the NeuraReport reporting engine.

    GOAL
    - Apply the user's natural-language instructions to the existing report template HTML.
    - Make only the requested changes; do not redesign the template unless explicitly asked.

    CONSTRAINTS
    - Preserve all dynamic tokens/placeholders exactly as written (examples: {token}, {{ token }}, {row_token}, etc.).
    - Preserve repeat markers and structural markers such as <!-- BEGIN:BLOCK_REPEAT ... --> / <!-- END:BLOCK_REPEAT -->.
    - Do not remove or rename IDs, classes, data-* attributes, or comments that look like implementation markers
      unless the user explicitly asks.
    - Keep the HTML self-contained: no external CSS/JS, no <script> tags, no external URLs.

    EDITING BEHAVIOUR
    - Work in-place on the provided HTML.
    - Prefer minimal structural changes: adjust text, styles, and small layout tweaks unless the instructions clearly
      request bigger changes.
    - If an instruction conflicts with token semantics (for example, replacing a token with fixed text) then only do so
      when the user explicitly asks.
    - Maintain valid HTML.

    OUTPUT FORMAT (STRICT JSON, no markdown fences, no commentary):
    {
      "updated_html": "<string>",          // full HTML after applying the instructions
      "summary": ["change 1", "change 2"]  // short, human-readable descriptions of the main changes
    }

    - Always return BOTH keys.
    - Ensure JSON is valid UTF-8 and properly escaped.
    """
).strip()


def build_template_edit_prompt(template_html: str, instructions: str, kind: str = "pdf") -> Dict[str, Any]:
    """
    Build a chat-completions payload for editing an existing template HTML using natural-language instructions.

    Args:
        template_html: The current HTML template content
        instructions: Natural-language editing instructions
        kind: Template kind — 'pdf' or 'excel'

    Returns a dict with:
        {
          "system": <system_prompt>,
          "messages": [ ... ],
          "version": TEMPLATE_EDIT_PROMPT_VERSION,
        }
    """
    system_prompt = TEMPLATE_EDIT_SYSTEM_PROMPT
    if kind == "excel":
        # _EXCEL_GUIDANCE is defined above in this file
        system_prompt = system_prompt + "\n\n" + _EXCEL_GUIDANCE

    payload: Dict[str, Any] = {
        "template_html": template_html or "",
        "instructions": instructions or "",
    }
    payload_json = json.dumps(payload, ensure_ascii=False, indent=2)
    user_text = "Apply the instructions in this JSON payload:\n" + payload_json

    messages: List[Dict[str, Any]] = [
        {
            "role": "system",
            "content": [{"type": "text", "text": system_prompt}],
        },
        {
            "role": "user",
            "content": [{"type": "text", "text": user_text}],
        },
    ]

    return {
        "system": system_prompt,
        "messages": messages,
        "version": TEMPLATE_EDIT_PROMPT_VERSION,
    }


__all__ = [
    "TEMPLATE_EDIT_PROMPT_VERSION",
    "TEMPLATE_EDIT_SYSTEM_PROMPT",
    "build_template_edit_prompt",
]


# PROMPT_REGISTRY

"""
Prompt versioning and registry service.

Stores prompts as YAML templates with:
- Version tracking and rollback
- Jinja2 template rendering
- A/B testing support via variant weights
- Runtime prompt serving to agents

Based on: bigscience-workshop/promptsource + promptslab/Promptify patterns.
"""
from typing import Any, Optional

import yaml
from jinja2 import Template

logger = logging.getLogger("neura.prompts.registry")


@dataclass
class PromptVersion:
    """A single version of a prompt template."""
    version: str
    template: str
    description: str = ""
    model: str = "qwen"
    temperature: float = 0.7
    max_tokens: int = 1024
    metadata: dict[str, Any] = field(default_factory=dict)

    def render(self, **kwargs) -> str:
        """Render the prompt template with variables."""
        return Template(self.template).render(**kwargs)


@dataclass
class PromptEntry:
    """A prompt with multiple versions."""
    name: str
    category: str
    description: str = ""
    versions: dict[str, PromptVersion] = field(default_factory=dict)
    active_version: str = "v1"

    @property
    def active(self) -> Optional[PromptVersion]:
        return self.versions.get(self.active_version)

    def render(self, version: Optional[str] = None, **kwargs) -> str:
        """Render the active (or specified) version of the prompt."""
        v = self.versions.get(version or self.active_version)
        if v is None:
            raise ValueError(f"Prompt version '{version or self.active_version}' not found for '{self.name}'")
        return v.render(**kwargs)


class PromptRegistry:
    """
    Central registry for versioned prompt templates.

    Loads prompts from YAML files in a directory structure:
        prompts/
            analysis/
                document_analysis.yaml
                chart_suggestion.yaml
            generation/
                report_generation.yaml
            agents/
                research.yaml
    """

    def __init__(self, prompts_dir: Optional[str] = None):
        self._prompts: dict[str, PromptEntry] = {}
        self._prompts_dir = prompts_dir or str(
            Path(__file__).parent / "registry"
        )

    def load_from_directory(self, directory: Optional[str] = None) -> int:
        """Load all prompt YAML files from the directory tree."""
        base = Path(directory or self._prompts_dir)
        if not base.exists():
            logger.warning("prompts_dir_missing", extra={"event": "prompts_dir_missing", "path": str(base)})
            return 0

        count = 0
        for yaml_file in base.rglob("*.yaml"):
            try:
                with open(yaml_file, "r", encoding="utf-8") as f:
                    data = yaml.safe_load(f)
                if not isinstance(data, dict):
                    continue

                name = data.get("name", yaml_file.stem)
                category = data.get("category", yaml_file.parent.name)

                entry = PromptEntry(
                    name=name,
                    category=category,
                    description=data.get("description", ""),
                    active_version=data.get("active_version", "v1"),
                )

                for ver_key, ver_data in data.get("versions", {}).items():
                    entry.versions[ver_key] = PromptVersion(
                        version=ver_key,
                        template=ver_data.get("template", ""),
                        description=ver_data.get("description", ""),
                        model=ver_data.get("model", "qwen"),
                        temperature=ver_data.get("temperature", 0.7),
                        max_tokens=ver_data.get("max_tokens", 1024),
                        metadata=ver_data.get("metadata", {}),
                    )

                self._prompts[name] = entry
                count += 1
            except Exception as exc:
                logger.warning("prompt_load_failed", extra={"event": "prompt_load_failed", "file": str(yaml_file), "error": str(exc)})

        logger.info("prompts_loaded", extra={"event": "prompts_loaded", "count": count})
        return count

    def register(self, name: str, template: str, version: str = "v1", **kwargs) -> None:
        """Register a prompt programmatically."""
        if name not in self._prompts:
            self._prompts[name] = PromptEntry(name=name, category=kwargs.get("category", "default"))
        self._prompts[name].versions[version] = PromptVersion(
            version=version, template=template, **{k: v for k, v in kwargs.items() if k != "category"},
        )
        if len(self._prompts[name].versions) == 1:
            self._prompts[name].active_version = version

    def get(self, name: str) -> Optional[PromptEntry]:
        return self._prompts.get(name)

    def render(self, name: str, version: Optional[str] = None, **kwargs) -> str:
        """Render a prompt by name."""
        entry = self._prompts.get(name)
        if entry is None:
            raise KeyError(f"Prompt '{name}' not found in registry")
        return entry.render(version=version, **kwargs)

    def list_prompts(self) -> list[dict[str, Any]]:
        """List all registered prompts."""
        return [
            {
                "name": e.name,
                "category": e.category,
                "description": e.description,
                "active_version": e.active_version,
                "versions": list(e.versions.keys()),
            }
            for e in sorted(self._prompts.values(), key=lambda e: e.name)
        ]


_registry: Optional[PromptRegistry] = None


def get_prompt_registry() -> PromptRegistry:
    """Get or create the singleton prompt registry."""
    global _registry
    if _registry is None:
        _registry = PromptRegistry()
        _registry.load_from_directory()
    return _registry


# =============================================================================
# RECOMMENDATION SERVICE (ported from recommendations/service.py)
# =============================================================================

class RecommendationService:
    """Service for template recommendations using AI."""

    def __init__(self):
        self._llm_client = None

    def _get_llm_client(self):
        if self._llm_client is None:
            from backend.app.services.llm import get_llm_client
            self._llm_client = get_llm_client()
        return self._llm_client

    def recommend_templates(
        self,
        connection_id: Optional[str] = None,
        schema_info: Optional[Dict[str, Any]] = None,
        context: Optional[str] = None,
        limit: int = 5,
        correlation_id: Optional[str] = None,
    ) -> List[Dict[str, Any]]:
        """Recommend templates based on connection schema and context."""
        _logger = logging.getLogger("neura.domain.recommendations")
        _logger.info("Generating template recommendations", extra={"correlation_id": correlation_id})

        from backend.app.repositories import state_store
        with state_store.transaction() as state:
            templates = state.get("templates", {})
        approved = [t for t in templates.values() if t.get("status") == "approved"]

        if not approved:
            return []

        if connection_id and not schema_info:
            try:
                from backend.app.repositories import resolve_connection_ref
                schema_info = resolve_connection_ref(connection_id)
            except Exception as e:
                _logger.warning("Failed to get schema for connection %s: %s", connection_id, e)

        template_catalog = [
            {"id": t.get("id"), "name": t.get("name"), "kind": t.get("kind"), "tags": t.get("tags", [])}
            for t in approved
        ]

        prompt = f"""Recommend templates from this catalog based on the user's needs.

TEMPLATE CATALOG:
{template_catalog}

"""
        if schema_info:
            tables = [t["name"] for t in schema_info.get("tables", [])] if isinstance(schema_info, dict) else []
            if tables:
                prompt += f"DATABASE TABLES: {', '.join(tables)}\n\n"
        if context:
            prompt += f"USER CONTEXT: {context}\n\n"

        prompt += f"""Return a JSON array of the top {limit} recommended templates:
[
  {{"template_id": "id", "score": 0.95, "reason": "Why this template matches"}}
]

Return ONLY the JSON array."""

        try:
            client = self._get_llm_client()
            response = client.complete(
                messages=[{"role": "user", "content": prompt}],
                description="template_recommendations",
                temperature=0.3,
            )
            content = response["choices"][0]["message"]["content"]
            json_match = re.search(r"\[[\s\S]*\]", content)
            if json_match:
                recommendations = json.loads(json_match.group())
                for rec in recommendations:
                    tid = rec.get("template_id")
                    if tid in templates:
                        rec["template"] = templates[tid]
                return recommendations[:limit]
        except Exception as exc:
            _logger.error(f"Recommendation generation failed: {exc}")

        sorted_templates = sorted(approved, key=lambda t: t.get("created_at", ""), reverse=True)
        return [{"template_id": t["id"], "template": t, "score": 0.5, "reason": "Recently created"} for t in sorted_templates[:limit]]

    def get_similar_templates(self, template_id: str, limit: int = 3) -> List[Dict[str, Any]]:
        """Get templates similar to a given template."""
        from backend.app.repositories import state_store
        with state_store.transaction() as state:
            templates = state.get("templates", {})
        target = templates.get(template_id)
        if not target:
            return []
        target_tags = set(target.get("tags", []))
        similar = []
        for tid, t in templates.items():
            if tid == template_id or t.get("status") != "approved":
                continue
            t_tags = set(t.get("tags", []))
            overlap = len(target_tags & t_tags)
            if overlap > 0:
                similar.append({"template": t, "score": overlap / max(len(target_tags), 1)})
        similar.sort(key=lambda x: x["score"], reverse=True)
        return similar[:limit]


# =============================================================================
# FEDERATION SERVICE (ported from federation/service.py)
# =============================================================================

class FederationService:
    """Service for cross-database federation operations."""

    def __init__(self):
        self._llm_client = None

    def _get_llm_client(self):
        if self._llm_client is None:
            from backend.app.services.llm import get_llm_client
            self._llm_client = get_llm_client()
        return self._llm_client

    def create_virtual_schema(self, request, correlation_id: Optional[str] = None):
        """Create a new virtual schema."""
        _logger = logging.getLogger("neura.domain.federation")
        _logger.info(f"Creating virtual schema: {request.name}", extra={"correlation_id": correlation_id})
        schema_id = str(uuid.uuid4())[:8]
        now = datetime.now(timezone.utc).replace(microsecond=0).isoformat()

        tables = []
        for conn_id in request.connection_ids:
            try:
                from backend.app.repositories import resolve_connection_ref
                schema = resolve_connection_ref(conn_id)
                for table in (schema.get("tables", []) if isinstance(schema, dict) else []):
                    tables.append({"connection_id": conn_id, "table_name": table["name"], "alias": f"{conn_id[:4]}_{table['name']}"})
            except Exception as exc:
                _logger.warning(f"Failed to get schema for {conn_id}: {exc}")

        virtual_schema = {
            "id": schema_id, "name": request.name, "description": getattr(request, 'description', ''),
            "connections": request.connection_ids, "tables": tables, "joins": [],
            "created_at": now, "updated_at": now,
        }

        from backend.app.repositories import state_store
        with state_store.transaction() as state:
            state.setdefault("virtual_schemas", {})[schema_id] = virtual_schema
        return virtual_schema

    def suggest_joins(self, connection_ids: List[str], correlation_id: Optional[str] = None) -> List[Dict[str, Any]]:
        """Suggest joins between tables in different connections using AI."""
        _logger = logging.getLogger("neura.domain.federation")
        _logger.info(f"Suggesting joins for {len(connection_ids)} connections", extra={"correlation_id": correlation_id})

        schemas = {}
        for conn_id in connection_ids:
            try:
                from backend.app.repositories import resolve_connection_ref
                schemas[conn_id] = resolve_connection_ref(conn_id)
            except Exception as exc:
                _logger.warning(f"Failed to get schema for {conn_id}: {exc}")

        if len(schemas) < 2:
            return []

        schema_desc = []
        for conn_id, schema in schemas.items():
            if not isinstance(schema, dict):
                continue
            tables_desc = []
            for table in schema.get("tables", []):
                cols = [f"{c['name']} ({c.get('type', 'TEXT')})" for c in table.get("columns", [])]
                tables_desc.append(f"  - {table['name']}: {', '.join(cols)}")
            schema_desc.append(f"Connection {conn_id}:\n" + "\n".join(tables_desc))

        prompt = f"""Analyze these database schemas and suggest joins between tables from different connections.

{chr(10).join(schema_desc)}

Return a JSON array of join suggestions:
[{{"left_connection_id":"conn1","left_table":"t1","left_column":"c1","right_connection_id":"conn2","right_table":"t2","right_column":"c2","confidence":0.9,"reason":"explanation"}}]

Return ONLY the JSON array."""

        try:
            client = self._get_llm_client()
            response = client.complete(messages=[{"role": "user", "content": prompt}], description="join_suggestion", temperature=0.0)
            content = response["choices"][0]["message"]["content"]
            json_match = re.search(r"\[[\s\S]*\]", content)
            if json_match:
                return json.loads(json_match.group())
        except Exception as exc:
            _logger.error(f"Join suggestion failed: {exc}")
        return []

    def list_virtual_schemas(self) -> List[Dict[str, Any]]:
        from backend.app.repositories import state_store
        with state_store.transaction() as state:
            return list(state.get("virtual_schemas", {}).values())

    def get_virtual_schema(self, schema_id: str) -> Optional[Dict[str, Any]]:
        from backend.app.repositories import state_store
        with state_store.transaction() as state:
            return state.get("virtual_schemas", {}).get(schema_id)

    def delete_virtual_schema(self, schema_id: str) -> bool:
        from backend.app.repositories import state_store
        with state_store.transaction() as state:
            schemas = state.get("virtual_schemas", {})
            if schema_id not in schemas:
                return False
            del schemas[schema_id]
        return True

    def execute_query(self, request, correlation_id: Optional[str] = None) -> Dict[str, Any]:
        """Execute a federated query across multiple databases."""
        _logger = logging.getLogger("neura.domain.federation")
        schema = self.get_virtual_schema(request.virtual_schema_id)
        if not schema:
            raise ValueError(f"Virtual schema {request.virtual_schema_id} not found")

        connections = schema.get("connections", [])
        if not connections:
            raise ValueError("Virtual schema has no connections")

        # Single connection: execute directly
        target_conn = connections[0]
        try:
            from backend.app.repositories import resolve_connection_ref
            result = resolve_connection_ref(target_conn)
            return {"columns": [], "rows": [], "row_count": 0, "schema_id": request.virtual_schema_id, "executed_on": [target_conn], "routing": "single"}
        except Exception as exc:
            raise ValueError(f"Federated query failed: {exc}")


# =============================================================================
# NL2SQL SERVICE (ported from nl2sql/service.py)
# =============================================================================

_TRAILING_SEMICOLONS_RE = re.compile(r";+\s*$")

def _strip_trailing_semicolons(sql: str) -> str:
    return _TRAILING_SEMICOLONS_RE.sub("", (sql or "").strip())

def _coerce_value(value):
    """Convert bytes and other non-JSON types to serializable formats."""
    if isinstance(value, (bytes, bytearray, memoryview)):
        return bytes(value).hex()
    try:
        if hasattr(value, "item"):
            return value.item()
    except Exception:
        pass
    return value


class NL2SQLService:
    """Service for natural language to SQL operations."""

    def __init__(self):
        self._text_to_sql = None

    def _get_text_to_sql(self):
        if self._text_to_sql is None:
            from backend.app.services.llm import get_llm_client, TextToSQL
            client = get_llm_client()
            self._text_to_sql = TextToSQL(client=client, dialect="sqlite")
        return self._text_to_sql

    def _resolve_connection(self, connection_id: str):
        try:
            from backend.app.repositories import resolve_db_path, verify_sqlite
            db_path = resolve_db_path(connection_id=connection_id, db_url=None, db_path=None)
            verify_sqlite(db_path)
            return db_path
        except Exception as exc:
            raise ValueError(f"Invalid or unreachable database connection: {exc}")

    def _get_schema_for_connection(self, db_path, tables=None):
        from backend.app.repositories import get_loader
        loader = get_loader(db_path)
        schema = {}
        table_names = tables if tables else loader.table_names()
        for table_name in table_names:
            columns = [{"name": col.get("name"), "type": col.get("type", "TEXT"), "description": ""} for col in loader.pragma_table_info(table_name)]
            sample_values = {}
            try:
                frame = loader.frame(table_name)
                if not frame.empty:
                    sample_rows = frame.head(3)
                    for col in columns:
                        col_name = col["name"]
                        if col_name in sample_rows.columns:
                            values = [_coerce_value(v) for v in sample_rows[col_name].tolist()]
                            if values:
                                sample_values[col_name] = values[:3]
            except Exception:
                pass
            schema[table_name] = {"columns": columns, "foreign_keys": loader.foreign_keys(table_name), "sample_values": sample_values}
        return schema

    def generate_sql(self, request, correlation_id: Optional[str] = None):
        """Generate SQL from a natural language question."""
        _logger = logging.getLogger("neura.domain.nl2sql")
        _logger.info(f"Generating SQL for: {request.question[:100]}...", extra={"correlation_id": correlation_id})
        db_path = self._resolve_connection(request.connection_id)
        schema = self._get_schema_for_connection(db_path, getattr(request, 'tables', None))
        if not schema:
            raise ValueError("No tables found in the database")

        t2sql = self._get_text_to_sql()
        t2sql._schemas.clear()
        t2sql.add_schemas_from_catalog(schema)
        result = t2sql.generate_sql(question=request.question, tables=getattr(request, 'tables', None), context=getattr(request, 'context', None))

        self._record_history(question=request.question, sql=result.sql, connection_id=request.connection_id, confidence=result.confidence, success=True)
        return {"sql": result.sql, "explanation": result.explanation, "confidence": result.confidence, "warnings": result.warnings, "original_question": request.question}

    def execute_query(self, request, correlation_id: Optional[str] = None):
        """Execute a SQL query and return results using DataFrames."""
        import time as _time
        _logger = logging.getLogger("neura.domain.nl2sql")
        db_path = self._resolve_connection(request.connection_id)
        from backend.app.repositories import sqlite_shim, ensure_connection_loaded
        ensure_connection_loaded(request.connection_id, db_path)
        sql_clean = _strip_trailing_semicolons(request.sql)
        limit = getattr(request, 'limit', 100)
        offset = getattr(request, 'offset', 0)

        started = _time.time()
        try:
            with sqlite_shim.connect(str(db_path)) as con:
                con.row_factory = sqlite_shim.Row
                total_count = None
                if getattr(request, 'include_total', False):
                    try:
                        total_count = con.execute(f"SELECT COUNT(*) as cnt FROM ({sql_clean}) AS subq").fetchone()["cnt"]
                    except Exception:
                        pass
                limited_sql = f"SELECT * FROM ({sql_clean}) AS subq LIMIT {limit} OFFSET {offset}"
                cur = con.execute(limited_sql)
                rows_raw = cur.fetchall()
                columns = [desc[0] for desc in cur.description] if cur.description else []
                rows = [{col: _coerce_value(row[i]) for i, col in enumerate(columns)} for row in rows_raw]
        except Exception as exc:
            raise ValueError(f"Query execution failed: {exc}")

        execution_time_ms = int((_time.time() - started) * 1000)
        truncated = (total_count > limit) if total_count is not None else (len(rows) >= limit)
        return {"columns": columns, "rows": rows, "row_count": len(rows), "total_count": total_count, "execution_time_ms": execution_time_ms, "truncated": truncated}

    def explain_query(self, sql: str, correlation_id: Optional[str] = None) -> str:
        t2sql = self._get_text_to_sql()
        return t2sql.explain_sql(sql)

    def save_query(self, request, correlation_id: Optional[str] = None):
        self._resolve_connection(request.connection_id)
        query_id = str(uuid.uuid4())
        now = datetime.now(timezone.utc).replace(microsecond=0).isoformat()
        saved_query = {"id": query_id, "name": request.name, "description": getattr(request, 'description', ''), "sql": request.sql, "connection_id": request.connection_id, "original_question": getattr(request, 'original_question', ''), "tags": getattr(request, 'tags', []) or [], "created_at": now, "updated_at": now, "run_count": 0}
        from backend.app.repositories import state_store
        state_store.save_query(saved_query)
        return saved_query

    def list_saved_queries(self, connection_id=None, tags=None):
        from backend.app.repositories import state_store
        queries = state_store.list_saved_queries()
        if connection_id:
            queries = [q for q in queries if q.get("connection_id") == connection_id]
        if tags:
            tag_set = set(tags)
            queries = [q for q in queries if tag_set.intersection(set(q.get("tags", [])))]
        return queries

    def get_saved_query(self, query_id: str):
        from backend.app.repositories import state_store
        return state_store.get_saved_query(query_id)

    def delete_saved_query(self, query_id: str) -> bool:
        from backend.app.repositories import state_store
        return state_store.delete_saved_query(query_id)

    def get_query_history(self, connection_id=None, limit=50):
        from backend.app.repositories import state_store
        history = state_store.get_query_history(limit=limit)
        if connection_id:
            history = [h for h in history if h.get("connection_id") == connection_id]
        return history

    def _record_history(self, question, sql, connection_id, confidence, success, error=None):
        from backend.app.repositories import state_store
        entry = {"id": str(uuid.uuid4())[:8], "question": question, "sql": sql, "connection_id": connection_id, "confidence": confidence, "success": success, "error": error, "created_at": datetime.now(timezone.utc).replace(microsecond=0).isoformat()}
        state_store.add_query_history(entry)


# =============================================================================
# EXPORT SERVICE (ported from export/service.py)
# =============================================================================

class ExportService:
    """Service for exporting documents to various formats."""

    def __init__(self):
        self._export_jobs: dict[str, dict] = {}
        self._embed_tokens: dict[str, dict] = {}

    async def export_to_pdf(self, content: bytes, options: dict[str, Any]) -> bytes:
        if options.get("pdfa_compliant"):
            return await self._convert_to_pdfa(content, options)
        return content

    async def _convert_to_pdfa(self, pdf_content: bytes, options: dict[str, Any]) -> bytes:
        try:
            import fitz
            doc = fitz.open(stream=pdf_content, filetype="pdf")
            metadata = doc.metadata
            metadata["format"] = "PDF/A-1b"
            doc.set_metadata(metadata)
            output = io.BytesIO()
            doc.save(output)
            doc.close()
            return output.getvalue()
        except ImportError:
            return pdf_content

    async def export_to_docx(self, content: str, options: dict[str, Any]) -> bytes:
        from docx import Document
        doc = Document()
        if options.get("title"):
            doc.add_heading(options["title"], 0)
        for para in content.split("\n\n"):
            if para.strip():
                doc.add_paragraph(para.strip())
        output = io.BytesIO()
        doc.save(output)
        return output.getvalue()

    async def export_to_pptx(self, content: str, options: dict[str, Any]) -> bytes:
        from pptx import Presentation
        prs = Presentation()
        title_slide = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide)
        slide.shapes.title.text = options.get("title", "Document Export")
        slide.placeholders[1].text = options.get("subtitle", "Generated by NeuraReport")
        for i, section in enumerate(content.split("\n\n")):
            if section.strip():
                s = prs.slides.add_slide(prs.slide_layouts[1])
                s.shapes.title.text = f"Section {i + 1}"
                s.placeholders[1].text_frame.text = section.strip()
        output = io.BytesIO()
        prs.save(output)
        return output.getvalue()

    async def export_to_markdown(self, content: str, options: dict[str, Any]) -> bytes:
        md = ""
        if options.get("include_frontmatter", True):
            md = f"---\ntitle: {options.get('title', 'Document')}\ndate: {datetime.now(timezone.utc).strftime('%Y-%m-%d')}\n---\n\n"
        md += content
        return md.encode("utf-8")

    async def export_to_html(self, content: str, options: dict[str, Any]) -> bytes:
        import html as html_mod
        title = html_mod.escape(options.get("title", "Document"))
        if options.get("standalone", True):
            html_content = f"""<!DOCTYPE html><html><head><meta charset="UTF-8"><title>{title}</title></head><body><h1>{title}</h1>{content}</body></html>"""
        else:
            html_content = content
        return html_content.encode("utf-8")

    async def create_export_job(self, document_id: str, format: str, options: dict[str, Any]) -> dict:
        job_id = str(uuid.uuid4())
        job = {"job_id": job_id, "document_id": document_id, "format": format, "options": options, "status": "pending", "created_at": datetime.now(timezone.utc)}
        self._export_jobs[job_id] = job
        return job

    async def bulk_export(self, document_ids: list[str], format: str, options: dict[str, Any]) -> dict:
        job_id = str(uuid.uuid4())
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
            for i, doc_id in enumerate(document_ids):
                zf.writestr(f"document_{i + 1}.{format}", f"Content of document {doc_id}")
        return {"job_id": job_id, "status": "completed", "file_size": zip_buffer.tell()}

    async def generate_embed_token(self, document_id: str, options: dict[str, Any]) -> dict:
        import secrets as _secrets
        token = _secrets.token_urlsafe(32)
        self._embed_tokens[token] = {"token": token, "document_id": document_id, "options": options}
        width = options.get("width", 800)
        height = options.get("height", 600)
        embed_url = f"/embed/{token}"
        return {"token": token, "embed_url": embed_url, "embed_code": f'<iframe src="{embed_url}" width="{width}" height="{height}" frameborder="0"></iframe>'}

    async def validate_embed_token(self, token: str):
        return self._embed_tokens.get(token)

    async def revoke_embed_token(self, token_id: str) -> bool:
        if token_id in self._embed_tokens:
            del self._embed_tokens[token_id]
            return True
        return False


class DistributionService:
    """Service for distributing documents to various channels."""

    def __init__(self):
        self._distribution_jobs: dict[str, dict] = {}

    async def send_email(self, document_id: str, recipients: list[str], subject: str, message: str, attachments=None) -> dict:
        job_id = str(uuid.uuid4())
        job = {"job_id": job_id, "channel": "email", "document_id": document_id, "recipients": recipients, "status": "sent", "sent_at": datetime.now(timezone.utc), "recipients_count": len(recipients)}
        self._distribution_jobs[job_id] = job
        return job

    async def send_to_slack(self, document_id: str, channel: str, message: Optional[str] = None, file_content: Optional[bytes] = None) -> dict:
        job_id = str(uuid.uuid4())
        status = "skipped"
        try:
            from slack_sdk import WebClient
            import os as _os
            token = _os.getenv("SLACK_BOT_TOKEN")
            if token:
                client = WebClient(token=token)
                if file_content:
                    client.files_upload_v2(channel=channel, content=file_content, title=f"Document {document_id}", initial_comment=message or "Document shared from NeuraReport")
                else:
                    client.chat_postMessage(channel=channel, text=message or f"Document {document_id} shared from NeuraReport")
                status = "sent"
        except ImportError:
            status = "error"
        except Exception:
            status = "error"
        return {"job_id": job_id, "channel": "slack", "document_id": document_id, "slack_channel": channel, "status": status}

    async def send_to_teams(self, document_id: str, webhook_url: str, title: Optional[str] = None, message: Optional[str] = None) -> dict:
        job_id = str(uuid.uuid4())
        try:
            import aiohttp
            payload = {"@type": "MessageCard", "@context": "http://schema.org/extensions", "themeColor": "0076D7", "summary": title or f"Document {document_id}", "sections": [{"activityTitle": title or "Document Shared", "text": message or f"Document {document_id} shared from NeuraReport"}]}
            async with aiohttp.ClientSession() as session:
                async with session.post(webhook_url, json=payload) as resp:
                    status = "sent" if resp.status == 200 else "error"
        except Exception:
            status = "error"
        return {"job_id": job_id, "channel": "teams", "document_id": document_id, "status": status}

    async def send_webhook(self, document_id: str, webhook_url: str, method: str = "POST", headers=None, payload=None) -> dict:
        job_id = str(uuid.uuid4())
        try:
            import aiohttp
            if payload is None:
                payload = {"document_id": document_id, "timestamp": datetime.now(timezone.utc).isoformat(), "source": "neura_report"}
            async with aiohttp.ClientSession() as session:
                async with session.request(method, webhook_url, json=payload, headers=headers) as resp:
                    status = "sent" if resp.status < 400 else "error"
        except Exception:
            status = "error"
        return {"job_id": job_id, "channel": "webhook", "document_id": document_id, "status": status}


export_service = ExportService()
distribution_service = DistributionService()


# =============================================================================
# FACT CHECKER (ported from validation/fact_checker.py)
# =============================================================================

@dataclass
class Claim:
    """A single factual claim extracted from text."""
    text: str
    source_sentence: str
    evidence: list = field(default_factory=list)
    verdict: str = "unverified"  # verified, refuted, unverified, unsupported
    confidence: float = 0.0
    reasoning: str = ""


@dataclass
class FactCheckResult:
    """Result of fact-checking an LLM output."""
    original_text: str
    claims: list
    overall_score: float = 0.0
    passed: bool = True
    pass_threshold: float = 0.6

    def to_dict(self) -> dict:
        return {
            "passed": self.passed,
            "overall_score": round(self.overall_score, 3),
            "total_claims": len(self.claims),
            "verified_claims": sum(1 for c in self.claims if c.verdict == "verified"),
            "refuted_claims": sum(1 for c in self.claims if c.verdict == "refuted"),
            "unsupported_claims": sum(1 for c in self.claims if c.verdict == "unsupported"),
            "claims": [{"text": c.text, "verdict": c.verdict, "confidence": round(c.confidence, 3), "reasoning": c.reasoning, "evidence_count": len(c.evidence)} for c in self.claims],
        }


class FactChecker:
    """Fact-checking pipeline for LLM outputs. 3-stage: decompose claims, retrieve evidence, verify."""

    def __init__(self, pass_threshold: float = 0.6):
        self.pass_threshold = pass_threshold

    async def decompose_claims(self, text: str) -> list:
        try:
            from backend.app.services.llm import get_llm_client
            import asyncio
            from functools import partial
            client = get_llm_client()
            loop = asyncio.get_event_loop()
            response = await loop.run_in_executor(None, partial(client.complete, messages=[{"role": "user", "content": f"Extract all factual claims from this text. Return each claim on a new line, prefixed with \"- \".\nOnly include verifiable factual statements.\n\nText:\n{text[:3000]}\n\nClaims:"}], description="fact_check_decompose_claims", temperature=0.0, max_tokens=1024))
            raw = response.get("choices", [{}])[0].get("message", {}).get("content", "")
            claims = []
            for line in raw.strip().split("\n"):
                line = line.strip().lstrip("- ").strip()
                if line and len(line) > 10:
                    claims.append(Claim(text=line, source_sentence=line))
            return claims[:20]
        except Exception:
            sentences = [s.strip() for s in re.split(r'[.!?]+', text) if s.strip() and len(s.strip()) > 15]
            return [Claim(text=s, source_sentence=s) for s in sentences[:20]]

    async def retrieve_evidence(self, claim, context_docs: list) -> list:
        evidence = []
        claim_lower = claim.text.lower()
        for i, doc in enumerate(context_docs):
            doc_lower = doc.lower()
            claim_words = set(claim_lower.split())
            doc_words = set(doc_lower.split())
            overlap = len(claim_words & doc_words) / max(len(claim_words), 1)
            if overlap > 0.3:
                sentences = [s.strip() for s in doc.split('.') if s.strip()]
                best_sentence = max(sentences, key=lambda s: len(set(s.lower().split()) & claim_words), default=doc[:200])
                evidence.append({"doc_index": i, "relevant_text": best_sentence[:300], "overlap_score": round(overlap, 3)})
        return sorted(evidence, key=lambda e: e["overlap_score"], reverse=True)[:3]

    async def verify_claim(self, claim, context_docs: list):
        evidence = await self.retrieve_evidence(claim, context_docs)
        claim.evidence = evidence
        if not evidence:
            claim.verdict = "unsupported"
            claim.confidence = 0.2
            claim.reasoning = "No supporting evidence found in context documents"
            return claim
        best_score = max(e["overlap_score"] for e in evidence)
        if best_score > 0.6:
            claim.verdict = "verified"
            claim.confidence = min(best_score * 1.2, 1.0)
            claim.reasoning = f"Strong evidence found (overlap: {best_score:.2f})"
        elif best_score > 0.4:
            claim.verdict = "verified"
            claim.confidence = best_score
            claim.reasoning = f"Moderate evidence found (overlap: {best_score:.2f})"
        else:
            claim.verdict = "unverified"
            claim.confidence = best_score
            claim.reasoning = f"Weak evidence (overlap: {best_score:.2f})"
        return claim

    async def check(self, text: str, context_docs: Optional[list] = None, context_doc_ids: Optional[list] = None) -> FactCheckResult:
        if context_docs is None:
            context_docs = []
        claims = await self.decompose_claims(text)
        if not claims:
            return FactCheckResult(original_text=text, claims=[], overall_score=1.0, passed=True, pass_threshold=self.pass_threshold)
        verified_claims = []
        for claim in claims:
            verified = await self.verify_claim(claim, context_docs)
            verified_claims.append(verified)
        score_map = {"verified": 1.0, "unverified": 0.5, "unsupported": 0.2, "refuted": 0.0}
        overall_score = sum(score_map.get(c.verdict, 0.5) for c in verified_claims) / len(verified_claims) if verified_claims else 1.0
        return FactCheckResult(original_text=text, claims=verified_claims, overall_score=overall_score, passed=overall_score >= self.pass_threshold, pass_threshold=self.pass_threshold)


# =============================================================================
# WRITING SERVICE (ported from ai/writing_service.py)
# =============================================================================

MAX_TEXT_CHARS = 100_000
MAX_TEXT_CHARS_EXPAND = 50_000


class WritingTone(str, Enum):
    PROFESSIONAL = "professional"
    CASUAL = "casual"
    FORMAL = "formal"
    FRIENDLY = "friendly"
    ACADEMIC = "academic"
    TECHNICAL = "technical"
    PERSUASIVE = "persuasive"
    CONCISE = "concise"


class GrammarIssue(BaseModel):
    start: int = Field(..., description="Start position in text")
    end: int = Field(..., description="End position in text")
    original: str = Field(..., description="Original text")
    suggestion: str = Field(..., description="Suggested correction")
    issue_type: str = Field(..., description="Type of issue")
    explanation: str = Field(..., description="Explanation of the issue")
    severity: str = Field(default="warning")


class GrammarCheckResult(BaseModel):
    issues: List[GrammarIssue] = Field(default_factory=list)
    corrected_text: str
    issue_count: int
    score: float = Field(..., ge=0, le=100)


class SummarizeResult(BaseModel):
    summary: str
    key_points: List[str] = Field(default_factory=list)
    word_count_original: int
    word_count_summary: int
    compression_ratio: float


class RewriteResult(BaseModel):
    rewritten_text: str
    tone: str
    changes_made: List[str] = Field(default_factory=list)


class ExpandResult(BaseModel):
    expanded_text: str
    sections_added: List[str] = Field(default_factory=list)
    word_count_original: int
    word_count_expanded: int


class TranslateResult(BaseModel):
    translated_text: str
    source_language: str
    target_language: str
    confidence: float = Field(default=1.0, ge=0, le=1)


class WritingServiceError(Exception):
    pass

class InputValidationError(WritingServiceError):
    pass

class LLMResponseError(WritingServiceError):
    pass

class LLMUnavailableError(WritingServiceError):
    pass


def _extract_json_writing(raw: str) -> dict:
    text = raw.strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?\s*\n?", "", text, count=1)
        text = re.sub(r"\n?```\s*$", "", text, count=1)
    parsed = json.loads(text)
    if not isinstance(parsed, dict):
        raise ValueError(f"Expected JSON object, got {type(parsed).__name__}")
    return parsed


def _validate_grammar_positions(issues: list, text_length: int) -> list:
    valid = []
    for issue in issues:
        start = max(0, min(issue.get("start", 0), text_length))
        end = max(start, min(issue.get("end", 0), text_length))
        issue["start"] = start
        issue["end"] = end
        valid.append(issue)
    return valid


class WritingService:
    """AI-powered writing assistance: grammar, summarize, rewrite, expand, translate."""

    def __init__(self):
        self._llm_client = None

    def _get_llm_client(self):
        if self._llm_client is None:
            from backend.app.services.llm import get_llm_client
            self._llm_client = get_llm_client()
        return self._llm_client

    async def _call_llm(self, system_prompt: str, user_prompt: str, max_tokens: int = 2000, description: str = "writing_service") -> str:
        import asyncio as _asyncio
        client = self._get_llm_client()
        messages = [{"role": "system", "content": system_prompt}, {"role": "user", "content": user_prompt}]
        try:
            response = await _asyncio.to_thread(client.complete, messages=messages, description=description, max_tokens=max_tokens)
        except RuntimeError as exc:
            if "temporarily unavailable" in str(exc).lower():
                raise LLMUnavailableError(str(exc)) from exc
            raise LLMResponseError(str(exc)) from exc
        except Exception as exc:
            raise LLMResponseError(f"LLM call failed: {exc}") from exc
        content = response.get("choices", [{}])[0].get("message", {}).get("content", "")
        if not content:
            raise LLMResponseError("LLM returned empty response")
        return content

    async def check_grammar(self, text: str, language: str = "en", strict: bool = False) -> GrammarCheckResult:
        stripped = text.strip()
        if not stripped:
            return GrammarCheckResult(issues=[], corrected_text=text, issue_count=0, score=100.0)
        if len(text) > MAX_TEXT_CHARS:
            raise InputValidationError(f"Text exceeds maximum length of {MAX_TEXT_CHARS:,} characters.")
        strict_note = " (be strict — flag all style issues)" if strict else ""
        system_prompt = f"You are an expert grammar checker for {language} text.{strict_note}\nRespond ONLY with JSON: {{\"issues\": [{{\"start\": N, \"end\": N, \"original\": \"\", \"suggestion\": \"\", \"issue_type\": \"\", \"explanation\": \"\", \"severity\": \"\"}}], \"corrected_text\": \"\", \"score\": N}}"
        raw = await self._call_llm(system_prompt, f"Check this text:\n\n{text}", description="grammar_check")
        result = _extract_json_writing(raw)
        raw_issues = result.get("issues", [])
        if not isinstance(raw_issues, list):
            raw_issues = []
        validated = _validate_grammar_positions(raw_issues, len(text))
        issues = []
        for d in validated:
            try:
                issues.append(GrammarIssue(**d))
            except Exception:
                pass
        return GrammarCheckResult(issues=issues, corrected_text=result.get("corrected_text", text), issue_count=len(issues), score=max(0.0, min(100.0, float(result.get("score", 100.0)))))

    async def summarize(self, text: str, max_length: Optional[int] = None, style: str = "bullet_points") -> SummarizeResult:
        stripped = text.strip()
        if not stripped:
            return SummarizeResult(summary="", key_points=[], word_count_original=0, word_count_summary=0, compression_ratio=1.0)
        if len(text) > MAX_TEXT_CHARS:
            raise InputValidationError(f"Text exceeds maximum length of {MAX_TEXT_CHARS:,} characters.")
        wco = len(text.split())
        length_inst = f"Keep under {max_length} words." if max_length else ""
        system_prompt = f"You are an expert summarizer. {length_inst}\nRespond ONLY with JSON: {{\"summary\": \"\", \"key_points\": []}}"
        raw = await self._call_llm(system_prompt, f"Summarize:\n\n{text}", description="summarize")
        result = _extract_json_writing(raw)
        summary = result.get("summary", "")
        wcs = len(summary.split()) if summary else 0
        return SummarizeResult(summary=summary, key_points=result.get("key_points", []), word_count_original=wco, word_count_summary=wcs, compression_ratio=wcs / wco if wco > 0 else 1.0)

    async def rewrite(self, text: str, tone: WritingTone = WritingTone.PROFESSIONAL, preserve_meaning: bool = True) -> RewriteResult:
        stripped = text.strip()
        if not stripped:
            return RewriteResult(rewritten_text=text, tone=tone.value, changes_made=[])
        if len(text) > MAX_TEXT_CHARS:
            raise InputValidationError(f"Text exceeds maximum length of {MAX_TEXT_CHARS:,} characters.")
        system_prompt = f"You are an expert writer. Rewrite to be {tone.value}.{'Preserve meaning.' if preserve_meaning else ''}\nRespond ONLY with JSON: {{\"rewritten_text\": \"\", \"changes_made\": []}}"
        raw = await self._call_llm(system_prompt, f"Rewrite:\n\n{text}", description="rewrite")
        result = _extract_json_writing(raw)
        return RewriteResult(rewritten_text=result.get("rewritten_text", text), tone=tone.value, changes_made=result.get("changes_made", []))

    async def expand(self, text: str, target_length: Optional[int] = None, add_examples: bool = False, add_details: bool = True) -> ExpandResult:
        stripped = text.strip()
        if not stripped:
            return ExpandResult(expanded_text=text, sections_added=[], word_count_original=0, word_count_expanded=0)
        if len(text) > MAX_TEXT_CHARS_EXPAND:
            raise InputValidationError(f"Text exceeds maximum length of {MAX_TEXT_CHARS_EXPAND:,} characters for expansion.")
        wco = len(text.split())
        instructions = []
        if add_examples: instructions.append("Include examples")
        if add_details: instructions.append("Add details")
        if target_length: instructions.append(f"Aim for ~{target_length} words")
        system_prompt = f"Expand the text. {', '.join(instructions)}.\nRespond ONLY with JSON: {{\"expanded_text\": \"\", \"sections_added\": []}}"
        raw = await self._call_llm(system_prompt, f"Expand:\n\n{text}", max_tokens=4000, description="expand")
        result = _extract_json_writing(raw)
        expanded = result.get("expanded_text", text)
        return ExpandResult(expanded_text=expanded, sections_added=result.get("sections_added", []), word_count_original=wco, word_count_expanded=len(expanded.split()))

    async def translate(self, text: str, target_language: str, source_language: Optional[str] = None, preserve_formatting: bool = True) -> TranslateResult:
        stripped = text.strip()
        if not stripped:
            return TranslateResult(translated_text=text, source_language=source_language or "unknown", target_language=target_language, confidence=1.0)
        if len(text) > MAX_TEXT_CHARS:
            raise InputValidationError(f"Text exceeds maximum length of {MAX_TEXT_CHARS:,} characters.")
        src = f"from {source_language}" if source_language else "(detect source)"
        system_prompt = f"Translate {src} to {target_language}.{'Preserve formatting.' if preserve_formatting else ''}\nRespond ONLY with JSON: {{\"translated_text\": \"\", \"source_language\": \"\", \"confidence\": N}}"
        raw = await self._call_llm(system_prompt, f"Translate:\n\n{text}", max_tokens=4000, description="translate")
        result = _extract_json_writing(raw)
        return TranslateResult(translated_text=result.get("translated_text", text), source_language=result.get("source_language", source_language or "auto"), target_language=target_language, confidence=max(0.0, min(1.0, float(result.get("confidence", 0.9)))))

    async def generate_content(self, prompt: str, context: Optional[str] = None, tone: WritingTone = WritingTone.PROFESSIONAL, max_length: Optional[int] = None) -> str:
        if not prompt.strip():
            raise InputValidationError("Prompt cannot be empty.")
        system_prompt = f"Generate {tone.value} content.{f' Max {max_length} words.' if max_length else ''}{f' Context: {context}' if context else ''}"
        return await self._call_llm(system_prompt, prompt, max_tokens=4000, description="generate_content")


writing_service = WritingService()


# =============================================================================
# SPREADSHEET AI SERVICE (ported from ai/spreadsheet_ai_service.py)
# =============================================================================

MAX_DATA_ROWS = 5_000
MAX_FORMULA_LENGTH = 5_000


class FormulaResult(BaseModel):
    formula: str
    explanation: str
    examples: List[str] = Field(default_factory=list)
    alternative_formulas: List[str] = Field(default_factory=list)


class DataCleaningSuggestion(BaseModel):
    column: str
    issue: str
    suggestion: str
    severity: str = "medium"
    affected_rows: int = 0
    auto_fixable: bool = False


class DataCleaningResult(BaseModel):
    suggestions: List[DataCleaningSuggestion] = Field(default_factory=list)
    quality_score: float = Field(..., ge=0, le=100)
    summary: str


class SpreadsheetAnomaly(BaseModel):
    location: str
    value: Any
    expected_range: str
    confidence: float = Field(..., ge=0, le=1)
    explanation: str
    anomaly_type: str


class AnomalyDetectionResult(BaseModel):
    anomalies: List[SpreadsheetAnomaly] = Field(default_factory=list)
    total_rows_analyzed: int
    anomaly_count: int
    summary: str


class PredictionColumn(BaseModel):
    column_name: str
    predictions: List[Any] = Field(default_factory=list)
    confidence_scores: List[float] = Field(default_factory=list)
    methodology: str
    accuracy_estimate: float = Field(..., ge=0, le=1)


class FormulaExplanation(BaseModel):
    formula: str
    summary: str
    step_by_step: List[str] = Field(default_factory=list)
    components: Dict[str, str] = Field(default_factory=dict)
    potential_issues: List[str] = Field(default_factory=list)


class SpreadsheetAIService:
    """AI-powered spreadsheet: NL-to-formula, data quality, anomaly detection, predictions."""

    def __init__(self):
        self._llm_client = None

    def _get_llm_client(self):
        if self._llm_client is None:
            from backend.app.services.llm import get_llm_client
            self._llm_client = get_llm_client()
        return self._llm_client

    async def _call_llm(self, system_prompt: str, user_prompt: str, max_tokens: int = 2000, description: str = "spreadsheet_ai") -> str:
        import asyncio as _asyncio
        client = self._get_llm_client()
        messages = [{"role": "system", "content": system_prompt}, {"role": "user", "content": user_prompt}]
        try:
            response = await _asyncio.to_thread(client.complete, messages=messages, description=description, max_tokens=max_tokens)
        except Exception as exc:
            raise LLMResponseError(f"LLM call failed: {exc}") from exc
        content = response.get("choices", [{}])[0].get("message", {}).get("content", "")
        if not content:
            raise LLMResponseError("LLM returned empty response")
        return content

    async def natural_language_to_formula(self, description: str, context: Optional[str] = None, spreadsheet_type: str = "excel") -> FormulaResult:
        if not description.strip():
            raise InputValidationError("Description cannot be empty.")
        ctx = f"\n\nContext:\n{context}" if context else ""
        system_prompt = f"Convert to {spreadsheet_type} formula.\nRespond ONLY with JSON: {{\"formula\": \"\", \"explanation\": \"\", \"examples\": [], \"alternative_formulas\": []}}"
        raw = await self._call_llm(system_prompt, f"Create formula for: {description}{ctx}", description="nl_to_formula")
        result = _extract_json_writing(raw)
        return FormulaResult(formula=result.get("formula", ""), explanation=result.get("explanation", ""), examples=result.get("examples", []), alternative_formulas=result.get("alternative_formulas", []))

    async def analyze_data_quality(self, data_sample: List[Dict[str, Any]], column_info: Optional[Dict[str, str]] = None) -> DataCleaningResult:
        if not data_sample:
            return DataCleaningResult(suggestions=[], quality_score=100.0, summary="No data provided")
        if len(data_sample) > MAX_DATA_ROWS:
            raise InputValidationError(f"Data exceeds max of {MAX_DATA_ROWS:,} rows.")
        preview = json.dumps(data_sample[:20], indent=2, default=str)
        col_ctx = f"\n\nExpected types:\n{json.dumps(column_info, indent=2)}" if column_info else ""
        system_prompt = "Analyze data quality.\nRespond ONLY with JSON: {\"suggestions\": [{\"column\": \"\", \"issue\": \"\", \"suggestion\": \"\", \"severity\": \"\", \"affected_rows\": N, \"auto_fixable\": bool}], \"quality_score\": N, \"summary\": \"\"}"
        raw = await self._call_llm(system_prompt, f"Analyze:\n\n{preview}{col_ctx}", description="data_quality")
        result = _extract_json_writing(raw)
        suggestions = []
        for s in result.get("suggestions", []):
            try: suggestions.append(DataCleaningSuggestion(**s))
            except Exception: pass
        return DataCleaningResult(suggestions=suggestions, quality_score=max(0.0, min(100.0, float(result.get("quality_score", 0)))), summary=result.get("summary", ""))

    async def detect_anomalies(self, data: List[Dict[str, Any]], columns_to_analyze: Optional[List[str]] = None, sensitivity: str = "medium") -> AnomalyDetectionResult:
        if not data:
            return AnomalyDetectionResult(anomalies=[], total_rows_analyzed=0, anomaly_count=0, summary="No data")
        if len(data) > MAX_DATA_ROWS:
            raise InputValidationError(f"Data exceeds max of {MAX_DATA_ROWS:,} rows.")
        preview = json.dumps(data[:50], indent=2, default=str)
        col_ctx = f"\nFocus on: {', '.join(columns_to_analyze)}" if columns_to_analyze else ""
        system_prompt = f"Detect anomalies (sensitivity: {sensitivity}).\nRespond ONLY with JSON: {{\"anomalies\": [{{\"location\": \"\", \"value\": \"\", \"expected_range\": \"\", \"confidence\": N, \"explanation\": \"\", \"anomaly_type\": \"\"}}], \"total_rows_analyzed\": N, \"summary\": \"\"}}"
        raw = await self._call_llm(system_prompt, f"Detect anomalies:\n\n{preview}{col_ctx}", description="anomaly_detection")
        result = _extract_json_writing(raw)
        anomalies = []
        for a in result.get("anomalies", []):
            try: anomalies.append(SpreadsheetAnomaly(**a))
            except Exception: pass
        return AnomalyDetectionResult(anomalies=anomalies, total_rows_analyzed=result.get("total_rows_analyzed", len(data)), anomaly_count=len(anomalies), summary=result.get("summary", ""))

    async def generate_predictive_column(self, data: List[Dict[str, Any]], target_description: str, based_on_columns: List[str]) -> PredictionColumn:
        if not data:
            raise InputValidationError("Data cannot be empty.")
        if len(data) > MAX_DATA_ROWS:
            raise InputValidationError(f"Data exceeds max of {MAX_DATA_ROWS:,} rows.")
        preview = json.dumps(data[:30], indent=2, default=str)
        system_prompt = "Generate predictions.\nRespond ONLY with JSON: {\"column_name\": \"\", \"predictions\": [], \"confidence_scores\": [], \"methodology\": \"\", \"accuracy_estimate\": N}"
        raw = await self._call_llm(system_prompt, f"Predict: {target_description}\nBased on: {', '.join(based_on_columns)}\n\n{preview}", max_tokens=4000, description="predictive_column")
        result = _extract_json_writing(raw)
        return PredictionColumn(column_name=result.get("column_name", "Predicted"), predictions=result.get("predictions", []), confidence_scores=result.get("confidence_scores", []), methodology=result.get("methodology", ""), accuracy_estimate=max(0.0, min(1.0, float(result.get("accuracy_estimate", 0)))))

    async def explain_formula(self, formula: str, context: Optional[str] = None) -> FormulaExplanation:
        if not formula.strip():
            raise InputValidationError("Formula cannot be empty.")
        ctx = f"\n\nContext: {context}" if context else ""
        system_prompt = "Explain this formula.\nRespond ONLY with JSON: {\"formula\": \"\", \"summary\": \"\", \"step_by_step\": [], \"components\": {}, \"potential_issues\": []}"
        raw = await self._call_llm(system_prompt, f"Explain: {formula}{ctx}", description="explain_formula")
        result = _extract_json_writing(raw)
        return FormulaExplanation(formula=formula, summary=result.get("summary", ""), step_by_step=result.get("step_by_step", []), components=result.get("components", {}), potential_issues=result.get("potential_issues", []))

    async def suggest_formulas(self, data_sample: List[Dict[str, Any]], analysis_goals: Optional[str] = None) -> List[FormulaResult]:
        if not data_sample:
            return []
        preview = json.dumps(data_sample[:10], indent=2, default=str)
        goals = f"\nGoals: {analysis_goals}" if analysis_goals else ""
        system_prompt = "Suggest useful formulas.\nRespond ONLY with JSON: {\"suggestions\": [{\"formula\": \"\", \"explanation\": \"\", \"examples\": [], \"alternative_formulas\": []}]}"
        raw = await self._call_llm(system_prompt, f"Suggest formulas:\n\n{preview}{goals}", description="suggest_formulas")
        result = _extract_json_writing(raw)
        suggestions = []
        for s in result.get("suggestions", []):
            try: suggestions.append(FormulaResult(**s))
            except Exception: pass
        return suggestions


spreadsheet_ai_service = SpreadsheetAIService()


# =============================================================================
# QUICKCHART CLIENT (ported from charts/quickchart.py)
# =============================================================================

@dataclass
class ChartConfig:
    """Configuration for a chart."""
    chart_type: str = "bar"
    labels: list = field(default_factory=list)
    datasets: list = field(default_factory=list)
    title: str = ""
    width: int = 500
    height: int = 300
    background_color: str = "white"
    device_pixel_ratio: float = 2.0


class QuickChartClient:
    """Client for generating chart images via QuickChart.io API."""

    BASE_URL = "https://quickchart.io/chart"

    def __init__(self, base_url: Optional[str] = None):
        self.base_url = base_url or self.BASE_URL

    def _build_chart_json(self, config: ChartConfig) -> dict:
        return {
            "type": config.chart_type,
            "data": {"labels": config.labels, "datasets": config.datasets},
            "options": {"title": {"display": bool(config.title), "text": config.title}, "responsive": True},
        }

    def get_chart_url(self, config: ChartConfig) -> str:
        import urllib.parse
        chart_json = json.dumps(self._build_chart_json(config))
        params = {"c": chart_json, "w": config.width, "h": config.height, "bkg": config.background_color, "devicePixelRatio": config.device_pixel_ratio}
        return f"{self.base_url}?{urllib.parse.urlencode(params)}"

    def get_chart_bytes(self, config: ChartConfig) -> bytes:
        import urllib.request
        url = self.get_chart_url(config)
        with urllib.request.urlopen(url, timeout=30) as resp:
            return resp.read()

    def get_short_url(self, config: ChartConfig) -> str:
        import urllib.request
        chart_json = json.dumps(self._build_chart_json(config))
        payload = json.dumps({"chart": chart_json, "width": config.width, "height": config.height}).encode()
        req = urllib.request.Request(f"{self.base_url}/create", data=payload, headers={"Content-Type": "application/json"}, method="POST")
        with urllib.request.urlopen(req, timeout=30) as resp:
            result = json.loads(resp.read())
            return result.get("url", "")


def create_bar_chart(labels: list, data: list, title: str = "", **kwargs) -> ChartConfig:
    return ChartConfig(chart_type="bar", labels=labels, datasets=[{"label": title, "data": data}], title=title, **kwargs)

def create_line_chart(labels: list, data: list, title: str = "", **kwargs) -> ChartConfig:
    return ChartConfig(chart_type="line", labels=labels, datasets=[{"label": title, "data": data}], title=title, **kwargs)

def create_pie_chart(labels: list, data: list, title: str = "", **kwargs) -> ChartConfig:
    return ChartConfig(chart_type="pie", labels=labels, datasets=[{"data": data}], title=title, **kwargs)


quickchart_client = QuickChartClient()

