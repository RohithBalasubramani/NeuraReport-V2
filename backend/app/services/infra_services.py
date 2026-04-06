"""Consolidated Infrastructure Services (Phase B6).

Merged from: ocr, render, v2_config, events, streaming,
vector, validation, connections, export.
"""
from __future__ import annotations

# mypy: ignore-errors
"""
Shared OCR extraction utility (GLM-OCR via Ollama).

Calls Ollama's native /api/generate endpoint (NOT the OpenAI-compatible /v1
endpoint) because Ollama vision models use a raw prompt template
({{ .Prompt }}) that /v1 doesn't handle correctly.

Used by:
- providers.py (_try_ocr_image) — when the text-only LLM agent receives images
- pdf_extractors.py (OCREngine) — fallback OCR for scanned PDFs
- invoice_parser.py — invoice text extraction
- receipt_scanner.py — receipt text extraction
- TemplateVerify.py (_extract_vision_text) — template verification OCR
"""

import base64
import logging
from pathlib import Path
from typing import Optional
from backend.app.common import strip_code_fences, utc_now, utc_now_iso

logger = logging.getLogger("neura.ocr")


def ocr_extract(
    image_bytes: bytes,
    *,
    prompt: str = "OCR the text in this image.",
    max_tokens: int = 4096,
    timeout: int = 180,
) -> Optional[str]:
    """Extract text from an image using GLM-OCR (or configured vision model) via Ollama.

    Default: GLM-OCR 0.9B — 94.62% on OmniDocBench, ~2GB VRAM.
    Controlled by VISION_LLM_MODEL env var.

    Returns extracted text, or None if vision is unavailable or extraction fails.
    """
    from backend.app.services.llm import get_llm_config

    config = get_llm_config()
    if not config.vision_enabled or not config.vision_model:
        return None

    try:
        import requests

        b64_image = base64.b64encode(image_bytes).decode()

        # Ollama native API base: strip /v1 suffix if present
        api_base = (config.vision_api_base or "http://localhost:11434/v1").rstrip("/")
        if api_base.endswith("/v1"):
            api_base = api_base[:-3]

        resp = requests.post(
            f"{api_base}/api/generate",
            json={
                "model": config.vision_model,
                "prompt": prompt,
                "images": [b64_image],
                "stream": False,
                "options": {"num_predict": max_tokens, "temperature": 0},
            },
            timeout=timeout,
        )
        resp.raise_for_status()
        text = resp.json().get("response", "")
        # Strip markdown bold markers some models add
        text = text.replace("**", "")

        if text and len(text.strip()) > 20:
            logger.info(
                "ocr_extracted",
                extra={"event": "ocr_extracted", "model": config.vision_model, "chars": len(text)},
            )
            return text.strip()
        return None

    except Exception as exc:
        logger.warning(
            "ocr_failed",
            extra={"event": "ocr_failed", "model": getattr(config, 'vision_model', '?'), "error": str(exc)},
        )
        return None


def ocr_extract_from_file(
    image_path: Path,
    **kwargs,
) -> Optional[str]:
    """Read an image file and extract text via the configured OCR model."""
    image_path = Path(image_path)
    if not image_path.exists():
        return None
    return ocr_extract(image_path.read_bytes(), **kwargs)


# Backward-compatible aliases (legacy code may still import these names)
deepseek_ocr_extract = ocr_extract
deepseek_ocr_extract_from_file = ocr_extract_from_file
# Preferred aliases
glm_ocr_extract = ocr_extract
glm_ocr_extract_from_file = ocr_extract_from_file


# ── Structured OCR Extraction ──────────────────────────────────────────

STRUCTURED_OCR_PROMPT = """\
You are a precision document analyzer examining a business report PDF page. Extract ALL visible text with careful attention to structure.

TASK: Analyze this report page and return your findings in these clearly separated sections.

== SCALAR FIELDS ==
List every standalone label-value pair visible in the header, footer, or margins of the report.
Format each as: LABEL: value
Examples of what to look for: report title, company name, plant name, batch number, date ranges, print date, page numbers, prepared by, approved by.
Include the EXACT text as printed -- do not paraphrase or normalize.

== COLUMN HEADERS ==
List every column header in the data table(s), left to right, exactly as printed.
Format as a numbered list:
1. First Column Header
2. Second Column Header
...
If there are multiple tables, separate them with TABLE_BREAK.
Preserve abbreviations, units in parentheses, and special characters exactly as shown.

== DATA SAMPLES ==
Copy the first 2-3 data rows from each table, preserving the column order from COLUMN HEADERS.
Format each row as pipe-separated values: value1 | value2 | value3 | ...
Include numeric formatting as shown (decimal places, thousands separators).

== LAYOUT NOTES ==
Report these observations:
- Number of distinct tables on the page
- Whether the page has repeating blocks (identical sections stacked vertically, like multiple batches)
- Whether there is a totals/summary row at the bottom of any table
- Position of header fields (top, left margin, right margin)
- Position of footer fields (bottom, below table)
- Any merged cells or spanning headers in the table

Do NOT add interpretation, commentary, or markdown formatting. Extract ONLY what is physically visible on the page.\
"""


VALIDATION_OCR_PROMPT = """\
Read ALL text visible in this generated report page. Focus specifically on:
1. Any text that looks like {placeholder} or {{token}} patterns — these are template tokens that should have been replaced with real data and indicate a data leak.
2. All column headers in any tables — list them left to right.
3. Whether data cells contain actual values or are empty/blank.
4. Any obviously broken formatting (overlapping text, cut-off content, misaligned columns).

Return the complete text exactly as shown, then add a line:
ISSUES: [list any {token} patterns, empty cells, or formatting problems found, or "none"]\
"""

import re as _re


def _normalize_to_snake_case(text: str) -> str:
    """Normalize a header label to lowercase_snake_case."""
    # Remove units in parentheses for the normalized form, keep for matching
    s = _re.sub(r"[()/%]", " ", text)
    s = _re.sub(r"[^A-Za-z0-9\s_]", "", s)
    s = s.strip().lower()
    s = _re.sub(r"\s+", "_", s)
    s = _re.sub(r"_+", "_", s)
    return s.strip("_")


def ocr_extract_structured(
    image_bytes: bytes,
    *,
    timeout: int = 180,
) -> dict:
    """Extract structured OCR from a report page image via GLM-OCR.

    Returns a dict with sections: scalar_fields, column_headers,
    data_samples, layout_notes, and raw_text.
    Falls back gracefully if parsing fails.
    """
    raw = ocr_extract(image_bytes, prompt=STRUCTURED_OCR_PROMPT, max_tokens=4096, timeout=timeout)
    if not raw:
        return {"raw_text": "", "sections": {}}

    result = {
        "version": "1.0",
        "extracted_at": utc_now_iso(),
        "raw_text": raw,
        "sections": {
            "scalar_fields": [],
            "column_headers": [],
            "data_samples": [],
            "layout_notes": {},
        },
    }

    try:
        _parse_structured_sections(raw, result["sections"])
    except Exception:
        logger.debug("structured_ocr_parse_failed", exc_info=True)

    return result


def _parse_structured_sections(raw: str, sections: dict) -> None:
    """Parse the 4 structured sections from GLM-OCR's response."""
    # Split by section markers
    scalar_match = _re.search(r"==\s*SCALAR FIELDS\s*==\s*\n(.*?)(?===|$)", raw, _re.DOTALL)
    header_match = _re.search(r"==\s*COLUMN HEADERS\s*==\s*\n(.*?)(?===|$)", raw, _re.DOTALL)
    sample_match = _re.search(r"==\s*DATA SAMPLES\s*==\s*\n(.*?)(?===|$)", raw, _re.DOTALL)
    layout_match = _re.search(r"==\s*LAYOUT NOTES\s*==\s*\n(.*?)(?===|$)", raw, _re.DOTALL)

    # Parse scalar fields: "LABEL: value"
    if scalar_match:
        for line in scalar_match.group(1).strip().splitlines():
            line = line.strip()
            if ":" in line and not line.startswith("-"):
                parts = line.split(":", 1)
                label = parts[0].strip()
                value = parts[1].strip() if len(parts) > 1 else ""
                if label:
                    sections["scalar_fields"].append({
                        "label": label,
                        "normalized": _normalize_to_snake_case(label),
                        "sample_value": value,
                    })

    # Parse column headers: numbered list "1. Header Text"
    if header_match:
        pos = 0
        for line in header_match.group(1).strip().splitlines():
            line = line.strip()
            if line.upper() == "TABLE_BREAK":
                continue
            num_match = _re.match(r"^\d+[\.\)]\s*(.+)$", line)
            if num_match:
                pos += 1
                text = num_match.group(1).strip()
                sections["column_headers"].append({
                    "position": pos,
                    "text": text,
                    "normalized": _normalize_to_snake_case(text),
                })

    # Parse data samples: pipe-separated rows
    if sample_match:
        row_num = 0
        for line in sample_match.group(1).strip().splitlines():
            line = line.strip()
            if "|" in line:
                row_num += 1
                values = [v.strip() for v in line.split("|") if v.strip()]
                sections["data_samples"].append({
                    "row": row_num,
                    "values": values,
                })

    # Parse layout notes: key observations
    if layout_match:
        layout_text = layout_match.group(1).strip()
        notes = sections["layout_notes"]
        # Extract specific observations
        notes["raw"] = layout_text
        if _re.search(r"(\d+)\s*(?:distinct\s+)?table", layout_text, _re.IGNORECASE):
            m = _re.search(r"(\d+)\s*(?:distinct\s+)?table", layout_text, _re.IGNORECASE)
            notes["table_count"] = int(m.group(1))
        notes["has_repeating_blocks"] = bool(
            _re.search(r"repeating|stacked|multiple.*batch", layout_text, _re.IGNORECASE)
            and not _re.search(r"no\s+repeating|does\s+not\s+have\s+repeating", layout_text, _re.IGNORECASE)
        )
        notes["has_totals_row"] = bool(
            _re.search(r"total|summary", layout_text, _re.IGNORECASE)
            and not _re.search(r"no\s+total|no\s+summary", layout_text, _re.IGNORECASE)
        )


def format_ocr_for_llm(structured: dict) -> str:
    """Format structured OCR for PDF_PAGE_TEXT injection in llm_call_1 (template generation).

    Produces a text block that tells Qwen exactly what headers, labels,
    and data exist on the PDF page — so it generates accurate token names.
    """
    sections = structured.get("sections", {})
    if not sections:
        return structured.get("raw_text", "")

    parts = []

    scalars = sections.get("scalar_fields", [])
    if scalars:
        parts.append("HEADER/FOOTER FIELDS:")
        for s in scalars:
            parts.append(f"  {s['label']}: {s.get('sample_value', '')}")

    headers = sections.get("column_headers", [])
    if headers:
        parts.append("\nTABLE COLUMN HEADERS (left to right):")
        for h in headers:
            parts.append(f"  {h['position']}. {h['text']}")

    samples = sections.get("data_samples", [])
    if samples:
        parts.append("\nSAMPLE DATA ROWS:")
        for row in samples[:3]:
            parts.append(f"  Row {row['row']}: {' | '.join(row['values'])}")

    layout = sections.get("layout_notes", {})
    if layout:
        notes_parts = []
        if layout.get("table_count"):
            notes_parts.append(f"{layout['table_count']} table(s)")
        if layout.get("has_totals_row"):
            notes_parts.append("has totals row")
        if layout.get("has_repeating_blocks"):
            notes_parts.append("has repeating blocks")
        if notes_parts:
            parts.append(f"\nLAYOUT: {', '.join(notes_parts)}")

    return "\n".join(parts) if parts else structured.get("raw_text", "")


def format_ocr_for_mapping(structured: dict) -> str:
    """Format structured OCR for the mapping prompt (auto_map_tokens, refine_mapping).

    Provides column headers with normalized snake_case names, scalar fields
    with sample values, and data samples — optimized for Qwen's token mapping.
    """
    sections = structured.get("sections", {})
    if not sections:
        return structured.get("raw_text", "")

    parts = []

    headers = sections.get("column_headers", [])
    if headers:
        parts.append("## Column Headers (from source PDF, left to right)")
        for h in headers:
            parts.append(f'  {h["position"]}. "{h["text"]}" → normalized: {h["normalized"]}')

    scalars = sections.get("scalar_fields", [])
    if scalars:
        parts.append("\n## Scalar Fields (header/footer labels)")
        for s in scalars:
            parts.append(f'  - {s["label"]}: "{s.get("sample_value", "?")}" → normalized: {s["normalized"]}')

    samples = sections.get("data_samples", [])
    if samples:
        parts.append("\n## Sample Data Rows")
        for row in samples[:2]:
            parts.append(f"  Row {row['row']}: {' | '.join(row['values'])}")

    layout = sections.get("layout_notes", {})
    if layout:
        notes = []
        if layout.get("table_count"):
            notes.append(f"{layout['table_count']} table(s)")
        if layout.get("has_totals_row"):
            notes.append("totals_row=yes")
        else:
            notes.append("totals_row=no")
        if layout.get("has_repeating_blocks"):
            notes.append("repeating_blocks=yes")
        else:
            notes.append("repeating_blocks=no")
        parts.append(f"\n## Layout: {', '.join(notes)}")

    return "\n".join(parts) if parts else structured.get("raw_text", "")


import os
from typing import Literal

import fitz  # PyMuPDF
from playwright.sync_api import sync_playwright

MM_PER_INCH = 25.4

A4_MM_W = 210.0
A4_MM_H = 297.0

def _a4_enforcing_css() -> str:
    # Ensures one exact A4 page with white background for both methods
    return """
    @page { size: A4; margin: 0 }
    html, body { margin: 0; padding: 0; background: #fff; }
    .page { width: 210mm; min-height: 297mm; box-sizing: border-box; background: #fff; }
    """

def _wait_for_fonts(page) -> None:
    # Best-effort wait to ensure webfonts loaded
    try:
        page.wait_for_function("document.fonts && document.fonts.status === 'loaded'", timeout=10000)
    except Exception:
        pass

def _html_to_pdf_bytes_with_playwright(html: str) -> bytes:
    with sync_playwright() as pw:
        browser = pw.chromium.launch()
        ctx = browser.new_context()
        page = ctx.new_page()
        page.set_content(html, wait_until="networkidle")
        page.emulate_media(media="print")
        page.add_style_tag(content=_a4_enforcing_css())
        _wait_for_fonts(page)
        pdf_bytes = page.pdf(
            format="A4",
            print_background=True,
            margin={"top": "0", "right": "0", "bottom": "0", "left": "0"},
        )
        ctx.close()
        browser.close()
        return pdf_bytes

def _rasterize_pdf_first_page_to_png(pdf_bytes: bytes, dpi: int) -> bytes:
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    page = doc.load_page(0)
    scale = dpi / 72.0  # PDF points are 72/in
    pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale), colorspace=fitz.csRGB, alpha=False)
    return pix.tobytes("png")

def _screenshot_element_to_png(html: str, selector: str, dpi: int) -> bytes:
    # Fallback if you need a direct DOM screenshot. Chromium may cap device_scale_factor.
    css_w = round(A4_MM_W / MM_PER_INCH * 96)  # ~794 px
    css_h = round(A4_MM_H / MM_PER_INCH * 96)  # ~1123 px
    dsf = dpi / 96.0  # 400/96 ≈ 4.1667 (may be clamped)
    with sync_playwright() as pw:
        browser = pw.chromium.launch()
        ctx = browser.new_context(
            device_scale_factor=dsf,
            viewport={"width": css_w, "height": css_h},
        )
        page = ctx.new_page()
        page.set_content(html, wait_until="networkidle")
        page.emulate_media(media="screen")
        page.add_style_tag(content=_a4_enforcing_css())
        _wait_for_fonts(page)
        locator = page.locator(selector)
        if locator.count() == 0:
            locator = page.locator("body")
        png = locator.screenshot(type="png")
        ctx.close()
        browser.close()
        return png

def rasterize_html_to_png(
    html: str,
    dpi: int = 400,
    method: Literal["pdf", "screenshot"] = "pdf",
    selector: str = ".page",
) -> bytes:
    """
    Returns a tightly-cropped A4 PNG at the requested DPI.
    Preferred: method='pdf' (uses print engine + PyMuPDF rasterize).
    Fallback: method='screenshot' crops selector ('.page').
    """
    if method == "pdf":
        pdf = _html_to_pdf_bytes_with_playwright(html)
        return _rasterize_pdf_first_page_to_png(pdf, dpi=dpi)
    return _screenshot_element_to_png(html, selector=selector, dpi=dpi)

def save_png(png_bytes: bytes, out_path: str) -> str:
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    with open(out_path, "wb") as fh:
        fh.write(png_bytes)
    return out_path

# mypy: ignore-errors
"""
Centralized V2 Configuration — Feature Flags + Thresholds.

Pydantic BaseSettings with environment variable overrides.
All feature flags default to False for safe incremental rollout.

Environment variable prefix: V2_
Example: V2_ENABLE_LANGGRAPH_PIPELINE=true
"""

import logging
from typing import Optional

from pydantic import Field
from pydantic_settings import BaseSettings

logger = logging.getLogger("neura.v2_config")

class V2Config(BaseSettings):
    """All V2 feature flags and tunable thresholds."""

    model_config = {"env_prefix": "V2_", "env_file": ".env", "extra": "ignore"}

    # ── Framework Feature Flags ──────────────────────────────────────
    enable_langgraph_pipeline: bool = Field(
        default=False,
        description="Route report generation through LangGraph state-graph pipeline",
    )
    enable_rag_augmentation: bool = Field(
        default=False,
        description="Inject RAG context into mapping, docqa, and report prompts",
    )
    enable_autogen_teams: bool = Field(
        default=False,
        description="Dispatch eligible agent tasks to AutoGen multi-agent teams",
    )
    enable_crewai_crews: bool = Field(
        default=False,
        description="Dispatch eligible agent tasks to CrewAI role-based crews",
    )
    enable_quality_loop: bool = Field(
        default=False,
        description="Wrap agent/report outputs in iterative quality-loop evaluation",
    )
    enable_conversation_memory: bool = Field(
        default=False,
        description="Persist per-session conversation context and entity tracking",
    )
    enable_dspy_signatures: bool = Field(
        default=False,
        description="Use DSPy compiled modules for structured LLM calls",
    )
    enable_semantic_cache: bool = Field(
        default=False,
        description="Enable L2 embedding-based semantic cache for LLM responses",
    )
    enable_sse_streaming: bool = Field(
        default=False,
        description="Stream pipeline stage progress via SSE to the frontend",
    )

    # ── Quality Loop Thresholds ──────────────────────────────────────
    quality_threshold: float = Field(
        default=0.7,
        ge=0.0,
        le=1.0,
        description="Minimum quality score to accept output without retry",
    )
    max_quality_iterations: int = Field(
        default=3,
        ge=1,
        le=10,
        description="Maximum quality-loop retry iterations",
    )
    quality_timeout_seconds: int = Field(
        default=120,
        ge=10,
        description="Quality loop timeout in seconds",
    )

    # Quality dimension weights (BFI 6-dimension pattern)
    quality_weight_completeness: float = Field(default=0.25)
    quality_weight_accuracy: float = Field(default=0.30)
    quality_weight_clarity: float = Field(default=0.20)
    quality_weight_relevance: float = Field(default=0.15)
    quality_weight_formatting: float = Field(default=0.10)

    # ── Semantic Cache ───────────────────────────────────────────────
    cache_l2_enabled: bool = Field(
        default=False,
        description="Enable L2 embedding-similarity cache layer",
    )
    cache_l2_threshold: float = Field(
        default=0.92,
        ge=0.5,
        le=1.0,
        description="Cosine similarity threshold for L2 cache hits",
    )
    cache_l2_max_size: int = Field(
        default=500,
        ge=10,
        description="Maximum number of entries in L2 cache",
    )

    # ── RAG Settings ─────────────────────────────────────────────────
    rag_top_k: int = Field(
        default=5,
        ge=1,
        le=20,
        description="Number of RAG retrieval results to inject",
    )
    rag_relevance_threshold: float = Field(
        default=0.5,
        ge=0.0,
        le=1.0,
        description="Minimum relevance score for RAG results",
    )

    # ── RL / Experience ──────────────────────────────────────────────
    rl_persist_to_jsonl: bool = Field(
        default=True,
        description="Persist Thompson Sampling experience to JSONL file",
    )
    rl_jsonl_path: str = Field(
        default="state/rl_experience.jsonl",
        description="Path for RL experience JSONL persistence",
    )

    # ── Pipeline ─────────────────────────────────────────────────────
    pipeline_checkpoint_enabled: bool = Field(
        default=True,
        description="Enable LangGraph checkpoint persistence for pipeline recovery",
    )
    pipeline_selective_retry: bool = Field(
        default=True,
        description="On quality failure, retry only failed stages (BFI pattern)",
    )

    # ── Agent Teams ──────────────────────────────────────────────────
    team_max_rounds: int = Field(
        default=5,
        ge=1,
        le=20,
        description="Maximum rounds for multi-agent team conversations",
    )
    team_timeout_seconds: int = Field(
        default=180,
        ge=30,
        description="Timeout for team-based execution",
    )

    def get_quality_weights(self) -> dict[str, float]:
        """Return quality dimension weights as a dict."""
        return {
            "completeness": self.quality_weight_completeness,
            "accuracy": self.quality_weight_accuracy,
            "clarity": self.quality_weight_clarity,
            "relevance": self.quality_weight_relevance,
            "formatting": self.quality_weight_formatting,
        }

    def log_active_flags(self) -> None:
        """Log which V2 flags are currently active."""
        flags = {
            k: v
            for k, v in self.model_dump().items()
            if k.startswith("enable_") and v is True
        }
        if flags:
            logger.info("V2 active flags: %s", ", ".join(flags.keys()))
        else:
            logger.info("V2: all feature flags disabled (safe mode)")

# ── Singleton ────────────────────────────────────────────────────────

_instance: Optional[V2Config] = None

def get_v2_config() -> V2Config:
    """Return the global V2Config singleton (lazy-initialized)."""
    global _instance
    if _instance is None:
        _instance = V2Config()
        _instance.log_active_flags()
    return _instance

def reset_v2_config() -> None:
    """Reset the singleton (for testing)."""
    global _instance
    _instance = None

# mypy: ignore-errors
"""
Typed EventBus — Pydantic-typed events with priority-ordered handlers.

Modelled after the BFI pipeline's EventBus pattern with:
- Pydantic event models for type safety
- Priority-ordered handler dispatch
- Scoped handlers (per-request lifecycle cleanup)
- Non-blocking publish (errors logged, not propagated)
"""

import logging
import time
import uuid
from dataclasses import dataclass, field
from typing import Any, Callable, Dict, List, Optional, Type

from pydantic import BaseModel, Field

logger = logging.getLogger("neura.events")

# ── Event Models ─────────────────────────────────────────────────────

class BaseEvent(BaseModel):
    """Base class for all typed events."""

    event_id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    timestamp: float = Field(default_factory=time.time)
    source: str = Field(default="unknown")

class PipelineStageEvent(BaseEvent):
    """Emitted when a pipeline stage changes status."""

    source: str = "pipeline"
    run_id: str = ""
    stage: str = ""
    status: str = ""  # started, completed, failed, retrying
    duration_ms: Optional[float] = None
    metadata: Dict[str, Any] = Field(default_factory=dict)
    error: Optional[str] = None

class QualityScoreEvent(BaseEvent):
    """Emitted when a quality evaluation completes."""

    source: str = "quality"
    entity_type: str = ""
    entity_id: str = ""
    overall_score: float = 0.0
    dimension_scores: Dict[str, float] = Field(default_factory=dict)
    iteration: int = 1
    accepted: bool = False

class FeedbackEvent(BaseEvent):
    """Emitted when user feedback is recorded."""

    source: str = "feedback"
    feedback_type: str = ""  # thumbs, rating, correction
    entity_type: str = ""
    entity_id: str = ""
    rating: Optional[float] = None
    thumbs_up: Optional[bool] = None
    comment: Optional[str] = None

class AgentTeamEvent(BaseEvent):
    """Emitted during multi-agent team execution."""

    source: str = "agent_team"
    team_name: str = ""
    agent_name: str = ""
    action: str = ""  # started, message, round_complete, finished
    round_num: int = 0
    content: Optional[str] = None
    metadata: Dict[str, Any] = Field(default_factory=dict)

# ── Handler Registration ─────────────────────────────────────────────

@dataclass
class HandlerEntry:
    """A registered event handler with priority and optional scope."""

    handler: Callable
    priority: int = 0
    scope: Optional[str] = None
    handler_id: str = field(default_factory=lambda: str(uuid.uuid4()))

class EventBus:
    """
    Typed event bus with priority-ordered handler dispatch.

    Features:
    - Subscribe handlers for specific event types
    - Priority ordering (lower = runs first)
    - Scoped handlers for per-request lifecycle cleanup
    - Non-blocking publish (errors logged, not propagated)
    """

    def __init__(self) -> None:
        self._handlers: Dict[Type[BaseEvent], List[HandlerEntry]] = {}
        self._global_handlers: List[HandlerEntry] = []

    def subscribe(
        self,
        event_type: Type[BaseEvent],
        handler: Callable,
        priority: int = 0,
        scope: Optional[str] = None,
    ) -> str:
        """Register a handler for a specific event type."""
        entry = HandlerEntry(
            handler=handler, priority=priority, scope=scope
        )
        if event_type not in self._handlers:
            self._handlers[event_type] = []
        self._handlers[event_type].append(entry)
        # Re-sort by priority
        self._handlers[event_type].sort(key=lambda h: h.priority)
        return entry.handler_id

    def subscribe_all(
        self,
        handler: Callable,
        priority: int = 0,
        scope: Optional[str] = None,
    ) -> str:
        """Register a handler that receives ALL event types."""
        entry = HandlerEntry(
            handler=handler, priority=priority, scope=scope
        )
        self._global_handlers.append(entry)
        self._global_handlers.sort(key=lambda h: h.priority)
        return entry.handler_id

    def unsubscribe(self, handler_id: str) -> bool:
        """Remove a handler by its ID."""
        for entries in self._handlers.values():
            for entry in entries:
                if entry.handler_id == handler_id:
                    entries.remove(entry)
                    return True
        for entry in self._global_handlers:
            if entry.handler_id == handler_id:
                self._global_handlers.remove(entry)
                return True
        return False

    def cleanup_scope(self, scope: str) -> int:
        """Remove all handlers registered under a given scope."""
        removed = 0
        for event_type in self._handlers:
            before = len(self._handlers[event_type])
            self._handlers[event_type] = [
                h for h in self._handlers[event_type] if h.scope != scope
            ]
            removed += before - len(self._handlers[event_type])

        before = len(self._global_handlers)
        self._global_handlers = [
            h for h in self._global_handlers if h.scope != scope
        ]
        removed += before - len(self._global_handlers)

        if removed:
            logger.debug("Cleaned up %d handlers for scope %s", removed, scope)
        return removed

    def publish(self, event: BaseEvent) -> None:
        """
        Dispatch an event to all matching handlers.

        Handlers are called synchronously in priority order.
        Errors are logged but do not propagate (non-blocking).
        """
        event_type = type(event)
        handlers = list(self._handlers.get(event_type, []))
        handlers.extend(self._global_handlers)
        handlers.sort(key=lambda h: h.priority)

        for entry in handlers:
            try:
                entry.handler(event)
            except Exception:
                logger.exception(
                    "Event handler %s failed for %s",
                    entry.handler_id,
                    event_type.__name__,
                )

    def handler_count(self, event_type: Optional[Type[BaseEvent]] = None) -> int:
        """Return the number of registered handlers."""
        if event_type:
            return len(self._handlers.get(event_type, []))
        total = sum(len(h) for h in self._handlers.values())
        total += len(self._global_handlers)
        return total

    def clear(self) -> None:
        """Remove all handlers."""
        self._handlers.clear()
        self._global_handlers.clear()

# ── Singleton ────────────────────────────────────────────────────────

_instance: Optional[EventBus] = None

def get_event_bus() -> EventBus:
    """Return the global EventBus singleton."""
    global _instance
    if _instance is None:
        _instance = EventBus()
    return _instance

# mypy: ignore-errors
"""
SSE Pipeline Bridge — Queue-based thread→async streaming bridge.

Modelled after the BFI pipeline's SSE pattern:
- Sync pipeline nodes call bridge.emit() from worker threads
- Async SSE endpoint consumes bridge.stream() as an AsyncIterator
- Uses threading.Queue + threading.Event for thread-safe coordination
- Supports stage lifecycle events, quality scores, and errors

Event schema:
    {
        "event": "stage_start"|"stage_complete"|"stage_retry"|"quality_score"|
                 "pipeline_start"|"pipeline_complete"|"error"|"heartbeat",
        "run_id": str,
        "stage": str | None,
        "data": dict,
        "timestamp": float
    }
"""

import asyncio
import json
import logging
import time
import uuid
from dataclasses import dataclass, field
from queue import Empty, Queue
from threading import Event
from typing import Any, AsyncIterator, Dict, Optional

logger = logging.getLogger("neura.streaming")

@dataclass
class SSEEvent:
    """Structured SSE event."""

    event: str
    run_id: str = ""
    stage: Optional[str] = None
    data: Dict[str, Any] = field(default_factory=dict)
    timestamp: float = field(default_factory=time.time)

    def to_dict(self) -> Dict[str, Any]:
        d = {
            "event": self.event,
            "run_id": self.run_id,
            "timestamp": self.timestamp,
        }
        if self.stage:
            d["stage"] = self.stage
        if self.data:
            d["data"] = self.data
        return d

    def to_sse_line(self) -> str:
        """Format as SSE data line."""
        return f"data: {json.dumps(self.to_dict())}\n\n"

class PipelineSSEBridge:
    """
    Queue-based bridge for streaming pipeline events via SSE.

    Thread-safe: emit() can be called from any thread.
    The async stream() method yields events for StreamingResponse.
    """

    def __init__(
        self,
        run_id: Optional[str] = None,
        queue_maxsize: int = 100,
        heartbeat_interval: float = 15.0,
    ) -> None:
        self.run_id = run_id or str(uuid.uuid4())
        self._queue: Queue[Optional[SSEEvent]] = Queue(maxsize=queue_maxsize)
        self._finished = Event()
        self._started_at = time.time()
        self._heartbeat_interval = heartbeat_interval
        self._stage_timings: Dict[str, float] = {}

    # ── Sync API (called from pipeline threads) ──────────────────────

    def emit(self, event_dict: Dict[str, Any]) -> None:
        """Emit a raw event dict from a sync pipeline context."""
        evt = SSEEvent(
            event=event_dict.get("event", "info"),
            run_id=self.run_id,
            stage=event_dict.get("stage"),
            data=event_dict.get("data", {}),
        )
        try:
            self._queue.put_nowait(evt)
        except Exception:
            logger.warning("SSE queue full, dropping event: %s", evt.event)

    def emit_stage_start(self, stage: str, **extra: Any) -> None:
        """Convenience: emit a stage_start event."""
        self._stage_timings[stage] = time.time()
        self.emit({
            "event": "stage_start",
            "stage": stage,
            "data": {"message": f"Starting {stage}", **extra},
        })

    def emit_stage_complete(self, stage: str, **extra: Any) -> None:
        """Convenience: emit a stage_complete event with duration."""
        started = self._stage_timings.pop(stage, self._started_at)
        duration_ms = (time.time() - started) * 1000
        self.emit({
            "event": "stage_complete",
            "stage": stage,
            "data": {"duration_ms": round(duration_ms, 1), **extra},
        })

    def emit_stage_retry(self, stage: str, attempt: int, reason: str = "") -> None:
        """Convenience: emit a stage_retry event."""
        self.emit({
            "event": "stage_retry",
            "stage": stage,
            "data": {"attempt": attempt, "reason": reason},
        })

    def emit_quality_score(
        self, score: float, dimensions: Optional[Dict[str, float]] = None, **extra: Any
    ) -> None:
        """Convenience: emit a quality_score event."""
        self.emit({
            "event": "quality_score",
            "data": {
                "score": score,
                "dimensions": dimensions or {},
                **extra,
            },
        })

    def emit_error(self, stage: str, error: str) -> None:
        """Convenience: emit an error event."""
        self.emit({
            "event": "error",
            "stage": stage,
            "data": {"error": error},
        })

    def finish(self, success: bool = True, summary: Optional[Dict[str, Any]] = None) -> None:
        """
        Signal pipeline completion.

        Emits a final pipeline_complete or pipeline_fail event,
        then puts a sentinel None to terminate the stream.
        """
        total_ms = (time.time() - self._started_at) * 1000
        event_type = "pipeline_complete" if success else "pipeline_fail"
        data = {"total_duration_ms": round(total_ms, 1)}
        if summary:
            data["summary"] = summary

        self.emit({"event": event_type, "data": data})
        # Sentinel to stop stream iteration
        self._queue.put(None)
        self._finished.set()

    # ── Async API (called from SSE endpoint) ─────────────────────────

    async def stream(self) -> AsyncIterator[str]:
        """
        Async generator that yields SSE-formatted strings.

        Use with FastAPI StreamingResponse:
            return StreamingResponse(
                bridge.stream(),
                media_type="text/event-stream"
            )
        """
        # Emit initial pipeline_start
        start_event = SSEEvent(
            event="pipeline_start",
            run_id=self.run_id,
            data={"message": "Pipeline started"},
        )
        yield start_event.to_sse_line()

        last_heartbeat = time.time()

        while not self._finished.is_set() or not self._queue.empty():
            try:
                # Non-blocking poll with small sleep for async cooperation
                evt = self._queue.get_nowait()
                if evt is None:
                    # Sentinel — pipeline is done
                    break
                yield evt.to_sse_line()
                last_heartbeat = time.time()
            except Empty:
                # Send heartbeat if idle too long (keeps connection alive)
                if time.time() - last_heartbeat > self._heartbeat_interval:
                    heartbeat = SSEEvent(
                        event="heartbeat",
                        run_id=self.run_id,
                    )
                    yield heartbeat.to_sse_line()
                    last_heartbeat = time.time()

                # Yield control to event loop
                await asyncio.sleep(0.1)

        # Drain any remaining events
        while not self._queue.empty():
            evt = self._queue.get_nowait()
            if evt is not None:
                yield evt.to_sse_line()

    @property
    def is_finished(self) -> bool:
        return self._finished.is_set()

    @property
    def elapsed_seconds(self) -> float:
        return time.time() - self._started_at

"""
Vector embedding and semantic search service.

Provides document embedding, vector storage (pgvector/Qdrant), and
semantic retrieval for the RAG pipeline.

Supports:
- sentence-transformers for local embeddings
- OpenAI embeddings API as fallback
- pgvector for PostgreSQL-native vector storage
- Qdrant for dedicated vector search

Based on: pgvector/pgvector-python + qdrant/qdrant-client patterns.
"""

import hashlib
import logging
import os
import threading
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional, Tuple

import numpy as np

logger = logging.getLogger("neura.vector")

# Configuration
DEFAULT_MODEL = os.getenv("NEURA_EMBEDDING_MODEL", "all-MiniLM-L6-v2")
DEFAULT_DIMENSION = 384  # all-MiniLM-L6-v2 output dimension
QDRANT_URL = os.getenv("NEURA_QDRANT_URL", "http://localhost:6333")
QDRANT_COLLECTION = os.getenv("NEURA_QDRANT_COLLECTION", "neurareport_docs")

@dataclass
class EmbeddingResult:
    """Result of an embedding operation."""
    text: str
    embedding: List[float]
    model: str
    token_count: int = 0

@dataclass
class SearchResult:
    """A single search result with score."""
    document_id: str
    text: str
    score: float
    metadata: Dict[str, Any] = field(default_factory=dict)

class EmbeddingService:
    """
    Generate embeddings using sentence-transformers (local) or OpenAI API.

    Prefers local models for privacy and cost; falls back to OpenAI if
    sentence-transformers is unavailable.
    """

    def __init__(self, model_name: str = DEFAULT_MODEL):
        self.model_name = model_name
        self._model = None
        self._lock = threading.Lock()
        self._use_openai = False

    def _load_model(self):
        """Lazy-load the embedding model."""
        if self._model is not None:
            return

        with self._lock:
            if self._model is not None:
                return
            try:
                from sentence_transformers import SentenceTransformer
                self._model = SentenceTransformer(self.model_name)
                logger.info(
                    "embedding_model_loaded",
                    extra={"event": "embedding_model_loaded", "model": self.model_name},
                )
            except ImportError:
                logger.warning(
                    "sentence_transformers_unavailable",
                    extra={"event": "sentence_transformers_unavailable"},
                )
                self._use_openai = True
            except Exception as exc:
                logger.warning(
                    "embedding_model_load_failed",
                    extra={"event": "embedding_model_load_failed", "error": str(exc)},
                )
                self._use_openai = True

    def embed_texts(self, texts: List[str]) -> List[EmbeddingResult]:
        """Embed a batch of texts."""
        if not texts:
            return []

        self._load_model()

        if self._use_openai:
            return self._embed_openai(texts)

        embeddings = self._model.encode(texts, normalize_embeddings=True)
        return [
            EmbeddingResult(
                text=text,
                embedding=emb.tolist(),
                model=self.model_name,
            )
            for text, emb in zip(texts, embeddings)
        ]

    def embed_text(self, text: str) -> EmbeddingResult:
        """Embed a single text."""
        results = self.embed_texts([text])
        return results[0] if results else EmbeddingResult(text=text, embedding=[], model=self.model_name)

    def _embed_openai(self, texts: List[str]) -> List[EmbeddingResult]:
        """Fallback to OpenAI embeddings API."""
        try:
            import openai
            client = openai.OpenAI()
            response = client.embeddings.create(
                input=texts,
                model="text-embedding-3-small",
            )
            return [
                EmbeddingResult(
                    text=texts[i],
                    embedding=item.embedding,
                    model="text-embedding-3-small",
                    token_count=response.usage.total_tokens // len(texts),
                )
                for i, item in enumerate(response.data)
            ]
        except Exception as exc:
            logger.error("openai_embedding_failed", extra={"event": "openai_embedding_failed", "error": str(exc)})
            return [
                EmbeddingResult(text=t, embedding=[], model="none")
                for t in texts
            ]

class VectorStore:
    """
    Vector storage and retrieval backend.

    Supports Qdrant (preferred) and pgvector (PostgreSQL extension).
    Falls back to in-memory numpy-based search if neither is available.
    """

    def __init__(self, backend: str = "auto"):
        self._backend = backend
        self._client = None
        self._memory_store: Dict[str, Tuple[List[float], Dict[str, Any]]] = {}
        self._lock = threading.Lock()
        self._resolved_backend = None

    def _resolve_backend(self):
        """Resolve which backend to use."""
        if self._resolved_backend:
            return

        if self._backend == "qdrant" or (self._backend == "auto"):
            try:
                from qdrant_client import QdrantClient
                from qdrant_client.models import Distance, VectorParams

                self._client = QdrantClient(url=QDRANT_URL)
                # Ensure collection exists
                collections = [c.name for c in self._client.get_collections().collections]
                if QDRANT_COLLECTION not in collections:
                    self._client.create_collection(
                        collection_name=QDRANT_COLLECTION,
                        vectors_config=VectorParams(
                            size=DEFAULT_DIMENSION,
                            distance=Distance.COSINE,
                        ),
                    )
                self._resolved_backend = "qdrant"
                logger.info("vector_store_backend", extra={"event": "vector_store_backend", "backend": "qdrant"})
                return
            except Exception as exc:
                if self._backend == "qdrant":
                    raise
                logger.info(f"Qdrant unavailable, trying pgvector: {exc}")

        if self._backend in ("pgvector", "auto"):
            try:
                self._resolved_backend = "pgvector"
                logger.info("vector_store_backend", extra={"event": "vector_store_backend", "backend": "pgvector"})
                return
            except Exception as exc:
                if self._backend == "pgvector":
                    raise
                logger.info(f"pgvector unavailable, using memory: {exc}")

        self._resolved_backend = "memory"
        logger.info("vector_store_backend", extra={"event": "vector_store_backend", "backend": "memory"})

    def upsert(self, document_id: str, embedding: List[float], metadata: Optional[Dict[str, Any]] = None) -> None:
        """Insert or update a vector."""
        self._resolve_backend()

        if self._resolved_backend == "qdrant":
            from qdrant_client.models import PointStruct
            point = PointStruct(
                id=self._hash_id(document_id),
                vector=embedding,
                payload={"document_id": document_id, **(metadata or {})},
            )
            self._client.upsert(collection_name=QDRANT_COLLECTION, points=[point])
        elif self._resolved_backend == "pgvector":
            self._pgvector_upsert(document_id, embedding, metadata)
        else:
            with self._lock:
                self._memory_store[document_id] = (embedding, metadata or {})

    def upsert_batch(self, items: List[Tuple[str, List[float], Dict[str, Any]]]) -> int:
        """Batch upsert vectors. Returns count of upserted items."""
        self._resolve_backend()

        if self._resolved_backend == "qdrant":
            from qdrant_client.models import PointStruct
            points = [
                PointStruct(
                    id=self._hash_id(doc_id),
                    vector=emb,
                    payload={"document_id": doc_id, **meta},
                )
                for doc_id, emb, meta in items
            ]
            self._client.upsert(collection_name=QDRANT_COLLECTION, points=points)
            return len(points)
        else:
            for doc_id, emb, meta in items:
                self.upsert(doc_id, emb, meta)
            return len(items)

    def search(self, query_embedding: List[float], top_k: int = 10) -> List[SearchResult]:
        """Search for similar vectors."""
        self._resolve_backend()

        if self._resolved_backend == "qdrant":
            results = self._client.search(
                collection_name=QDRANT_COLLECTION,
                query_vector=query_embedding,
                limit=top_k,
            )
            return [
                SearchResult(
                    document_id=hit.payload.get("document_id", ""),
                    text=hit.payload.get("text", ""),
                    score=hit.score,
                    metadata={k: v for k, v in hit.payload.items() if k not in ("document_id", "text")},
                )
                for hit in results
            ]
        elif self._resolved_backend == "pgvector":
            return self._pgvector_search(query_embedding, top_k)
        else:
            return self._memory_search(query_embedding, top_k)

    def delete(self, document_id: str) -> bool:
        """Delete a vector by document ID."""
        self._resolve_backend()

        if self._resolved_backend == "qdrant":
            from qdrant_client.models import PointIdsList
            self._client.delete(
                collection_name=QDRANT_COLLECTION,
                points_selector=PointIdsList(points=[self._hash_id(document_id)]),
            )
            return True
        elif self._resolved_backend == "memory":
            with self._lock:
                return self._memory_store.pop(document_id, None) is not None
        return False

    def count(self) -> int:
        """Get the number of stored vectors."""
        self._resolve_backend()

        if self._resolved_backend == "qdrant":
            info = self._client.get_collection(QDRANT_COLLECTION)
            return info.points_count
        elif self._resolved_backend == "memory":
            return len(self._memory_store)
        return 0

    def _memory_search(self, query_embedding: List[float], top_k: int) -> List[SearchResult]:
        """In-memory cosine similarity search."""
        if not self._memory_store:
            return []

        query = np.array(query_embedding, dtype=np.float32)
        query_norm = np.linalg.norm(query)
        if query_norm == 0:
            return []
        query = query / query_norm

        scored = []
        for doc_id, (emb, meta) in self._memory_store.items():
            vec = np.array(emb, dtype=np.float32)
            vec_norm = np.linalg.norm(vec)
            if vec_norm == 0:
                continue
            similarity = float(np.dot(query, vec / vec_norm))
            scored.append((doc_id, similarity, meta))

        scored.sort(key=lambda x: x[1], reverse=True)

        return [
            SearchResult(
                document_id=doc_id,
                text=meta.get("text", ""),
                score=score,
                metadata={k: v for k, v in meta.items() if k != "text"},
            )
            for doc_id, score, meta in scored[:top_k]
        ]

    def _pgvector_upsert(self, document_id: str, embedding: List[float], metadata: Optional[Dict[str, Any]]) -> None:
        """Upsert using pgvector extension."""
        # This requires the pgvector extension and appropriate table setup
        # Implementation deferred to Alembic migration for table creation
        logger.warning("pgvector_upsert_not_implemented")

    def _pgvector_search(self, query_embedding: List[float], top_k: int) -> List[SearchResult]:
        """Search using pgvector extension."""
        logger.warning("pgvector_search_not_implemented")
        return []

    @staticmethod
    def _hash_id(document_id: str) -> int:
        """Convert string ID to integer for Qdrant."""
        return int(hashlib.md5(document_id.encode()).hexdigest()[:16], 16)

# Module-level singletons
_embedding_service: Optional[EmbeddingService] = None
_vector_store: Optional[VectorStore] = None

def get_embedding_service() -> EmbeddingService:
    global _embedding_service
    if _embedding_service is None:
        _embedding_service = EmbeddingService()
    return _embedding_service

def get_vector_store() -> VectorStore:
    global _vector_store
    if _vector_store is None:
        _vector_store = VectorStore()
    return _vector_store

"""
Hallucination detection and fact-checking pipeline.

3-stage verification:
1. Claim decomposition: break LLM output into individual claims
2. Evidence retrieval: search document store for supporting evidence
3. Claim verification: score each claim against evidence

Based on: OpenFactCheck + Exa hallucination detector patterns.
"""
import asyncio
import re
import logging
from dataclasses import dataclass, field
from functools import partial
from typing import Any, Optional

logger = logging.getLogger("neura.validation.factcheck")

@dataclass
class Claim:
    """A single factual claim extracted from text."""
    text: str
    source_sentence: str
    evidence: list[dict[str, Any]] = field(default_factory=list)
    verdict: str = "unverified"  # verified, refuted, unverified, unsupported
    confidence: float = 0.0
    reasoning: str = ""

@dataclass
class FactCheckResult:
    """Result of fact-checking an LLM output."""
    original_text: str
    claims: list[Claim]
    overall_score: float = 0.0  # 0.0 = all hallucinated, 1.0 = all verified
    passed: bool = True
    pass_threshold: float = 0.6

    def to_dict(self) -> dict:
        return {
            "passed": self.passed,
            "overall_score": round(self.overall_score, 3),
            "total_claims": len(self.claims),
            "verified_claims": sum(1 for c in self.claims if c.verdict == "verified"),
            "refuted_claims": sum(1 for c in self.claims if c.verdict == "refuted"),
            "unsupported_claims": sum(1 for c in self.claims if c.verdict == "unsupported"),
            "claims": [
                {
                    "text": c.text,
                    "verdict": c.verdict,
                    "confidence": round(c.confidence, 3),
                    "reasoning": c.reasoning,
                    "evidence_count": len(c.evidence),
                }
                for c in self.claims
            ],
        }

class FactChecker:
    """
    Fact-checking pipeline for LLM outputs.

    Usage:
        checker = FactChecker()
        result = await checker.check(llm_output, context_docs=[...])
        if not result.passed:
            # Flag response as potentially unreliable
    """

    def __init__(self, pass_threshold: float = 0.6):
        self.pass_threshold = pass_threshold

    async def decompose_claims(self, text: str) -> list[Claim]:
        """Extract individual factual claims from text using LLM."""
        try:
            from backend.app.services.llm import get_llm_client
            client = get_llm_client()

            loop = asyncio.get_event_loop()
            response = await loop.run_in_executor(
                None,
                partial(
                    client.complete,
                    messages=[{
                        "role": "user",
                        "content": f"""Extract all factual claims from this text. Return each claim on a new line, prefixed with "- ".
Only include verifiable factual statements, not opinions or qualifiers.

Text:
{text[:3000]}

Claims:"""
                    }],
                    description="fact_check_decompose_claims",
                    temperature=0.0,
                    max_tokens=1024,
                ),
            )

            raw = (
                response.get("choices", [{}])[0]
                .get("message", {})
                .get("content", "")
            )
            claims = []
            for line in raw.strip().split("\n"):
                line = line.strip().lstrip("- ").strip()
                if line and len(line) > 10:
                    claims.append(Claim(text=line, source_sentence=line))

            return claims[:20]  # Cap at 20 claims

        except Exception:
            # Fallback: split by sentences
            sentences = [s.strip() for s in re.split(r'[.!?]+', text) if s.strip() and len(s.strip()) > 15]
            return [Claim(text=s, source_sentence=s) for s in sentences[:20]]

    async def retrieve_evidence(self, claim: Claim, context_docs: list[str]) -> list[dict[str, Any]]:
        """Search context documents for evidence. Uses vector search when available, falls back to keyword overlap."""
        # Try vector-based evidence retrieval
        try:
            from backend.app.services.infra_services import EmbeddingPipeline
            from backend.app.services.infra_services import PgVectorStore
            from backend.app.services.config import get_settings
            settings = get_settings()
            if "postgresql" in settings.database_url:
                pipeline = EmbeddingPipeline()
                query_embedding = await pipeline.embed_query(claim.text)
                # Use pgvector similarity search
                store = PgVectorStore(settings.database_url)
                from backend.app.services.db.engine import get_session_factory
                async with get_session_factory()() as session:
                    results = await store.search_similar(session, query_embedding, top_k=3)
                    if results:
                        return [
                            {
                                "doc_index": r.chunk_index,
                                "relevant_text": r.content[:300],
                                "overlap_score": round(r.similarity, 3),
                                "source": r.source,
                            }
                            for r in results
                        ]
        except Exception as exc:
            logger.debug("vector_evidence_fallback", extra={"event": "vector_evidence_fallback", "error": str(exc)})

        # Fallback: keyword overlap scoring
        evidence = []
        claim_lower = claim.text.lower()

        for i, doc in enumerate(context_docs):
            # Simple keyword overlap scoring
            doc_lower = doc.lower()
            claim_words = set(claim_lower.split())
            doc_words = set(doc_lower.split())
            overlap = len(claim_words & doc_words) / max(len(claim_words), 1)

            if overlap > 0.3:
                # Find the most relevant sentence
                sentences = [s.strip() for s in doc.split('.') if s.strip()]
                best_sentence = max(
                    sentences,
                    key=lambda s: len(set(s.lower().split()) & claim_words),
                    default=doc[:200],
                )
                evidence.append({
                    "doc_index": i,
                    "relevant_text": best_sentence[:300],
                    "overlap_score": round(overlap, 3),
                })

        return sorted(evidence, key=lambda e: e["overlap_score"], reverse=True)[:3]

    async def verify_claim(self, claim: Claim, context_docs: list[str]) -> Claim:
        """Verify a single claim against context documents."""
        evidence = await self.retrieve_evidence(claim, context_docs)
        claim.evidence = evidence

        if not evidence:
            claim.verdict = "unsupported"
            claim.confidence = 0.2
            claim.reasoning = "No supporting evidence found in context documents"
            return claim

        # Use best evidence overlap as confidence proxy
        best_score = max(e["overlap_score"] for e in evidence)

        if best_score > 0.6:
            claim.verdict = "verified"
            claim.confidence = min(best_score * 1.2, 1.0)
            claim.reasoning = f"Strong evidence found (overlap: {best_score:.2f})"
        elif best_score > 0.4:
            claim.verdict = "verified"
            claim.confidence = best_score
            claim.reasoning = f"Moderate evidence found (overlap: {best_score:.2f})"
        else:
            claim.verdict = "unverified"
            claim.confidence = best_score
            claim.reasoning = f"Weak evidence (overlap: {best_score:.2f})"

        return claim

    async def check(
        self,
        text: str,
        context_docs: Optional[list[str]] = None,
        context_doc_ids: Optional[list[str]] = None,
    ) -> FactCheckResult:
        """Full fact-checking pipeline."""
        if context_docs is None:
            context_docs = []

        # 1. Decompose into claims
        claims = await self.decompose_claims(text)

        if not claims:
            return FactCheckResult(
                original_text=text, claims=[], overall_score=1.0, passed=True,
                pass_threshold=self.pass_threshold,
            )

        # 2 & 3. Retrieve evidence and verify each claim
        verified_claims = []
        for claim in claims:
            verified = await self.verify_claim(claim, context_docs)
            verified_claims.append(verified)

        # 4. Compute overall score
        if verified_claims:
            score_map = {"verified": 1.0, "unverified": 0.5, "unsupported": 0.2, "refuted": 0.0}
            total = sum(score_map.get(c.verdict, 0.5) for c in verified_claims)
            overall_score = total / len(verified_claims)
        else:
            overall_score = 1.0

        passed = overall_score >= self.pass_threshold

        result = FactCheckResult(
            original_text=text,
            claims=verified_claims,
            overall_score=overall_score,
            passed=passed,
            pass_threshold=self.pass_threshold,
        )

        logger.info("fact_check_completed", extra={
            "event": "fact_check_completed",
            "total_claims": len(verified_claims),
            "overall_score": round(overall_score, 3),
            "passed": passed,
        })

        return result

import logging
import re
import time
from pathlib import Path
from urllib.parse import urlparse

from backend.app.repositories import ConnectionRepository
from backend.app.schemas import ConnectionResponse, ConnectionTestRequest, ConnectionUpsertRequest
from backend.app.utils import AppError

logger = logging.getLogger(__name__)

# Supported database types and their default ports
SUPPORTED_DB_TYPES = {
    "sqlite": {"port": None, "requires_auth": False},
    "postgresql": {"port": 5432, "requires_auth": True},
    "mysql": {"port": 3306, "requires_auth": True},
    "mssql": {"port": 1433, "requires_auth": True},
    "mariadb": {"port": 3306, "requires_auth": True},
}

def _parse_db_url(db_url: str) -> dict:
    """Parse a database URL into components."""
    if not db_url:
        return {}

    # Handle SQLite special case
    if db_url.startswith("sqlite"):
        match = re.match(r"sqlite(?:3)?:///(.+)", db_url)
        if match:
            return {"db_type": "sqlite", "database": match.group(1)}
        return {"db_type": "sqlite", "database": db_url}

    try:
        parsed = urlparse(db_url)
        db_type = parsed.scheme.lower()
        # Normalize postgres -> postgresql
        if db_type == "postgres":
            db_type = "postgresql"
        return {
            "db_type": db_type,
            "host": parsed.hostname or "localhost",
            "port": parsed.port,
            "database": parsed.path.lstrip("/") if parsed.path else "",
            "username": parsed.username,
            "password": parsed.password,
        }
    except Exception:
        return {}

class ConnectionService:
    def __init__(self, repo: ConnectionRepository | None = None):
        self.repo = repo or ConnectionRepository()

    def list(self, correlation_id: str | None = None):
        """List all connections via the repository."""
        return self.repo.list()

    def _resolve_and_verify(
        self,
        *,
        connection_id: str | None,
        db_url: str | None,
        db_path: str | None,
        db_type: str | None = None,
        verify: bool = True,
    ) -> Path | str:
        """Resolve and optionally verify a database connection.

        For SQLite, returns a Path. For other databases, returns the connection URL.
        """
        # Determine database type from URL or explicit parameter
        parsed = _parse_db_url(db_url) if db_url else {}
        detected_type = db_type or parsed.get("db_type") or "sqlite"

        if detected_type == "sqlite":
            # SQLite: verify file path
            try:
                path = self.repo.resolve_path(connection_id=connection_id, db_url=db_url, db_path=db_path)
                if verify:
                    self.repo.verify(path)
                return path
            except Exception as exc:
                raise AppError(
                    code="invalid_database",
                    message="Invalid or unreachable database",
                    detail=None,
                    status_code=400,
                )
        else:
            # Other databases: validate URL format and test connection
            if not db_url:
                raise AppError(
                    code="invalid_database",
                    message="Database URL is required for non-SQLite databases",
                    status_code=400,
                )
            if verify:
                self._verify_network_database(db_url, detected_type)
            return db_url

    def _verify_network_database(self, db_url: str, db_type: str) -> None:
        """Verify a network database connection."""
        parsed = _parse_db_url(db_url)
        host = parsed.get("host", "localhost")
        port = parsed.get("port")

        if not port:
            port = SUPPORTED_DB_TYPES.get(db_type, {}).get("port", 5432)

        # Basic connectivity check using socket
        import socket
        sock = None
        try:
            sock = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
            sock.settimeout(5)
            result = sock.connect_ex((host, int(port)))
            if result != 0:
                raise AppError(
                    code="connection_failed",
                    message=f"Cannot connect to {db_type} at {host}:{port}",
                    status_code=503,
                )
        except socket.error as exc:
            raise AppError(
                code="connection_failed",
                message=f"Network error connecting to {db_type}",
                detail=None,
                status_code=503,
            )
        finally:
            if sock:
                sock.close()

    def test(self, payload: ConnectionTestRequest, correlation_id: str | None = None) -> dict:
        started = time.time()

        # Parse the URL to determine database type
        parsed = _parse_db_url(payload.db_url) if payload.db_url else {}
        db_type = payload.db_type or parsed.get("db_type") or "sqlite"

        if db_type not in SUPPORTED_DB_TYPES:
            raise AppError(
                code="unsupported_database",
                message=f"Database type '{db_type}' is not supported. Supported: {', '.join(SUPPORTED_DB_TYPES.keys())}",
                status_code=400,
            )

        # Resolve and verify the connection
        connection_ref = self._resolve_and_verify(
            connection_id=None,
            db_url=payload.db_url,
            db_path=payload.database,
            db_type=db_type,
            verify=True,
        )

        latency_ms = int((time.time() - started) * 1000)

        # Build connection config
        if db_type == "sqlite":
            resolved = Path(connection_ref).resolve()
            cfg = {
                "db_type": db_type,
                "database": str(resolved),
                "db_url": payload.db_url,
                "name": resolved.name,
                "status": "connected",
                "latency_ms": latency_ms,
            }
            display_name = resolved.name
            normalized = {"db_type": db_type, "database": str(resolved)}
        else:
            cfg = {
                "db_type": db_type,
                "database": parsed.get("database", ""),
                "host": parsed.get("host", "localhost"),
                "port": parsed.get("port") or SUPPORTED_DB_TYPES[db_type]["port"],
                "db_url": payload.db_url,
                "name": f"{db_type}://{parsed.get('host', 'localhost')}/{parsed.get('database', '')}",
                "status": "connected",
                "latency_ms": latency_ms,
            }
            display_name = f"{parsed.get('host', 'localhost')}:{parsed.get('database', '')}"
            normalized = {"db_type": db_type, "host": parsed.get("host"), "database": parsed.get("database")}

        connection_id = self.repo.save(cfg)
        self.repo.record_ping(connection_id, status="connected", detail="Connected", latency_ms=latency_ms)

        return {
            "ok": True,
            "details": f"Connected ({display_name})",
            "latency_ms": latency_ms,
            "connection_id": connection_id,
            "normalized": normalized,
            "correlation_id": correlation_id,
        }

    def upsert(self, payload: ConnectionUpsertRequest, correlation_id: str | None = None) -> ConnectionResponse:
        # Parse to determine type
        parsed = _parse_db_url(payload.db_url) if payload.db_url else {}
        db_type = payload.db_type or parsed.get("db_type") or "sqlite"

        connection_ref = self._resolve_and_verify(
            connection_id=payload.id,
            db_url=payload.db_url,
            db_path=payload.database,
            db_type=db_type,
            verify=False,
        )

        if db_type == "sqlite":
            db_path_str = str(connection_ref)
            name = payload.name or Path(connection_ref).name
        else:
            db_path_str = payload.db_url or str(connection_ref)
            name = payload.name or f"{db_type}://{parsed.get('host', 'localhost')}/{parsed.get('database', '')}"

        record = self.repo.upsert(
            conn_id=payload.id,
            name=name,
            db_type=db_type,
            database_path=db_path_str,
            secret_payload={"db_url": payload.db_url, "database": db_path_str},
            status=payload.status,
            latency_ms=payload.latency_ms,
            tags=payload.tags,
        )
        if payload.status:
            self.repo.record_ping(record["id"], status=payload.status, detail=None, latency_ms=payload.latency_ms)
        return ConnectionResponse(
            id=record["id"],
            name=record["name"],
            db_type=record["db_type"],
            database_path=Path(db_path_str) if db_type == "sqlite" else None,
            status=record.get("status") or "unknown",
            latency_ms=record.get("latency_ms"),
        )

    def delete(self, connection_id: str) -> None:
        if not self.repo.delete(connection_id):
            raise AppError(code="connection_not_found", message="Connection not found", status_code=404)

    def _execute_health_query(self, db_path: str, db_type: str) -> tuple[bool, str]:
        """Execute a simple query to verify database is actually working."""
        if db_type == "sqlite":
            import sqlite3
            conn = None
            try:
                conn = sqlite3.connect(db_path, timeout=5)
                cursor = conn.cursor()
                cursor.execute("SELECT 1")
                cursor.fetchone()
                return True, "Query executed successfully"
            except sqlite3.Error as e:
                logger.warning("sqlite_connection_test_failed", extra={"error": str(e)})
                return False, "SQLite connection test failed"
            finally:
                if conn:
                    conn.close()
        elif db_type in ("postgresql", "postgres"):
            from sqlalchemy import create_engine, text as sa_text
            engine = None
            try:
                engine = create_engine(db_path, connect_args={"connect_timeout": 5})
                with engine.connect() as conn:
                    conn.execute(sa_text("SELECT 1"))
                return True, "Query executed successfully"
            except Exception as e:
                logger.warning("postgres_connection_test_failed", extra={"error": str(e)})
                return False, "PostgreSQL connection test failed"
            finally:
                if engine:
                    engine.dispose()
        else:
            return True, "Port accessible"

    def healthcheck(self, connection_id: str, correlation_id: str | None = None) -> dict:
        """Verify a saved connection is still accessible and can execute queries."""
        secrets = self.repo.get_secrets(connection_id)
        if not secrets:
            raise AppError(code="connection_not_found", message="Connection not found", status_code=404)

        db_path = secrets.get("database_path") or secrets.get("database")
        db_url = secrets.get("db_url")
        if not db_path and not db_url:
            raise AppError(
                code="invalid_connection",
                message="Connection has no database path or URL",
                status_code=400,
            )

        # Determine database type
        parsed = _parse_db_url(db_url) if db_url else {}
        db_type = parsed.get("db_type") or "sqlite"

        started = time.time()

        # First check file/network accessibility
        try:
            if db_type == "sqlite":
                self.repo.verify(Path(db_path))
            else:
                self._verify_network_database(db_url, db_type)
        except Exception as exc:
            latency_ms = int((time.time() - started) * 1000)
            self.repo.record_ping(connection_id, status="error", detail=str(exc), latency_ms=latency_ms)
            raise AppError(
                code="connection_failed",
                message="Database connection failed",
                detail=None,
                status_code=503,
            )

        # Then execute a test query to verify database is actually working
        success, message = self._execute_health_query(db_path or db_url, db_type)
        latency_ms = int((time.time() - started) * 1000)

        if not success:
            self.repo.record_ping(connection_id, status="error", detail=message, latency_ms=latency_ms)
            raise AppError(
                code="connection_failed",
                message="Database query failed",
                detail=message,
                status_code=503,
            )

        self.repo.record_ping(connection_id, status="connected", detail="Health check passed", latency_ms=latency_ms)
        return {
            "status": "connected",
            "latency_ms": latency_ms,
            "connection_id": connection_id,
            "correlation_id": correlation_id,
        }

"""Export Service.

Handles document export to various formats and distribution.
"""

import html as html_mod
import io
import logging
import secrets
import uuid
import zipfile
from datetime import datetime, timezone
from typing import Any, Optional

from backend.app.services.config import get_settings

logger = logging.getLogger(__name__)


class ExportService:
    """Service for exporting documents to various formats."""

    def __init__(self):
        self._export_jobs: dict[str, dict] = {}
        self._embed_tokens: dict[str, dict] = {}
        base_root = get_settings().uploads_root
        self._uploads_dir = base_root / "exports"
        self._uploads_dir.mkdir(parents=True, exist_ok=True)

    async def export_to_pdf(
        self,
        content: bytes,
        options: dict[str, Any],
    ) -> bytes:
        """Export content to PDF format."""
        # If already PDF, apply options
        if options.get("pdfa_compliant"):
            return await self._convert_to_pdfa(content, options)
        return content

    async def _convert_to_pdfa(
        self,
        pdf_content: bytes,
        options: dict[str, Any],
    ) -> bytes:
        """Convert PDF to PDF/A format for archival."""
        try:
            import fitz
            doc = fitz.open(stream=pdf_content, filetype="pdf")
            # Add PDF/A metadata
            metadata = doc.metadata
            metadata["format"] = "PDF/A-1b"
            doc.set_metadata(metadata)
            output = io.BytesIO()
            doc.save(output)
            doc.close()
            return output.getvalue()
        except ImportError:
            logger.warning("PyMuPDF not available, returning original PDF")
            return pdf_content

    async def export_to_docx(
        self,
        content: str,
        options: dict[str, Any],
    ) -> bytes:
        """Export content to Word DOCX format."""
        try:
            from docx import Document
            from docx.shared import Inches, Pt

            doc = Document()

            # Add title if provided
            if options.get("title"):
                doc.add_heading(options["title"], 0)

            # Add content paragraphs
            for para in content.split("\n\n"):
                if para.strip():
                    doc.add_paragraph(para.strip())

            # Save to bytes
            output = io.BytesIO()
            doc.save(output)
            return output.getvalue()

        except ImportError:
            logger.warning("python-docx not available")
            raise ImportError("python-docx is required for DOCX export")

    async def export_to_pptx(
        self,
        content: str,
        options: dict[str, Any],
    ) -> bytes:
        """Export content to PowerPoint PPTX format."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt

            prs = Presentation()

            # Add title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = options.get("title", "Document Export")
            subtitle.text = options.get("subtitle", "Generated by NeuraReport")

            # Add content slides
            bullet_slide_layout = prs.slide_layouts[1]
            sections = content.split("\n\n")

            for i, section in enumerate(sections):
                if not section.strip():
                    continue

                slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]

                title_shape.text = f"Section {i + 1}"

                tf = body_shape.text_frame
                tf.text = section.strip()

            # Save to bytes
            output = io.BytesIO()
            prs.save(output)
            return output.getvalue()

        except ImportError:
            logger.warning("python-pptx not available")
            raise ImportError("python-pptx is required for PPTX export")

    async def export_to_epub(
        self,
        content: str,
        options: dict[str, Any],
    ) -> bytes:
        """Export content to ePub format."""
        try:
            from ebooklib import epub

            book = epub.EpubBook()

            # Set metadata
            book.set_identifier(str(uuid.uuid4()))
            book.set_title(options.get("title", "Document Export"))
            book.set_language(options.get("language", "en"))

            if options.get("author"):
                book.add_author(options["author"])

            # Create chapter
            chapter = epub.EpubHtml(
                title="Content",
                file_name="content.xhtml",
                lang="en",
            )
            safe_title = html_mod.escape(options.get('title', 'Document'))
            chapter.content = f"<html><body><h1>{safe_title}</h1>{content}</body></html>"
            book.add_item(chapter)

            # Add navigation
            book.toc = [epub.Link("content.xhtml", "Content", "content")]
            book.add_item(epub.EpubNcx())
            book.add_item(epub.EpubNav())

            book.spine = ["nav", chapter]

            # Write to bytes
            output = io.BytesIO()
            epub.write_epub(output, book)
            return output.getvalue()

        except ImportError:
            logger.warning("ebooklib not available")
            raise ImportError("ebooklib is required for ePub export")

    @staticmethod
    def _escape_latex(text: str) -> str:
        """Escape LaTeX special characters to prevent injection."""
        replacements = [
            ('\\', '\\textbackslash{}'),
            ('{', '\\{'), ('}', '\\}'),
            ('$', '\\$'), ('#', '\\#'), ('%', '\\%'),
            ('&', '\\&'), ('^', '\\textasciicircum{}'),
            ('_', '\\_'), ('~', '\\textasciitilde{}'),
        ]
        for old, new in replacements:
            text = text.replace(old, new)
        return text

    async def export_to_latex(
        self,
        content: str,
        options: dict[str, Any],
    ) -> bytes:
        """Export content to LaTeX format."""
        doc_class = self._escape_latex(options.get("document_class", "article"))
        title = self._escape_latex(options.get("title", "Document"))
        author = self._escape_latex(options.get("author", ""))

        latex_content = f"""\\documentclass{{{doc_class}}}
\\usepackage[utf8]{{inputenc}}
\\usepackage{{graphicx}}
\\usepackage{{hyperref}}

\\title{{{title}}}
\\author{{{author}}}
\\date{{\\today}}

\\begin{{document}}

\\maketitle

{content}

\\end{{document}}
"""
        return latex_content.encode("utf-8")

    async def export_to_markdown(
        self,
        content: str,
        options: dict[str, Any],
    ) -> bytes:
        """Export content to Markdown format."""
        flavor = options.get("flavor", "gfm")
        include_frontmatter = options.get("include_frontmatter", True)

        md_content = ""

        if include_frontmatter:
            md_content = f"""---
title: {options.get('title', 'Document')}
author: {options.get('author', '')}
date: {utc_now().strftime('%Y-%m-%d')}
---

"""

        md_content += content

        return md_content.encode("utf-8")

    async def export_to_html(
        self,
        content: str,
        options: dict[str, Any],
    ) -> bytes:
        """Export content to HTML format."""
        title = html_mod.escape(options.get("title", "Document"))
        standalone = options.get("standalone", True)

        if standalone:
            html_content = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>
        body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; max-width: 800px; margin: 0 auto; padding: 20px; }}
        h1 {{ color: #333; }}
        p {{ line-height: 1.6; }}
    </style>
</head>
<body>
    <h1>{title}</h1>
    {content}
</body>
</html>
"""
        else:
            html_content = content

        return html_content.encode("utf-8")

    async def create_export_job(
        self,
        document_id: str,
        format: str,
        options: dict[str, Any],
    ) -> dict:
        """Create an export job."""
        job_id = str(uuid.uuid4())
        now = utc_now()

        job = {
            "job_id": job_id,
            "document_id": document_id,
            "format": format,
            "options": options,
            "status": "pending",
            "created_at": now,
            "completed_at": None,
            "download_url": None,
            "file_size": None,
            "error": None,
        }

        self._export_jobs[job_id] = job

        # Persist to state store
        try:
            from backend.app.repositories import state_store
            with state_store.transaction() as state:
                jobs = state.get("export_jobs", {})
                jobs[job_id] = job
                state["export_jobs"] = jobs
        except Exception as e:
            logger.warning(f"Failed to persist export job: {e}")

        return job

    async def get_export_job(self, job_id: str) -> Optional[dict]:
        """Get export job status."""
        job = self._export_jobs.get(job_id)
        if job is not None:
            return job
        try:
            from backend.app.repositories import state_store
            with state_store.transaction() as state:
                return state.get("export_jobs", {}).get(job_id)
        except Exception:
            return None

    async def bulk_export(
        self,
        document_ids: list[str],
        format: str,
        options: dict[str, Any],
    ) -> dict:
        """Export multiple documents as a ZIP file."""
        job_id = str(uuid.uuid4())
        now = utc_now()

        # Create ZIP file in memory
        zip_buffer = io.BytesIO()

        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
            for i, doc_id in enumerate(document_ids):
                # Add placeholder files
                filename = f"document_{i + 1}.{format}"
                zf.writestr(filename, f"Content of document {doc_id}")

        zip_content = zip_buffer.getvalue()

        # Save ZIP file
        zip_filename = f"export_{job_id}.zip"
        zip_path = self._uploads_dir / zip_filename
        zip_path.write_bytes(zip_content)

        job = {
            "job_id": job_id,
            "document_ids": document_ids,
            "format": format,
            "status": "completed",
            "created_at": now,
            "completed_at": utc_now(),
            "download_url": f"/uploads/exports/{zip_filename}",
            "file_size": len(zip_content),
        }

        self._export_jobs[job_id] = job

        return job

    async def generate_embed_token(
        self,
        document_id: str,
        options: dict[str, Any],
    ) -> dict:
        """Generate an embed token for a document."""
        token = secrets.token_urlsafe(32)
        now = utc_now()

        embed_data = {
            "token": token,
            "document_id": document_id,
            "options": options,
            "created_at": now,
            "expires_at": None,  # Could add expiration
        }

        self._embed_tokens[token] = embed_data

        # Generate embed URL and code
        embed_url = f"/embed/{token}"
        width = options.get("width", 800)
        height = options.get("height", 600)

        embed_code = f'<iframe src="{embed_url}" width="{width}" height="{height}" frameborder="0"></iframe>'

        return {
            "token": token,
            "embed_url": embed_url,
            "embed_code": embed_code,
            "expires_at": embed_data["expires_at"],
        }

    async def validate_embed_token(self, token: str) -> Optional[dict]:
        """Validate an embed token."""
        return self._embed_tokens.get(token)

    async def revoke_embed_token(self, token_id: str) -> bool:
        """Revoke an embed token. Returns True if the token existed and was revoked."""
        if token_id in self._embed_tokens:
            del self._embed_tokens[token_id]
            return True
        return False

    async def list_embed_tokens(self, document_id: str) -> list[dict]:
        """List all embed tokens for a given document."""
        return [
            token_data
            for token_data in self._embed_tokens.values()
            if token_data.get("document_id") == document_id
        ]

    async def print_document(
        self,
        document_id: str,
        printer_id: Optional[str] = None,
        copies: int = 1,
        options: Optional[dict[str, Any]] = None,
    ) -> dict:
        """Send a document to a printer."""
        job_id = str(uuid.uuid4())
        now = utc_now()

        job = {
            "job_id": job_id,
            "document_id": document_id,
            "printer_id": printer_id or "default",
            "copies": copies,
            "options": options or {},
            "status": "queued",
            "created_at": now,
        }

        self._export_jobs[job_id] = job
        logger.info(f"Print job {job_id} queued for document {document_id}")
        return job

    async def list_printers(self) -> list[dict]:
        """List available printers."""
        # In a real implementation this would query network/system printers.
        return [
            {
                "printer_id": "default",
                "name": "Default Printer",
                "status": "ready",
                "location": "Local",
            },
        ]

    async def list_export_jobs(
        self,
        status: Optional[str] = None,
        format: Optional[str] = None,
        limit: int = 50,
        offset: int = 0,
    ) -> dict:
        """List export jobs with optional filtering."""
        jobs = list(self._export_jobs.values())

        if status:
            jobs = [j for j in jobs if j.get("status") == status]
        if format:
            jobs = [j for j in jobs if j.get("format") == format]

        total = len(jobs)
        jobs = jobs[offset : offset + limit]

        return {"jobs": jobs, "total": total, "limit": limit, "offset": offset}

    async def cancel_export_job(self, job_id: str) -> Optional[dict]:
        """Cancel an export job. Returns the updated job or None if not found."""
        job = self._export_jobs.get(job_id)
        if job is None:
            return None
        if job.get("status") in ("completed", "failed"):
            return job  # Cannot cancel a finished job; caller should check status.
        job["status"] = "cancelled"
        job["completed_at"] = utc_now()
        return job

class DistributionService:
    """Service for distributing documents to various channels."""

    def __init__(self):
        self._distribution_jobs: dict[str, dict] = {}

    async def send_email(
        self,
        document_id: str,
        recipients: list[str],
        subject: str,
        message: str,
        attachments: list[bytes] = None,
    ) -> dict:
        """Send document via email."""
        job_id = str(uuid.uuid4())
        now = utc_now()

        # Email sending would be implemented here
        # Using existing mailer service if available

        job = {
            "job_id": job_id,
            "channel": "email",
            "document_id": document_id,
            "recipients": recipients,
            "status": "sent",
            "sent_at": now,
            "recipients_count": len(recipients),
        }

        self._distribution_jobs[job_id] = job

        logger.info(f"Email sent to {len(recipients)} recipients")

        return job

    async def send_to_slack(
        self,
        document_id: str,
        channel: str,
        message: Optional[str] = None,
        file_content: Optional[bytes] = None,
    ) -> dict:
        """Send document to Slack channel."""
        job_id = str(uuid.uuid4())
        now = utc_now()

        try:
            from slack_sdk import WebClient

            # Would need SLACK_BOT_TOKEN environment variable
            import os
            token = os.getenv("SLACK_BOT_TOKEN")

            if token:
                client = WebClient(token=token)

                if file_content:
                    client.files_upload_v2(
                        channel=channel,
                        content=file_content,
                        title=f"Document {document_id}",
                        initial_comment=message or "Document shared from NeuraReport",
                    )
                else:
                    client.chat_postMessage(
                        channel=channel,
                        text=message or f"Document {document_id} shared from NeuraReport",
                    )

                status = "sent"
            else:
                status = "skipped"
                logger.warning("Slack token not configured")

        except ImportError:
            status = "error"
            logger.warning("slack_sdk not available")
        except Exception as e:
            status = "error"
            logger.error(f"Slack send failed: {e}")

        job = {
            "job_id": job_id,
            "channel": "slack",
            "document_id": document_id,
            "slack_channel": channel,
            "status": status,
            "sent_at": now,
        }

        self._distribution_jobs[job_id] = job

        return job

    async def send_to_teams(
        self,
        document_id: str,
        webhook_url: str,
        title: Optional[str] = None,
        message: Optional[str] = None,
    ) -> dict:
        """Send document to Microsoft Teams via webhook."""
        job_id = str(uuid.uuid4())
        now = utc_now()

        try:
            import aiohttp

            payload = {
                "@type": "MessageCard",
                "@context": "http://schema.org/extensions",
                "themeColor": "0076D7",
                "summary": title or f"Document {document_id}",
                "sections": [
                    {
                        "activityTitle": title or "Document Shared",
                        "text": message or f"Document {document_id} shared from NeuraReport",
                    }
                ],
            }

            async with aiohttp.ClientSession() as session:
                async with session.post(webhook_url, json=payload) as resp:
                    if resp.status == 200:
                        status = "sent"
                    else:
                        status = "error"
                        logger.warning(f"Teams webhook returned {resp.status}")

        except Exception as e:
            status = "error"
            logger.error(f"Teams send failed: {e}")

        job = {
            "job_id": job_id,
            "channel": "teams",
            "document_id": document_id,
            "status": status,
            "sent_at": now,
        }

        self._distribution_jobs[job_id] = job

        return job

    async def send_webhook(
        self,
        document_id: str,
        webhook_url: str,
        method: str = "POST",
        headers: dict[str, str] = None,
        payload: dict = None,
    ) -> dict:
        """Send document data via webhook."""
        job_id = str(uuid.uuid4())
        now = utc_now()

        try:
            import aiohttp

            if payload is None:
                payload = {
                    "document_id": document_id,
                    "timestamp": now.isoformat(),
                    "source": "neura_report",
                }

            async with aiohttp.ClientSession() as session:
                async with session.request(
                    method,
                    webhook_url,
                    json=payload,
                    headers=headers,
                ) as resp:
                    status = "sent" if resp.status < 400 else "error"
                    response_status = resp.status

        except Exception as e:
            status = "error"
            response_status = None
            logger.error(f"Webhook send failed: {e}")

        job = {
            "job_id": job_id,
            "channel": "webhook",
            "document_id": document_id,
            "webhook_url": webhook_url,
            "status": status,
            "response_status": response_status,
            "sent_at": now,
        }

        self._distribution_jobs[job_id] = job

        return job

    async def publish_to_portal(
        self,
        document_id: str,
        portal_path: str,
        options: dict[str, Any],
    ) -> dict:
        """Publish document to portal."""
        job_id = str(uuid.uuid4())
        now = utc_now()

        # Portal publishing would be implemented here
        portal_url = f"/portal/{portal_path}/{document_id}"

        job = {
            "job_id": job_id,
            "channel": "portal",
            "document_id": document_id,
            "portal_path": portal_path,
            "portal_url": portal_url,
            "public": options.get("public", False),
            "status": "published",
            "published_at": now,
        }

        self._distribution_jobs[job_id] = job

        return job

# Singleton instances
export_service = ExportService()
distribution_service = DistributionService()

"""
pgvector document store for semantic search and RAG.

Uses PostgreSQL's pgvector extension for:
- Cosine similarity search with HNSW indexing
- Hybrid search (vector + full-text via RRF)
- Source attribution for RAG outputs

Based on: pgvector/pgvector patterns + pgvector-python
"""

import logging
import json
from typing import Any, Optional
from dataclasses import dataclass

logger = logging.getLogger("neura.vectorstore")

@dataclass
class SearchResult:
    """A single search result with source attribution."""
    chunk_id: int
    doc_id: str
    chunk_index: int
    content: str
    source: str
    metadata: dict[str, Any]
    similarity: float

    def to_citation(self, index: int) -> dict:
        return {
            "index": index,
            "source": self.source,
            "content_preview": self.content[:200],
            "similarity": round(self.similarity, 4),
            "doc_id": self.doc_id,
        }

class PgVectorStore:
    """
    Vector store using PostgreSQL + pgvector extension.

    Schema:
    - document_chunks table with vector(dimensions) column
    - HNSW index for approximate nearest neighbor search
    - Full-text search via tsvector for hybrid retrieval
    """

    def __init__(self, connection_string: str, embedding_dim: int | None = None):
        from backend.app.services.config import get_settings

        self.connection_string = connection_string
        self.embedding_dim = int(embedding_dim or get_settings().embedding_dim)

    async def ensure_schema(self, session) -> None:
        """Create pgvector extension and document_chunks table if not exists."""
        await session.execute("CREATE EXTENSION IF NOT EXISTS vector")
        await session.execute(f"""
            CREATE TABLE IF NOT EXISTS document_chunks (
                id BIGSERIAL PRIMARY KEY,
                doc_id VARCHAR(255) NOT NULL,
                chunk_index INTEGER NOT NULL,
                content TEXT NOT NULL,
                embedding vector({self.embedding_dim}),
                source VARCHAR(500),
                metadata JSONB DEFAULT '{{}}',
                created_at TIMESTAMPTZ DEFAULT NOW(),
                UNIQUE(doc_id, chunk_index)
            )
        """)
        await session.execute("""
            CREATE INDEX IF NOT EXISTS idx_chunks_doc_id ON document_chunks (doc_id)
        """)
        # Full-text search index
        await session.execute("""
            CREATE INDEX IF NOT EXISTS idx_chunks_fts
            ON document_chunks USING gin(to_tsvector('english', content))
        """)
        await session.commit()
        logger.info("vectorstore_schema_ensured", extra={"event": "vectorstore_schema_ensured"})

    async def create_hnsw_index(self, session, m: int = 16, ef_construction: int = 64) -> None:
        """Create HNSW index. Call AFTER initial data load for best performance."""
        await session.execute(f"""
            CREATE INDEX IF NOT EXISTS idx_chunks_embedding_cosine
            ON document_chunks
            USING hnsw (embedding vector_cosine_ops)
            WITH (m = {m}, ef_construction = {ef_construction})
        """)
        await session.commit()
        logger.info("hnsw_index_created", extra={"event": "hnsw_index_created", "m": m, "ef_construction": ef_construction})

    async def upsert_chunks(self, session, chunks: list[dict[str, Any]]) -> int:
        """Upsert document chunks with embeddings."""
        count = 0
        for chunk in chunks:
            await session.execute(
                """
                INSERT INTO document_chunks (doc_id, chunk_index, content, embedding, source, metadata)
                VALUES (:doc_id, :chunk_index, :content, :embedding, :source, :metadata::jsonb)
                ON CONFLICT (doc_id, chunk_index)
                DO UPDATE SET content = EXCLUDED.content, embedding = EXCLUDED.embedding,
                             source = EXCLUDED.source, metadata = EXCLUDED.metadata
                """,
                {
                    "doc_id": chunk["doc_id"],
                    "chunk_index": chunk["chunk_index"],
                    "content": chunk["content"],
                    "embedding": str(chunk["embedding"]),
                    "source": chunk.get("source", ""),
                    "metadata": json.dumps(chunk.get("metadata", {})),
                },
            )
            count += 1
        await session.commit()
        return count

    async def search_similar(
        self, session, query_embedding: list[float], top_k: int = 10,
        source_filter: Optional[str] = None, ef_search: int = 100,
    ) -> list[SearchResult]:
        """Cosine similarity search with optional source filtering."""
        await session.execute(f"SET LOCAL hnsw.ef_search = {ef_search}")

        query_vec = str(query_embedding)

        if source_filter:
            rows = await session.execute(
                """
                SELECT id, doc_id, chunk_index, content, source, metadata,
                       1 - (embedding <=> :vec) AS similarity
                FROM document_chunks WHERE source = :source
                ORDER BY embedding <=> :vec LIMIT :top_k
                """,
                {"vec": query_vec, "source": source_filter, "top_k": top_k},
            )
        else:
            rows = await session.execute(
                """
                SELECT id, doc_id, chunk_index, content, source, metadata,
                       1 - (embedding <=> :vec) AS similarity
                FROM document_chunks ORDER BY embedding <=> :vec LIMIT :top_k
                """,
                {"vec": query_vec, "top_k": top_k},
            )

        return [
            SearchResult(
                chunk_id=row.id, doc_id=row.doc_id, chunk_index=row.chunk_index,
                content=row.content, source=row.source,
                metadata=row.metadata if isinstance(row.metadata, dict) else {},
                similarity=float(row.similarity),
            )
            for row in rows.fetchall()
        ]

    async def hybrid_search(
        self, session, query_embedding: list[float], query_text: str,
        top_k: int = 10, vector_weight: float = 0.7,
    ) -> list[SearchResult]:
        """Reciprocal Rank Fusion hybrid search (vector + full-text)."""
        k = 60  # RRF constant
        query_vec = str(query_embedding)

        rows = await session.execute(
            """
            WITH vector_results AS (
                SELECT id, doc_id, chunk_index, content, source, metadata,
                       ROW_NUMBER() OVER (ORDER BY embedding <=> :vec) AS vector_rank
                FROM document_chunks ORDER BY embedding <=> :vec LIMIT :limit
            ),
            text_results AS (
                SELECT id, doc_id, chunk_index, content, source, metadata,
                       ROW_NUMBER() OVER (
                           ORDER BY ts_rank_cd(to_tsvector('english', content), plainto_tsquery('english', :query)) DESC
                       ) AS text_rank
                FROM document_chunks
                WHERE to_tsvector('english', content) @@ plainto_tsquery('english', :query)
                LIMIT :limit
            )
            SELECT COALESCE(v.id, t.id) AS id,
                   COALESCE(v.doc_id, t.doc_id) AS doc_id,
                   COALESCE(v.chunk_index, t.chunk_index) AS chunk_index,
                   COALESCE(v.content, t.content) AS content,
                   COALESCE(v.source, t.source) AS source,
                   COALESCE(v.metadata, t.metadata) AS metadata,
                   COALESCE(:vw / (:k + v.vector_rank), 0) + COALESCE((1 - :vw) / (:k + t.text_rank), 0) AS similarity
            FROM vector_results v FULL OUTER JOIN text_results t ON v.id = t.id
            ORDER BY similarity DESC LIMIT :top_k
            """,
            {"vec": query_vec, "query": query_text, "limit": top_k * 2, "vw": vector_weight, "k": k, "top_k": top_k},
        )

        return [
            SearchResult(
                chunk_id=row.id, doc_id=row.doc_id, chunk_index=row.chunk_index,
                content=row.content, source=row.source,
                metadata=row.metadata if isinstance(row.metadata, dict) else {},
                similarity=float(row.similarity),
            )
            for row in rows.fetchall()
        ]

    async def delete_document(self, session, doc_id: str) -> int:
        """Delete all chunks for a document."""
        result = await session.execute(
            "DELETE FROM document_chunks WHERE doc_id = :doc_id",
            {"doc_id": doc_id},
        )
        await session.commit()
        return result.rowcount

"""
Embedding pipeline: chunk text, generate embeddings, store in vector DB.

Uses sentence-transformers for local embedding generation (no API key needed).
Falls back to TF-IDF vector hashing when sentence-transformers is unavailable.
"""
import hashlib
import logging
import math
import re
from typing import Any, Optional

logger = logging.getLogger("neura.vectorstore.embedding")

# Lazy-loaded sentence-transformers model
_st_model = None
_st_model_name: str = ""

def _get_sentence_transformer(model_name: str = "all-MiniLM-L6-v2"):
    """Lazy-load a sentence-transformers model."""
    global _st_model, _st_model_name
    if _st_model is not None and _st_model_name == model_name:
        return _st_model
    try:
        from sentence_transformers import SentenceTransformer
        _st_model = SentenceTransformer(model_name)
        _st_model_name = model_name
        logger.info("sentence_transformer_loaded", extra={"event": "sentence_transformer_loaded", "model": model_name})
        return _st_model
    except ImportError:
        logger.warning("sentence_transformers_not_installed", extra={"event": "sentence_transformers_not_installed"})
        return None

def _tfidf_hash_embedding(text: str, dim: int = 384) -> list[float]:
    """Generate a deterministic pseudo-embedding via token hashing (fallback)."""
    tokens = re.findall(r"\b\w+\b", text.lower())
    vec = [0.0] * dim
    for token in tokens:
        h = int(hashlib.md5(token.encode()).hexdigest(), 16)
        idx = h % dim
        sign = 1.0 if (h // dim) % 2 == 0 else -1.0
        vec[idx] += sign
    # L2 normalize
    norm = math.sqrt(sum(v * v for v in vec)) or 1.0
    return [v / norm for v in vec]

class EmbeddingPipeline:
    """Generate embeddings using sentence-transformers (local, no API key)."""

    def __init__(self, model: str | None = None, embedding_dim: int | None = None):
        from backend.app.services.config import get_settings

        settings = get_settings()
        self.model = model or settings.embedding_model
        self.embedding_dim = int(embedding_dim or settings.embedding_dim)

    def chunk_text(self, text: str, chunk_size: int = 512, chunk_overlap: int = 50) -> list[str]:
        """Split text into overlapping chunks."""
        chunks = []
        start = 0
        while start < len(text):
            end = start + chunk_size
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            start = end - chunk_overlap
        return chunks

    async def generate_embeddings(self, texts: list[str]) -> list[list[float]]:
        """Generate embeddings using sentence-transformers or TF-IDF fallback."""
        model = _get_sentence_transformer(self.model)
        if model is not None:
            embeddings = model.encode(texts, show_progress_bar=False)
            return [emb.tolist() for emb in embeddings]

        # Fallback: TF-IDF hash embeddings
        logger.debug("using_tfidf_fallback", extra={"event": "using_tfidf_fallback"})
        return [_tfidf_hash_embedding(t, self.embedding_dim) for t in texts]

    async def process_document(
        self, doc_id: str, content: str, source: str,
        metadata: Optional[dict[str, Any]] = None,
    ) -> list[dict]:
        """Full pipeline: chunk -> embed -> return records for storage."""
        chunks = self.chunk_text(content)
        if not chunks:
            return []

        embeddings = await self.generate_embeddings(chunks)

        records = []
        for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
            records.append({
                "doc_id": doc_id,
                "chunk_index": i,
                "content": chunk,
                "embedding": embedding,
                "source": source,
                "metadata": metadata or {},
            })

        logger.info("document_embedded", extra={
            "event": "document_embedded", "doc_id": doc_id,
            "chunks": len(records), "model": self.model,
        })
        return records

    async def embed_query(self, query: str) -> list[float]:
        """Embed a single query string."""
        embeddings = await self.generate_embeddings([query])
        return embeddings[0]

"""
RAG pipeline: retrieve relevant chunks, generate answer with source citations.

Based on: Haystack PromptBuilder + AnswerBuilder + source attribution patterns.
"""
import re
import logging
from typing import Any, Optional
from dataclasses import dataclass

logger = logging.getLogger("neura.vectorstore.rag")

RAG_PROMPT_TEMPLATE = """You are a helpful research assistant. Answer the question using ONLY the
provided context documents. For every claim, cite the source using [N] where N is the document index.

If the context does not contain enough information, say "I don't have enough information" and explain what is missing.

Context documents:
{context}

Question: {query}

Answer (with citations):"""

@dataclass
class RAGResponse:
    """RAG response with source attribution."""
    answer: str
    sources: list[dict[str, Any]]
    referenced_sources: list[dict[str, Any]]
    confidence: float = 0.0

    def to_dict(self) -> dict:
        return {
            "answer": self.answer,
            "sources": self.sources,
            "referenced_sources": self.referenced_sources,
            "confidence": self.confidence,
        }

class RAGPipeline:
    """End-to-end RAG: embed query -> retrieve -> generate with citations."""

    def __init__(self, vector_store, embedding_pipeline, llm_model: str = "qwen"):
        self.vector_store = vector_store
        self.embedding_pipeline = embedding_pipeline
        self.llm_model = llm_model

    async def query(
        self, session, question: str, top_k: int = 5,
        source_filter: Optional[str] = None, use_hybrid: bool = True,
    ) -> RAGResponse:
        """Run a RAG query with source attribution."""
        # 1. Embed query
        query_embedding = await self.embedding_pipeline.embed_query(question)

        # 2. Retrieve relevant chunks
        if use_hybrid:
            results = await self.vector_store.hybrid_search(session, query_embedding, question, top_k=top_k)
        else:
            results = await self.vector_store.search_similar(session, query_embedding, top_k=top_k, source_filter=source_filter)

        if not results:
            return RAGResponse(answer="No relevant documents found.", sources=[], referenced_sources=[])

        # 3. Build context with numbered sources
        context_parts = []
        sources = []
        for i, result in enumerate(results, 1):
            context_parts.append(f"[{i}] (Source: {result.source})\n{result.content}\n")
            sources.append(result.to_citation(i))

        context = "\n".join(context_parts)
        prompt = RAG_PROMPT_TEMPLATE.format(context=context, query=question)

        # 4. Generate answer using centralized LLM client
        from backend.app.services.llm import get_llm_client
        client = get_llm_client()

        response = client.complete(
            messages=[{"role": "user", "content": prompt}],
            model=self.llm_model,
            description="rag_generate_answer",
            temperature=0.1,
            max_tokens=1024,
        )
        answer = (
            response.get("choices", [{}])[0]
            .get("message", {})
            .get("content", "")
        )

        # 5. Extract referenced sources from [N] patterns
        referenced_indices = set(int(m) for m in re.findall(r"\[(\d+)\]", answer))
        referenced_sources = [s for s in sources if s["index"] in referenced_indices]

        # 6. Compute confidence (average similarity of referenced sources)
        if referenced_sources:
            confidence = sum(s["similarity"] for s in referenced_sources) / len(referenced_sources)
        else:
            confidence = sum(s["similarity"] for s in sources) / len(sources) if sources else 0.0

        return RAGResponse(
            answer=answer,
            sources=sources,
            referenced_sources=referenced_sources,
            confidence=round(confidence, 4),
        )

"""Service utilities — core helpers + LLM utils."""

# mypy: ignore-errors
"""
Core utility functions - merged from validation, zip_tools, render, html, mailer,
artifacts, lock, tokens, text, and context modules.
"""

import contextvars
import hashlib
import json
import logging
import mimetypes
import os
import re
import shutil
import smtplib
import ssl
import zipfile
from collections.abc import Iterable, Mapping
from contextlib import contextmanager
from copy import deepcopy
from dataclasses import dataclass
from datetime import datetime, timezone
from email.message import EmailMessage
from itertools import zip_longest
from pathlib import Path
from typing import Any, Dict, Generator, List, Optional, Sequence, Tuple

try:
    from jsonschema import Draft7Validator  # type: ignore
except ImportError:  # pragma: no cover
    Draft7Validator = None  # type: ignore

try:
    import bleach
    from bleach.css_sanitizer import CSSSanitizer
except ImportError:  # pragma: no cover
    bleach = None  # type: ignore
    CSSSanitizer = None  # type: ignore

try:
    from filelock import FileLock, Timeout as FileLockTimeout
except ImportError:  # pragma: no cover
    FileLock = None  # type: ignore
    FileLockTimeout = None  # type: ignore

try:
    from PIL import Image  # type: ignore
except ImportError:  # pragma: no cover
    Image = None  # type: ignore

try:
    from playwright.sync_api import sync_playwright  # type: ignore
except ImportError:  # pragma: no cover
    sync_playwright = None  # type: ignore

from backend.app.utils import write_json_atomic, write_text_atomic

logger = logging.getLogger("neura.utils")

_correlation_id: contextvars.ContextVar[Optional[str]] = contextvars.ContextVar("correlation_id", default=None)

def set_correlation_id(value: Optional[str]) -> None:
    _correlation_id.set(value)

def get_correlation_id() -> Optional[str]:
    return _correlation_id.get()

_DOUBLE_BRACE_PATTERN = re.compile(r"\{\{\s*([^{}]+?)\s*\}\}")
_TOKEN_PATTERN = re.compile(r"\{\{[^{}]+\}\}|\{[^{}]+\}")
_TOKEN_NAME_PATTERN = re.compile(r"^[A-Za-z0-9_][A-Za-z0-9_\-\.]*$")

def normalize_token_braces(text: str) -> str:
    """
    Convert double-braced placeholders like ``{{token}}`` into single-braced
    ``{token}`` while leaving existing single-braced tokens untouched.
    Non-token usages (e.g., empty braces) are returned unchanged.
    """
    if not text:
        return text

    def _replace(match: re.Match[str]) -> str:
        inner = match.group(1).strip()
        if not inner:
            return match.group(0)
        return "{" + inner + "}"

    return _DOUBLE_BRACE_PATTERN.sub(_replace, text)

def extract_tokens(html: str) -> List[str]:
    """
    Return a list of token names found in the HTML text.
    Tokens wrapped in either `{token}` or `{{ token }}` are recognized.
    """
    if not html:
        return []
    normalized = normalize_token_braces(html)
    tokens: list[str] = []
    for match in _TOKEN_PATTERN.findall(normalized):
        token = match.strip()
        if token.startswith("{{") and token.endswith("}}"):
            token = token[2:-2]
        elif token.startswith("{") and token.endswith("}"):
            token = token[1:-1]
        token = token.strip()
        if not token:
            continue
        if not _TOKEN_NAME_PATTERN.match(token):
            continue
        tokens.append(token)
    return tokens

__all__ = ["normalize_token_braces", "extract_tokens"]


def extract_json_object(text: str) -> dict | None:
    """Extract a JSON object from LLM text that may contain prose before/after.

    Strategy:
    1. Try strip_code_fences first (handles ```json ... ```)
    2. Try json.loads on the result
    3. Fallback: walk through the string finding each '{' and try to parse
       from that position (handles prose with stray braces like {token_name})
    """
    if not text:
        return None

    cleaned = strip_code_fences(text)

    # Fast path: cleaned text is already valid JSON
    try:
        obj = json.loads(cleaned)
        if isinstance(obj, dict):
            return obj
    except (json.JSONDecodeError, ValueError):
        pass

    # Walk through finding each '{' and try to parse a complete JSON object
    search_in = cleaned if cleaned != text.strip() else text
    pos = 0
    while True:
        idx = search_in.find("{", pos)
        if idx == -1:
            break
        try:
            obj, end = json.JSONDecoder().raw_decode(search_in, idx)
            if isinstance(obj, dict):
                return obj
        except (json.JSONDecodeError, ValueError):
            pass
        pos = idx + 1

    return None

__all__ = ["strip_code_fences", "extract_json_object"]

logger = logging.getLogger("neura.artifacts")

MANIFEST_NAME = "artifact_manifest.json"
MANIFEST_SCHEMA_VERSION = "1.0"

def compute_checksums(files: Mapping[str, Path]) -> dict[str, str]:
    checksums: dict[str, str] = {}
    for name, path in files.items():
        if not path.exists():
            continue
        h = hashlib.sha256()
        with path.open("rb") as handle:
            for chunk in iter(lambda: handle.read(65536), b""):
                h.update(chunk)
        checksums[name] = h.hexdigest()
    return checksums

def write_artifact_manifest(
    template_dir: Path,
    *,
    step: str,
    files: Mapping[str, Path],
    inputs: Iterable[str],
    correlation_id: str | None = None,
) -> Path:
    """
    Persist artifact manifest alongside template artifacts.
    """
    template_dir = template_dir.resolve()
    manifest_path = template_dir / MANIFEST_NAME

    existing_payload: dict[str, Any] = {}
    if manifest_path.exists():
        try:
            existing_payload = json.loads(manifest_path.read_text(encoding="utf-8"))
        except Exception:
            logger.exception(
                "manifest_read_failed",
                extra={"event": "manifest_read_failed", "path": str(manifest_path)},
            )
            existing_payload = {}

    existing_files = dict(existing_payload.get("files") or {})
    existing_checksums = dict(existing_payload.get("file_checksums") or {})
    existing_inputs = list(existing_payload.get("input_refs") or [])

    def _store_path(path: Path) -> str:
        try:
            return str(path.resolve().relative_to(template_dir))
        except ValueError:
            return str(path)

    merged_files = existing_files | {name: _store_path(path) for name, path in files.items()}
    new_checksums = compute_checksums(files)
    merged_checksums = existing_checksums | new_checksums

    input_list = list(inputs)
    merged_inputs: list[str] = []
    for item in existing_inputs + input_list:
        if item not in merged_inputs:
            merged_inputs.append(item)

    payload = {
        "schema_version": MANIFEST_SCHEMA_VERSION,
        "produced_at": datetime.now(timezone.utc).isoformat(),
        "step": step,
        "files": merged_files,
        "file_checksums": merged_checksums,
        "input_refs": merged_inputs,
        "correlation_id": correlation_id,
        "pid": os.getpid(),
    }
    write_json_atomic(manifest_path, payload, indent=2, ensure_ascii=False, sort_keys=True, step="artifact_manifest")
    logger.info(
        "artifact_manifest_written",
        extra={
            "event": "artifact_manifest_written",
            "template_dir": str(template_dir),
            "step": step,
            "correlation_id": correlation_id,
        },
    )
    return manifest_path

def load_manifest(template_dir: Path) -> dict | None:
    path = template_dir / MANIFEST_NAME
    if not path.exists():
        return None
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        logger.exception(
            "manifest_load_failed",
            extra={"event": "manifest_load_failed", "path": str(path)},
        )
        return None

ALLOWED_TAGS = {
    "html",
    "head",
    "body",
    "meta",
    "title",
    "link",
    "style",
    "div",
    "span",
    "section",
    "article",
    "header",
    "footer",
    "main",
    "table",
    "thead",
    "tbody",
    "tfoot",
    "tr",
    "th",
    "td",
    "colgroup",
    "col",
    "p",
    "h1",
    "h2",
    "h3",
    "h4",
    "h5",
    "h6",
    "ul",
    "ol",
    "li",
    "strong",
    "em",
    "b",
    "i",
    "u",
    "br",
    "hr",
    "img",
    "canvas",
    "svg",
}

_REPEAT_COMMENT_RE = re.compile(r"^\s*(BEGIN:BLOCK_REPEAT\b.*|END:BLOCK_REPEAT\b.*)\s*$", re.IGNORECASE)

ALLOWED_ATTRS = {
    "*": {
        "class",
        "style",
        "id",
        "colspan",
        "rowspan",
        "align",
        "valign",
        "width",
        "height",
        "data-title",
        "data-index",
        "data-name",
        "data-value",
        "data-label",
    },
    "img": {"src", "alt"},
    "meta": {"charset"},
    "link": {"rel", "href"},
    "svg": {"viewbox", "xmlns"},
    "canvas": {"width", "height"},
}

_CSS_SANITIZER = CSSSanitizer()
_COMMENT_RE = re.compile(r"<!--(.*?)-->", re.DOTALL)
_STYLE_ATTR_RE = re.compile(r'style="([^"]*?)"')

# Matches CSS rules that apply position:fixed to footer-like selectors.
# Rewrites them to use normal document flow to prevent overlap with content.
_FIXED_FOOTER_RE = re.compile(
    r"(\.footer[-\w]*|#report-footer|footer)\s*\{([^}]*?)position\s*:\s*fixed\b([^}]*?)\}",
    re.IGNORECASE | re.DOTALL,
)
_FIXED_POS_PROPS_RE = re.compile(
    r"\b(bottom|left|right)\s*:\s*[^;]+;?",
    re.IGNORECASE,
)

def _bleach_attributes() -> Dict[str, List[str]]:
    attrs: Dict[str, List[str]] = {}
    for tag, allowed in ALLOWED_ATTRS.items():
        attrs[tag] = sorted(allowed)
    return attrs

def _filter_comments(html: str) -> str:
    def _replace(match: re.Match[str]) -> str:
        data = match.group(1)
        if _REPEAT_COMMENT_RE.match(data):
            return f"<!--{data}-->"
        return ""

    return _COMMENT_RE.sub(_replace, html)

def _fix_fixed_footers(html: str) -> str:
    """Rewrite ``position: fixed`` footer rules to normal document flow.

    LLM-generated templates sometimes use ``position: fixed`` on footer
    elements.  This causes the footer to overlap table content on long
    reports.  We convert it to a normal-flow footer with a ``@media print``
    rule that uses ``position: fixed`` only in print context (where the
    browser paginates correctly).
    """

    def _rewrite(match: re.Match[str]) -> str:
        selector = match.group(1)
        before = match.group(2)
        after = match.group(3)

        # Remove position:fixed and bottom/left/right from the main rule
        body = before + after
        body = re.sub(r"position\s*:\s*fixed\s*;?", "", body)
        body = _FIXED_POS_PROPS_RE.sub("", body)
        # Clean up stray semicolons, whitespace, and blank lines
        body = re.sub(r";\s*;", ";", body)
        body = re.sub(r"\n\s*\n", "\n", body)
        body = body.strip().strip(";").strip()

        # Build clean flow rule + print-only fixed rule
        if body:
            flow_rule = f"{selector} {{\n  {body};\n  margin-top: 4mm;\n  padding-top: 2mm;\n}}"
        else:
            flow_rule = f"{selector} {{\n  margin-top: 4mm;\n  padding-top: 2mm;\n}}"
        print_rule = (
            f"@media print {{\n"
            f"  {selector} {{ position: fixed; bottom: 0; left: 0; right: 0; }}\n"
            f"}}"
        )
        return f"{flow_rule}\n{print_rule}"

    return _FIXED_FOOTER_RE.sub(_rewrite, html)

def sanitize_html(html: str) -> str:
    cleaned = bleach.clean(
        html or "",
        tags=sorted(ALLOWED_TAGS),
        attributes=_bleach_attributes(),
        protocols=["http", "https", "data"],
        strip=True,
        strip_comments=False,
        css_sanitizer=_CSS_SANITIZER,
    )
    cleaned = _filter_comments(cleaned)

    def _strip_style(match: re.Match[str]) -> str:
        value = match.group(1)
        value = re.sub(r"[;\s]+$", "", value)
        return f'style="{value}"'

    cleaned = _STYLE_ATTR_RE.sub(_strip_style, cleaned)
    cleaned = _fix_fixed_footers(cleaned)
    return cleaned

logger = logging.getLogger("neura.lock")

class TemplateLockError(RuntimeError):
    """Raised when a template lock cannot be acquired."""

    def __init__(self, message: str, lock_holder: Optional[str] = None):
        super().__init__(message)
        self.lock_holder = lock_holder

def _locks_enabled() -> bool:
    disable_flag = os.getenv("NEURA_DISABLE_LOCKS", "").lower()
    if disable_flag in {"1", "true", "yes"}:
        return False
    enable_flag = os.getenv("NEURA_LOCKS_ENABLED", "").lower()
    if enable_flag in {"1", "true", "yes"}:
        return True
    if os.getenv("PYTEST_CURRENT_TEST"):
        return False
    return True

def _acquire_lock(lock_path: Path, *, timeout: float, poll_interval: float) -> FileLock:
    lock = FileLock(str(lock_path))
    lock.acquire(timeout=timeout, poll_interval=poll_interval)
    return lock

class _EagerTemplateLock:
    """Context manager that acquires a file lock eagerly (at construction time).

    Unlike @contextmanager generators which defer to __enter__, this acquires
    the lock immediately so callers can catch TemplateLockError at the call site.
    """

    def __init__(self, lock: FileLock):
        self._lock = lock

    def __enter__(self):
        return self

    def __exit__(self, *args):
        self._lock.release()

def _clean_stale_lock_if_dead(lock_path: Path) -> bool:
    """Remove a lock file if it's older than 30 minutes (likely from a crashed process)."""
    if not lock_path.exists():
        return False
    try:
        import time as _time
        age = _time.time() - lock_path.stat().st_mtime
        if age > 1800:  # 30 minutes
            lock_path.unlink(missing_ok=True)
            logger.warning("stale_lock_removed_by_age", extra={
                "lock": str(lock_path), "age_seconds": int(age),
            })
            return True
    except Exception:
        logger.debug("stale_lock_check_failed", exc_info=True)
    return False


def acquire_template_lock(
    template_dir: Path,
    name: str,
    correlation_id: str | None = None,
    timeout: float = 900.0,
) -> _EagerTemplateLock:
    """Acquire a template lock and return a context manager for releasing it."""
    if not _locks_enabled():
        @contextmanager
        def _noop():
            yield
        return _noop()

    lock_path = Path(template_dir) / f".lock.{name}"
    _clean_stale_lock_if_dead(lock_path)
    holder = f"pid={os.getpid()}"
    if correlation_id:
        holder = f"{holder},corr={correlation_id}"

    logger.info(
        "lock_acquiring",
        extra={
            "event": "lock_acquiring",
            "lock": str(lock_path),
            "holder": holder,
            "correlation_id": correlation_id,
        },
    )

    try:
        lock = _acquire_lock(lock_path, timeout=timeout, poll_interval=0.1)
    except FileLockTimeout as exc:
        raise TemplateLockError(
            f"Failed to acquire template lock '{name}' within {timeout}s.",
            lock_holder="unknown",
        ) from exc

    return _EagerTemplateLock(lock)

@contextmanager
def try_acquire_template_lock(
    template_dir: Path,
    name: str,
    correlation_id: str | None = None,
    timeout: float = 5.0,
) -> Generator[bool, None, None]:
    """
    Non-blocking version that yields True if lock acquired, False otherwise.
    Does not raise an exception on failure.
    """
    if not _locks_enabled():
        yield True
        return

    lock_path = Path(template_dir) / f".lock.{name}"
    holder = f"pid={os.getpid()}"
    if correlation_id:
        holder = f"{holder},corr={correlation_id}"

    try:
        lock = _acquire_lock(lock_path, timeout=timeout, poll_interval=0.1)
        acquired = True
    except FileLockTimeout:
        acquired = False

    try:
        yield acquired
    finally:
        if acquired:
            lock.release()

logger = logging.getLogger("neura.render")

_A4_MM: Tuple[float, float] = (210.0, 297.0)
_MM_PER_INCH = 25.4

def _page_viewport(page_size: str, dpi: int) -> Tuple[int, int]:
    """
    Compute integer viewport (width, height) for the requested page size at the given DPI.
    Currently supports A4 portrait which is our pipeline default.
    """
    size = page_size.upper()
    if size != "A4":
        raise ValueError(f"Unsupported page_size '{page_size}'. Only 'A4' is currently supported.")
    width_px = int(round((_A4_MM[0] / _MM_PER_INCH) * dpi))
    height_px = int(round((_A4_MM[1] / _MM_PER_INCH) * dpi))
    return width_px, height_px

def _ensure_dimensions(path: Path, target: Tuple[int, int]) -> None:
    """
    Pad/crop the rendered screenshot to exactly match the target dimensions.
    """
    if Image is None:  # pragma: no cover
        logger.debug(
            "render_image_adjust_skip",
            extra={"event": "render_image_adjust_skip", "path": str(path)},
        )
        return
    with Image.open(path) as img:
        if img.size == target:
            return
        width, height = target
        if img.width > width or img.height > height:
            img = img.crop((0, 0, width, height))
        if img.width < width or img.height < height:
            canvas = Image.new("RGB", target, "white")
            canvas.paste(img, (0, 0))
            img = canvas
        img.save(path)
        logger.info(
            "render_image_adjusted",
            extra={
                "event": "render_image_adjusted",
                "path": str(path),
                "target_width": width,
                "target_height": height,
            },
        )

def _ensure_playwright_browsers_path() -> None:
    """
    Patch in the system-level Playwright browser cache when packaging omits it.
    """
    if os.environ.get("PLAYWRIGHT_BROWSERS_PATH"):
        return
    local_app = os.getenv("LOCALAPPDATA")
    if not local_app:
        return
    candidate = Path(local_app) / "ms-playwright"
    if candidate.exists():
        os.environ["PLAYWRIGHT_BROWSERS_PATH"] = str(candidate)

def render_html_to_png(
    html_path: Path,
    out_png_path: Path,
    *,
    page_size: str = "A4",
    dpi: int = 400,
    wait_until: str = "networkidle",
) -> None:
    """
    Render HTML to a PNG using Playwright's sync API.
    Ensures deterministic viewport sizing so that SSIM comparisons are meaningful.
    """
    if sync_playwright is None:  # pragma: no cover
        raise RuntimeError(
            "playwright is required for HTML rendering. Install with `pip install playwright` "
            "and run `playwright install chromium`."
        )

    _ensure_playwright_browsers_path()

    html_path = Path(html_path).resolve()
    out_png_path = Path(out_png_path).resolve()
    out_png_path.parent.mkdir(parents=True, exist_ok=True)

    viewport = _page_viewport(page_size, dpi)

    with sync_playwright() as pw:
        browser = pw.chromium.launch()
        try:
            page = browser.new_page(
                viewport={"width": viewport[0], "height": viewport[1]},
                device_scale_factor=1,
            )
            page.goto(f"file://{html_path}", wait_until=wait_until)
            page.screenshot(path=str(out_png_path), full_page=True)
        finally:
            browser.close()

    _ensure_dimensions(out_png_path, viewport)

    logger.info(
        "html_rendered_to_png",
        extra={
            "event": "html_rendered_to_png",
            "html_path": str(html_path),
            "png_path": str(out_png_path),
            "page_size": page_size,
            "dpi": dpi,
        },
    )

__all__ = ["render_html_to_png"]

logger = logging.getLogger("neura.mailer")

def _env_bool(name: str, default: bool = True) -> bool:
    raw = os.getenv(name)
    if raw is None:
        return default
    return raw.strip().lower() in {"1", "true", "yes", "on"}

def _env_int(name: str, default: int) -> int:
    raw = os.getenv(name)
    if raw is None:
        return default
    try:
        return int(raw)
    except ValueError:
        logger.warning("invalid_mail_port", extra={"event": "invalid_mail_port", "env": name, "value": raw})
        return default

@dataclass(frozen=True)
class MailerConfig:
    host: str | None
    port: int
    username: str | None
    password: str | None
    sender: str | None
    use_tls: bool
    enabled: bool

def _load_from_state_store() -> dict | None:
    """Try to load SMTP settings from the persistent state store."""
    try:
        from backend.app.services.config import get_user_preferences
        prefs = get_user_preferences()
        smtp = prefs.get("smtp")
        if smtp and isinstance(smtp, dict) and smtp.get("host"):
            return smtp
    except Exception:
        pass
    return None

def _load_mailer_config() -> MailerConfig:
    # Priority: state store (UI-configured) > environment variables
    stored = _load_from_state_store()
    if stored:
        host = stored.get("host")
        sender = stored.get("sender")
        username = stored.get("username")
        password = stored.get("password")
        port = int(stored.get("port") or 587)
        use_tls = bool(stored.get("use_tls", True))
        source = "state_store"
    else:
        host = os.getenv("NEURA_MAIL_HOST")
        sender = os.getenv("NEURA_MAIL_SENDER")
        username = os.getenv("NEURA_MAIL_USERNAME")
        password = os.getenv("NEURA_MAIL_PASSWORD")
        port = _env_int("NEURA_MAIL_PORT", 587)
        use_tls = _env_bool("NEURA_MAIL_USE_TLS", True)
        source = "env"

    enabled = bool(host and sender)
    if not enabled:
        logger.info(
            "mail_disabled",
            extra={"event": "mail_disabled", "reason": "host_or_sender_missing", "host": bool(host), "sender": bool(sender), "source": source},
        )
    else:
        logger.info("mail_configured", extra={"event": "mail_configured", "source": source, "host": host})
    return MailerConfig(
        host=host,
        port=port,
        username=username,
        password=password,
        sender=sender,
        use_tls=use_tls,
        enabled=enabled,
    )

# Lazy-loaded: None means "not yet loaded".  At import time the state store
# may not be initialised yet (e.g. PyInstaller desktop app), so we defer the
# first load to the moment the config is actually needed.  This guarantees
# stored SMTP settings are picked up even after a cold restart.
MAILER_CONFIG: MailerConfig | None = None

def _get_config() -> MailerConfig:
    """Return the current config, lazy-loading on first access."""
    global MAILER_CONFIG
    if MAILER_CONFIG is None:
        MAILER_CONFIG = _load_mailer_config()
    return MAILER_CONFIG

def refresh_mailer_config() -> MailerConfig:
    global MAILER_CONFIG
    MAILER_CONFIG = _load_mailer_config()
    return MAILER_CONFIG

def _normalize_recipients(recipients: Iterable[str] | None) -> list[str]:
    if not recipients:
        return []
    seen: set[str] = set()
    normalized: list[str] = []
    for raw in recipients:
        text = str(raw or "").strip()
        if not text or text in seen:
            continue
        seen.add(text)
        normalized.append(text)
    return normalized

def send_report_email(
    *,
    to_addresses: Sequence[str],
    subject: str,
    body: str,
    attachments: Sequence[Path] | None = None,
) -> bool:
    config = _get_config()
    recipients = _normalize_recipients(to_addresses)
    if not recipients:
        return False
    if not config.enabled or not config.host or not config.sender:
        logger.warning(
            "mail_not_configured",
            extra={"event": "mail_not_configured", "recipients": len(recipients)},
        )
        return False

    message = EmailMessage()
    message["From"] = config.sender
    message["To"] = ", ".join(recipients)
    message["Subject"] = subject
    message.set_content(body)

    if attachments:
        for path in attachments:
            if not path:
                continue
            try:
                resolved = Path(path).resolve(strict=True)
            except (FileNotFoundError, OSError):
                logger.warning(
                    "mail_attachment_missing",
                    extra={"event": "mail_attachment_missing", "path": str(path)},
                )
                continue
            mime_type, encoding = mimetypes.guess_type(str(resolved))
            maintype = "application"
            subtype = "octet-stream"
            if mime_type and "/" in mime_type:
                maintype, subtype = mime_type.split("/", 1)
            try:
                data = resolved.read_bytes()
            except OSError:
                logger.warning(
                    "mail_attachment_read_failed",
                    extra={"event": "mail_attachment_read_failed", "path": str(resolved)},
                )
                continue
            message.add_attachment(data, maintype=maintype, subtype=subtype, filename=resolved.name)

    try:
        if config.use_tls:
            with smtplib.SMTP(config.host, config.port, timeout=15) as client:
                context = ssl.create_default_context()
                client.starttls(context=context)
                if config.username and config.password:
                    client.login(config.username, config.password)
                client.send_message(message)
        else:
            with smtplib.SMTP(config.host, config.port, timeout=15) as client:
                if config.username and config.password:
                    client.login(config.username, config.password)
                client.send_message(message)
    except Exception:
        logger.exception(
            "mail_send_failed",
            extra={"event": "mail_send_failed", "recipients": len(recipients)},
        )
        return False

    logger.info(
        "mail_sent",
        extra={"event": "mail_sent", "recipients": len(recipients)},
    )
    return True

def _is_safe_member(member: zipfile.ZipInfo) -> bool:
    name = member.filename.replace("\\", "/")
    if name.startswith(("/", "\\")):
        return False
    parts = [p for p in name.split("/") if p and p not in (".", "..")]
    if not parts:
        return False
    if ":" in parts[0]:
        return False
    return not any(p in ("..",) for p in parts)

def _is_within_dir(base_dir: Path, target: Path) -> bool:
    try:
        target.resolve().relative_to(base_dir)
        return True
    except Exception:
        return False

def detect_zip_root(members: Iterable[str]) -> Optional[str]:
    roots = []
    for m in members:
        p = m.replace("\\", "/")
        if p.endswith("/"):
            p = p[:-1]
        parts = [seg for seg in p.split("/") if seg]
        if not parts:
            continue
        roots.append(parts[0])
    if not roots:
        return None
    first = roots[0]
    if all(r == first for r in roots):
        return first
    return None

def create_zip_from_dir(src_dir: Path, dest_zip: Path, *, include_root: bool = True) -> Path:
    src_dir = Path(src_dir).resolve()
    dest_zip = Path(dest_zip).resolve()
    dest_zip.parent.mkdir(parents=True, exist_ok=True)

    with zipfile.ZipFile(dest_zip, mode="w", compression=zipfile.ZIP_DEFLATED) as zf:
        base_prefix = src_dir.name if include_root else ""
        for path in src_dir.rglob("*"):
            if path.is_dir():
                continue
            # Skip lock files
            if path.name.startswith(".lock."):
                continue
            rel = path.relative_to(src_dir)
            arcname = f"{base_prefix}/{rel.as_posix()}" if base_prefix else rel.as_posix()
            zf.write(path, arcname=arcname)
    return dest_zip

def extract_zip_to_dir(
    zip_path: Path,
    dest_dir: Path,
    *,
    strip_root: bool = True,
    max_entries: int | None = None,
    max_uncompressed_bytes: int | None = None,
    max_file_bytes: int | None = None,
) -> Tuple[Path, str]:
    """
    Extract zip into dest_dir safely. Returns (extracted_root, root_name_in_zip).
    If strip_root is True and the zip has a single top-level folder, the contents of that
    folder are extracted directly under dest_dir. Otherwise, members are extracted with their
    original paths.
    """
    dest_dir = Path(dest_dir).resolve()
    dest_dir.mkdir(parents=True, exist_ok=True)
    base_dir = dest_dir.resolve()

    with zipfile.ZipFile(zip_path, mode="r") as zf:
        members = list(zf.infolist())
        file_members = [member for member in members if not member.is_dir()]

        unsafe_members = [member.filename for member in members if not _is_safe_member(member)]
        if unsafe_members:
            sample = ", ".join(unsafe_members[:5])
            suffix = f" (+{len(unsafe_members) - 5} more)" if len(unsafe_members) > 5 else ""
            raise ValueError(f"Zip contains unsafe members: {sample}{suffix}")

        if max_entries is not None and len(file_members) > max_entries:
            raise ValueError(f"Zip contains {len(file_members)} files, exceeds limit {max_entries}")

        if max_uncompressed_bytes is not None:
            total_uncompressed = sum(member.file_size for member in file_members)
            if total_uncompressed > max_uncompressed_bytes:
                raise ValueError(
                    f"Zip uncompressed size {total_uncompressed} bytes exceeds limit {max_uncompressed_bytes}"
                )

        if max_file_bytes is not None:
            for member in file_members:
                if member.file_size > max_file_bytes:
                    raise ValueError(
                        f"Zip member {member.filename} size {member.file_size} exceeds limit {max_file_bytes}"
                    )

        root = detect_zip_root(m.filename for m in members)
        for member in members:
            name = member.filename.replace("\\", "/")
            if strip_root and root and name.startswith(root + "/"):
                rel = name[len(root) + 1 :]
            else:
                rel = name
            if not rel:
                continue
            target = dest_dir / rel
            if not _is_within_dir(base_dir, target):
                raise ValueError(f"Zip member {member.filename} would extract outside {base_dir}")
            if member.is_dir():
                target.mkdir(parents=True, exist_ok=True)
                continue
            target.parent.mkdir(parents=True, exist_ok=True)
            with zf.open(member, "r") as src, open(target, "wb") as dst:
                shutil.copyfileobj(src, dst, length=1024 * 1024)
    return dest_dir, (root or "")

SCHEMA_DIR = Path(__file__).resolve().parents[2] / "schemas"
JSON_SCHEMA_DIR = Path(__file__).resolve().parent / "json_schemas"

class SchemaValidationError(ValueError):
    pass

def _run_jsonschema_validator(validator, data: Any, label: str) -> None:
    """Run a jsonschema validator, raising SchemaValidationError on first error."""
    if validator is None:
        raise RuntimeError(f"jsonschema is required to validate {label} payloads.")
    errors = sorted(validator.iter_errors(data), key=lambda err: list(err.path))
    if errors:
        err = errors[0]
        path = ".".join(str(p) for p in err.path)
        location = f" at {path}" if path else ""
        raise SchemaValidationError(f"{label} validation error{location}: {err.message}")

_DIRECT_COLUMN_RE = re.compile(r"^\s*(?P<table>[A-Za-z_][\w]*)\s*\.\s*(?P<column>[A-Za-z_][\w]*)\s*$")

def _infer_parent_table_from_mapping(mapping: Mapping[str, Any] | None) -> str | None:
    if not isinstance(mapping, Mapping):
        return None
    for expr in mapping.values():
        if not isinstance(expr, str):
            continue
        match = _DIRECT_COLUMN_RE.match(expr.strip())
        if not match:
            continue
        table_name = match.group("table").strip(' "`[]')
        if table_name and not table_name.lower().startswith("params"):
            return table_name
    return None

def _load_schema(name: str) -> dict:
    path = SCHEMA_DIR / name
    return json.loads(path.read_text(encoding="utf-8"))

MAPPING_SCHEMA = _load_schema("mapping_pdf_labels.schema.json")
CONTRACT_SCHEMA = _load_schema("contract.schema.json")

def _coerce_join_block(data: Mapping[str, Any]) -> dict[str, Any] | None:
    if not isinstance(data, Mapping):
        return None
    join_raw = data.get("join")
    join = dict(join_raw) if isinstance(join_raw, Mapping) else {}
    parent_table = str(join.get("parent_table") or "").strip()
    parent_key = str(join.get("parent_key") or "").strip()
    if not parent_table:
        inferred = _infer_parent_table_from_mapping(data.get("mapping"))
        if inferred:
            parent_table = inferred
    if not parent_table:
        return None
    if not parent_key:
        parent_key = "__rowid__"
    child_table = str(join.get("child_table") or "").strip()
    child_key = str(join.get("child_key") or "").strip()
    join.update(
        {
            "parent_table": parent_table,
            "parent_key": parent_key,
            "child_table": child_table,
            "child_key": child_key,
        }
    )
    return join

if Draft7Validator is not None:
    _MAPPING_INLINE_V4_SCHEMA = json.loads(
        (JSON_SCHEMA_DIR / "mapping_inline_v4.schema.json").read_text(encoding="utf-8")
    )
    _MAPPING_INLINE_V4_VALIDATOR = Draft7Validator(_MAPPING_INLINE_V4_SCHEMA)
    _LLM_CALL_3_5_SCHEMA = json.loads((JSON_SCHEMA_DIR / "llm_call_3_5.schema.json").read_text(encoding="utf-8"))
    _LLM_CALL_3_5_VALIDATOR = Draft7Validator(_LLM_CALL_3_5_SCHEMA)
    _CONTRACT_V2_SCHEMA = json.loads((JSON_SCHEMA_DIR / "contract_v2.schema.json").read_text(encoding="utf-8"))
    _CONTRACT_V2_VALIDATOR = Draft7Validator(_CONTRACT_V2_SCHEMA)
    _CONTRACT_V2_OPTIONAL_JOIN_SCHEMA = deepcopy(_CONTRACT_V2_SCHEMA)
    required_fields = _CONTRACT_V2_OPTIONAL_JOIN_SCHEMA.get("required")
    if isinstance(required_fields, list) and "join" in required_fields:
        _CONTRACT_V2_OPTIONAL_JOIN_SCHEMA["required"] = [field for field in required_fields if field != "join"]
    join_schema = _CONTRACT_V2_OPTIONAL_JOIN_SCHEMA.get("properties", {}).get("join")
    if isinstance(join_schema, dict):
        join_required = join_schema.get("required")
        if isinstance(join_required, list):
            join_schema["required"] = [field for field in join_required if field in ("parent_table", "parent_key")]
    _CONTRACT_V2_OPTIONAL_JOIN_VALIDATOR = Draft7Validator(_CONTRACT_V2_OPTIONAL_JOIN_SCHEMA)
    _STEP5_REQUIREMENTS_SCHEMA = json.loads(
        (JSON_SCHEMA_DIR / "step5_requirements.schema.json").read_text(encoding="utf-8")
    )
    _STEP5_REQUIREMENTS_VALIDATOR = Draft7Validator(_STEP5_REQUIREMENTS_SCHEMA)
    _GENERATOR_SQL_PACK_SCHEMA = json.loads(
        (JSON_SCHEMA_DIR / "generator_sql_pack.schema.json").read_text(encoding="utf-8")
    )
    _GENERATOR_SQL_PACK_VALIDATOR = Draft7Validator(_GENERATOR_SQL_PACK_SCHEMA)
    _GENERATOR_OUTPUT_SCHEMAS_SCHEMA = json.loads(
        (JSON_SCHEMA_DIR / "generator_output_schemas.schema.json").read_text(encoding="utf-8")
    )
    _GENERATOR_OUTPUT_SCHEMAS_VALIDATOR = Draft7Validator(_GENERATOR_OUTPUT_SCHEMAS_SCHEMA)
    _GENERATOR_LLM_RESPONSE_SCHEMA = json.loads(
        (JSON_SCHEMA_DIR / "generator_llm_response.schema.json").read_text(encoding="utf-8")
    )
    _GENERATOR_LLM_RESPONSE_VALIDATOR = Draft7Validator(_GENERATOR_LLM_RESPONSE_SCHEMA)
else:  # pragma: no cover - optional dependency missing
    _MAPPING_INLINE_V4_SCHEMA = None
    _MAPPING_INLINE_V4_VALIDATOR = None
    _LLM_CALL_3_5_SCHEMA = None
    _LLM_CALL_3_5_VALIDATOR = None
    _CONTRACT_V2_SCHEMA = None
    _CONTRACT_V2_VALIDATOR = None
    _CONTRACT_V2_OPTIONAL_JOIN_SCHEMA = None
    _CONTRACT_V2_OPTIONAL_JOIN_VALIDATOR = None
    _STEP5_REQUIREMENTS_SCHEMA = None
    _STEP5_REQUIREMENTS_VALIDATOR = None
    _GENERATOR_SQL_PACK_SCHEMA = None
    _GENERATOR_SQL_PACK_VALIDATOR = None
    _GENERATOR_OUTPUT_SCHEMAS_SCHEMA = None
    _GENERATOR_OUTPUT_SCHEMAS_VALIDATOR = None
    _GENERATOR_LLM_RESPONSE_SCHEMA = None
    _GENERATOR_LLM_RESPONSE_VALIDATOR = None

def validate_mapping_schema(data: Any) -> None:
    if not isinstance(data, list):
        raise SchemaValidationError("mapping must be a list")
    for idx, item in enumerate(data):
        if not isinstance(item, dict):
            raise SchemaValidationError(f"mapping[{idx}] must be an object")
        for key in ("header", "placeholder", "mapping"):
            if key not in item or not isinstance(item[key], str) or not item[key].strip():
                raise SchemaValidationError(f"mapping[{idx}].{key} must be a non-empty string")

def _stringify_scalar(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, str):
        return value.strip()
    return str(value).strip()

def _as_sequence(value: Any) -> list[Any]:
    if value is None:
        return []
    if isinstance(value, (list, tuple, set)):
        return list(value)
    if isinstance(value, Iterable) and not isinstance(value, (str, bytes, dict)):
        return list(value)
    return [value]

def _flatten_over_value(value: Any) -> list[str]:
    if value is None:
        return []
    if isinstance(value, dict):
        flattened: list[str] = []
        for sub_key, sub_val in value.items():
            for item in _flatten_over_value(sub_val):
                entry = f"{sub_key}:{item}" if item else str(sub_key)
                entry = entry.strip()
                if entry:
                    flattened.append(entry)
        return flattened
    items: list[str] = []
    for item in _as_sequence(value):
        if isinstance(item, dict):
            items.extend(_flatten_over_value(item))
            continue
        entry = _stringify_scalar(item)
        if entry:
            items.append(entry)
    return items

def _normalize_string_list(values: Any) -> list[str]:
    items = _as_sequence(values)
    normalized: list[str] = []
    seen: set[str] = set()
    for item in items:
        if isinstance(item, str):
            text = item.strip()
        else:
            text = str(item or "").strip()
        if not text or text in seen:
            continue
        seen.add(text)
        normalized.append(text)
    return normalized

def _normalize_hint_value(raw_hint: Any) -> dict[str, Any]:
    if not isinstance(raw_hint, dict):
        op = _stringify_scalar(raw_hint) or "UNKNOWN"
        return {"op": op, "over": []}

    op = _stringify_scalar(raw_hint.get("op")) or "UNKNOWN"
    over_entries: list[str] = []

    over_entries.extend(_flatten_over_value(raw_hint.get("over")))

    if "over_a" in raw_hint or "over_b" in raw_hint:
        seq_a = [_stringify_scalar(item) for item in _as_sequence(raw_hint.get("over_a"))]
        seq_b = [_stringify_scalar(item) for item in _as_sequence(raw_hint.get("over_b"))]
        for a, b in zip_longest(seq_a, seq_b, fillvalue=""):
            a = a.strip()
            b = b.strip()
            if a and b:
                over_entries.append(f"{a} - {b}")
            elif a:
                over_entries.append(a)
            elif b:
                over_entries.append(f"- {b}")

    for key in ("num_ref", "den_ref", "formula"):
        val = raw_hint.get(key)
        text = _stringify_scalar(val)
        if text:
            over_entries.append(f"{key}:{text}")

    allowed_keys = {"op", "over", "over_a", "over_b", "num_ref", "den_ref", "formula"}
    for extra_key, extra_value in raw_hint.items():
        if extra_key in allowed_keys or extra_value is None:
            continue
        if isinstance(extra_value, dict):
            for sub_key, sub_value in extra_value.items():
                for item in _flatten_over_value({sub_key: sub_value}):
                    over_entries.append(f"{extra_key}.{item}")
        else:
            for item in _flatten_over_value(extra_value):
                over_entries.append(f"{extra_key}:{item}")

    cleaned: list[str] = []
    seen: set[str] = set()
    for entry in over_entries:
        normalized = entry.strip()
        if not normalized or normalized in seen:
            continue
        seen.add(normalized)
        cleaned.append(normalized)

    return {"op": op, "over": cleaned}

def normalize_mapping_inline_payload(payload: Any) -> Any:
    if not isinstance(payload, dict):
        return payload

    # LLMs sometimes return a top-level "constants" key for tokens they
    # consider static.  Merge those entries into "mapping" so the pipeline
    # handles them correctly, then drop the key before schema validation.
    constants = payload.pop("constants", None)
    if isinstance(constants, dict):
        mapping = payload.setdefault("mapping", {})
        for key, value in constants.items():
            if key not in mapping:
                mapping[str(key)] = str(value)

    meta = payload.get("meta")
    if not isinstance(meta, dict):
        return payload

    hints = meta.get("hints")
    if not isinstance(hints, dict):
        return payload

    normalized: dict[str, dict[str, Any]] = {}
    for key, value in hints.items():
        normalized[str(key)] = _normalize_hint_value(value)

    meta["hints"] = normalized
    return payload

def validate_contract_schema(data: Any) -> None:
    if not isinstance(data, dict):
        raise SchemaValidationError("contract must be an object")
    if "literals" not in data:
        data["literals"] = {}
    coerced_join = _coerce_join_block(data)
    if coerced_join is not None:
        data["join"] = coerced_join
    required = CONTRACT_SCHEMA["required"]
    for key in required:
        if key not in data:
            raise SchemaValidationError(f"contract missing key '{key}'")
    if not isinstance(data["mapping"], dict):
        raise SchemaValidationError("contract.mapping must be an object")
    if not isinstance(data["join"], dict):
        raise SchemaValidationError("contract.join must be an object")
    join = data["join"]
    for key in ("parent_table", "parent_key"):
        value = join.get(key)
        if not isinstance(value, str) or not value.strip():
            raise SchemaValidationError(f"contract.join.{key} must be a non-empty string")

    child_table_raw = join.get("child_table")
    if child_table_raw is not None and not isinstance(child_table_raw, str):
        raise SchemaValidationError("contract.join.child_table must be a string or null")
    child_key_raw = join.get("child_key")
    if child_key_raw is not None and not isinstance(child_key_raw, str):
        raise SchemaValidationError("contract.join.child_key must be a string or null")

    child_table_text = child_table_raw.strip() if isinstance(child_table_raw, str) else ""
    child_key_text = child_key_raw.strip() if isinstance(child_key_raw, str) else ""
    if child_table_text and not child_key_text:
        raise SchemaValidationError("contract.join.child_key must be a non-empty string when child_table is provided")
    for key in ("date_columns", "totals", "literals"):
        if not isinstance(data[key], dict):
            raise SchemaValidationError(f"contract.{key} must be an object")
    for key in ("header_tokens", "row_tokens", "row_order"):
        arr = data.get(key)
        if not isinstance(arr, list) or not all(isinstance(item, str) for item in arr):
            raise SchemaValidationError(f"contract.{key} must be an array of strings")

def validate_mapping_inline_v4(data: Any) -> None:
    """Validate LLM Call 3 output against the mapping_inline_v4 schema."""
    data = normalize_mapping_inline_payload(data)
    _run_jsonschema_validator(_MAPPING_INLINE_V4_VALIDATOR, data, "mapping_inline_v4")

def validate_llm_call_3_5(data: Any) -> None:
    """Validate LLM Call 3.5 response against the schema."""
    _run_jsonschema_validator(_LLM_CALL_3_5_VALIDATOR, data, "llm_call_3_5")

def validate_contract_v2(data: Any, *, require_join: bool = True) -> None:
    """Validate contract.json produced by LLM Call 4."""
    validator = _CONTRACT_V2_VALIDATOR if (require_join or _CONTRACT_V2_OPTIONAL_JOIN_VALIDATOR is None) else _CONTRACT_V2_OPTIONAL_JOIN_VALIDATOR
    _run_jsonschema_validator(validator, data, "contract_v2")

    reshape_rules = data.get("reshape_rules")
    if isinstance(reshape_rules, list) and reshape_rules:
        column_rule_found = False
        for idx, rule in enumerate(reshape_rules):
            if not isinstance(rule, Mapping):
                continue
            columns = rule.get("columns")
            if columns is None:
                continue
            if not isinstance(columns, list) or not columns:
                raise SchemaValidationError(
                    f"contract.reshape_rules[{idx}].columns must be a non-empty array when provided"
                )
            for col_idx, column in enumerate(columns):
                if not isinstance(column, Mapping):
                    raise SchemaValidationError(f"contract.reshape_rules[{idx}].columns[{col_idx}] must be an object")
                alias = column.get("as")
                if not isinstance(alias, str) or not alias.strip():
                    raise SchemaValidationError(
                        f"contract.reshape_rules[{idx}].columns[{col_idx}].as must be a non-empty string"
                    )
            column_rule_found = True
        if not column_rule_found:
            raise SchemaValidationError("contract.reshape_rules must include at least one rule with column definitions")

    join = data.get("join")
    if isinstance(join, dict):
        parent_table = join.get("parent_table")
        if not isinstance(parent_table, str) or not parent_table.strip():
            raise SchemaValidationError("contract.join.parent_table must be a non-empty string")
        parent_key = join.get("parent_key")
        if not isinstance(parent_key, str) or not parent_key.strip():
            raise SchemaValidationError("contract.join.parent_key must be a non-empty string")

        child_table = join.get("child_table")
        child_key = join.get("child_key")
        child_table_text = child_table.strip() if isinstance(child_table, str) else ""
        if child_table_text and (not isinstance(child_key, str) or not child_key.strip()):
            raise SchemaValidationError(
                "contract.join.child_key must be a non-empty string when child_table is provided"
            )

def validate_step5_requirements(data: Any) -> None:
    """Validate step5_requirements.json produced by LLM Call 4."""
    _run_jsonschema_validator(_STEP5_REQUIREMENTS_VALIDATOR, data, "step5_requirements")

def validate_generator_sql_pack(data: Any) -> None:
    """Validate the sql_pack section returned by LLM Call 5."""
    _run_jsonschema_validator(_GENERATOR_SQL_PACK_VALIDATOR, data, "generator_sql_pack")
    dialect = str(data.get("dialect") or "").strip().lower()
    if not dialect or dialect == "sqlite":
        raise SchemaValidationError(
            "generator_sql_pack.dialect must be 'duckdb' or 'postgres' (SQLite SQL is no longer supported)"
        )
    if dialect not in {"duckdb", "postgres"}:
        raise SchemaValidationError(f"generator_sql_pack.dialect '{data.get('dialect')}' is not supported")

def validate_generator_output_schemas(data: Any) -> None:
    """Validate the output_schemas section returned by LLM Call 5."""
    _run_jsonschema_validator(_GENERATOR_OUTPUT_SCHEMAS_VALIDATOR, data, "generator_output_schemas")

def validate_generator_llm_response(data: Any) -> None:
    """Validate the full LLM Call 5 response payload."""
    _run_jsonschema_validator(_GENERATOR_LLM_RESPONSE_VALIDATOR, data, "generator_llm_response")

    if "key_tokens" in data:
        tokens = data.get("key_tokens")
        if not isinstance(tokens, list):
            raise SchemaValidationError("generator_llm_response.key_tokens must be an array of strings")
        cleaned = _normalize_string_list(tokens)
        if len(cleaned) != len(tokens):
            raise SchemaValidationError("generator_llm_response.key_tokens must contain unique, non-empty strings")
        for idx, token in enumerate(tokens):
            if not isinstance(token, str):
                raise SchemaValidationError(f"generator_llm_response.key_tokens[{idx}] must be a string")
            if token.strip() != token:
                raise SchemaValidationError(
                    f"generator_llm_response.key_tokens[{idx}] must not contain leading or trailing whitespace"
                )

# ── LLM Utilities ──
logger = logging.getLogger("neura.llm")

_LOG_PATH_ENV = os.getenv("LLM_RAW_OUTPUT_PATH")
if _LOG_PATH_ENV:
    _RAW_OUTPUT_PATH = Path(_LOG_PATH_ENV).expanduser()
else:
    _RAW_OUTPUT_PATH = Path(__file__).resolve().parents[3] / "llm_raw_outputs.md"

_RAW_OUTPUT_LOCK = threading.Lock()

def _llm_repair_json(
    raw_text: str,
    llm_client: Any,
    *,
    expect_type: str = "dict",
) -> Any | None:
    """Ask the LLM to extract/repair JSON from its own malformed output.

    Returns the parsed result on success, or None on failure.
    This is a last-resort fallback after all regex/brace strategies fail.
    """
    truncated = raw_text[:4000]
    type_hint = "JSON object (dict)" if expect_type == "dict" else "JSON array (list)"
    repair_prompt = (
        f"The following text was supposed to contain a {type_hint} but "
        "JSON parsing failed. Extract or repair the JSON and return ONLY "
        f"the valid {type_hint}, nothing else — no markdown, no explanation:\n\n"
        + truncated
    )
    try:
        from backend.app.services.llm import LLMClient

        if isinstance(llm_client, LLMClient):
            resp = llm_client.complete(
                messages=[{"role": "user", "content": repair_prompt}],
                description="json_self_repair",
            )
            repaired_text = resp if isinstance(resp, str) else getattr(resp, "content", str(resp))
        else:
            # Fallback: try calling directly if it quacks like a client
            resp = llm_client.complete(
                messages=[{"role": "user", "content": repair_prompt}],
                description="json_self_repair",
            )
            repaired_text = resp if isinstance(resp, str) else getattr(resp, "content", str(resp))

        result = json.loads(repaired_text.strip())
        logger.info("json_self_repair succeeded for %s", expect_type)
        return result
    except Exception as exc:
        logger.debug("json_self_repair failed: %s", exc)
        return None

def extract_json_from_llm_response(
    content: str,
    default: Optional[Dict[str, Any]] = None,
    *,
    llm_client: Any = None,
) -> Dict[str, Any]:
    """Extract JSON from an LLM response that may contain markdown code blocks."""
    if default is None:
        default = {}

    if not content or not content.strip():
        return default

    cleaned = content.strip()

    # Handle ```json ... ``` or ``` ... ``` code blocks
    # Match both with and without 'json' language tag
    json_block_match = re.search(r"```(?:json)?\s*\n?([\s\S]*?)\n?```", cleaned)
    if json_block_match:
        cleaned = json_block_match.group(1).strip()
    elif cleaned.startswith("```"):
        # Fallback: manually strip opening/closing fences
        parts = cleaned.split("```", 2)
        if len(parts) >= 2:
            cleaned = parts[1].strip()
            # Remove 'json' language tag if present at start
            if cleaned.startswith("json"):
                cleaned = cleaned[4:].strip()

    # Try direct parse first
    try:
        result = json.loads(cleaned)
        if isinstance(result, dict):
            return result
        if isinstance(result, list):
            logger.debug("extract_json: expected dict, got list — wrapping as {data: [...]}")
            return {"data": result}
        # Scalar value — not useful as dict
        logger.debug("extract_json: expected dict, got scalar %s — returning default", type(result).__name__)
        return default
    except json.JSONDecodeError:
        pass

    # Try to find JSON object in the content
    # Look for outermost balanced braces
    start_idx = cleaned.find("{")
    if start_idx != -1:
        depth = 0
        in_string = False
        escape_next = False

        for i, char in enumerate(cleaned[start_idx:], start_idx):
            if escape_next:
                escape_next = False
                continue
            if char == "\\":
                escape_next = True
                continue
            if char == '"' and not escape_next:
                in_string = not in_string
                continue
            if in_string:
                continue
            if char == "{":
                depth += 1
            elif char == "}":
                depth -= 1
                if depth == 0:
                    json_str = cleaned[start_idx:i + 1]
                    try:
                        return json.loads(json_str)
                    except json.JSONDecodeError:
                        break

    # Try regex patterns as fallback
    for pattern in [r"\{[\s\S]*\}", r"\[[\s\S]*\]"]:
        match = re.search(pattern, cleaned)
        if match:
            try:
                result = json.loads(match.group())
                if isinstance(result, dict):
                    return result
                return {"data": result}
            except json.JSONDecodeError:
                continue

    # Last resort: ask the LLM to self-repair the JSON
    if llm_client is not None:
        repaired = _llm_repair_json(content, llm_client, expect_type="dict")
        if isinstance(repaired, dict):
            return repaired
        if isinstance(repaired, list):
            logger.debug("json_self_repair: expected dict, got list — wrapping")
            return {"data": repaired}

    logger.debug(f"Failed to extract JSON from LLM response: {content[:200]}...")
    return default

def extract_json_array_from_llm_response(
    content: str,
    default: Optional[list] = None,
    *,
    llm_client: Any = None,
) -> list:
    """Extract a JSON array from an LLM response."""
    if default is None:
        default = []

    if not content or not content.strip():
        return default

    cleaned = content.strip()

    # Handle code blocks
    json_block_match = re.search(r"```(?:json)?\s*\n?([\s\S]*?)\n?```", cleaned)
    if json_block_match:
        cleaned = json_block_match.group(1).strip()

    # Try direct parse
    try:
        result = json.loads(cleaned)
        if isinstance(result, list):
            return result
    except json.JSONDecodeError:
        pass

    # Try to find array in content
    match = re.search(r"\[[\s\S]*\]", cleaned)
    if match:
        try:
            result = json.loads(match.group())
            if isinstance(result, list):
                return result
        except json.JSONDecodeError:
            pass

    # Last resort: ask the LLM to self-repair the JSON
    if llm_client is not None:
        repaired = _llm_repair_json(content, llm_client, expect_type="list")
        if isinstance(repaired, list):
            return repaired

    return default

def validate_llm_json(
    raw: str,
    required_keys: set[str] | None = None,
) -> dict:
    """Parse LLM response, strip fences, validate required keys.

    Raises ValueError on failure (empty, non-dict, missing keys).
    """
    from ..utils.core import strip_code_fences

    if not raw or not raw.strip():
        raise ValueError("LLM returned empty response")

    cleaned = strip_code_fences(raw).strip()

    # Find JSON object if response has surrounding text
    if cleaned and not cleaned.startswith("{"):
        start = cleaned.find("{")
        end = cleaned.rfind("}")
        if start >= 0 and end > start:
            cleaned = cleaned[start : end + 1]

    try:
        result = json.loads(cleaned)
    except json.JSONDecodeError as exc:
        raise ValueError(f"LLM response is not valid JSON: {exc}") from exc

    if isinstance(result, list):
        logger.warning("validate_llm_json: expected dict, got list — wrapping")
        result = {"items": result}
    elif not isinstance(result, dict):
        raise ValueError(f"LLM response is {type(result).__name__}, expected dict")

    if required_keys:
        missing = required_keys - set(result.keys())
        if missing:
            raise ValueError(f"LLM response missing required keys: {missing}")

    return result

def _coerce_jsonable(value: Any) -> Any:
    """Best-effort conversion of LLM responses to JSON-serialisable data."""
    if isinstance(value, (str, int, float, bool)) or value is None:
        return value

    if isinstance(value, dict):
        return {str(k): _coerce_jsonable(v) for k, v in value.items()}

    if isinstance(value, (list, tuple, set)):
        return [_coerce_jsonable(v) for v in value]

    for attr in ("model_dump", "to_dict", "dict"):
        method = getattr(value, attr, None)
        if callable(method):
            try:
                return _coerce_jsonable(method())
            except Exception:
                continue

    json_method = getattr(value, "model_dump_json", None)
    if callable(json_method):
        try:
            return _coerce_jsonable(json.loads(json_method()))
        except Exception:
            pass

    return repr(value)

def _append_raw_output(description: str, response: Any) -> None:
    """Append the raw LLM response to a Markdown log file."""
    timestamp = datetime.now(timezone.utc).isoformat(timespec="seconds") + "Z"
    entry = _coerce_jsonable(response)

    try:
        with _RAW_OUTPUT_LOCK:
            _RAW_OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
            with _RAW_OUTPUT_PATH.open("a", encoding="utf-8") as handle:
                handle.write(f"## {timestamp} - {description}\n\n")
                handle.write("```json\n")
                handle.write(json.dumps(entry, indent=2))
                handle.write("\n```\n\n")
    except Exception as exc:  # pragma: no cover - logging must not break execution
        logger.debug(
            "llm_raw_output_log_failed",
            extra={"event": "llm_raw_output_log_failed"},
            exc_info=(type(exc), exc, exc.__traceback__),
        )

# Public alias for use by other modules
append_raw_llm_output = _append_raw_output

class DictAsObject:
    """Wrapper that allows dict access via both dict keys and object attributes."""

    def __init__(self, data: Any):
        if isinstance(data, dict):
            for key, value in data.items():
                setattr(self, key, DictAsObject(value) if isinstance(value, (dict, list)) else value)
            self._data = data
        elif isinstance(data, list):
            self._data = [DictAsObject(item) if isinstance(item, (dict, list)) else item for item in data]
        else:
            self._data = data

    def __getitem__(self, key):
        if isinstance(self._data, dict):
            value = self._data[key]
            return DictAsObject(value) if isinstance(value, (dict, list)) else value
        elif isinstance(self._data, list):
            value = self._data[key]
            return DictAsObject(value) if isinstance(value, (dict, list)) else value
        raise TypeError(f"Cannot index {type(self._data)}")

    def __iter__(self):
        if isinstance(self._data, list):
            for item in self._data:
                yield DictAsObject(item) if isinstance(item, (dict, list)) else item
        elif isinstance(self._data, dict):
            yield from self._data
        else:
            raise TypeError(f"Cannot iterate {type(self._data)}")

    def __len__(self):
        return len(self._data)

    def get(self, key, default=None):
        if isinstance(self._data, dict):
            value = self._data.get(key, default)
            return DictAsObject(value) if isinstance(value, (dict, list)) else value
        return default

    def __repr__(self):
        return f"DictAsObject({self._data!r})"

def _is_llm_retriable(exc: Exception) -> bool:
    """Classify whether an LLM error is transient and worth retrying."""
    err_msg = str(exc).lower()
    # Rate limit and capacity errors from LLM APIs are always retriable
    if any(p in err_msg for p in ("rate limit", "429", "503", "overloaded", "capacity")):
        return True
    # Check for generic transient indicators
    if any(p in err_msg for p in ("timeout", "connection", "temporary", "502")):
        return True
    return False


def call_chat_completion(
    client: Any,
    *,
    model: str,
    messages: Iterable[Dict[str, Any]],
    description: str,
    timeout: float | None = None,
    **kwargs: Any,
) -> Any:
    """Execute a chat completion using Claude Code CLI."""
    # Use the unified LLM client
    from backend.app.services.llm import get_llm_client, get_adaptive_timeout, LLMClient

    if isinstance(client, LLMClient):
        llm_client = client
    else:
        llm_client = get_llm_client()

    # Convert messages to list if needed
    messages_list = list(messages)

    # Enforce timeout via kwargs if not already set
    if timeout is None:
        timeout = get_adaptive_timeout(description)
    if "timeout" not in kwargs:
        kwargs["timeout"] = timeout

    response = llm_client.complete(
        messages=messages_list,
        model=model,
        description=description,
        **kwargs,
    )

    # Wrap dict response to support both dict and attribute access
    return DictAsObject(response) if isinstance(response, dict) else response

async def call_chat_completion_async(
    client: Any,
    *,
    model: str,
    messages: Iterable[Dict[str, Any]],
    description: str,
    timeout: float | None = None,
    **kwargs: Any,
) -> Any:
    """
    Async wrapper for call_chat_completion using Claude Code CLI.
    """
    import asyncio

    return await asyncio.to_thread(
        call_chat_completion,
        client,
        model=model,
        messages=messages,
        description=description,
        timeout=timeout,
        **kwargs,
    )

def _get_prompt_library():
    """Lazy accessor to avoid circular import with ai_services."""
    from backend.app.services.ai_services import PROMPT_LIBRARY
    return PROMPT_LIBRARY

def load_prompt(key: str, replacements: Dict[str, str] | None = None) -> str:
    """
    Load a prompt by key from the in-memory prompt library and optionally replace tokens.
    """
    pl = _get_prompt_library()
    if key not in pl:
        available = ", ".join(sorted(pl))
        raise KeyError(f"Prompt '{key}' not found. Available keys: {available}")

    prompt = pl[key]
    if replacements:
        for needle, value in replacements.items():
            prompt = prompt.replace(needle, value)
    return prompt

def available_prompts() -> Dict[str, str]:
    """
    Return a copy of the prompt library (key -> prompt text).
    """
    return dict(_get_prompt_library())

